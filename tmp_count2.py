from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

path = r'test_output.pptx'
prs = Presentation(path)
slide_count = len(prs.slides)
charts = 0
pictures = 0
text_summary = []

for idx, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        st = getattr(shape, 'shape_type', None)
        if st == MSO_SHAPE_TYPE.CHART:
            charts += 1
        if st == MSO_SHAPE_TYPE.PICTURE:
            pictures += 1

out = {
    'slide_count': slide_count,
    'chart_count': charts,
    'picture_count': pictures,
}
print(json.dumps(out, ensure_ascii=False, indent=2))