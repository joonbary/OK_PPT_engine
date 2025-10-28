"""간단한 PPTX 검사"""
from pptx import Presentation

file_path = r"C:\Users\apro\Downloads\generated_presentation_3ab66d39-8588-4944-a108-603a862f9f3d_iter2.pptx"
prs = Presentation(file_path)

print(f"Total slides: {len(prs.slides)}")
print("=" * 60)

# Check first 3 slides in detail
slide_count = 0
for slide in prs.slides:
    slide_count += 1
    if slide_count > 3:
        break
    print(f"\nSlide {slide_count}:")
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text[:100]}")
    
    # Get all text content
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
            text = shape.text.strip()
            if text:
                texts.append(text)
    
    if texts:
        print("Content preview:")
        for text in texts[:2]:
            print(f"  {text[:100]}...")
    
    print(f"Total text length: {sum(len(t) for t in texts)} chars")

# Check for patterns
all_text = ""
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            all_text += shape.text + " "

# Check for mock patterns
mock_patterns = ["100.0", "Mock", "PLACEHOLDER", "[keyword]"]
mock_found = []
for pattern in mock_patterns:
    if pattern in all_text:
        mock_found.append(pattern)

print("\n" + "=" * 60)
if mock_found:
    print(f"MOCK PATTERNS FOUND: {', '.join(mock_found)}")
else:
    print("NO MOCK PATTERNS FOUND - CONTENT IS AI GENERATED")

# Check language
if any(ord(c) >= 0xAC00 and ord(c) <= 0xD7AF for c in all_text):
    print("Language: Korean detected")
elif all(ord(c) < 128 for c in all_text if c.isalpha()):
    print("Language: English only")
else:
    print("Language: Mixed or other")