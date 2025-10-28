from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

path = r'test_output.pptx'
prs = Presentation(path)
slide_count = len(prs.slides)
charts = 0
text_summary = []

for idx, slide in enumerate(prs.slides, start=1):
    title = None
    if slide.shapes.title is not None and hasattr(slide.shapes.title, 'text_frame'):
        title = slide.shapes.title.text_frame.text
    texts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            tf = shape.text_frame
            if tf:
                buf = [p.text.strip() for p in tf.paragraphs if p.text]
                if buf:
                    texts.append(" | ".join(buf))
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                charts += 1
        except Exception:
            pass
    text_summary.append({
        'slide': idx,
        'title': (title or '').strip(),
        'first_text': (texts[0] if texts else '')[:120]
    })

out = {
    'slide_count': slide_count,
    'chart_count': charts,
    'slides': text_summary[:6]
}
print(json.dumps(out, ensure_ascii=False, indent=2))