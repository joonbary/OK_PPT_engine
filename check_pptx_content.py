"""
PPTX 파일 내용 검사 스크립트
생성된 PPT의 실제 내용을 확인합니다.
"""

from pptx import Presentation
import os

def analyze_pptx(file_path):
    """PPTX 파일 내용 분석"""
    if not os.path.exists(file_path):
        print(f"파일을 찾을 수 없습니다: {file_path}")
        return
    
    print(f"분석 중: {file_path}")
    print("=" * 80)
    
    # 프레젠테이션 열기
    prs = Presentation(file_path)
    
    print(f"총 슬라이드 수: {len(prs.slides)}")
    print("=" * 80)
    
    # Mock 패턴 감지
    mock_patterns = ["100.0", "[키워드]", "Mock", "템플릿", "PLACEHOLDER", "높은 수준"]
    mock_found = False
    
    # 각 슬라이드 분석
    for idx, slide in enumerate(prs.slides, 1):
        print(f"\n슬라이드 {idx}:")
        print("-" * 40)
        
        # 제목 추출
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text
            print(f"제목: {title}")
            
            # Mock 패턴 확인
            for pattern in mock_patterns:
                if pattern in title:
                    print(f"  [WARNING] Mock 패턴 감지: '{pattern}'")
                    mock_found = True
        
        # 내용 추출
        content_items = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    content_items.append(text[:200] + "..." if len(text) > 200 else text)
                    
                    # Mock 패턴 확인
                    for pattern in mock_patterns:
                        if pattern in text:
                            print(f"  [WARNING] Mock 패턴 감지: '{pattern}'")
                            mock_found = True
        
        if content_items:
            print("내용:")
            for item in content_items[:3]:  # 처음 3개만 출력
                print(f"  - {item}")
            if len(content_items) > 3:
                print(f"  ... 외 {len(content_items) - 3}개 항목")
        
        # 차트/테이블 확인
        charts = 0
        tables = 0
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                charts += 1
            if hasattr(shape, 'has_table') and shape.has_table:
                tables += 1
        
        if charts > 0:
            print(f"차트: {charts}개")
        if tables > 0:
            print(f"테이블: {tables}개")
    
    print("\n" + "=" * 80)
    print("분석 완료")
    
    if mock_found:
        print("[FAIL] Mock 데이터 패턴이 감지되었습니다!")
    else:
        print("[PASS] Mock 데이터 패턴이 발견되지 않았습니다.")
    
    # 콘텐츠 품질 평가
    total_text_length = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                total_text_length += len(shape.text)
    
    avg_text_per_slide = total_text_length / len(prs.slides) if prs.slides else 0
    print(f"\n평균 텍스트 길이: {avg_text_per_slide:.0f}자/슬라이드")
    
    if avg_text_per_slide < 50:
        print("[WARNING] 콘텐츠가 너무 적습니다.")
    elif avg_text_per_slide > 500:
        print("[WARNING] 콘텐츠가 너무 많습니다.")
    else:
        print("[PASS] 적절한 콘텐츠 양입니다.")

if __name__ == "__main__":
    # 분석할 파일 경로
    file_path = r"C:\Users\apro\Downloads\generated_presentation_3ab66d39-8588-4944-a108-603a862f9f3d_iter2.pptx"
    analyze_pptx(file_path)