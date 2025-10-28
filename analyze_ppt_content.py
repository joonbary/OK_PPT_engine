"""
PPT 내용 분석 스크립트
"""

from pptx import Presentation
import json

def analyze_ppt():
    # PPT 파일 열기
    ppt = Presentation('generated_from_test_document_bb5ab511.pptx')
    
    print('='*60)
    print('Generated PPT Content Analysis')
    print('='*60)
    
    # Mock 콘텐츠 키워드
    mock_keywords = [
        'Strategic approach',
        'Key insights',
        'Performance metrics',
        'Growth trajectory',
        'Market dynamics',
        'Implementation roadmap'
    ]
    
    # AI 콘텐츠 키워드 (test_document.md에서)
    ai_keywords = [
        '8,500',  # 시장 규모
        'Tesla',
        'BYD',
        'ROI',
        '180%',
        '전기차',
        'Phase 1',
        'Phase 2',
        '3.5%',  # 목표 점유율
        '150억'  # 투자 규모
    ]
    
    # 슬라이드 분석
    all_text = []
    for i, slide in enumerate(ppt.slides, 1):
        print(f'\nSlide {i}:')
        print('-'*40)
        
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
                    all_text.append(text)
        
        if slide_texts:
            # 제목 출력
            title = slide_texts[0]
            print(f'Title: {title[:50]}...' if len(title) > 50 else f'Title: {title}')
            
            # 내용 미리보기
            if len(slide_texts) > 1:
                print('Has content: Yes')
                # 첫 번째 내용 항목만 표시
                content_preview = slide_texts[1][:50]
                print(f'First content: {content_preview}...')
        else:
            print('[Empty slide]')
    
    # 콘텐츠 타입 판별
    print('\n' + '='*60)
    print('Content Type Analysis')
    print('='*60)
    
    # 전체 텍스트 결합
    full_text = ' '.join(all_text).lower()
    
    # Mock 키워드 카운트
    mock_count = 0
    mock_found = []
    for keyword in mock_keywords:
        if keyword.lower() in full_text:
            mock_count += 1
            mock_found.append(keyword)
    
    # AI 키워드 카운트
    ai_count = 0
    ai_found = []
    for keyword in ai_keywords:
        if keyword.lower() in full_text or keyword in full_text:
            ai_count += 1
            ai_found.append(keyword)
    
    print(f'\nMock Keywords Found: {mock_count}/{len(mock_keywords)}')
    if mock_found:
        for kw in mock_found[:5]:  # 최대 5개만 표시
            print(f'  - {kw}')
    
    print(f'\nAI Keywords Found: {ai_count}/{len(ai_keywords)}')
    if ai_found:
        for kw in ai_found[:5]:  # 최대 5개만 표시
            print(f'  - {kw}')
    
    # 판정
    print('\n' + '='*60)
    print('VERDICT:')
    print('='*60)
    
    if ai_count > mock_count:
        print('[SUCCESS] This appears to be AI-generated content!')
        print(f'AI keywords: {ai_count}, Mock keywords: {mock_count}')
    elif mock_count > ai_count:
        print('[WARNING] This appears to be MOCK content!')
        print(f'Mock keywords: {mock_count}, AI keywords: {ai_count}')
        print('\nThe AI content is NOT being properly integrated.')
        print('The PPT is still showing template/mock data.')
    else:
        print('[UNCLEAR] Cannot determine content type')
        print(f'Mock keywords: {mock_count}, AI keywords: {ai_count}')
    
    # 슬라이드 수 확인
    print(f'\nTotal slides: {len(ppt.slides)}')
    
    # 첫 번째 슬라이드 상세 분석
    print('\n' + '='*60)
    print('First Slide Detailed Analysis:')
    print('='*60)
    
    if ppt.slides:
        first_slide = ppt.slides[0]
        for shape in first_slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                # 긴 텍스트는 잘라서 표시
                if len(text) > 100:
                    print(f'Text: {text[:100]}...')
                else:
                    print(f'Text: {text}')

if __name__ == "__main__":
    analyze_ppt()