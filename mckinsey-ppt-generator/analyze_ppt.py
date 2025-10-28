from pptx import Presentation
import os
import sys

# UTF-8 인코딩 설정
sys.stdout.reconfigure(encoding='utf-8')

# 가장 최근 파일 찾기
ppt_dir = 'output/generated_presentations'
files = [f for f in os.listdir(ppt_dir) if f.endswith('.pptx')]
if not files:
    print('No PPTX files found')
    exit()

latest_file = max(files, key=lambda f: os.path.getmtime(os.path.join(ppt_dir, f)))
file_path = os.path.join(ppt_dir, latest_file)

print(f'=== 분석 파일: {latest_file} ===')
print(f'파일 크기: {os.path.getsize(file_path):,} bytes')
print()

# PPT 열기
prs = Presentation(file_path)

print(f'총 슬라이드 수: {len(prs.slides)}')
print('='*50)

# 각 슬라이드 내용 확인
for i, slide in enumerate(prs.slides, 1):
    print(f'\n[슬라이드 {i}]')
    
    # 제목 찾기
    title = None
    if slide.shapes.title:
        title = slide.shapes.title.text
        print(f'제목: {title}')
    
    # 텍스트 내용 확인
    text_content = []
    chart_count = 0
    table_count = 0
    image_count = 0
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
            text = shape.text.strip()
            if text and len(text) > 0:
                # 긴 텍스트는 축약
                if len(text) > 200:
                    text = text[:200] + '...'
                text_content.append(text)
        
        # 차트/테이블 확인
        if shape.has_chart:
            chart_count += 1
        if shape.has_table:
            table_count += 1
        if hasattr(shape, 'image'):
            image_count += 1
    
    if text_content:
        print('내용:')
        for j, text in enumerate(text_content[:3], 1):  # 최대 3개 텍스트만
            print(f'  {j}. {text}')
    
    # 시각 요소 수
    if chart_count > 0:
        print(f'차트: {chart_count}개')
    if table_count > 0:
        print(f'테이블: {table_count}개')
    if image_count > 0:
        print(f'이미지: {image_count}개')
    
    print('-'*40)

print('\n=== 품질 평가 ===')
slide_count = len(prs.slides)
has_content = sum(1 for slide in prs.slides if any(shape.text for shape in slide.shapes if hasattr(shape, 'text')))
content_ratio = has_content / slide_count * 100 if slide_count > 0 else 0

print(f'슬라이드 수: {slide_count}')
print(f'내용이 있는 슬라이드: {has_content} ({content_ratio:.1f}%)')

# 실제 콘텐츠 vs 템플릿 확인
template_keywords = ['템플릿', 'PLACEHOLDER', 'INSERT', 'TODO', '여기에', '추가', '핵심 메시지를 명확하게']
actual_content_slides = 0
mockup_slides = 0

for slide in prs.slides:
    slide_text = ' '.join(shape.text for shape in slide.shapes if hasattr(shape, 'text') and shape.text)
    if slide_text:
        # 템플릿 키워드 체크
        if any(keyword.lower() in slide_text.lower() for keyword in template_keywords):
            mockup_slides += 1
        else:
            actual_content_slides += 1

print(f'실제 콘텐츠 슬라이드: {actual_content_slides}')
print(f'모의 데이터 슬라이드: {mockup_slides}')
print(f'품질 점수: {actual_content_slides / slide_count * 100:.1f}%' if slide_count > 0 else '0%')

# 콘텐츠 깊이 분석
print('\n=== 콘텐츠 상세 분석 ===')
total_text_length = 0
slides_with_data = 0
slides_with_long_content = 0

for slide in prs.slides:
    slide_text_length = sum(len(shape.text) for shape in slide.shapes if hasattr(shape, 'text') and shape.text)
    total_text_length += slide_text_length
    
    if slide_text_length > 50:
        slides_with_data += 1
    if slide_text_length > 200:
        slides_with_long_content += 1

avg_text_per_slide = total_text_length / slide_count if slide_count > 0 else 0
print(f'평균 텍스트 길이: {avg_text_per_slide:.0f} 문자/슬라이드')
print(f'내용이 충실한 슬라이드 (50자 이상): {slides_with_data}개')
print(f'상세 내용 슬라이드 (200자 이상): {slides_with_long_content}개')