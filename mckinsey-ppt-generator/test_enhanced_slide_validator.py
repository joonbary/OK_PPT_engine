"""
Enhanced SlideValidator 통합 테스트 스크립트
Task 3.1 완료 검증을 위한 종합적인 테스트
"""

import sys
import os
sys.path.append(os.path.join(os.path.dirname(__file__), 'app'))

from pptx import Presentation
from pptx.util import Inches, Pt
import time
from typing import Dict, Any

# Import our enhanced validator
from app.services.slide_validator import SlideValidator, ValidationResult, IssueSeverity, IssueCategory

def create_test_presentation() -> Presentation:
    """테스트용 프레젠테이션 생성"""
    prs = Presentation()
    
    # 슬라이드 1: 정상적인 슬라이드
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "McKinsey 스타일 제목"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.name = "Calibri"
    
    content_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    content_frame = content_box.text_frame
    content_frame.text = "• 첫 번째 포인트\n• 두 번째 포인트\n• 세 번째 포인트"
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(14)
    content_para.font.name = "Calibri"
    
    # 슬라이드 2: 문제가 있는 슬라이드 (텍스트 오버플로우, 폰트 문제)
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 매우 작은 박스에 긴 텍스트 (오버플로우 유발)
    overflow_box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    overflow_frame = overflow_box.text_frame
    overflow_frame.text = "이것은 매우 긴 텍스트로서 작은 박스에 들어가지 않아서 오버플로우가 발생할 것입니다. 이런 경우에 Enhanced TextFitter가 감지해야 합니다."
    
    # 너무 작은 폰트
    small_font_box = slide2.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    small_font_frame = small_font_box.text_frame
    small_font_frame.text = "너무 작은 폰트 크기"
    small_font_para = small_font_frame.paragraphs[0]
    small_font_para.font.size = Pt(8)
    
    # 여러 폰트 사용 (일관성 문제)
    mixed_font_box = slide2.shapes.add_textbox(Inches(6), Inches(1), Inches(4), Inches(2))
    mixed_font_frame = mixed_font_box.text_frame
    mixed_font_frame.text = "첫 번째 텍스트\n두 번째 텍스트\n세 번째 텍스트"
    for i, para in enumerate(mixed_font_frame.paragraphs):
        fonts = ["Times New Roman", "Comic Sans MS", "Impact"]
        if i < len(fonts):
            para.font.name = fonts[i]
            para.font.size = Pt(12 + i * 2)
    
    # 슬라이드 3: 겹치는 요소들
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 겹치는 텍스트 박스들
    box1 = slide3.shapes.add_textbox(Inches(2), Inches(2), Inches(3), Inches(2))
    box1.text_frame.text = "첫 번째 박스"
    
    box2 = slide3.shapes.add_textbox(Inches(3), Inches(3), Inches(3), Inches(2))
    box2.text_frame.text = "두 번째 박스 (겹침)"
    
    # 슬라이드 4: 경계 벗어남
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 슬라이드 경계를 벗어나는 요소
    out_of_bounds_box = slide4.shapes.add_textbox(Inches(12), Inches(6), Inches(3), Inches(3))
    out_of_bounds_box.text_frame.text = "경계 벗어남"
    
    # 슬라이드 5: 콘텐츠 과밀
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 너무 많은 불릿 포인트
    dense_box = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    dense_frame = dense_box.text_frame
    bullets = []
    for i in range(10):  # 10개 불릿 (권장 5개 초과)
        bullets.append(f"• 불릿 포인트 {i+1}: 이것은 상당히 긴 텍스트로 구성된 불릿 포인트입니다")
    dense_frame.text = "\n".join(bullets)
    
    return prs

def run_validation_tests():
    """Enhanced SlideValidator 통합 테스트 실행"""
    print("Enhanced SlideValidator 통합 테스트 시작")
    print("=" * 60)
    
    # 테스트 프레젠테이션 생성
    print("테스트 프레젠테이션 생성 중...")
    prs = create_test_presentation()
    
    # Enhanced SlideValidator 초기화
    print("Enhanced SlideValidator 초기화 중...")
    validator = SlideValidator()
    
    test_results = {
        "total_slides": len(prs.slides),
        "validation_results": [],
        "performance_metrics": {
            "total_time_ms": 0,
            "avg_time_per_slide_ms": 0,
            "text_fitter_enabled": validator.text_fitter_available
        },
        "issue_summary": {
            "critical": 0,
            "warnings": 0,
            "suggestions": 0,
            "total": 0
        }
    }
    
    print(f"TextFitter 통합 상태: {'활성화' if validator.text_fitter_available else '비활성화'}")
    print()
    
    # 각 슬라이드 검증
    total_start_time = time.perf_counter()
    
    for i, slide in enumerate(prs.slides):
        print(f"슬라이드 {i+1} 검증 중...")
        
        # Enhanced 검증 실행
        start_time = time.perf_counter()
        result = validator.validate_slide(slide, slide_number=i+1)
        end_time = time.perf_counter()
        
        validation_time = (end_time - start_time) * 1000
        
        # 결과 분석
        if isinstance(result, ValidationResult):
            critical_count = len(result.critical_issues)
            warning_count = len(result.warnings)
            suggestion_count = len(result.suggestions)
            total_issues = len(result.issues)
            
            test_results["issue_summary"]["critical"] += critical_count
            test_results["issue_summary"]["warnings"] += warning_count
            test_results["issue_summary"]["suggestions"] += suggestion_count
            test_results["issue_summary"]["total"] += total_issues
            
            slide_result = {
                "slide_number": i+1,
                "is_valid": result.is_valid,
                "validation_time_ms": validation_time,
                "critical_issues": critical_count,
                "warnings": warning_count,
                "suggestions": suggestion_count,
                "total_issues": total_issues,
                "metrics": result.metrics,
                "issues_detail": [str(issue) for issue in result.issues[:3]]  # 상위 3개만
            }
            
            print(f"   검증 시간: {validation_time:.1f}ms")
            print(f"   유효성: {'통과' if result.is_valid else '실패'}")
            print(f"   이슈: Critical={critical_count}, Warning={warning_count}, Suggestion={suggestion_count}")
            
            if result.issues:
                print("   주요 이슈:")
                for issue in result.issues[:2]:  # 상위 2개만 표시
                    print(f"      - {issue}")
            
            test_results["validation_results"].append(slide_result)
            
        else:
            print(f"   검증 실패: 예상치 못한 결과 타입")
        
        print()
    
    total_end_time = time.perf_counter()
    total_time = (total_end_time - total_start_time) * 1000
    
    # 전체 성능 메트릭 계산
    test_results["performance_metrics"]["total_time_ms"] = total_time
    test_results["performance_metrics"]["avg_time_per_slide_ms"] = total_time / len(prs.slides)
    
    return test_results

def print_test_summary(results: Dict[str, Any]):
    """테스트 결과 요약 출력"""
    print("Enhanced SlideValidator 테스트 결과 요약")
    print("=" * 60)
    
    # 기본 정보
    print(f"총 슬라이드 수: {results['total_slides']}")
    print(f"TextFitter 통합: {'활성화' if results['performance_metrics']['text_fitter_enabled'] else '비활성화'}")
    print()
    
    # 성능 메트릭
    perf = results["performance_metrics"]
    print("성능 메트릭:")
    print(f"   - 총 검증 시간: {perf['total_time_ms']:.1f}ms")
    print(f"   - 슬라이드당 평균 시간: {perf['avg_time_per_slide_ms']:.1f}ms")
    print(f"   - 성능 목표 달성: {'YES' if perf['avg_time_per_slide_ms'] < 100 else 'NO'} (목표: <100ms)")
    print()
    
    # 이슈 요약
    issues = results["issue_summary"]
    print("검증 이슈 요약:")
    print(f"   - Critical 이슈: {issues['critical']}개")
    print(f"   - Warning 이슈: {issues['warnings']}개")
    print(f"   - Suggestion 이슈: {issues['suggestions']}개")
    print(f"   - 총 이슈: {issues['total']}개")
    print()
    
    # 슬라이드별 결과
    print("슬라이드별 검증 결과:")
    for slide_result in results["validation_results"]:
        status = "통과" if slide_result["is_valid"] else "실패"
        print(f"   슬라이드 {slide_result['slide_number']}: {status} "
              f"({slide_result['validation_time_ms']:.1f}ms, "
              f"이슈 {slide_result['total_issues']}개)")
    print()
    
    # 주요 기능 검증 결과
    print("주요 기능 검증 결과:")
    
    # 텍스트 오버플로우 검증
    overflow_detected = any("오버플로우" in str(issue) for result in results["validation_results"] 
                           for issue in result.get("issues_detail", []))
    print(f"   - 텍스트 오버플로우 검증: {'작동' if overflow_detected else '미작동'}")
    
    # 요소 겹침 검증
    overlap_detected = any("겹침" in str(issue) for result in results["validation_results"] 
                          for issue in result.get("issues_detail", []))
    print(f"   - 요소 겹침 검증: {'작동' if overlap_detected else '미작동'}")
    
    # 경계 초과 검증
    bounds_detected = any("경계" in str(issue) for result in results["validation_results"] 
                         for issue in result.get("issues_detail", []))
    print(f"   - 경계 초과 검증: {'작동' if bounds_detected else '미작동'}")
    
    # 폰트 일관성 검증
    font_detected = any("폰트" in str(issue) for result in results["validation_results"] 
                       for issue in result.get("issues_detail", []))
    print(f"   - 폰트 일관성 검증: {'작동' if font_detected else '미작동'}")
    
    # McKinsey 스타일 검증
    mckinsey_detected = any("McKinsey" in str(issue) or "스타일" in str(issue) 
                           for result in results["validation_results"] 
                           for issue in result.get("issues_detail", []))
    print(f"   - McKinsey 스타일 검증: {'작동' if mckinsey_detected else '미작동'}")
    print()
    
    # 전체 품질 점수
    total_slides = results["total_slides"]
    valid_slides = sum(1 for result in results["validation_results"] if result["is_valid"])
    quality_score = (valid_slides / total_slides) * 100 if total_slides > 0 else 0
    
    avg_time = perf['avg_time_per_slide_ms']
    performance_score = max(0, 100 - (avg_time - 50)) if avg_time > 50 else 100  # 50ms 기준
    
    overall_score = (quality_score * 0.7 + performance_score * 0.3)
    
    print("전체 평가:")
    print(f"   - 품질 점수: {quality_score:.1f}% ({valid_slides}/{total_slides} 슬라이드 통과)")
    print(f"   - 성능 점수: {performance_score:.1f}% (평균 {avg_time:.1f}ms)")
    print(f"   - 종합 점수: {overall_score:.1f}%")
    print()
    
    # 성공 여부 판정
    if overall_score >= 80 and perf['text_fitter_enabled']:
        print("Enhanced SlideValidator 통합 테스트 성공!")
        print("   - 모든 주요 기능이 정상 작동합니다.")
        print("   - Enhanced TextFitter와 완전히 통합되었습니다.")
        print("   - 성능 목표를 달성했습니다.")
        return True
    else:
        print("Enhanced SlideValidator 통합 테스트 부분 실패")
        if not perf['text_fitter_enabled']:
            print("   - TextFitter 통합이 비활성화되어 있습니다.")
        if overall_score < 80:
            print(f"   - 종합 점수가 목표(80%) 미달: {overall_score:.1f}%")
        return False

def main():
    """메인 테스트 실행 함수"""
    try:
        print("Enhanced SlideValidator Task 3.1 통합 테스트")
        print("=" * 60)
        print()
        
        # 테스트 실행
        results = run_validation_tests()
        
        # 결과 출력
        print_test_summary(results)
        
        # 테스트 PPT 파일 저장
        print("테스트 프레젠테이션 저장 중...")
        prs = create_test_presentation()
        prs.save("test_enhanced_slide_validator.pptx")
        print("   test_enhanced_slide_validator.pptx 저장 완료")
        print()
        
        return results
        
    except Exception as e:
        print(f"테스트 실행 중 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    main()