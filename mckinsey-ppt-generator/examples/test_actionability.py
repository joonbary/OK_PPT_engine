"""
Actionability 점수 향상 테스트
"""

import asyncio
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from app.services.content_generator import ContentGenerator
from app.services.quality_controller import QualityController


async def test_actionability():
    """실행 가능성 테스트"""
    print("=" * 60)
    print("Actionability 테스트 시작...")
    print("=" * 60)
    
    # ContentGenerator 생성
    generator = ContentGenerator()
    
    # 실행 가능한 권고사항이 많은 문서
    document = """
    ### 2024년 사업 성과 및 2025년 전략 권고
    
    ## 현황
    매출 1,200억원으로 전년 대비 25% 성장했습니다.
    신제품이 전체 매출의 70% 기여했습니다.
    
    ## 전략적 권고사항 (우선순위)
    
    1순위: 신제품 R&D 투자를 200억원에서 300억원으로 50% 확대해야 합니다.
    - 기대효과: 2025년 신제품 매출 30% 증가 가능
    - 실행시기: 2025년 1분기 내 
    - 담당: R&D 본부
    
    2순위: 동남아 3개국(베트남, 태국, 인도네시아) 긴급 진출이 필요합니다.
    - 기대효과: 해외 매출 40% 증가 가능
    - 실행시기: 2025년 상반기
    - 투자규모: 100억원
    
    3순위: 디지털 전환 투자로 운영 효율 20% 개선을 실행해야 합니다.
    - AI 기반 수요예측 시스템 즉시 도입
    - 모바일 앱 2025년 2월까지 리뉴얼 완료
    - 예상 비용 절감: 연간 50억원
    
    최우선 과제: 향후 3개월 내 신제품 3종 출시로 시장 선점이 핵심입니다.
    """
    
    # 프레젠테이션 생성
    prs = await generator.generate(
        document=document,
        num_slides=8,
        target_audience="executive",
        presentation_purpose="recommendation"
    )
    
    # 품질 평가
    controller = QualityController(target_score=0.85)
    score = controller.evaluate(prs)
    
    print(f"\n품질 점수 결과:")
    print(f"총점: {score.total:.3f}")
    print(f"- Clarity: {score.clarity:.3f}")
    print(f"- Insight: {score.insight:.3f}")
    print(f"- Structure: {score.structure:.3f}")
    print(f"- Visual: {score.visual:.3f}")
    print(f"- Actionability: {score.actionability:.3f} {'[LOW]' if score.actionability < 0.7 else '[OK]'}")
    
    # Actionability 세부 분석
    if score.details.get("actionability"):
        details = score.details["actionability"]
        print(f"\nActionability 세부 분석:")
        print(f"- 실행 가능한 슬라이드: {details.get('actionable_slides', 0)}/{details.get('total_slides', 0)}")
        print(f"- 정량화된 슬라이드: {details.get('quantified_slides', 0)}/{details.get('total_slides', 0)}")
        print(f"- 우선순위 명시 슬라이드: {details.get('prioritized_slides', 0)}/{details.get('total_slides', 0)}")
    
    # 슬라이드별 내용 확인
    print(f"\n슬라이드별 콘텐츠 분석:")
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        content = ""
        
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text
        
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        content += text + " "
        
        # 실행 가능성 키워드 체크
        action_keywords = ["권고", "제안", "실행", "추진", "필요", "해야", "시행"]
        has_action = any(keyword in title + content for keyword in action_keywords)
        
        priority_keywords = ["우선", "핵심", "중요", "긴급", "1순위", "최우선"]
        has_priority = any(keyword in title + content for keyword in priority_keywords)
        
        quantified = any(char.isdigit() for char in content)
        
        status = []
        if has_action:
            status.append("ACTION")
        if has_priority:
            status.append("PRIORITY")
        if quantified:
            status.append("QUANTIFIED")
        
        print(f"\n슬라이드 {i}: {title[:50]}...")
        print(f"  상태: {' + '.join(status) if status else '[NONE]'}")
        if not status:
            print(f"  내용 샘플: {content[:100]}...")
    
    # 개선 제안
    if score.actionability < 0.7:
        print(f"\n[개선 필요]")
        print("1. 각 슬라이드에 구체적 실행 방안 추가")
        print("2. 정량적 목표와 기한 명시")
        print("3. 우선순위와 담당 조직 지정")
        print("4. InsightLadder Level 4 (Action) 활성화 필요")
    
    return score.actionability


if __name__ == "__main__":
    actionability_score = asyncio.run(test_actionability())
    
    if actionability_score >= 0.7:
        print(f"\n[SUCCESS] Actionability 점수 충분: {actionability_score:.3f}")
    else:
        print(f"\n[FAILED] Actionability 점수 부족: {actionability_score:.3f}")
        print("ContentGenerator의 generate_text_content 메서드 개선 필요")