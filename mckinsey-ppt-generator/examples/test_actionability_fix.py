"""
Actionability 개선 검증
"""

import asyncio
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from app.services.content_generator import ContentGenerator
from app.services.quality_controller import QualityController
from app.models.workflow_models import SlideGenerationSpec


async def test_actionability_enhancement():
    """Actionability 강화 테스트"""
    
    # Given: 간단한 슬라이드 스펙
    slide_specs = [
        SlideGenerationSpec(
            slide_number=1,
            title="전략적 권고사항",
            content_type="conclusion",
            content={
                "bullets": [
                    "신제품 R&D 투자 확대",
                    "해외 시장 진출",
                    "운영 효율화"
                ]
            },
            layout_type="title_and_content"
        ),
        SlideGenerationSpec(
            slide_number=2,
            title="실행 계획",
            content_type="recommendation",
            content={
                "bullets": [
                    "디지털 전환 추진",
                    "고객 만족도 개선",
                    "비용 절감"
                ],
                "next_steps": [
                    "1분기 내 실행",
                    "투자 계획 수립",
                    "성과 모니터링"
                ]
            },
            layout_type="title_and_content"
        ),
        SlideGenerationSpec(
            slide_number=3,
            title="시장 분석",
            content_type="text",
            content={
                "body": "시장이 연간 15% 성장하고 있습니다. 신규 진입 기회가 있습니다.",
                "bullets": [
                    "아시아 시장 성장률 25%",
                    "유럽 시장 안정화",
                    "북미 시장 포화"
                ]
            },
            layout_type="title_and_content"
        )
    ]
    
    # When: ContentGenerator로 PPT 생성
    generator = ContentGenerator()
    prs = Presentation()
    
    for spec in slide_specs:
        await generator._create_slide(prs, spec)
    
    # Then: QualityController로 평가
    controller = QualityController()
    score = controller.evaluate(prs)
    
    print("\n" + "=" * 60)
    print("Actionability 개선 후 품질 점수")
    print("=" * 60)
    print(f"Actionability: {score.actionability:.3f}")
    print(f"Clarity: {score.clarity:.3f}")
    print(f"Insight: {score.insight:.3f}")
    print(f"Structure: {score.structure:.3f}")
    print(f"Visual: {score.visual:.3f}")
    print(f"Total Score: {score.total:.3f}")
    print(f"Target (0.85): {'[SUCCESS]' if score.total >= 0.85 else '[FAILED]'}")
    
    # 세부 분석
    details = score.details.get("actionability", {})
    print(f"\n상세:")
    print(f"  - 실행 가능한 슬라이드: {details.get('actionable_slides', 0)}/{details.get('total_slides', 0)}")
    print(f"  - 정량화 슬라이드: {details.get('quantified_slides', 0)}/{details.get('total_slides', 0)}")
    print(f"  - 우선순위 명시: {details.get('prioritized_slides', 0)}/{details.get('total_slides', 0)}")
    
    # 슬라이드 내용 출력
    print("\n" + "=" * 60)
    print("생성된 슬라이드 내용")
    print("=" * 60)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n슬라이드 {i+1}:")
        if slide.shapes.title:
            print(f"  제목: {slide.shapes.title.text}")
        
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        print(f"    • {para.text}")
    
    return score.actionability, score.total


async def test_full_generation():
    """전체 생성 테스트"""
    print("\n" + "=" * 60)
    print("전체 문서 생성 테스트")
    print("=" * 60)
    
    document = """
    ## 2025년 전략 권고사항
    
    우리 회사는 2024년 매출 1,200억원을 달성했습니다.
    이는 전년 대비 25% 성장한 수치입니다.
    
    ### 핵심 권고사항
    1. 신제품 R&D 투자를 확대하여 시장 선점
    2. 동남아 3개국 진출로 해외 매출 확대
    3. 디지털 전환으로 운영 효율 개선
    
    ### 실행 계획
    - 1분기: R&D 투자 계획 수립
    - 2분기: 해외 진출 파트너십 체결
    - 3분기: 디지털 시스템 구축
    - 4분기: 성과 평가 및 조정
    """
    
    # ContentGenerator로 생성
    generator = ContentGenerator()
    prs = await generator.generate(
        document=document,
        num_slides=8,
        target_audience="executive",
        presentation_purpose="recommendation"
    )
    
    # 품질 평가
    controller = QualityController(target_score=0.85)
    score = controller.evaluate(prs)
    
    print(f"\n품질 점수:")
    print(f"  Actionability: {score.actionability:.3f}")
    print(f"  Total: {score.total:.3f}")
    print(f"  Pass (0.85): {score.passed}")
    
    return score.total


if __name__ == "__main__":
    # 단순 테스트
    actionability, total = asyncio.run(test_actionability_enhancement())
    
    if actionability >= 0.80:
        print(f"\n[SUCCESS] Actionability 개선 확인: {actionability:.3f}")
        
        # 전체 테스트
        full_score = asyncio.run(test_full_generation())
        
        if full_score >= 0.85:
            print(f"\n[SUCCESS] 목표 달성! Total Score: {full_score:.3f}")
        else:
            print(f"\n[WARNING] 목표 미달. Total Score: {full_score:.3f}")
    else:
        print(f"\n[FAILED] Actionability 개선 필요: {actionability:.3f}")