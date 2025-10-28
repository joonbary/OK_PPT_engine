import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

"""
최종 품질 점수 0.85 달성 검증
"""

import logging
import asyncio
from pptx import Presentation
from app.services.content_generator import ContentGenerator
from app.services.quality_controller import QualityController

# 로깅 설정
logging.basicConfig(level=logging.INFO)


async def test_final_quality():
    """최종 품질 테스트"""
    
    print("\n" + "=" * 70)
    print("최종 품질 점수 0.85 달성 테스트")
    print("=" * 70)
    
    # Given: 결론/권고 슬라이드 스펙
    slide_specs = [
        {
            "type": "title",
            "title": "2025년 전략 방향",
            "content": ["전략적 권고사항 및 실행 계획"]
        },
        {
            "type": "conclusion",
            "title": "핵심 전략 권고",
            "content": [
                "신제품 R&D 투자 확대",
                "해외 시장 진출 가속화",
                "운영 효율성 개선"
            ]
        },
        {
            "type": "recommendation",
            "title": "우선 실행 과제",
            "content": [
                "디지털 전환 추진",
                "고객 경험 개선",
                "비용 구조 최적화"
            ]
        },
        {
            "type": "action",
            "title": "분기별 실행 계획",
            "content": [
                "Q1: 신제품 개발 착수",
                "Q2: 해외 시장 진출 준비",
                "Q3: 운영 시스템 개선"
            ]
        }
    ]
    
    # When: ContentGenerator로 PPT 생성
    print("\n[1/3] 슬라이드 생성 중...")
    generator = ContentGenerator()
    prs = Presentation()
    
    for i, spec in enumerate(slide_specs):
        print(f"  - 슬라이드 {i+1}/{len(slide_specs)} 생성")
        slide = await generator._create_slide(prs, spec)
    
    # Then: 슬라이드 내용 확인
    print("\n[2/3] 생성된 슬라이드 내용 확인")
    print("-" * 70)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n슬라이드 {i+1}:")
        
        if slide.shapes.title:
            print(f"  제목: {slide.shapes.title.text}")
        
        # 본문 콘텐츠 추출
        content_found = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        print(f"    - {para.text.encode('utf-8')}")
                        content_found = True
        
        if not content_found:
            print("    (본문 없음)")
    
    # 품질 평가
    print("\n[3/3] 품질 평가 중...")
    print("-" * 70)
    
    controller = QualityController(target_score=0.85)
    score = controller.evaluate(prs)
    
    # 결과 출력
    print("\n" + "=" * 70)
    print("품질 평가 결과")
    print("=" * 70)
    
    print(f"\n 종합 점수: {score.total:.3f} {'PASS' if score.passed else 'FAIL'}")
    print(f"   목표: 0.850 (차이: {score.total - 0.85:+.3f})")
    
    print(f"\n 세부 점수:")
    print(f"   1. Clarity (명확성):        {score.clarity:.3f} {'PASS' if score.clarity >= 0.85 else 'WARN'}")
    print(f"   2. Insight (인사이트):       {score.insight:.3f} {'PASS' if score.insight >= 0.85 else 'WARN'}")
    print(f"   3. Structure (구조):        {score.structure:.3f} {'PASS' if score.structure >= 0.80 else 'WARN'}")
    print(f"   4. Visual (시각):           {score.visual:.3f} {'PASS' if score.visual >= 0.70 else 'WARN'}")
    print(f"   5. Actionability (실행성):  {score.actionability:.3f} {'PASS' if score.actionability >= 0.80 else 'WARN'}")
    
    # Actionability 상세
    action_details = score.details.get("actionability", {})
    print(f"\n Actionability 상세:")
    print(f"   - 실행 가능: {action_details.get('actionable_slides', 0)}/{action_details.get('total_slides', 0)} 슬라이드")
    print(f"   - 정량화: {action_details.get('quantified_slides', 0)}/{action_details.get('total_slides', 0)} 슬라이드")
    print(f"   - 우선순위: {action_details.get('prioritized_slides', 0)}/{action_details.get('total_slides', 0)} 슬라이드")
    
    # 최종 판정
    print("\n" + "=" * 70)
    if score.passed:
        print("성공! 품질 점수 0.85 달성!")
    else:
        print("개선 필요:")
        if score.clarity < 0.85:
            print(f"   - Clarity: {0.85 - score.clarity:.3f}점 부족")
        if score.insight < 0.85:
            print(f"   - Insight: {0.85 - score.insight:.3f}점 부족")
        if score.structure < 0.80:
            print(f"   - Structure: {0.80 - score.structure:.3f}점 부족")
        if score.visual < 0.70:
            print(f"   - Visual: {0.70 - score.visual:.3f}점 부족")
        if score.actionability < 0.80:
            print(f"   - Actionability: {0.80 - score.actionability:.3f}점 부족")
    print("=" * 70)
    
    return score.passed


if __name__ == "__main__":
    success = asyncio.run(test_final_quality())
    exit(0 if success else 1)