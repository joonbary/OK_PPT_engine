"""
레이아웃 시스템 테스트 스크립트
이 스크립트를 실행하여 구현된 기능을 테스트할 수 있습니다.
"""

import os
import sys
from pathlib import Path

# 프로젝트 루트를 Python 경로에 추가
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))

from pptx import Presentation
from pptx.util import Inches
import asyncio

from app.services.content_analyzer import ContentAnalyzer
from app.services.layout_library import LayoutLibrary
from app.services.layout_applier import LayoutApplier
from app.services.text_fitter import TextFitter
from app.services.slide_validator import SlideValidator
from app.services.slide_fixer import SlideFixer
from app.agents.design_agent import DesignAgent


def test_content_analyzer():
    """ContentAnalyzer 테스트"""
    print("\n" + "="*60)
    print("1. ContentAnalyzer 테스트")
    print("="*60)
    
    analyzer = ContentAnalyzer()
    
    # 테스트 케이스들
    test_cases = [
        {
            "title": "전략 계획",
            "content": """• 시장 분석 완료
• 경쟁사 벤치마킹
• 고객 니즈 파악
• 실행 계획 수립""",
            "expected": "list"
        },
        {
            "title": "Before vs After 비교",
            "content": """기존 방식 vs 새로운 방식
- Before: 전통적 접근
- After: 혁신적 접근""",
            "expected": "comparison"
        },
        {
            "title": "간단한 소개",
            "content": "McKinsey & Company",
            "expected": "title_only"
        }
    ]
    
    for i, case in enumerate(test_cases, 1):
        analysis = analyzer.analyze_slide_content(case["content"], case["title"])
        print(f"\n테스트 {i}: {case['title']}")
        print(f"  - 감지된 타입: {analysis['content_type']}")
        print(f"  - 추천 레이아웃: {analysis['recommended_layout']}")
        print(f"  - 복잡도: {analysis['complexity']}")
        print(f"  - 텍스트 밀도: {analysis['text_density']}")
        
        if analysis['content_type'] == case['expected']:
            print(f"  [SUCCESS] 예상한 타입 '{case['expected']}'과 일치")
        else:
            print(f"  [FAIL] 예상 '{case['expected']}', 실제 '{analysis['content_type']}'")


def test_text_fitter():
    """TextFitter 테스트"""
    print("\n" + "="*60)
    print("2. TextFitter 테스트")
    print("="*60)
    
    fitter = TextFitter()
    
    # 긴 텍스트 테스트
    long_text = "이것은 매우 긴 텍스트입니다. " * 20
    
    # 작은 박스에 맞추기
    result = fitter.fit_text_to_box(
        long_text,
        box_width=3.0,  # 3 inches
        box_height=2.0,  # 2 inches
        initial_font_size=14
    )
    
    print(f"\n긴 텍스트 fitting 결과:")
    print(f"  - 원본 길이: {len(long_text)}자")
    print(f"  - 조정된 폰트 크기: {result['adjusted_font_size']}pt")
    print(f"  - 텍스트 맞음: {result['fits']}")
    print(f"  - 텍스트 축약됨: {result['truncated']}")
    
    if result['truncated']:
        print(f"  - 조정된 텍스트 미리보기: {result['adjusted_text'][:50]}...")
    
    # 한글 텍스트 줄바꿈 테스트
    korean_text = "맥킨지앤컴퍼니는 전세계적으로 가장 권위있는 경영컨설팅 회사입니다"
    wrapped = fitter.smart_line_break(korean_text, 20)
    
    print(f"\n한글 텍스트 줄바꿈:")
    print(f"  원본: {korean_text}")
    print(f"  줄바꿈 결과:")
    for line in wrapped.split('\n'):
        print(f"    {line}")


def test_layout_applier():
    """LayoutApplier 테스트"""
    print("\n" + "="*60)
    print("3. LayoutApplier 테스트")
    print("="*60)
    
    # 프레젠테이션 생성
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    applier = LayoutApplier()
    
    # 불릿 리스트 레이아웃 적용
    content = {
        "title": "주요 성과 지표",
        "bullets": [
            "매출 성장률 25% 달성",
            "시장 점유율 15% 확대",
            "고객 만족도 92% 기록",
            "운영 효율성 30% 개선",
            "신규 시장 3개 진출"
        ]
    }
    
    result = applier.apply_layout(slide, "bullet_list", content)
    
    print(f"\n불릿 리스트 레이아웃 적용 결과:")
    print(f"  - 적용된 레이아웃: {result['layout_applied']}")
    print(f"  - 배치된 요소 수: {result['elements_placed']}")
    print(f"  - 텍스트 조정됨: {result['text_adjusted']}")
    
    if result['issues']:
        print(f"  - 발생한 이슈:")
        for issue in result['issues']:
            print(f"    • {issue}")
    else:
        print(f"  [SUCCESS] 이슈 없이 성공적으로 적용됨")
    
    # PPT 파일 저장
    output_file = "test_output_applier.pptx"
    prs.save(output_file)
    print(f"\n  [SAVED] 테스트 PPT 저장됨: {output_file}")


def test_slide_validator_and_fixer():
    """SlideValidator와 SlideFixer 테스트"""
    print("\n" + "="*60)
    print("4. SlideValidator & SlideFixer 테스트")
    print("="*60)
    
    # 문제가 있는 슬라이드 생성
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 의도적으로 문제 생성
    # 1. 텍스트 오버플로우
    text_box1 = slide.shapes.add_textbox(
        Inches(1), Inches(1), 
        Inches(1), Inches(0.5)  # 매우 작은 박스
    )
    text_box1.text_frame.text = "이것은 매우 긴 텍스트입니다. " * 20
    
    # 2. 슬라이드 경계 벗어남
    text_box2 = slide.shapes.add_textbox(
        Inches(12), Inches(1),  # 오른쪽 경계 근처
        Inches(3), Inches(1)
    )
    text_box2.text_frame.text = "경계를 벗어난 텍스트"
    
    # 검증
    validator = SlideValidator()
    validation_result = validator.validate_slide(slide)
    
    print(f"\n검증 결과:")
    print(f"  - 유효함: {validation_result['is_valid']}")
    print(f"  - 발견된 이슈: {len(validation_result['issues'])}개")
    for issue in validation_result['issues']:
        print(f"    [ISSUE] {issue}")
    print(f"  - 경고: {len(validation_result['warnings'])}개")
    for warning in validation_result['warnings']:
        print(f"    [WARNING] {warning}")
    
    # 자동 수정
    if not validation_result['is_valid']:
        print(f"\n자동 수정 시도...")
        fixer = SlideFixer()
        fix_result = fixer.fix_slide(slide, validation_result)
        
        print(f"  - 적용된 수정: {len(fix_result['fixes_applied'])}개")
        for fix in fix_result['fixes_applied']:
            print(f"    [FIXED] {fix}")
        
        if fix_result['fixes_failed']:
            print(f"  - 실패한 수정: {len(fix_result['fixes_failed'])}개")
            for fail in fix_result['fixes_failed']:
                print(f"    [FAILED] {fail}")
        
        # 재검증
        print(f"\n재검증 결과:")
        revalidation = validator.validate_slide(slide)
        print(f"  - 유효함: {revalidation['is_valid']}")
        print(f"  - 남은 이슈: {len(revalidation['issues'])}개")
    
    # 저장
    output_file = "test_output_fixed.pptx"
    prs.save(output_file)
    print(f"\n  [SAVED] 수정된 PPT 저장됨: {output_file}")


async def test_design_agent():
    """DesignAgent 통합 테스트"""
    print("\n" + "="*60)
    print("5. DesignAgent 통합 테스트")
    print("="*60)
    
    agent = DesignAgent()
    
    # 테스트 슬라이드 데이터
    slides_data = [
        {
            "title": "디지털 전환 전략",
            "content": [
                "클라우드 마이그레이션 완료",
                "AI/ML 도입으로 효율성 30% 향상",
                "디지털 채널 매출 비중 50% 달성",
                "데이터 기반 의사결정 체계 구축"
            ]
        },
        {
            "title": "시장 분석: Before vs After",
            "content": [
                "기존: 전통적 오프라인 중심",
                "기존: 수동적 고객 대응",
                "신규: 옴니채널 통합 전략",
                "신규: 예측 기반 선제적 대응"
            ]
        },
        {
            "title": "2024년 로드맵",
            "content": "전사적 디지털 전환 완성"
        }
    ]
    
    # 실제 프레젠테이션 생성
    prs = Presentation()
    for _ in slides_data:
        prs.slides.add_slide(prs.slide_layouts[6])
    
    # DesignAgent 처리
    input_data = {
        "slides": slides_data,
        "presentation": prs,
        "style": "mckinsey_professional"
    }
    
    result = await agent.process(input_data)
    
    print(f"\nDesignAgent 처리 결과:")
    print(f"\n📊 요약:")
    summary = result.get("summary", {})
    print(f"  - 총 슬라이드: {summary.get('total_slides', 0)}")
    print(f"  - 유효 슬라이드: {summary.get('valid_slides', 0)}")
    print(f"  - 발견된 이슈: {summary.get('total_issues_found', 0)}")
    print(f"  - 적용된 수정: {summary.get('total_fixes_applied', 0)}")
    print(f"  - 품질 점수: {summary.get('overall_quality_score', 0)}/100")
    
    print(f"\n📋 슬라이드별 분석:")
    for slide in result.get("processed_slides", []):
        print(f"\n  슬라이드 {slide['slide_number']}:")
        print(f"    - 선택된 레이아웃: {slide.get('selected_layout', 'N/A')}")
        
        analysis = slide.get('content_analysis', {})
        print(f"    - 콘텐츠 타입: {analysis.get('content_type', 'N/A')}")
        print(f"    - 복잡도: {analysis.get('complexity', 'N/A')}")
        
        validation = slide.get('validation', {})
        if validation:
            print(f"    - 검증 통과: {'[OK]' if validation.get('is_valid') else '[FAIL]'}")
    
    print(f"\n💡 개선 권고사항:")
    for i, recommendation in enumerate(result.get("recommendations", []), 1):
        print(f"  {i}. {recommendation}")
    
    # PPT 저장
    output_file = "test_output_design_agent.pptx"
    prs.save(output_file)
    print(f"\n  [SAVED] 최종 PPT 저장됨: {output_file}")


def run_all_tests():
    """모든 테스트 실행"""
    print("\n" + "="*60)
    print(" 맥킨지 PPT 레이아웃 시스템 테스트")
    print("="*60)
    
    try:
        # 1. ContentAnalyzer 테스트
        test_content_analyzer()
        
        # 2. TextFitter 테스트
        test_text_fitter()
        
        # 3. LayoutApplier 테스트
        test_layout_applier()
        
        # 4. Validator & Fixer 테스트
        test_slide_validator_and_fixer()
        
        # 5. DesignAgent 통합 테스트
        print("\n비동기 테스트 실행 중...")
        asyncio.run(test_design_agent())
        
        print("\n" + "="*60)
        print(" [SUCCESS] 모든 테스트 완료!")
        print("="*60)
        print("\n생성된 테스트 파일들:")
        print("  - test_output_applier.pptx: LayoutApplier 테스트 결과")
        print("  - test_output_fixed.pptx: SlideFixer 수정 결과")
        print("  - test_output_design_agent.pptx: 통합 시스템 결과")
        
    except Exception as e:
        print(f"\n[ERROR] 테스트 중 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    # 현재 디렉토리 확인
    print(f"현재 디렉토리: {os.getcwd()}")
    
    # 테스트 실행
    run_all_tests()
    
    print("\n테스트가 완료되었습니다!")
    print("생성된 PPT 파일들을 확인해보세요.")