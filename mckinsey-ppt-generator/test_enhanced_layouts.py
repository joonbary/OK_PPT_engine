#!/usr/bin/env python3
"""
Test script for enhanced layout system - without pytest dependency
"""

import sys
import os

# Add the app directory to the Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'app'))

from services.layout_library import LayoutLibrary
from services.content_analyzer import ContentAnalyzer
from services.layout_applier import LayoutApplier
from pptx import Presentation


def test_layout_availability():
    """Test that all layouts are available including new ones"""
    print("Testing layout availability...")
    
    library = LayoutLibrary()
    expected_layouts = [
        # Original 8 layouts
        "title_slide", "single_column", "bullet_list", "two_column",
        "three_column", "image_text", "matrix_2x2", "table_layout",
        # New 7 layouts
        "timeline", "process_flow", "pyramid", "dashboard_grid",
        "quote_highlight", "split_screen", "agenda_toc"
    ]
    
    for layout_name in expected_layouts:
        layout = library.get_layout(layout_name)
        assert layout is not None, f"Layout {layout_name} not found"
        assert "elements" in layout, f"Layout {layout_name} missing elements"
        assert "name" in layout, f"Layout {layout_name} missing name"
        assert "use_cases" in layout, f"Layout {layout_name} missing use_cases"
        print(f"  [OK] {layout_name}: {layout['name']}")
    
    print(f"[OK] All {len(expected_layouts)} layouts available")


def test_keyword_based_selection():
    """Test enhanced keyword-based layout selection"""
    print("\nTesting keyword-based layout selection...")
    
    library = LayoutLibrary()
    analysis = {"content_type": "list", "bullet_count": 4, "complexity": "medium"}
    
    test_cases = [
        ("project timeline with milestones", "Project Roadmap", "timeline"),
        ("workflow process steps", "Process Guide", "process_flow"),
        ("organization hierarchy levels", "Company Structure", "pyramid"),
        ("KPI dashboard metrics", "Performance Dashboard", "dashboard_grid"),
        ("customer testimonial quote", "Client Feedback", "quote_highlight"),
        ("split comparison 50/50", "Side by Side", "split_screen"),
        ("meeting agenda items", "Meeting Schedule", "agenda_toc")
    ]
    
    for content_text, title, expected_layout in test_cases:
        selected = library.select_layout_for_content(analysis, content_text, title)
        assert selected == expected_layout, f"Expected {expected_layout}, got {selected} for '{content_text}'"
        print(f"  [OK] '{content_text}' -> {selected}")
    
    print("[OK] Keyword-based selection working correctly")


def test_layout_complexity_scoring():
    """Test layout complexity scoring system"""
    print("\nTesting layout complexity scoring...")
    
    library = LayoutLibrary()
    analysis = {"bullet_count": 3, "text_density": "medium", "complexity": "simple"}
    
    # Simple layouts should have low complexity
    simple_score = library.calculate_layout_complexity_score("title_slide", analysis)
    assert simple_score < 0.3, f"Title slide should be simple, got {simple_score}"
    print(f"  [OK] title_slide complexity: {simple_score:.2f} (simple)")
    
    # Complex layouts should have high complexity
    complex_score = library.calculate_layout_complexity_score("dashboard_grid", analysis)
    assert complex_score > 0.7, f"Dashboard should be complex, got {complex_score}"
    print(f"  [OK] dashboard_grid complexity: {complex_score:.2f} (complex)")
    
    print("[OK] Complexity scoring working correctly")


def test_enhanced_content_analysis():
    """Test enhanced content analysis with keyword detection"""
    print("\nTesting enhanced content analysis...")
    
    analyzer = ContentAnalyzer()
    
    test_cases = [
        ("프로젝트 타임라인과 주요 마일스톤", "Project Timeline", "timeline"),
        ("다음 단계의 워크플로우 프로세스", "Process Guide", "process_flow"),
        ("KPI 대시보드 성과 지표", "Performance Dashboard", "dashboard_grid"),
        ('"This is amazing" - John Smith', "Customer Quote", "quote_highlight"),
        ("1. Opening 2. Update 3. Review", "Meeting Agenda", "agenda_toc")
    ]
    
    for content_text, title, expected_layout in test_cases:
        result = analyzer.analyze_slide_content(content_text, title)
        recommended = result["recommended_layout"]
        assert recommended == expected_layout, f"Expected {expected_layout}, got {recommended} for '{content_text}'"
        print(f"  [OK] '{content_text}' -> {recommended}")
    
    print("[OK] Enhanced content analysis working correctly")


def test_layout_applier_integration():
    """Test that layout applier can handle new layouts"""
    print("\nTesting layout applier integration...")
    
    applier = LayoutApplier()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout
    
    # Test timeline layout
    timeline_content = {
        "title": "Project Timeline",
        "milestone_1": "Planning",
        "milestone_2": "Development",
        "milestone_3": "Testing",
        "milestone_4": "Launch"
    }
    
    result = applier.apply_layout(slide, "timeline", timeline_content)
    assert result["layout_applied"] == "timeline", "Timeline layout not applied correctly"
    assert result["elements_placed"] > 0, "No elements placed for timeline"
    print(f"  [OK] timeline layout: {result['elements_placed']} elements placed")
    
    # Test dashboard layout
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[6])
    dashboard_content = {
        "title": "Performance Dashboard",
        "kpi_1": {"title": "Revenue", "value": "$1.2M", "description": "Monthly"},
        "kpi_2": {"title": "Users", "value": "10,500", "description": "Active"}
    }
    
    result2 = applier.apply_layout(slide2, "dashboard_grid", dashboard_content)
    assert result2["layout_applied"] == "dashboard_grid", "Dashboard layout not applied correctly"
    assert result2["elements_placed"] > 0, "No elements placed for dashboard"
    print(f"  [OK] dashboard_grid layout: {result2['elements_placed']} elements placed")
    
    print("[OK] Layout applier integration working correctly")


def main():
    """Run all tests"""
    print("Testing Enhanced PPT Layout Library")
    print("=" * 50)
    
    try:
        test_layout_availability()
        test_keyword_based_selection()
        test_layout_complexity_scoring()
        test_enhanced_content_analysis()
        test_layout_applier_integration()
        
        print("\n" + "=" * 50)
        print("ALL TESTS PASSED! Enhanced layout library is working correctly.")
        print("\nSummary:")
        print("  - 15 layouts available (8 original + 7 new)")
        print("  - Keyword-based selection implemented")
        print("  - Complexity scoring functional")
        print("  - Enhanced content analysis working")
        print("  - Layout applier supports all new layouts")
        
    except Exception as e:
        print(f"\nTEST FAILED: {str(e)}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    sys.exit(main())