"""
Integration tests for the layout optimization system.
Tests ContentAnalyzer, LayoutApplier, SlideValidator, and SlideFixer.
"""

import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.content_analyzer import ContentAnalyzer
from app.services.layout_library import LayoutLibrary
from app.services.layout_applier import LayoutApplier
from app.services.slide_validator import SlideValidator
from app.services.slide_fixer import SlideFixer
from app.services.text_fitter import TextFitter
from app.agents.design_agent import DesignAgent


class TestContentAnalyzer:
    """Test ContentAnalyzer functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.analyzer = ContentAnalyzer()
    
    def test_detect_bullet_list(self):
        """Test bullet list detection"""
        markdown = """
        - Point 1
        - Point 2
        - Point 3
        - Point 4
        """
        result = self.analyzer.analyze_slide_content(markdown, "Test Title")
        
        assert result["content_type"] == "list"
        assert result["bullet_count"] == 4
        assert result["recommended_layout"] == "bullet_list"
    
    def test_detect_comparison(self):
        """Test comparison content detection"""
        markdown = "## Before vs After comparison"
        result = self.analyzer.analyze_slide_content(markdown, "Comparison Title")
        
        assert result["content_type"] == "comparison"
        assert result["recommended_layout"] == "two_column"
    
    def test_detect_matrix(self):
        """Test matrix/table detection"""
        markdown = """
        | Header 1 | Header 2 |
        |----------|----------|
        | Cell 1   | Cell 2   |
        """
        result = self.analyzer.analyze_slide_content(markdown)
        
        assert result["content_type"] == "matrix"
        assert result["has_table"] == True
    
    def test_text_density_calculation(self):
        """Test text density calculation"""
        short_text = "Short text"
        medium_text = " ".join(["Word"] * 100)
        long_text = " ".join(["Word"] * 200)
        
        assert self.analyzer.calculate_text_density(short_text) == "low"
        assert self.analyzer.calculate_text_density(medium_text) == "medium"
        assert self.analyzer.calculate_text_density(long_text) == "high"
    
    def test_korean_text_analysis(self):
        """Test Korean text analysis"""
        korean_text = """
        • 첫 번째 포인트
        • 두 번째 포인트
        • 세 번째 포인트
        """
        result = self.analyzer.analyze_slide_content(korean_text)
        
        assert result["content_type"] == "list"
        assert result["bullet_count"] == 3
    
    def test_enhanced_layout_detection(self):
        """Test enhanced layout detection with keywords"""
        # Test timeline detection
        timeline_text = "프로젝트 타임라인과 주요 마일스톤을 보여주는 일정"
        result = self.analyzer.analyze_slide_content(timeline_text, "Project Timeline")
        assert result["recommended_layout"] == "timeline"
        
        # Test process flow detection  
        process_text = "다음 단계의 워크플로우 프로세스를 따라주세요"
        result = self.analyzer.analyze_slide_content(process_text, "Process Guide")
        assert result["recommended_layout"] == "process_flow"
        
        # Test dashboard detection
        dashboard_text = "KPI 대시보드 성과 지표 모니터링"
        result = self.analyzer.analyze_slide_content(dashboard_text, "Performance Dashboard")
        assert result["recommended_layout"] == "dashboard_grid"
        
        # Test quote detection
        quote_text = 'John said "This is an amazing product"'
        result = self.analyzer.analyze_slide_content(quote_text, "Customer Testimonial")
        assert result["recommended_layout"] == "quote_highlight"
        
        # Test agenda detection
        agenda_text = """
        1. Opening remarks
        2. Status update
        3. Budget review
        4. Next steps
        """
        result = self.analyzer.analyze_slide_content(agenda_text, "Meeting Agenda")
        assert result["recommended_layout"] == "agenda_toc"


class TestLayoutLibrary:
    """Test LayoutLibrary functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.library = LayoutLibrary()
    
    def test_layout_availability(self):
        """Test that all layouts are available including new ones"""
        expected_layouts = [
            # Original 8 layouts
            "title_slide",
            "single_column", 
            "bullet_list",
            "two_column",
            "three_column",
            "image_text",
            "matrix_2x2",
            "table_layout",
            # New 7 layouts
            "timeline",
            "process_flow",
            "pyramid",
            "dashboard_grid",
            "quote_highlight",
            "split_screen",
            "agenda_toc"
        ]
        
        for layout_name in expected_layouts:
            layout = self.library.get_layout(layout_name)
            assert layout is not None
            assert "elements" in layout
            assert "name" in layout
            assert "use_cases" in layout
    
    def test_layout_selection(self):
        """Test automatic layout selection"""
        # Test bullet list selection
        analysis = {
            "content_type": "list",
            "bullet_count": 4,
            "complexity": "medium"
        }
        selected = self.library.select_layout_for_content(analysis)
        assert selected == "bullet_list"
        
        # Test two column selection for comparison
        analysis = {
            "content_type": "comparison",
            "complexity": "medium"
        }
        selected = self.library.select_layout_for_content(analysis)
        assert selected == "two_column"
    
    def test_text_limits(self):
        """Test text length limits for layouts"""
        limits = self.library.get_text_limits_for_layout("bullet_list")
        
        assert "headline" in limits or "default" in limits
        assert isinstance(limits.get("bullet", limits.get("default", 100)), int)
    
    def test_keyword_based_layout_selection(self):
        """Test enhanced keyword-based layout selection"""
        # Test timeline keywords
        analysis = {"content_type": "list", "bullet_count": 4, "complexity": "medium"}
        selected = self.library.select_layout_for_content(
            analysis, "project timeline with milestones", "Project Roadmap"
        )
        assert selected == "timeline"
        
        # Test process flow keywords
        selected = self.library.select_layout_for_content(
            analysis, "workflow process steps", "Process Guide"
        )
        assert selected == "process_flow"
        
        # Test pyramid keywords
        selected = self.library.select_layout_for_content(
            analysis, "organization hierarchy levels", "Company Structure"
        )
        assert selected == "pyramid"
        
        # Test dashboard keywords
        selected = self.library.select_layout_for_content(
            analysis, "KPI dashboard metrics", "Performance Dashboard"
        )
        assert selected == "dashboard_grid"
        
        # Test quote keywords
        selected = self.library.select_layout_for_content(
            analysis, "customer testimonial quote", "Client Feedback"
        )
        assert selected == "quote_highlight"
        
        # Test split screen keywords
        selected = self.library.select_layout_for_content(
            analysis, "split comparison 50/50", "Side by Side"
        )
        assert selected == "split_screen"
        
        # Test agenda keywords
        selected = self.library.select_layout_for_content(
            analysis, "meeting agenda items", "Meeting Schedule"
        )
        assert selected == "agenda_toc"
    
    def test_layout_complexity_scoring(self):
        """Test layout complexity scoring system"""
        analysis = {"bullet_count": 3, "text_density": "medium", "complexity": "simple"}
        
        # Simple layouts should have low complexity
        simple_score = self.library.calculate_layout_complexity_score("title_slide", analysis)
        assert simple_score < 0.3
        
        # Complex layouts should have high complexity
        complex_score = self.library.calculate_layout_complexity_score("dashboard_grid", analysis)
        assert complex_score > 0.7
        
        # Complexity should increase with more bullets
        high_bullet_analysis = {"bullet_count": 10, "text_density": "high", "complexity": "complex"}
        increased_score = self.library.calculate_layout_complexity_score("bullet_list", high_bullet_analysis)
        base_score = self.library.calculate_layout_complexity_score("bullet_list", analysis)
        assert increased_score > base_score
    
    def test_layout_fallback_chains(self):
        """Test layout fallback chain functionality"""
        # Test timeline fallback
        fallbacks = self.library.get_layout_fallback_chain("timeline")
        assert "process_flow" in fallbacks
        assert "single_column" in fallbacks
        
        # Test pyramid fallback
        fallbacks = self.library.get_layout_fallback_chain("pyramid")
        assert "matrix_2x2" in fallbacks
        
        # Test unknown layout fallback
        fallbacks = self.library.get_layout_fallback_chain("unknown_layout")
        assert fallbacks == ["single_column"]
    
    def test_layout_compatibility_validation(self):
        """Test layout compatibility validation"""
        # Compatible scenario
        analysis = {"bullet_count": 3, "text_density": "medium", "complexity": "simple"}
        validation = self.library.validate_layout_compatibility("bullet_list", analysis)
        assert validation["compatible"] == True
        assert validation["compatibility_score"] > 0.8
        
        # Incompatible scenario
        high_density_analysis = {"bullet_count": 2, "text_density": "high", "complexity": "simple"}
        validation = self.library.validate_layout_compatibility("title_slide", high_density_analysis)
        assert validation["compatible"] == False
        assert len(validation["issues"]) > 0
    
    def test_layout_variants(self):
        """Test layout variant generation"""
        variants = self.library.get_layout_variants("bullet_list")
        
        assert "standard" in variants
        assert "compact" in variants
        assert "spacious" in variants
        
        # Verify variant differences
        standard = variants["standard"]
        compact = variants["compact"]
        
        assert standard != compact  # Should be different
        assert len(compact["elements"]) == len(standard["elements"])  # Same structure


class TestTextFitter:
    """Test TextFitter functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.fitter = TextFitter()
    
    def test_text_fitting(self):
        """Test text fitting to box"""
        text = "This is a test text that needs to fit in a box"
        result = self.fitter.fit_text_to_box(
            text,
            box_width=4.0,  # 4 inches
            box_height=2.0,  # 2 inches
            initial_font_size=14
        )
        
        assert "adjusted_font_size" in result
        assert "adjusted_text" in result
        assert "fits" in result
        assert result["adjusted_font_size"] >= 10  # Minimum font size
    
    def test_text_truncation(self):
        """Test text truncation with ellipsis"""
        long_text = "This is a very long text " * 20
        truncated = self.fitter.truncate_with_ellipsis(long_text, 50)
        
        assert len(truncated) <= 50
        assert truncated.endswith("...")
    
    def test_smart_line_break(self):
        """Test smart line breaking"""
        text = "This is a long sentence that needs to be broken into multiple lines"
        result = self.fitter.smart_line_break(text, 20)
        
        lines = result.split('\n')
        for line in lines:
            assert len(line) <= 25  # Allow some flexibility for word boundaries
    
    def test_korean_text_truncation(self):
        """Test Korean text truncation"""
        korean_text = "이것은 매우 긴 한국어 텍스트입니다" * 5
        truncated = self.fitter.truncate_with_ellipsis(korean_text, 30)
        
        assert len(truncated) <= 30
        assert truncated.endswith("...")


class TestSlideValidator:
    """Test SlideValidator functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.validator = SlideValidator()
        self.presentation = Presentation()
        self.slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[1]
        )
    
    def test_empty_slide_detection(self):
        """Test empty slide detection"""
        result = self.validator.validate_slide(self.slide)
        
        # Empty slide should have a warning
        assert any("빈 슬라이드" in warning for warning in result.get("warnings", []))
    
    def test_text_overflow_detection(self):
        """Test text overflow detection"""
        # Add a small text box with lots of text
        text_box = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(1), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = "Very long text " * 50
        
        result = self.validator.validate_slide(self.slide)
        
        # Should detect potential overflow
        metrics = result.get("metrics", {})
        assert metrics.get("text_overflow_count", 0) > 0 or len(result.get("issues", [])) > 0
    
    def test_font_consistency_check(self):
        """Test font consistency validation"""
        # Add text with different fonts
        for i in range(5):
            text_box = self.slide.shapes.add_textbox(
                Inches(1), Inches(i), Inches(2), Inches(0.5)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"Text {i}"
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10 + i * 2)  # Different sizes
        
        result = self.validator.validate_slide(self.slide)
        
        # Should warn about font inconsistency
        assert result.get("metrics", {}).get("font_consistency") == False or \
               any("폰트" in warning for warning in result.get("warnings", []))


class TestSlideFixer:
    """Test SlideFixer functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.fixer = SlideFixer()
        self.presentation = Presentation()
        self.slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[1]
        )
    
    def test_text_overflow_fix(self):
        """Test fixing text overflow"""
        # Create overflow situation
        text_box = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = "Very long text that will definitely overflow " * 10
        
        validation_result = {
            "is_valid": False,
            "issues": ["텍스트 오버플로우 감지: 1개 텍스트박스"]
        }
        
        result = self.fixer.fix_slide(self.slide, validation_result)
        
        assert len(result["fixes_applied"]) > 0
        # Text should be truncated or font size reduced
        assert len(text_frame.text) <= len("Very long text that will definitely overflow " * 10)
    
    def test_out_of_bounds_fix(self):
        """Test fixing out of bounds shapes"""
        # Create shape outside slide bounds
        text_box = self.slide.shapes.add_textbox(
            Inches(12), Inches(1), Inches(3), Inches(1)
        )
        
        validation_result = {
            "is_valid": False,
            "issues": ["슬라이드 경계 벗어남: 1개 요소"]
        }
        
        result = self.fixer.fix_slide(self.slide, validation_result)
        
        assert len(result["fixes_applied"]) > 0
        # Shape should be moved within bounds
        assert text_box.left < Inches(13.33)


class TestLayoutApplier:
    """Test LayoutApplier functionality"""
    
    def setup_method(self):
        """Set up test fixtures"""
        self.applier = LayoutApplier()
        self.presentation = Presentation()
        self.slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
    
    def test_bullet_list_application(self):
        """Test applying bullet list layout"""
        content = {
            "title": "Test Bullet List",
            "bullets": [
                "First point",
                "Second point",
                "Third point"
            ]
        }
        
        result = self.applier.apply_layout(self.slide, "bullet_list", content)
        
        assert result["layout_applied"] == "bullet_list"
        assert result["elements_placed"] > 0
        assert len(result["issues"]) == 0
    
    def test_two_column_application(self):
        """Test applying two column layout"""
        content = {
            "title": "Comparison",
            "left_header": "Before",
            "right_header": "After",
            "left_content": ["Old point 1", "Old point 2"],
            "right_content": ["New point 1", "New point 2"]
        }
        
        result = self.applier.apply_layout(self.slide, "two_column", content)
        
        assert result["layout_applied"] == "two_column"
        assert result["elements_placed"] >= 3  # Title + 2 columns minimum
    
    def test_text_adjustment(self):
        """Test that text is adjusted when applied"""
        content = {
            "title": "Very long title " * 10,
            "bullets": ["Very long bullet point " * 20 for _ in range(10)]
        }
        
        result = self.applier.apply_layout(self.slide, "bullet_list", content)
        
        assert result["text_adjusted"] == True
        # Should have truncated some content
        assert any("Truncated" in issue or "truncated" in issue 
                  for issue in result.get("issues", []))
    
    def test_timeline_layout_application(self):
        """Test applying timeline layout"""
        content = {
            "title": "Project Timeline",
            "milestone_1": "Planning",
            "milestone_1_detail": "Initial planning phase",
            "milestone_2": "Development", 
            "milestone_2_detail": "Core development",
            "milestone_3": "Testing",
            "milestone_3_detail": "Quality assurance",
            "milestone_4": "Launch",
            "milestone_4_detail": "Product launch"
        }
        
        result = self.applier.apply_layout(self.slide, "timeline", content)
        
        assert result["layout_applied"] == "timeline"
        assert result["elements_placed"] >= 5  # Title + 4 milestones minimum
        assert len(result["issues"]) == 0
    
    def test_process_flow_layout_application(self):
        """Test applying process flow layout"""
        content = {
            "title": "Workflow Process",
            "step_1": "Initiate",
            "step_2": "Analyze",
            "step_3": "Design",
            "step_4": "Implement",
            "step_5": "Review"
        }
        
        result = self.applier.apply_layout(self.slide, "process_flow", content)
        
        assert result["layout_applied"] == "process_flow"
        assert result["elements_placed"] >= 6  # Title + 5 steps
    
    def test_pyramid_layout_application(self):
        """Test applying pyramid layout"""
        content = {
            "title": "Organization Structure",
            "level_1": "CEO",
            "level_2_left": "VP Sales",
            "level_2_right": "VP Engineering", 
            "level_3_1": "Sales Rep 1",
            "level_3_2": "Sales Rep 2",
            "level_3_3": "Engineer 1",
            "level_3_4": "Engineer 2"
        }
        
        result = self.applier.apply_layout(self.slide, "pyramid", content)
        
        assert result["layout_applied"] == "pyramid"
        assert result["elements_placed"] >= 4  # Title + levels
    
    def test_dashboard_grid_layout_application(self):
        """Test applying dashboard grid layout"""
        content = {
            "title": "Performance Dashboard",
            "kpi_1": {"title": "Revenue", "value": "$1.2M", "description": "Monthly revenue"},
            "kpi_2": {"title": "Users", "value": "10,500", "description": "Active users"},
            "kpi_3": {"title": "Growth", "value": "15%", "description": "Month over month"},
            "kpi_4": {"title": "Satisfaction", "value": "4.8/5", "description": "Customer rating"}
        }
        
        result = self.applier.apply_layout(self.slide, "dashboard_grid", content)
        
        assert result["layout_applied"] == "dashboard_grid"
        assert result["elements_placed"] >= 5  # Title + KPIs
    
    def test_quote_highlight_layout_application(self):
        """Test applying quote highlight layout"""
        content = {
            "title": "Customer Testimonial",
            "quote": "This product has revolutionized our workflow and increased productivity by 300%",
            "attribution": "John Smith",
            "context": "CEO, Tech Corp"
        }
        
        result = self.applier.apply_layout(self.slide, "quote_highlight", content)
        
        assert result["layout_applied"] == "quote_highlight"
        assert result["elements_placed"] >= 3  # Title + quote + attribution
    
    def test_split_screen_layout_application(self):
        """Test applying split screen layout"""
        content = {
            "title": "Before vs After",
            "left_content": ["Old process was manual", "Time consuming", "Error prone"],
            "right_content": ["New process is automated", "Fast execution", "High accuracy"]
        }
        
        result = self.applier.apply_layout(self.slide, "split_screen", content)
        
        assert result["layout_applied"] == "split_screen"
        assert result["elements_placed"] >= 3  # Title + 2 panels
    
    def test_agenda_toc_layout_application(self):
        """Test applying agenda/TOC layout"""
        content = {
            "title": "Meeting Agenda",
            "agenda_item_1": "Welcome and introductions",
            "agenda_item_2": "Project status update",
            "agenda_item_3": "Budget review",
            "agenda_item_4": "Next steps planning",
            "agenda_item_5": "Q&A and closing",
            "time_info": "Total: 60 minutes"
        }
        
        result = self.applier.apply_layout(self.slide, "agenda_toc", content)
        
        assert result["layout_applied"] == "agenda_toc"
        assert result["elements_placed"] >= 6  # Title + 5 agenda items


class TestDesignAgent:
    """Test DesignAgent integration"""
    
    @pytest.mark.asyncio
    async def test_design_agent_process(self):
        """Test DesignAgent full processing pipeline"""
        agent = DesignAgent()
        
        # Mock slide data
        slides_data = [
            {
                "title": "Executive Summary",
                "content": [
                    "Key finding 1",
                    "Key finding 2",
                    "Key finding 3"
                ]
            },
            {
                "title": "Comparison Analysis",
                "content": [
                    "Before: Old approach",
                    "After: New approach"
                ]
            }
        ]
        
        input_data = {
            "slides": slides_data,
            "style": "mckinsey_professional"
        }
        
        result = await agent.process(input_data)
        
        assert "processed_slides" in result
        assert "quality_metrics" in result
        assert "recommendations" in result
        assert len(result["processed_slides"]) == 2
        
        # Check first slide processing
        slide1 = result["processed_slides"][0]
        assert slide1["selected_layout"] == "bullet_list"
        assert "content_analysis" in slide1
        assert "design_recommendations" in slide1
        
        # Check second slide processing
        slide2 = result["processed_slides"][1]
        # Should detect comparison and recommend two_column
        assert slide2["content_analysis"]["content_type"] in ["comparison", "list"]
    
    @pytest.mark.asyncio
    async def test_design_agent_with_presentation(self):
        """Test DesignAgent with actual presentation object"""
        agent = DesignAgent()
        presentation = Presentation()
        
        # Add test slides
        slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide1.shapes.title.text = "Test Slide"
        
        # Process presentation
        result = await agent.process_presentation(presentation)
        
        assert "processed_slides" in result
        assert len(result["processed_slides"]) == 1
        assert "quality_metrics" in result
        assert result["quality_metrics"]["total_slides"] == 1


class TestIntegration:
    """End-to-end integration tests"""
    
    def test_full_pipeline(self):
        """Test complete pipeline from analysis to fixing"""
        # Initialize components
        analyzer = ContentAnalyzer()
        library = LayoutLibrary()
        applier = LayoutApplier()
        validator = SlideValidator()
        fixer = SlideFixer()
        
        # Create test presentation
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        
        # Test content
        markdown_content = """
        - Important point about strategy
        - Critical success factors to consider
        - Implementation timeline and milestones
        - Expected outcomes and benefits
        - Risk mitigation strategies
        """
        
        # Step 1: Analyze content
        analysis = analyzer.analyze_slide_content(markdown_content, "Strategic Plan")
        assert analysis["content_type"] == "list"
        
        # Step 2: Select layout
        layout_name = library.select_layout_for_content(analysis)
        assert layout_name == "bullet_list"
        
        # Step 3: Apply layout
        content = {
            "title": "Strategic Plan",
            "bullets": [line.strip("- ").strip() 
                       for line in markdown_content.strip().split("\n") 
                       if line.strip()]
        }
        
        apply_result = applier.apply_layout(slide, layout_name, content)
        assert apply_result["elements_placed"] > 0
        
        # Step 4: Validate
        validation = validator.validate_slide(slide)
        
        # Step 5: Fix if needed
        if not validation["is_valid"]:
            fix_result = fixer.fix_slide(slide, validation)
            assert "fixes_applied" in fix_result
            
            # Re-validate
            revalidation = validator.validate_slide(slide)
            # Should be better after fixes
            assert len(revalidation["issues"]) <= len(validation["issues"])
    
    def test_korean_content_pipeline(self):
        """Test pipeline with Korean content"""
        # Initialize components
        analyzer = ContentAnalyzer()
        applier = LayoutApplier()
        
        # Create test presentation
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        
        # Korean content
        korean_content = """
        • 전략적 목표 설정
        • 핵심 성공 요인 분석
        • 실행 계획 수립
        • 성과 측정 지표
        """
        
        # Analyze
        analysis = analyzer.analyze_slide_content(korean_content, "전략 계획")
        
        # Apply
        content = {
            "title": "전략 계획",
            "bullets": [line.strip("•").strip() 
                       for line in korean_content.strip().split("\n") 
                       if line.strip()]
        }
        
        layout_name = analysis["recommended_layout"]
        result = applier.apply_layout(slide, layout_name, content)
        
        assert result["elements_placed"] > 0
        assert result["layout_applied"] in ["bullet_list", "single_column"]


if __name__ == "__main__":
    pytest.main([__file__, "-v"])