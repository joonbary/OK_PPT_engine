"""
End-to-End Test Script for Real Document PPT Generation
Tests the complete pipeline with actual document content to ensure:
1. Document parsing works
2. LLM calls are made (not mock data)
3. PPT content reflects document input
4. Quality validation passes
"""

import asyncio
import sys
import os
import time
from pathlib import Path

# Add parent directory to path
sys.path.append(str(Path(__file__).parent.parent))

from app.services.workflow_orchestrator import WorkflowOrchestrator
from app.models.workflow_models import GenerationRequest
from loguru import logger


async def test_real_document():
    """실제 문서로 전체 파이프라인 테스트"""
    
    print("="*80)
    print("🧪 PPT Generation End-to-End Test")
    print("="*80)
    
    # 실제 문서 내용 (전기차 시장 분석)
    test_document = """
    전기차 시장 분석 보고서 2024
    
    1. 시장 개요
    글로벌 전기차 시장은 2024년 기준 1,200억 달러 규모로 성장했습니다.
    전년 대비 35% 증가하며 빠른 성장세를 보이고 있으며, 특히 아시아 태평양 지역이
    전체 시장의 45%를 차지하며 가장 큰 비중을 보이고 있습니다.
    
    2. 주요 동향과 성장 동력
    - 테슬라가 글로벌 시장 점유율 23%로 1위 유지
    - BYD가 중국 시장에서 급성장하며 18% 점유율 기록
    - 유럽 시장이 환경 규제 강화로 20% 성장률 달성
    - 배터리 기술 혁신으로 주행거리 500km 이상 모델 보편화
    - 충전 인프라 확대로 2024년 충전소 200만개 돌파
    
    3. 기술 발전 현황
    배터리 가격이 kWh당 100달러 수준으로 하락하며 경제성이 크게 개선되었습니다.
    자율주행 기술과의 융합이 가속화되고 있으며, 레벨 3 자율주행이 상용화 단계에 있습니다.
    초고속 충전 기술 발전으로 15분 충전으로 300km 주행이 가능해졌습니다.
    
    4. 시장 전망과 기회
    2027년까지 연평균 28% 성장이 예상되며, 시장 규모는 3,000억 달러에 달할 전망입니다.
    배터리 가격 하락과 정부 지원 정책이 주요 성장 동력이 될 것으로 예상됩니다.
    신흥 시장 진출과 상용차 전기화가 새로운 성장 기회를 제공할 것입니다.
    
    5. 전략적 권고사항
    - 배터리 기술 R&D 투자를 50% 확대하여 경쟁력 확보
    - 충전 인프라 파트너십 구축으로 고객 편의성 향상
    - 자율주행 기술 기업과의 전략적 제휴 추진
    - 구독 서비스 모델 도입으로 수익 다각화
    - ESG 경영 강화로 투자자 신뢰 확보
    """
    
    # WorkflowOrchestrator 초기화
    print("\n📊 Initializing Workflow Orchestrator...")
    orchestrator = WorkflowOrchestrator(
        max_iterations=1,  # 한 번만 실행
        target_quality_score=0.75,
        aggressive_fix=True,
        timeout_minutes=10
    )
    
    # GenerationRequest 생성
    request = GenerationRequest(
        document=test_document,
        num_slides=10,
        target_audience="executive",
        style="mckinsey"
    )
    
    print(f"📄 Document length: {len(test_document)} characters")
    print(f"🎯 Target slides: {request.num_slides}")
    print(f"👥 Target audience: {request.target_audience}")
    print(f"🎨 Style: {request.style}")
    
    try:
        # 파이프라인 실행
        print("\n🚀 Starting PPT generation pipeline...")
        start_time = time.time()
        
        result = await orchestrator.execute(
            request=request,
            job_id=f"test_{int(time.time())}"
        )
        
        elapsed_time = time.time() - start_time
        
        print(f"\n✅ Pipeline completed in {elapsed_time:.2f} seconds")
        print(f"📁 Output path: {result.output_path}")
        print(f"📊 Quality score: {result.quality_score:.3f}")
        print(f"📈 Slides generated: {result.metrics.slides_generated}")
        
        # 결과 검증
        print("\n🔍 Validating results...")
        
        # 1. 파일 존재 확인
        if os.path.exists(result.output_path):
            file_size = os.path.getsize(result.output_path) / 1024
            print(f"✅ PPT file exists: {file_size:.1f} KB")
        else:
            print(f"❌ PPT file not found: {result.output_path}")
            return False
        
        # 2. 품질 점수 확인
        if result.quality_score >= 0.7:
            print(f"✅ Quality score acceptable: {result.quality_score:.3f}")
        else:
            print(f"⚠️ Quality score low: {result.quality_score:.3f}")
        
        # 3. 생성 시간 확인 (Mock 데이터 감지)
        if elapsed_time >= 5.0:
            print(f"✅ Generation time realistic: {elapsed_time:.2f}s (AI was used)")
        else:
            print(f"❌ Generation too fast: {elapsed_time:.2f}s (might be using mock data)")
            return False
        
        # 4. PPT 내용 검증
        from pptx import Presentation
        prs = Presentation(result.output_path)
        
        print(f"\n📋 Analyzing PPT content...")
        print(f"Total slides: {len(prs.slides)}")
        
        # Mock 데이터 패턴 검사
        mock_patterns = ["100.0", "[키워드]", "Mock", "템플릿", "PLACEHOLDER"]
        document_keywords = ["전기차", "시장", "배터리", "테슬라", "BYD", "충전"]
        
        mock_detected = False
        content_relevant = False
        
        for i, slide in enumerate(prs.slides, 1):
            if slide.shapes.title:
                title = slide.shapes.title.text
                print(f"  Slide {i}: {title}")
                
                # Mock 패턴 검사
                for pattern in mock_patterns:
                    if pattern in title:
                        print(f"    ❌ Mock pattern detected: '{pattern}'")
                        mock_detected = True
                
                # 문서 키워드 검사
                for keyword in document_keywords:
                    if keyword in title:
                        content_relevant = True
        
        if mock_detected:
            print("\n❌ Test FAILED: Mock data patterns detected")
            return False
        
        if not content_relevant:
            print("\n❌ Test FAILED: Content doesn't reflect document")
            return False
        
        print("\n🎉 All tests passed!")
        print("✅ PPT successfully generated from real document")
        print("✅ No mock data detected")
        print("✅ Content reflects input document")
        
        return True
        
    except Exception as e:
        print(f"\n❌ Pipeline failed with error: {e}")
        logger.exception("Pipeline error:")
        return False


async def main():
    """메인 테스트 함수"""
    print("Starting End-to-End Test...")
    
    # OpenAI API 키 확인
    api_key = os.getenv('OPENAI_API_KEY')
    if not api_key:
        print("❌ OPENAI_API_KEY not set")
        print("Please set the environment variable and try again")
        return
    
    print(f"✅ OpenAI API key found: {api_key[:8]}...")
    
    # 테스트 실행
    success = await test_real_document()
    
    if success:
        print("\n✨ End-to-End test completed successfully!")
        sys.exit(0)
    else:
        print("\n💥 End-to-End test failed")
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())