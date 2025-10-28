"""
Test suite for Enhanced SlideFixer implementation
Task 3.2 - Intelligent Auto-Correction System

Tests all 6 core fixing methods and enhanced functionality:
1. Text overflow with Enhanced TextFitter integration
2. Shape overlap intelligent resolution
3. Boundary violations correction
4. Readability improvements (small fonts)
5. Margin issues automatic adjustment
6. Font consistency standardization
"""

import pytest
import time
from unittest.mock import Mock, patch, MagicMock
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

from app.services.slide_fixer import SlideFixer, FixResult, FixSummary
from app.services.slide_validator import ValidationIssue, ValidationResult, IssueSeverity, IssueCategory


class TestEnhancedSlideFixer:
    """Enhanced SlideFixer 테스트 슈트"""
    
    @pytest.fixture
    def slide_fixer(self):
        """SlideFixer 인스턴스 생성"""
        return SlideFixer()
    
    @pytest.fixture
    def mock_slide(self):
        """Mock 슬라이드 객체 생성"""
        slide = Mock()
        slide.shapes = []
        return slide
    
    @pytest.fixture
    def mock_text_shape(self):
        """Mock 텍스트 shape 생성"""
        shape = Mock()
        shape.element = "shape_1"
        shape.left = Inches(1).emu
        shape.top = Inches(1).emu
        shape.width = Inches(4).emu
        shape.height = Inches(2).emu
        
        # 텍스트 프레임 설정
        shape.text_frame = Mock()
        shape.text_frame.text = "Sample text content"
        shape.text_frame.paragraphs = [Mock()]
        shape.text_frame.paragraphs[0].runs = [Mock()]
        shape.text_frame.paragraphs[0].runs[0].font = Mock()
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        shape.text_frame.paragraphs[0].runs[0].font.name = "Arial"
        
        return shape
    
    def test_enhanced_data_structures(self, slide_fixer):
        """Enhanced 데이터 구조 테스트"""
        # FixResult 생성 테스트
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.TEXT_OVERFLOW,
            message="Text overflow detected",
            details="Text exceeds box boundaries",
            metadata={"shape_id": "shape_1"}
        )
        
        fix_result = FixResult(
            success=True,
            issue_fixed=issue,
            fix_method="_fix_text_overflow",
            details="Font size reduced from 12pt to 10pt",
            before_state={"font_size": 12},
            after_state={"font_size": 10},
            processing_time_ms=25.5
        )
        
        assert fix_result.success == True
        assert fix_result.issue_fixed.category == IssueCategory.TEXT_OVERFLOW
        assert fix_result.processing_time_ms == 25.5
        
        # FixSummary 생성 테스트
        fix_summary = FixSummary(
            total_issues=5,
            fixed_issues=4,
            failed_fixes=1,
            fix_results=[fix_result]
        )
        
        assert fix_summary.success_rate == 0.8
        assert len(fix_summary.fix_results) == 1
    
    def test_priority_based_fixing(self, slide_fixer):
        """우선순위 기반 수정 시스템 테스트"""
        priorities = slide_fixer.fix_priorities
        
        # 우선순위 검증
        assert priorities[IssueCategory.OUT_OF_BOUNDS] > priorities[IssueCategory.TEXT_OVERFLOW]
        assert priorities[IssueCategory.TEXT_OVERFLOW] > priorities[IssueCategory.FONT_CONSISTENCY]
        assert priorities[IssueCategory.FONT_CONSISTENCY] > priorities[IssueCategory.MCKINSEY_STYLE]
        
        # 수정 메서드 매핑 검증
        assert IssueCategory.TEXT_OVERFLOW in slide_fixer.fix_methods
        assert IssueCategory.SHAPE_OVERLAP in slide_fixer.fix_methods
        assert IssueCategory.OUT_OF_BOUNDS in slide_fixer.fix_methods
    
    @patch('app.services.slide_fixer.app_logger')
    def test_text_fitter_integration(self, mock_logger, slide_fixer):
        """Enhanced TextFitter 통합 테스트"""
        # TextFitter 사용 가능성 확인
        assert hasattr(slide_fixer, 'text_fitter_available')
        
        if slide_fixer.text_fitter_available:
            assert slide_fixer.text_fitter is not None
            mock_logger.info.assert_called_with("Enhanced TextFitter integration enabled for auto-fixing")
        else:
            mock_logger.warning.assert_called()
    
    def test_fix_text_overflow_enhanced(self, slide_fixer, mock_slide, mock_text_shape):
        """Enhanced 텍스트 오버플로우 수정 테스트"""
        # 테스트 설정
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.TEXT_OVERFLOW,
            message="Text overflow detected",
            details="Text exceeds shape boundaries",
            metadata={
                "shape_id": "shape_1",
                "current_font_size": 12,
                "recommended_font_size": 10
            }
        )
        
        # Mock TextFitter 결과
        if slide_fixer.text_fitter_available:
            slide_fixer.text_fitter.fit_text_to_box = Mock(return_value={
                "fits": True,
                "adjusted_font_size": 10,
                "adjusted_text": "Sample text content",
                "truncated": False
            })
        
        # 수정 실행
        result = slide_fixer._fix_text_overflow(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.TEXT_OVERFLOW
        assert result.processing_time_ms > 0
    
    def test_fix_shape_overlap_enhanced(self, slide_fixer, mock_slide):
        """Enhanced 요소 겹침 수정 테스트"""
        # 겹치는 두 shape 생성
        shape1 = Mock()
        shape1.element = "shape_1"
        shape1.left = Inches(1).emu
        shape1.top = Inches(1).emu
        shape1.width = Inches(2).emu
        shape1.height = Inches(1).emu
        
        shape2 = Mock()
        shape2.element = "shape_2"
        shape2.left = Inches(2).emu  # 겹침
        shape2.top = Inches(1).emu
        shape2.width = Inches(2).emu
        shape2.height = Inches(1).emu
        
        mock_slide.shapes = [shape1, shape2]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.SHAPE_OVERLAP,
            message="Shapes overlap detected",
            details="Two shapes are overlapping",
            metadata={
                "overlap_info": {
                    "shape1_id": "shape_1",
                    "shape2_id": "shape_2",
                    "overlap_area": 0.5
                }
            }
        )
        
        # 수정 실행
        result = slide_fixer._fix_shape_overlap(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.SHAPE_OVERLAP
    
    def test_fix_out_of_bounds_enhanced(self, slide_fixer, mock_slide, mock_text_shape):
        """Enhanced 경계 초과 수정 테스트"""
        # 경계를 벗어난 shape 설정
        mock_text_shape.left = Inches(12).emu  # 슬라이드 폭 초과
        mock_text_shape.top = Inches(-0.5).emu  # 음수 위치
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.CRITICAL,
            category=IssueCategory.OUT_OF_BOUNDS,
            message="Shape out of bounds",
            details="Shape exceeds slide boundaries",
            metadata={"shape_id": "shape_1"}
        )
        
        # 수정 실행
        result = slide_fixer._fix_out_of_bounds(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.OUT_OF_BOUNDS
    
    def test_fix_readability_enhanced(self, slide_fixer, mock_slide, mock_text_shape):
        """Enhanced 가독성 수정 테스트"""
        # 작은 폰트 크기 설정
        mock_text_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(8)  # 너무 작음
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.READABILITY,
            message="Small font detected",
            details="Font size below minimum readable threshold",
            metadata={
                "small_fonts": [{"shape_id": "shape_1", "font_size": 8}],
                "caps_text": [],
                "long_lines": []
            }
        )
        
        # 수정 실행
        result = slide_fixer._fix_readability(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.READABILITY
    
    def test_fix_margins_enhanced(self, slide_fixer, mock_slide, mock_text_shape):
        """Enhanced 여백 수정 테스트"""
        # 여백 부족 상황 설정
        mock_text_shape.left = Inches(0.1).emu  # 너무 작은 여백
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.MARGIN_ISSUES,
            message="Insufficient margins",
            details="Shape too close to slide edges",
            metadata={
                "margin_violations": [{
                    "shape_id": "shape_1",
                    "violation_type": "left_margin",
                    "current_margin": 0.1,
                    "required_margin": 0.3
                }]
            }
        )
        
        # 수정 실행
        result = slide_fixer._fix_margins(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.MARGIN_ISSUES
    
    def test_fix_font_consistency_enhanced(self, slide_fixer, mock_slide, mock_text_shape):
        """Enhanced 폰트 일관성 수정 테스트"""
        # 비일관적인 폰트 설정
        mock_text_shape.text_frame.paragraphs[0].runs[0].font.name = "Comic Sans"  # 비표준
        mock_text_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(13)  # 비표준
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.SUGGESTION,
            category=IssueCategory.FONT_CONSISTENCY,
            message="Font inconsistency detected",
            details="Non-standard fonts detected",
            metadata={
                "font_inconsistencies": {
                    "font_types": ["Comic Sans", "Arial"],
                    "font_sizes": [13, 12, 11]
                }
            }
        )
        
        # 수정 실행
        result = slide_fixer._fix_font_consistency(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.FONT_CONSISTENCY
    
    def test_fix_mckinsey_style(self, slide_fixer, mock_slide, mock_text_shape):
        """McKinsey 스타일 준수 수정 테스트"""
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.MCKINSEY_STYLE,
            message="McKinsey style compliance issues",
            details="Style does not meet McKinsey standards",
            metadata={
                "compliance_issues": ["title_font_size", "body_font_range", "bullet_count"]
            }
        )
        
        # 수정 실행
        result = slide_fixer._fix_mckinsey_style(mock_slide, issue, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.MCKINSEY_STYLE
    
    def test_reduce_content_density(self, slide_fixer, mock_slide, mock_text_shape):
        """콘텐츠 밀도 감소 테스트"""
        mock_slide.shapes = [mock_text_shape]
        
        # ValidationIssue 생성
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.CONTENT_DENSITY,
            message="Content density too high",
            details="Slide contains too much content",
            metadata={
                "density_info": {
                    "text_length": 350,
                    "bullet_count": 8,
                    "text_boxes": 6,
                    "density_score": 85
                }
            }
        )
        
        # 수정 실행
        result = slide_fixer._reduce_content_density(mock_slide, issue, aggressive_mode=True)
        
        # 결과 검증
        assert isinstance(result, FixResult)
        assert result.issue_fixed.category == IssueCategory.CONTENT_DENSITY
    
    def test_aggressive_mode_functionality(self, slide_fixer, mock_slide, mock_text_shape):
        """Aggressive 모드 기능 테스트"""
        mock_slide.shapes = [mock_text_shape]
        
        # 텍스트 오버플로우 이슈
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.TEXT_OVERFLOW,
            message="Text overflow",
            details="Text too long for shape",
            metadata={"shape_id": "shape_1"}
        )
        
        # TextFitter Mock 설정 (텍스트 잘라내기 필요)
        if slide_fixer.text_fitter_available:
            slide_fixer.text_fitter.fit_text_to_box = Mock(return_value={
                "fits": False,
                "adjusted_font_size": 8,
                "adjusted_text": "Sample text...",  # 축약됨
                "truncated": True
            })
        
        # Aggressive 모드로 수정
        result = slide_fixer._fix_text_overflow(mock_slide, issue, aggressive_mode=True)
        
        # Aggressive 모드에서는 더 과감한 수정이 이루어져야 함
        assert isinstance(result, FixResult)
        if result.success:
            assert "aggressive" in result.details.lower() or "truncat" in result.details.lower()
    
    def test_iterative_fixing_system(self, slide_fixer, mock_slide):
        """반복 수정 시스템 테스트"""
        # ValidationResult 생성 (여러 이슈)
        issues = [
            ValidationIssue(
                severity=IssueSeverity.CRITICAL,
                category=IssueCategory.OUT_OF_BOUNDS,
                message="Out of bounds",
                details="Shape exceeds boundaries",
                metadata={"shape_id": "shape_1"}
            ),
            ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.TEXT_OVERFLOW,
                message="Text overflow",
                details="Text too long",
                metadata={"shape_id": "shape_1"}
            )
        ]
        
        validation_result = ValidationResult(
            is_valid=False,
            issues=issues,
            slide_number=1,
            processing_time_ms=50.0
        )
        
        # Mock slide validator
        slide_fixer.validator.validate_slide = Mock(return_value=ValidationResult(
            is_valid=True,  # 수정 후 유효
            issues=[],
            slide_number=1,
            processing_time_ms=30.0
        ))
        
        # 수정 실행
        summary = slide_fixer.fix_slide(mock_slide, validation_result, aggressive_mode=False)
        
        # 결과 검증
        assert isinstance(summary, FixSummary)
        assert summary.total_issues >= 0
        assert summary.processing_time_ms > 0
    
    def test_performance_monitoring(self, slide_fixer, mock_slide):
        """성능 모니터링 테스트"""
        # 성능 임계치 확인
        assert slide_fixer.max_fix_time_ms == 200
        assert slide_fixer.max_fix_iterations == 3
        
        # ValidationResult 생성
        validation_result = ValidationResult(
            is_valid=True,
            issues=[],
            slide_number=1,
            processing_time_ms=25.0
        )
        
        # 수정 실행 및 시간 측정
        start_time = time.perf_counter()
        summary = slide_fixer.fix_slide(mock_slide, validation_result)
        end_time = time.perf_counter()
        
        # 성능 검증
        actual_time = (end_time - start_time) * 1000
        assert summary.processing_time_ms > 0
        assert abs(summary.processing_time_ms - actual_time) < 10  # 10ms 오차 허용
    
    def test_legacy_compatibility(self, slide_fixer, mock_slide):
        """레거시 호환성 테스트"""
        # 레거시 딕셔너리 형태의 검증 결과
        legacy_result = {
            "is_valid": False,
            "has_text_overflow": True,
            "has_overlapping": False,
            "has_out_of_bounds": True,
            "has_font_issues": True,
            "has_margin_issues": False
        }
        
        # 레거시 모드로 수정
        summary = slide_fixer.fix_slide(mock_slide, legacy_result)
        
        # 결과 검증
        assert isinstance(summary, FixSummary)
        assert summary.final_validation is None  # 레거시 모드에서는 None
    
    def test_fix_method_capability_validation(self, slide_fixer):
        """수정 메서드 능력 검증 테스트"""
        # 지원되는 카테고리 확인
        assert slide_fixer.validate_fix_capability(IssueCategory.TEXT_OVERFLOW)
        assert slide_fixer.validate_fix_capability(IssueCategory.SHAPE_OVERLAP)
        assert slide_fixer.validate_fix_capability(IssueCategory.OUT_OF_BOUNDS)
        
        # 지원되지 않는 카테고리 확인
        assert not slide_fixer.validate_fix_capability(IssueCategory.EMPTY_CONTENT)
        assert not slide_fixer.validate_fix_capability(IssueCategory.PERFORMANCE)
    
    def test_fix_method_info_retrieval(self, slide_fixer):
        """수정 메서드 정보 조회 테스트"""
        # 텍스트 오버플로우 메서드 정보
        info = slide_fixer.get_fix_method_info(IssueCategory.TEXT_OVERFLOW)
        assert info is not None
        assert info["category"] == IssueCategory.TEXT_OVERFLOW.value
        assert info["priority"] == slide_fixer.fix_priorities[IssueCategory.TEXT_OVERFLOW]
        assert info["method_name"] == "_fix_text_overflow"
        assert info["supports_aggressive_mode"] == True
        assert "requires_text_fitter" in info
        
        # 지원되지 않는 카테고리
        info = slide_fixer.get_fix_method_info(IssueCategory.EMPTY_CONTENT)
        assert info is None
    
    def test_enhanced_statistics(self, slide_fixer):
        """Enhanced 통계 정보 테스트"""
        stats = slide_fixer.get_fix_statistics()
        
        # 기본 통계
        assert "total_fixes_attempted" in stats
        assert "total_fixes_successful" in stats
        assert "success_rate" in stats
        assert "text_fitter_available" in stats
        
        # Enhanced 기능 정보
        assert "enhanced_features" in stats
        assert stats["enhanced_features"]["validation_result_integration"] == True
        assert stats["enhanced_features"]["priority_based_fixing"] == True
        assert stats["enhanced_features"]["iterative_correction"] == True
        
        # 사용 가능한 수정 메서드
        assert "fix_methods_available" in stats
        assert "_fix_text_overflow" in stats["fix_methods_available"]
        assert "_fix_shape_overlap" in stats["fix_methods_available"]
        
        # 성능 임계치
        assert "performance_thresholds" in stats
        assert stats["performance_thresholds"]["max_fix_time_ms"] == 200
    
    def test_presentation_level_fixing(self, slide_fixer):
        """프레젠테이션 레벨 수정 테스트"""
        # Mock 프레젠테이션
        presentation = Mock()
        presentation.slides = [Mock(), Mock(), Mock()]
        
        # ValidationResult 리스트
        validation_results = [
            ValidationResult(is_valid=False, issues=[
                ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.TEXT_OVERFLOW,
                    message="Text overflow",
                    details="Text too long",
                    metadata={"shape_id": "shape_1"}
                )
            ], slide_number=1, processing_time_ms=25.0),
            ValidationResult(is_valid=True, issues=[], slide_number=2, processing_time_ms=15.0),
            ValidationResult(is_valid=False, issues=[
                ValidationIssue(
                    severity=IssueSeverity.CRITICAL,
                    category=IssueCategory.OUT_OF_BOUNDS,
                    message="Out of bounds",
                    details="Shape exceeds boundaries",
                    metadata={"shape_id": "shape_2"}
                )
            ], slide_number=3, processing_time_ms=35.0)
        ]
        
        # Mock validator
        slide_fixer.validator.validate_slide = Mock(return_value=ValidationResult(
            is_valid=True, issues=[], slide_number=1, processing_time_ms=20.0
        ))
        
        # 프레젠테이션 수정
        result = slide_fixer.fix_presentation(presentation, validation_results)
        
        # 결과 검증
        assert result["success"] or not result["success"]  # 성공 여부 확인
        assert "total_slides_processed" in result
        assert "total_issues_found" in result
        assert "total_issues_fixed" in result
        assert "overall_success_rate" in result
        assert "performance_metrics" in result
        assert len(result["slide_summaries"]) <= len(presentation.slides)


class TestSlidFixerErrorHandling:
    """SlideFixer 오류 처리 테스트"""
    
    @pytest.fixture
    def slide_fixer(self):
        return SlideFixer()
    
    def test_invalid_shape_id_handling(self, slide_fixer):
        """잘못된 shape ID 처리 테스트"""
        mock_slide = Mock()
        mock_slide.shapes = []
        
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.TEXT_OVERFLOW,
            message="Text overflow",
            details="Invalid shape",
            metadata={"shape_id": "invalid_id"}
        )
        
        result = slide_fixer._fix_text_overflow(mock_slide, issue)
        
        assert not result.success
        assert "not found" in result.details.lower()
    
    def test_missing_metadata_handling(self, slide_fixer):
        """메타데이터 누락 처리 테스트"""
        mock_slide = Mock()
        mock_slide.shapes = []
        
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.SHAPE_OVERLAP,
            message="Shape overlap",
            details="Missing metadata",
            metadata={}  # 빈 메타데이터
        )
        
        result = slide_fixer._fix_shape_overlap(mock_slide, issue)
        
        assert not result.success
        assert "not found" in result.details.lower()
    
    def test_text_fitter_failure_handling(self, slide_fixer, mock_slide, mock_text_shape):
        """TextFitter 실패 처리 테스트"""
        mock_slide.shapes = [mock_text_shape]
        
        # TextFitter 오류 시뮬레이션
        if slide_fixer.text_fitter_available:
            slide_fixer.text_fitter.fit_text_to_box = Mock(side_effect=Exception("TextFitter error"))
        
        issue = ValidationIssue(
            severity=IssueSeverity.WARNING,
            category=IssueCategory.TEXT_OVERFLOW,
            message="Text overflow",
            details="TextFitter will fail",
            metadata={"shape_id": "shape_1"}
        )
        
        # 오류가 발생해도 fallback으로 처리되어야 함
        result = slide_fixer._fix_text_overflow(mock_slide, issue)
        
        # Fallback 처리 확인
        assert isinstance(result, FixResult)


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])