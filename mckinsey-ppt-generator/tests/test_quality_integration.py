
"""
품질 점수 0.85 달성 통합 테스트
HeadlineGenerator + InsightLadder + QualityController
"""

import pytest
from pptx import Presentation
from app.services.workflow_orchestrator import WorkflowOrchestrator
from app.services.quality_controller import QualityController # Added
from app.models.workflow_models import GenerationRequest # Corrected import


class TestQualityIntegration:
    """품질 통합 테스트"""
    
    @pytest.mark.asyncio
    async def test_full_pipeline_quality_score(self):
        """
        전체 파이프라인 품질 점수 0.85 달성 테스트
        """
        # Given: 테스트 문서
        document = """
        우리 회사의 2024년 매출은 1000억원으로 전년 대비 20% 증가했습니다.
        이는 신제품 출시가 70% 기여했으며, 업계 평균 5% 대비 4배 빠른 성장입니다.
        향후 신제품 라인을 3개로 확대하여 2025년 30% 추가 성장을 목표로 합니다.
        
        주요 경쟁사 대비 우리의 시장 점유율은 2배 높으며,
        고객 만족도는 90%로 업계 최고 수준입니다.
        
        전략적 권고사항:
        1. 신제품 R&D 투자 50% 확대
        2. 해외 시장 진출 가속화
        3. 디지털 전환 강화로 효율성 20% 개선
        """
        
        request = GenerationRequest(
            document=document,
            num_slides=10,
            target_audience="executive",
            template_name="mckinsey_standard" # Changed from style
        )
        
        # When: 워크플로우 실행
        orchestrator = WorkflowOrchestrator(
            max_iterations=3,
            target_quality_score=0.85,
            aggressive_fix=True
        )
        
        response = await orchestrator.execute(request)
        
        # Then: 품질 점수 검증
        assert response.success, "워크플로우 실행 실패"
        assert response.quality_score >= 0.85, f"품질 점수 부족: {response.quality_score}"
        
        # 세부 점수 검증
        quality_details = response.quality_breakdown
        assert quality_details.clarity >= 0.75, "명확성 부족"
        assert quality_details.insight >= 0.80, "인사이트 부족"
        assert quality_details.structure >= 0.80, "구조 부족"
        assert quality_details.visual >= 0.70, "시각적 품질 부족"
        assert quality_details.actionability >= 0.80, "실행가능성 부족"
    
    @pytest.mark.asyncio
    async def test_headline_quality_improvement(self):
        """
        HeadlineGenerator를 통한 명확성 개선 테스트
        """
        # Given: 단순한 제목을 가진 슬라이드
        from app.services.content_generator import ContentGenerator
        from app.models.workflow_models import SlideGenerationSpec # Added import for SlideGenerationSpec
        
        generator = ContentGenerator()
        
        slide_spec = SlideGenerationSpec( # Changed to instance of SlideGenerationSpec
            slide_number=1,
            title="매출 분석",
            content_type="text", # Assuming "text" is appropriate for this content
            content={"body": "2024년 매출 1000억원"} # Content as dict
        )
        
        # When: 헤드라인 생성
        enhanced_title = generator._generate_mckinsey_headline(slide_spec)
        
        # Then: So What 테스트 통과
        from app.services.headline_generator import SoWhatTester
        tester = SoWhatTester()
        result = tester.test(enhanced_title)
        
        assert result["passed"], f"So What 테스트 실패: {result['issues']}"
        assert result["score"] >= 0.7, f"헤드라인 점수 부족: {result['score']}"
    
    @pytest.mark.asyncio
    async def test_insight_ladder_improvement(self):
        """
        InsightLadder를 통한 인사이트 개선 테스트
        """
        # Given: 데이터
        from app.services.insight_ladder import InsightLadder
        
        data = {
            "metric": "매출",
            "value": 1000,
            "previous_value": 800,
            "benchmark": 900,
            "period": "2024년",
            "unit": "억원",
            "drivers": {"신제품": 70, "기존제품": 30}
        }
        
        # When: 4단계 인사이트 생성
        ladder = InsightLadder()
        insights = ladder.climb(data)
        
        # Then: Level 4 달성
        assert len(insights) == 4, "4단계 인사이트 미생성"
        assert insights[-1].level.value == 4, "Level 4 미달성"
        assert insights[-1].confidence >= 0.7, "신뢰도 부족"
        
        # 모든 인사이트에 정량화 포함
        assert all(any(char.isdigit() for char in insight.statement) for insight in insights), \
            "모든 인사이트에 정량화 부족"
    
    def test_quality_controller_scoring(self):
        """
        QualityController 점수 계산 테스트
        """
        # Given: 테스트 프레젠테이션 생성
        prs = Presentation()
        
        # 고품질 슬라이드 추가
        for i in range(10):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # McKinsey 수준 제목
            title = slide.shapes.title
            title.text = f"신제품 출시로 매출 {i*10}% 증가하여 시장 선점 기회 확보"
            
            # 4단계 인사이트 본문
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                body.text = f"""
                • 전년 대비 {i*10}% 증가, 업계 평균 대비 2배
                • 신제품이 매출의 70% 기여
                • 신제품 라인 확대로 30% 추가 성장 가능
                """
        
        # When: 품질 평가
        controller = QualityController(target_score=0.85)
        score = controller.evaluate(prs)
        
        # Then: 목표 점수 달성
        assert score.total >= 0.85, f"품질 점수 부족: {score.total}"
        assert score.passed, "품질 기준 미달"
        
        # 세부 점수 확인
        print(f"\n=== 품질 점수 ===")
        print(f"명확성: {score.clarity:.3f}")
        print(f"인사이트: {score.insight:.3f}")
        print(f"구조: {score.structure:.3f}")
        print(f"시각: {score.visual:.3f}")
        print(f"실행가능성: {score.actionability:.3f}")
        print(f"총점: {score.total:.3f}")
