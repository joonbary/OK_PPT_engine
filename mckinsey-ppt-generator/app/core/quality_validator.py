"""
McKinsey Quality Validator - Fix #4
품질 검증 및 자동 수정 시스템
"""

from typing import Dict, List, Any, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import logging

logger = logging.getLogger(__name__)


class McKinseyQualityValidator:
    """McKinsey 품질 기준 검증기 및 자동 수정기"""
    
    # McKinsey 표준 색상
    MCKINSEY_COLORS = {
        'primary': RGBColor(0, 118, 168),      # McKinsey Blue
        'secondary': RGBColor(244, 118, 33),   # Orange
        'text': RGBColor(83, 86, 90),          # Dark Gray
        'light_gray': RGBColor(217, 217, 217), # Light Gray
        'white': RGBColor(255, 255, 255),      # White
    }
    
    # 표준 폰트
    STANDARD_FONT = 'Arial'
    
    # 슬라이드 크기 제한
    SLIDE_WIDTH = Inches(13.33)  # 16:9 표준
    SLIDE_HEIGHT = Inches(7.5)
    
    def __init__(self):
        self.validation_rules = {
            'style_checks': True,
            'layout_checks': True,
            'text_checks': True,
            'chart_checks': True,
            'auto_fix': True
        }
    
    def validate_presentation(self, prs: Presentation, auto_fix: bool = True) -> Dict[str, Any]:
        """
        프레젠테이션 전체 품질 검증 및 자동 수정
        
        Args:
            prs: PowerPoint 프레젠테이션 객체
            auto_fix: 자동 수정 여부
            
        Returns:
            검증 결과 및 수정 내역
        """
        logger.info("🔍 McKinsey 품질 검증 시작")
        
        issues = {
            'style_violations': [],
            'layout_issues': [],
            'text_problems': [],
            'chart_errors': []
        }
        
        fixes_applied = {
            'style_fixes': 0,
            'layout_fixes': 0,
            'text_fixes': 0,
            'chart_fixes': 0
        }
        
        # 각 슬라이드 검증
        for slide_idx, slide in enumerate(prs.slides):
            slide_issues = self.validate_slide(slide, slide_idx)
            
            # 이슈 누적
            for category, problems in slide_issues.items():
                issues[category].extend(problems)
            
            # 자동 수정 적용
            if auto_fix:
                slide_fixes = self.auto_fix_slide(slide, slide_idx, slide_issues)
                for fix_type, count in slide_fixes.items():
                    fixes_applied[fix_type] += count
        
        # 품질 점수 계산
        total_issues = sum(len(v) for v in issues.values())
        quality_score = max(0, 1.0 - (total_issues * 0.03))  # 더 엄격한 기준
        
        # 전체 프레젠테이션 일관성 검증
        consistency_issues = self.validate_consistency(prs)
        issues['consistency_issues'] = consistency_issues
        
        result = {
            'quality_score': quality_score,
            'issues': issues,
            'fixes_applied': fixes_applied,
            'passed': quality_score >= 0.85,
            'total_issues': total_issues,
            'total_fixes': sum(fixes_applied.values())
        }
        
        logger.info(f"✅ 품질 검증 완료: 점수 {quality_score:.3f}, 이슈 {total_issues}개, 수정 {sum(fixes_applied.values())}개")
        
        return result
    
    def validate_slide(self, slide, slide_idx: int) -> Dict[str, List]:
        """개별 슬라이드 검증"""
        issues = {
            'style_violations': [],
            'layout_issues': [],
            'text_problems': [],
            'chart_errors': []
        }
        
        # 1. 스타일 검증
        style_issues = self._check_style_compliance(slide, slide_idx)
        issues['style_violations'].extend(style_issues)
        
        # 2. 레이아웃 검증
        layout_issues = self._check_layout_compliance(slide, slide_idx)
        issues['layout_issues'].extend(layout_issues)
        
        # 3. 텍스트 문제 검증
        text_issues = self._check_text_problems(slide, slide_idx)
        issues['text_problems'].extend(text_issues)
        
        # 4. 차트 검증
        chart_issues = self._check_chart_quality(slide, slide_idx)
        issues['chart_errors'].extend(chart_issues)
        
        return issues
    
    def _check_style_compliance(self, slide, slide_idx: int) -> List[Dict]:
        """스타일 준수 여부 검증"""
        violations = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    # 폰트 검증
                    if paragraph.font.name and paragraph.font.name != self.STANDARD_FONT:
                        violations.append({
                            'slide': slide_idx,
                            'shape': shape_idx,
                            'paragraph': para_idx,
                            'type': 'wrong_font',
                            'current': paragraph.font.name,
                            'expected': self.STANDARD_FONT,
                            'severity': 'medium'
                        })
                    
                    # 색상 검증
                    if paragraph.font.color.type == 1:  # SOLID COLOR
                        try:
                            current_rgb = paragraph.font.color.rgb
                            valid_colors = list(self.MCKINSEY_COLORS.values())
                            
                            if current_rgb not in valid_colors:
                                violations.append({
                                    'slide': slide_idx,
                                    'shape': shape_idx,
                                    'paragraph': para_idx,
                                    'type': 'wrong_color',
                                    'current': f'RGB({current_rgb})',  # RGBColor 객체는 직접 문자열로 변환
                                    'severity': 'medium'
                                })
                        except Exception as color_error:
                            # 색상 접근 오류는 무시하고 계속 진행
                            pass
                    
                    # 폰트 크기 검증
                    if paragraph.font.size:
                        size_pt = paragraph.font.size.pt
                        if size_pt < 10 or size_pt > 48:
                            violations.append({
                                'slide': slide_idx,
                                'shape': shape_idx,
                                'paragraph': para_idx,
                                'type': 'font_size_out_of_range',
                                'current': f'{size_pt}pt',
                                'severity': 'low'
                            })
        
        return violations
    
    def _check_layout_compliance(self, slide, slide_idx: int) -> List[Dict]:
        """레이아웃 준수 여부 검증"""
        issues = []
        
        # 제목 존재 여부 (첫 번째 슬라이드 제외하고 모든 슬라이드에 제목 필요)
        has_title = False
        title_shape = None
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                # 상단 1인치 이내에 텍스트가 있으면 제목으로 간주
                if shape.top < Inches(1):
                    has_title = True
                    title_shape = shape
                    break
        
        if not has_title and slide_idx > 0:
            issues.append({
                'slide': slide_idx,
                'type': 'missing_title',
                'severity': 'high'
            })
        
        # 제목 위치 검증
        if title_shape and title_shape.left > Inches(1):
            issues.append({
                'slide': slide_idx,
                'type': 'title_misaligned',
                'current_left': title_shape.left,
                'severity': 'medium'
            })
        
        # 여백 검증
        margin_issues = self._check_margins(slide, slide_idx)
        issues.extend(margin_issues)
        
        return issues
    
    def _check_text_problems(self, slide, slide_idx: int) -> List[Dict]:
        """텍스트 관련 문제 검증"""
        problems = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame'):
                # 텍스트 오버플로우 검증
                if shape.left + shape.width > self.SLIDE_WIDTH:
                    problems.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'overflow_right',
                        'current_right': shape.left + shape.width,
                        'max_right': self.SLIDE_WIDTH,
                        'severity': 'high'
                    })
                
                if shape.top + shape.height > self.SLIDE_HEIGHT:
                    problems.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'overflow_bottom',
                        'current_bottom': shape.top + shape.height,
                        'max_bottom': self.SLIDE_HEIGHT,
                        'severity': 'high'
                    })
                
                # 텍스트가 너무 작은 영역에 있는지 검증
                if shape.width < Inches(1) and len(shape.text_frame.text) > 50:
                    problems.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'text_area_too_small',
                        'severity': 'medium'
                    })
                
                # 빈 텍스트 박스 검증
                if not shape.text_frame.text.strip():
                    problems.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'empty_text_box',
                        'severity': 'low'
                    })
        
        return problems
    
    def _check_chart_quality(self, slide, slide_idx: int) -> List[Dict]:
        """차트 품질 검증"""
        errors = []
        
        # 차트 검색 및 검증
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 3:  # Chart type
                try:
                    chart = shape.chart
                    
                    # 차트 제목 검증
                    if not chart.has_title or not chart.chart_title.text_frame.text.strip():
                        errors.append({
                            'slide': slide_idx,
                            'shape': shape_idx,
                            'type': 'chart_missing_title',
                            'severity': 'medium'
                        })
                    
                    # 차트 크기 검증
                    if shape.width < Inches(3) or shape.height < Inches(2):
                        errors.append({
                            'slide': slide_idx,
                            'shape': shape_idx,
                            'type': 'chart_too_small',
                            'current_size': f'{shape.width} x {shape.height}',
                            'severity': 'medium'
                        })
                        
                except Exception as e:
                    errors.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'chart_access_error',
                        'error': str(e),
                        'severity': 'high'
                    })
        
        return errors
    
    def _check_margins(self, slide, slide_idx: int) -> List[Dict]:
        """여백 검증"""
        issues = []
        min_margin = Inches(0.5)
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame'):
                # 좌측 여백
                if shape.left < min_margin:
                    issues.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'insufficient_left_margin',
                        'current': shape.left,
                        'min_required': min_margin,
                        'severity': 'medium'
                    })
                
                # 상단 여백 (제목 제외)
                if shape.top < min_margin and shape.top > Inches(0.1):
                    issues.append({
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'type': 'insufficient_top_margin',
                        'current': shape.top,
                        'min_required': min_margin,
                        'severity': 'low'
                    })
        
        return issues
    
    def validate_consistency(self, prs: Presentation) -> List[Dict]:
        """프레젠테이션 전체 일관성 검증"""
        issues = []
        
        # 폰트 일관성 검증
        fonts_used = set()
        colors_used = set()
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.font.name:
                            fonts_used.add(paragraph.font.name)
                        if paragraph.font.color.type == 1:
                            try:
                                rgb = paragraph.font.color.rgb
                                colors_used.add(f'RGB({rgb})')  # RGBColor 객체는 직접 문자열로 변환
                            except Exception:
                                pass  # 색상 접근 오류는 무시하고 계속 진행
        
        # 폰트가 2개 이상 사용되면 일관성 문제
        if len(fonts_used) > 2:
            issues.append({
                'type': 'inconsistent_fonts',
                'fonts_found': list(fonts_used),
                'severity': 'medium'
            })
        
        # 색상이 너무 많이 사용되면 일관성 문제
        if len(colors_used) > 5:
            issues.append({
                'type': 'too_many_colors',
                'colors_found': list(colors_used),
                'severity': 'medium'
            })
        
        return issues
    
    def auto_fix_slide(self, slide, slide_idx: int, issues: Dict[str, List]) -> Dict[str, int]:
        """슬라이드 자동 수정"""
        fixes_applied = {
            'style_fixes': 0,
            'layout_fixes': 0,
            'text_fixes': 0,
            'chart_fixes': 0
        }
        
        # 스타일 자동 수정
        style_fixes = self._auto_fix_style_violations(slide, slide_idx, issues.get('style_violations', []))
        fixes_applied['style_fixes'] += style_fixes
        
        # 레이아웃 자동 수정
        layout_fixes = self._auto_fix_layout_issues(slide, slide_idx, issues.get('layout_issues', []))
        fixes_applied['layout_fixes'] += layout_fixes
        
        # 텍스트 문제 자동 수정
        text_fixes = self._auto_fix_text_problems(slide, slide_idx, issues.get('text_problems', []))
        fixes_applied['text_fixes'] += text_fixes
        
        # 차트 문제 자동 수정
        chart_fixes = self._auto_fix_chart_errors(slide, slide_idx, issues.get('chart_errors', []))
        fixes_applied['chart_fixes'] += chart_fixes
        
        return fixes_applied
    
    def _auto_fix_style_violations(self, slide, slide_idx: int, violations: List[Dict]) -> int:
        """스타일 위반 자동 수정"""
        fixes = 0
        
        for violation in violations:
            if violation['type'] == 'wrong_font':
                try:
                    shape = slide.shapes[violation['shape']]
                    paragraph = shape.text_frame.paragraphs[violation['paragraph']]
                    paragraph.font.name = self.STANDARD_FONT
                    fixes += 1
                    logger.debug(f"✅ Fixed font in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to fix font: {e}")
            
            elif violation['type'] == 'wrong_color':
                try:
                    shape = slide.shapes[violation['shape']]
                    paragraph = shape.text_frame.paragraphs[violation['paragraph']]
                    # 기본 텍스트 색상으로 수정
                    paragraph.font.color.rgb = self.MCKINSEY_COLORS['text']
                    fixes += 1
                    logger.debug(f"✅ Fixed color in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to fix color: {e}")
        
        return fixes
    
    def _auto_fix_layout_issues(self, slide, slide_idx: int, issues: List[Dict]) -> int:
        """레이아웃 문제 자동 수정"""
        fixes = 0
        
        for issue in issues:
            if issue['type'] == 'title_misaligned':
                try:
                    # 제목을 좌측 정렬로 수정
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.top < Inches(1):
                            shape.left = Inches(0.5)
                            fixes += 1
                            logger.debug(f"✅ Fixed title alignment in slide {slide_idx}")
                            break
                except Exception as e:
                    logger.warning(f"Failed to fix title alignment: {e}")
        
        return fixes
    
    def _auto_fix_text_problems(self, slide, slide_idx: int, problems: List[Dict]) -> int:
        """텍스트 문제 자동 수정"""
        fixes = 0
        
        for problem in problems:
            if problem['type'] == 'overflow_right':
                try:
                    shape = slide.shapes[problem['shape']]
                    # 너비를 슬라이드 경계 내로 조정
                    max_width = self.SLIDE_WIDTH - shape.left - Inches(0.5)
                    if max_width > Inches(1):
                        shape.width = max_width
                        fixes += 1
                        logger.debug(f"✅ Fixed right overflow in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to fix right overflow: {e}")
            
            elif problem['type'] == 'overflow_bottom':
                try:
                    shape = slide.shapes[problem['shape']]
                    # 높이를 슬라이드 경계 내로 조정
                    max_height = self.SLIDE_HEIGHT - shape.top - Inches(0.5)
                    if max_height > Inches(0.5):
                        shape.height = max_height
                        fixes += 1
                        logger.debug(f"✅ Fixed bottom overflow in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to fix bottom overflow: {e}")
            
            elif problem['type'] == 'empty_text_box':
                try:
                    # 빈 텍스트 박스 제거
                    shape = slide.shapes[problem['shape']]
                    slide.shapes._spTree.remove(shape._element)
                    fixes += 1
                    logger.debug(f"✅ Removed empty text box in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to remove empty text box: {e}")
        
        return fixes
    
    def _auto_fix_chart_errors(self, slide, slide_idx: int, errors: List[Dict]) -> int:
        """차트 오류 자동 수정"""
        fixes = 0
        
        for error in errors:
            if error['type'] == 'chart_too_small':
                try:
                    shape = slide.shapes[error['shape']]
                    # 최소 크기로 조정
                    min_width = Inches(4)
                    min_height = Inches(3)
                    
                    if shape.width < min_width:
                        shape.width = min_width
                    if shape.height < min_height:
                        shape.height = min_height
                    
                    fixes += 1
                    logger.debug(f"✅ Fixed chart size in slide {slide_idx}")
                except Exception as e:
                    logger.warning(f"Failed to fix chart size: {e}")
        
        return fixes
    
    def generate_quality_report(self, validation_result: Dict[str, Any]) -> str:
        """품질 검증 보고서 생성"""
        report = []
        report.append("=== McKinsey PPT 품질 검증 보고서 ===")
        report.append(f"품질 점수: {validation_result['quality_score']:.3f}")
        report.append(f"검증 통과: {'✅ PASS' if validation_result['passed'] else '❌ FAIL'}")
        report.append(f"총 이슈: {validation_result['total_issues']}개")
        report.append(f"자동 수정: {validation_result['total_fixes']}개")
        report.append("")
        
        # 이슈별 상세 내역
        issues = validation_result['issues']
        for category, problems in issues.items():
            if problems:
                report.append(f"[{category.upper()}] {len(problems)}개 이슈:")
                for problem in problems[:5]:  # 최대 5개만 표시
                    severity = problem.get('severity', 'unknown')
                    report.append(f"  - {problem['type']} (심각도: {severity})")
                if len(problems) > 5:
                    report.append(f"  ... 외 {len(problems) - 5}개")
                report.append("")
        
        # 수정 내역
        fixes = validation_result['fixes_applied']
        report.append("자동 수정 내역:")
        for fix_type, count in fixes.items():
            if count > 0:
                report.append(f"  - {fix_type}: {count}개")
        
        return "\n".join(report)