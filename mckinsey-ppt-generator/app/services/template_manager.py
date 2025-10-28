"""
Template management system for PowerPoint presentations.

This module handles template creation, management, and application for
generating professional presentations with consistent styling and layouts.
"""

import os
import json
import logging
from typing import Dict, List, Optional, Any, Tuple
from pathlib import Path
from dataclasses import dataclass
from datetime import datetime

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.core.config import settings
from app.models.presentation import Template
from app.services.mckinsey_styles import (
    McKinseyStyles, 
    SlideLayoutType, 
    create_mckinsey_template
)

# Configure logging
logger = logging.getLogger(__name__)


@dataclass
class TemplateConfig:
    """Template configuration data structure"""
    name: str
    description: str
    category: str
    colors: Dict[str, Any]
    fonts: Dict[str, Any]
    layouts: Dict[str, Any]
    chart_style: Dict[str, Any]
    table_style: Dict[str, Any]
    metadata: Dict[str, Any] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            "name": self.name,
            "description": self.description,
            "category": self.category,
            "colors": self.colors,
            "fonts": self.fonts,
            "layouts": self.layouts,
            "chart_style": self.chart_style,
            "table_style": self.table_style,
            "metadata": self.metadata or {}
        }


class TemplateManager:
    """Manages PowerPoint templates and their application"""
    
    def __init__(self, db: Session, templates_dir: str = None):
        """
        Initialize template manager
        
        Args:
            db: Database session
            templates_dir: Directory for storing template files
        """
        self.db = db
        self.mckinsey_styles = McKinseyStyles()
        
        # Set up templates directory
        if templates_dir:
            self.templates_dir = Path(templates_dir)
        else:
            self.templates_dir = Path("templates/ppt")
        
        self.templates_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize default templates
        self._initialize_default_templates()
    
    def _initialize_default_templates(self) -> None:
        """Initialize default templates if they don't exist"""
        try:
            # Check if McKinsey template exists
            mckinsey_template = self.get_template_by_name("McKinsey Professional")
            if not mckinsey_template:
                self.create_mckinsey_template()
                logger.info("Created default McKinsey template")
                
        except Exception as e:
            logger.error(f"Error initializing default templates: {e}")
    
    def create_template(self, config: TemplateConfig) -> Optional[Template]:
        """
        Create a new template
        
        Args:
            config: Template configuration
            
        Returns:
            Created template or None if failed
        """
        try:
            # Check if template with same name exists
            existing = self.get_template_by_name(config.name)
            if existing:
                raise ValueError(f"Template '{config.name}' already exists")
            
            # Create template record
            template = Template(
                name=config.name,
                description=config.description,
                category=config.category,
                config=config.to_dict(),
                usage_count=0,
                rating=0
            )
            
            # Save to database
            self.db.add(template)
            self.db.commit()
            self.db.refresh(template)
            
            # Create template file
            template_file_path = self._create_template_file(template)
            if template_file_path:
                template.template_file = str(template_file_path)
                self.db.commit()
            
            logger.info(f"Created template: {config.name}")
            return template
            
        except Exception as e:
            logger.error(f"Error creating template: {e}")
            self.db.rollback()
            return None
    
    def create_mckinsey_template(self) -> Optional[Template]:
        """Create the default McKinsey template"""
        try:
            mckinsey_config_data = create_mckinsey_template()
            
            config = TemplateConfig(
                name=mckinsey_config_data["name"],
                description=mckinsey_config_data["description"],
                category="Professional",
                colors=mckinsey_config_data["colors"],
                fonts=mckinsey_config_data["fonts"],
                layouts=mckinsey_config_data["layouts"],
                chart_style=mckinsey_config_data["chart_style"],
                table_style=mckinsey_config_data["table_style"],
                metadata={
                    "created_by": "system",
                    "version": "1.0",
                    "is_default": True
                }
            )
            
            return self.create_template(config)
            
        except Exception as e:
            logger.error(f"Error creating McKinsey template: {e}")
            return None
    
    def get_template_by_name(self, name: str) -> Optional[Template]:
        """Get template by name"""
        try:
            return self.db.query(Template).filter(Template.name == name).first()
        except Exception as e:
            logger.error(f"Error getting template by name: {e}")
            return None
    
    def get_template_by_id(self, template_id: str) -> Optional[Template]:
        """Get template by ID"""
        try:
            return self.db.query(Template).filter(Template.id == template_id).first()
        except Exception as e:
            logger.error(f"Error getting template by ID: {e}")
            return None
    
    def list_templates(self, category: str = None) -> List[Template]:
        """
        List available templates
        
        Args:
            category: Filter by category (optional)
            
        Returns:
            List of templates
        """
        try:
            query = self.db.query(Template)
            
            if category:
                query = query.filter(Template.category == category)
            
            return query.order_by(Template.name).all()
            
        except Exception as e:
            logger.error(f"Error listing templates: {e}")
            return []
    
    def get_template_categories(self) -> List[str]:
        """Get all template categories"""
        try:
            result = self.db.query(Template.category).distinct().all()
            return [row[0] for row in result if row[0]]
        except Exception as e:
            logger.error(f"Error getting template categories: {e}")
            return []
    
    def update_template(self, template_id: str, updates: Dict[str, Any]) -> bool:
        """
        Update template configuration
        
        Args:
            template_id: Template ID
            updates: Updates to apply
            
        Returns:
            True if successful, False otherwise
        """
        try:
            template = self.get_template_by_id(template_id)
            if not template:
                return False
            
            # Update allowed fields
            allowed_fields = ['description', 'category', 'config', 'rating']
            for field, value in updates.items():
                if field in allowed_fields:
                    setattr(template, field, value)
            
            self.db.commit()
            logger.info(f"Updated template: {template.name}")
            return True
            
        except Exception as e:
            logger.error(f"Error updating template: {e}")
            self.db.rollback()
            return False
    
    def delete_template(self, template_id: str) -> bool:
        """
        Delete a template
        
        Args:
            template_id: Template ID
            
        Returns:
            True if successful, False otherwise
        """
        try:
            template = self.get_template_by_id(template_id)
            if not template:
                return False
            
            # Don't delete default templates
            if template.config.get("metadata", {}).get("is_default", False):
                logger.warning(f"Cannot delete default template: {template.name}")
                return False
            
            # Delete template file if exists
            if template.template_file:
                template_path = Path(template.template_file)
                if template_path.exists():
                    template_path.unlink()
            
            # Delete preview if exists
            if template.preview_path:
                preview_path = Path(template.preview_path)
                if preview_path.exists():
                    preview_path.unlink()
            
            # Delete from database
            self.db.delete(template)
            self.db.commit()
            
            logger.info(f"Deleted template: {template.name}")
            return True
            
        except Exception as e:
            logger.error(f"Error deleting template: {e}")
            self.db.rollback()
            return False
    
    def apply_template_to_presentation(self, presentation: Presentation, 
                                     template: Template) -> bool:
        """
        Apply template styling to a presentation
        
        Args:
            presentation: PowerPoint presentation
            template: Template to apply
            
        Returns:
            True if successful, False otherwise
        """
        try:
            config = template.config
            
            # Apply master slide styling
            self._apply_master_styling(presentation, config)
            
            # Apply styling to existing slides
            for slide in presentation.slides:
                self._apply_slide_styling(slide, config)
            
            logger.info(f"Applied template '{template.name}' to presentation")
            return True
            
        except Exception as e:
            logger.error(f"Error applying template: {e}")
            return False
    
    def get_layout_config(self, template: Template, 
                         layout_type: str) -> Dict[str, Any]:
        """
        Get layout configuration for a specific slide type
        
        Args:
            template: Template
            layout_type: Type of layout
            
        Returns:
            Layout configuration
        """
        try:
            layouts = template.config.get("layouts", {})
            return layouts.get(layout_type, layouts.get("content", {}))
        except Exception as e:
            logger.error(f"Error getting layout config: {e}")
            return {}
    
    def increment_usage_count(self, template_id: str) -> None:
        """Increment template usage count"""
        try:
            template = self.get_template_by_id(template_id)
            if template:
                template.usage_count += 1
                self.db.commit()
        except Exception as e:
            logger.error(f"Error incrementing usage count: {e}")
    
    def _create_template_file(self, template: Template) -> Optional[Path]:
        """Create a template file with sample slides"""
        try:
            # Create a sample presentation with the template
            prs = Presentation()
            
            # Apply template styling
            self._apply_master_styling(prs, template.config)
            
            # Create sample slides for each layout type
            self._create_sample_slides(prs, template.config)
            
            # Save template file
            template_filename = f"{template.name.replace(' ', '_').lower()}.pptx"
            template_path = self.templates_dir / template_filename
            
            prs.save(str(template_path))
            
            return template_path
            
        except Exception as e:
            logger.error(f"Error creating template file: {e}")
            return None
    
    def _apply_master_styling(self, presentation: Presentation, 
                            config: Dict[str, Any]) -> None:
        """Apply master slide styling"""
        try:
            # Apply theme colors and fonts to master slides
            # Note: python-pptx has limited master slide modification capabilities
            # This is a simplified implementation
            pass
            
        except Exception as e:
            logger.error(f"Error applying master styling: {e}")
    
    def _apply_slide_styling(self, slide: Slide, config: Dict[str, Any]) -> None:
        """Apply template styling to a slide"""
        try:
            # Apply background color if specified
            colors = config.get("colors", {})
            if "background" in colors:
                # Set slide background (simplified)
                pass
            
            # Apply styling to text boxes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    self._apply_text_styling(shape.text_frame, config)
                    
        except Exception as e:
            logger.error(f"Error applying slide styling: {e}")
    
    def _apply_text_styling(self, text_frame, config: Dict[str, Any]) -> None:
        """Apply text styling from template config"""
        try:
            fonts = config.get("fonts", {})
            colors = config.get("colors", {})
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        font = run.font
                        
                        # Apply font
                        if "primary" in fonts:
                            font.name = fonts["primary"]
                        
                        # Apply color
                        if "text" in colors:
                            # Note: colors are stored as RGBColor objects
                            # This is a simplified implementation
                            pass
                            
        except Exception as e:
            logger.error(f"Error applying text styling: {e}")
    
    def _create_sample_slides(self, presentation: Presentation, 
                            config: Dict[str, Any]) -> None:
        """Create sample slides for template preview"""
        try:
            # Title slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            title = slide.shapes.title
            if title:
                title.text = "Sample Title Slide"
            
            # Content slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title = slide.shapes.title
            if title:
                title.text = "Sample Content Slide"
            
            # Apply styling to all slides
            for slide in presentation.slides:
                self._apply_slide_styling(slide, config)
                
        except Exception as e:
            logger.error(f"Error creating sample slides: {e}")
    
    def export_template_config(self, template_id: str, 
                             export_path: str) -> bool:
        """
        Export template configuration to JSON file
        
        Args:
            template_id: Template ID
            export_path: Path to save exported config
            
        Returns:
            True if successful, False otherwise
        """
        try:
            template = self.get_template_by_id(template_id)
            if not template:
                return False
            
            export_data = {
                "name": template.name,
                "description": template.description,
                "category": template.category,
                "config": template.config,
                "exported_at": datetime.now().isoformat(),
                "version": "1.0"
            }
            
            with open(export_path, 'w', encoding='utf-8') as f:
                json.dump(export_data, f, indent=2, ensure_ascii=False)
            
            logger.info(f"Exported template config: {export_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error exporting template config: {e}")
            return False
    
    def import_template_config(self, import_path: str) -> Optional[Template]:
        """
        Import template configuration from JSON file
        
        Args:
            import_path: Path to JSON config file
            
        Returns:
            Created template or None if failed
        """
        try:
            with open(import_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            config = TemplateConfig(
                name=data["name"],
                description=data["description"],
                category=data["category"],
                colors=data["config"]["colors"],
                fonts=data["config"]["fonts"],
                layouts=data["config"]["layouts"],
                chart_style=data["config"]["chart_style"],
                table_style=data["config"]["table_style"],
                metadata=data["config"].get("metadata", {})
            )
            
            template = self.create_template(config)
            logger.info(f"Imported template: {data['name']}")
            return template
            
        except Exception as e:
            logger.error(f"Error importing template config: {e}")
            return None


def get_default_template_config() -> TemplateConfig:
    """Get default template configuration"""
    mckinsey_config = create_mckinsey_template()
    
    return TemplateConfig(
        name=mckinsey_config["name"],
        description=mckinsey_config["description"],
        category="Professional",
        colors=mckinsey_config["colors"],
        fonts=mckinsey_config["fonts"],
        layouts=mckinsey_config["layouts"],
        chart_style=mckinsey_config["chart_style"],
        table_style=mckinsey_config["table_style"]
    )