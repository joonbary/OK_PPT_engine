"""
DesignApplicator (UTF-8)

Applies a clean McKinsey-like style to a presentation:
- White background, slide numbers on non-title slides
- Title slide improved to prevent overflow (word_wrap + dynamic font)
- Basic typography (Arial): title/body sizes and colors
"""

import logging
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)


class DesignApplicator:
    COLORS = {
        "blue": RGBColor(0, 118, 168),
        "orange": RGBColor(244, 118, 33),
        "gray": RGBColor(83, 86, 90),
        "white": RGBColor(255, 255, 255),
    }

    MCKINSEY_FONTS = {
        "title": {"name": "Arial", "size": Pt(28), "bold": True},
        "body": {"name": "Arial", "size": Pt(12), "bold": False},
    }

    MCKINSEY_SPACING = {
        "line_spacing": 1.2,
    }

    def __init__(self) -> None:
        self.logger = logging.getLogger(__name__)
        self.stats = {"slides_processed": 0}
        # Slide type -> layout name mapping (for future layout engines)
        self.SLIDE_TYPE_LAYOUTS = {
            "Title": "title_slide",
            "Executive Summary": "dual_header",
            "Market Analysis": "three_column",
            "Strategic Options": "matrix",
            "Recommendations": "action_list",
            "Internal Analysis": "dual_header",
            "Challenges": "problem_list",
            "Impact Analysis": "waterfall",
        }

    def _get_layout(self, prs, layout_name: str):
        # Placeholder: in this styling-only applicator we simply return blank
        return prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else 0]

    def _apply_three_column_layout(self, slide, data: dict):
        # Placeholder stub – styling engine only
        return

    def _apply_matrix_layout(self, slide, data: dict):
        # Placeholder stub – styling engine only
        return

    def _select_layout(self, slide_type: str, content: dict):
        """
        Select a layout name by slide_type. Logging added for traceability.
        This method is a no-op shim for compatibility with callers that expect
        a layout decision step before visual styling.
        """
        try:
            logger.info(f"Selecting layout for slide_type: {slide_type}")
            layout_mapping = {
                "title": "title_slide",
                "executive_summary": "dual_header",
                "market_analysis": "three_column",
                "strategy": "matrix",
                "financial": "waterfall",
                "conclusion": "title_slide",
            }
            layout_name = layout_mapping.get(slide_type, "dual_header")
            logger.info(f"Layout selected: {layout_name}")
            return layout_name
        except Exception as e:
            logger.warning(f"Layout selection failed: {e}; falling back to dual_header")
            return "dual_header"

    async def apply(self, presentation):
        logger.info(f"Applying McKinsey style to {len(presentation.slides)} slides")
        for idx, slide in enumerate(presentation.slides):
            await self.apply_mckinsey_style_to_slide(slide, idx, is_title_slide=(idx == 0))
        return presentation

    async def apply_mckinsey_style_to_slide(self, slide, slide_idx: int, is_title_slide: bool = False):
        self._set_background(slide)
        if not is_title_slide and slide_idx > 0:
            self._add_slide_number(slide, slide_idx)

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                await self._apply_text_style(shape, is_title_slide)

        if is_title_slide:
            await self._apply_title_slide_style_improved(slide)

        self.stats["slides_processed"] += 1

    def _set_background(self, slide):
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS["white"]
        except Exception as e:
            logger.debug(f"Could not set background: {e}")

    def _add_slide_number(self, slide, slide_idx: int):
        try:
            box = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(0.8), Inches(0.3))
            tf = box.text_frame
            tf.text = str(slide_idx)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = self.COLORS["gray"]
        except Exception as e:
            logger.debug(f"Could not add slide number: {e}")

    async def _apply_text_style(self, shape, is_title: bool):
        try:
            tf = shape.text_frame
            for i, paragraph in enumerate(tf.paragraphs):
                text_type = "title" if (is_title and i == 0) else "body"
                font_cfg = self.MCKINSEY_FONTS[text_type]
                # 문단 정렬/간격
                try:
                    paragraph.line_spacing = self.MCKINSEY_SPACING["line_spacing"]
                except Exception:
                    pass
                # 런 단위 폰트 일관 적용 (한국어 글리프 우선)
                try:
                    for run in paragraph.runs:
                        for fname in ['Noto Sans CJK KR', 'Noto Sans KR', 'NanumGothic', '맑은 고딕', font_cfg["name"], 'Arial']:
                            try:
                                run.font.name = fname
                                break
                            except Exception:
                                continue
                        run.font.size = font_cfg["size"]
                        run.font.bold = font_cfg["bold"]
                        run.font.color.rgb = self.COLORS["blue"] if text_type == "title" else self.COLORS["gray"]
                except Exception:
                    # 문단 폰트로 폴백
                    try:
                        for fname in ['Noto Sans CJK KR', 'Noto Sans KR', 'NanumGothic', '맑은 고딕', font_cfg["name"], 'Arial']:
                            try:
                                paragraph.font.name = fname
                                break
                            except Exception:
                                continue
                        paragraph.font.size = font_cfg["size"]
                        paragraph.font.bold = font_cfg["bold"]
                        paragraph.font.color.rgb = self.COLORS["blue"] if text_type == "title" else self.COLORS["gray"]
                    except Exception:
                        pass
        except Exception as e:
            logger.debug(f"Could not apply text style: {e}")

    async def _apply_title_slide_style_improved(self, slide):
        try:
            # pick the largest text shape with non-empty text
            title_shape = None
            max_area = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    try:
                        area = int(shape.width) * int(shape.height)
                    except Exception:
                        area = 0
                    if area > max_area:
                        max_area = area
                        title_shape = shape
            if not title_shape:
                return

            tf = title_shape.text_frame
            try:
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

            title_text = tf.text or ""
            n = len(title_text)
            if n < 30:
                fs = Pt(28)
            elif n < 50:
                fs = Pt(24)
            elif n < 80:
                fs = Pt(20)
            else:
                fs = Pt(18)

            for p in tf.paragraphs:
                p.font.name = "Arial"
                p.font.size = fs
                p.font.bold = True
                p.font.color.rgb = self.COLORS["blue"]
                p.alignment = PP_ALIGN.CENTER
                try:
                    p.line_spacing = self.MCKINSEY_SPACING["line_spacing"]
                except Exception:
                    pass
        except Exception as e:
            logger.debug(f"Could not apply improved title style: {e}")

    def get_design_stats(self):
        return self.stats
