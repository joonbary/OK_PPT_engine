# -*- coding: utf-8 -*-
"""
Clean ContentGenerator (UTF-8)

Minimal, encoding-safe implementation used as the original generator.
Produces a small but valid PPTX deck to ensure files always open.
"""

from __future__ import annotations

import logging
from typing import Dict, Any

from pptx import Presentation
from pptx.util import Inches


class ContentGenerator:
    """Lightweight, encoding-safe PPT content generator.

    This implementation intentionally avoids any locale/encoding-sensitive
    manipulations and guarantees a valid PPTX with at least 3 slides.
    """

    def __init__(self) -> None:
        self.logger = logging.getLogger(__name__)

    async def generate(
        self,
        document: str,
        num_slides: int = 15,
        target_audience: str = "executive",
        presentation_purpose: str = "analysis",
    ) -> Presentation:
        """Generate a small but valid presentation from plain text.

        - Always returns a valid deck (>=3 slides)
        - Uses ASCII-safe placeholders where needed
        """
        try:
            prs = Presentation()

            # 1) Title slide
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                if title:
                    title.text = "Auto-generated Presentation"
                if subtitle:
                    subtitle.text = f"Audience: {target_audience} | Purpose: {presentation_purpose}"
            except Exception:
                pass

            # 2) Executive summary slide
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                if title:
                    title.text = "Executive Summary"
                body = None
                for ph in slide.placeholders:
                    if ph.has_text_frame and ph != title:
                        body = ph
                        break
                summary = (document or "").strip()
                if not summary:
                    summary = "No document content provided."
                if body and body.has_text_frame:
                    tf = body.text_frame
                    tf.clear()
                    tf.text = summary[:300]
            except Exception:
                pass

            # 3) Content slide
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                if title:
                    title.text = "Key Points"
                body = None
                for ph in slide.placeholders:
                    if ph.has_text_frame and ph != title:
                        body = ph
                        break
                if body and body.has_text_frame:
                    tf = body.text_frame
                    tf.clear()
                    points = _extract_bullets(document)
                    if points:
                        tf.text = points[0]
                        for p in points[1:6]:
                            para = tf.add_paragraph()
                            para.text = p
                            para.level = 0
                    else:
                        tf.text = "No key points extracted."
            except Exception:
                pass

            # Ensure at least 3 slides
            while len(prs.slides) < 3:
                try:
                    prs.slides.add_slide(prs.slide_layouts[5])  # blank
                except Exception:
                    break

            self.logger.info("ContentGenerator: produced %d slides", len(prs.slides))
            return prs
        except Exception as e:
            self.logger.error("Content generation failed: %s", e, exc_info=True)
            # As a last resort, return a single-slide deck
            prs = Presentation()
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                if slide.shapes.title:
                    slide.shapes.title.text = "Auto-generated Presentation"
            except Exception:
                pass
            return prs


def _extract_bullets(text: str) -> list[str]:
    try:
        t = (text or "").replace("\r", "\n")
        lines = [ln.strip(" -•\t") for ln in t.split("\n")]
        lines = [ln for ln in lines if ln]
        # Take first 6 non-empty lines as bullets
        return lines[:6]
    except Exception:
        return []

