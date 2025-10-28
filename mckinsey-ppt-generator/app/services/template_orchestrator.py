# -*- coding: utf-8 -*-
"""
TemplateOrchestrator (UTF-8)

Selects and applies slide templates (layout patterns) using simple
heuristics with optional AI customization (not enabled by default here).
"""

from typing import Dict, Any
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from app.services.templates_ext import (
    TitleBodyTemplate,
    SectionDividerTemplate,
    ProsConsTemplate,
    BeforeAfterTemplate,
    TimelineTemplate,
    KPICardsTemplate,
)


class _TemplateBase:
    name: str = "base"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        raise NotImplementedError


class TitleSlideTemplate(_TemplateBase):
    name = "title_slide"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        t = (spec.get("title") or spec.get("headline") or "").strip() or "Presentation"
        sub = (spec.get("subtitle") or spec.get("supporting_detail") or "").strip()
        tbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = tbox.text_frame
        tf.text = t
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Arial'
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 118, 168)
        if sub:
            sbox = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.8))
            sf = sbox.text_frame
            sf.text = sub[:100]
            sf.word_wrap = True
            for p in sf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(83, 86, 90)


class FullTextTemplate(_TemplateBase):
    name = "full_text"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        items = []
        c = spec.get("content")
        if isinstance(c, list):
            items = [str(x) for x in c if str(x).strip()][:10]
        elif isinstance(c, dict):
            items = [str(v) for v in c.values() if str(v).strip()][:10]
        elif ("pros" in t) or ("cons" in t):
            chosen = self.templates["pros_cons"]
        elif ("before" in t) or ("after" in t) or ("비교" in t) or ("비교" in t):
            chosen = self.templates["before_after"]
        elif ("timeline" in t) or ("?ε??" in t) or ("roadmap" in t):
            chosen = self.templates["timeline"]
        elif "kpi" in t:
            chosen = self.templates["kpi_cards"]
        elif ("section" in t) or ("????" in t):
            chosen = self.templates["section_divider"]
        elif ("pros" in t) or ("cons" in t):
            chosen = self.templates["pros_cons"]
        elif ("before" in t) or ("after" in t) or ("비교" in t) or ("비교" in t):
            chosen = self.templates["before_after"]
        elif ("timeline" in t) or ("?ε??" in t) or ("roadmap" in t):
            chosen = self.templates["timeline"]
        elif "kpi" in t:
            chosen = self.templates["kpi_cards"]
        elif ("section" in t) or ("????" in t):
            chosen = self.templates["section_divider"]
        elif ("pros" in t) or ("cons" in t):
            chosen = self.templates["pros_cons"]
        elif ("before" in t) or ("after" in t) or ("전" in t) or ("후" in t):
            chosen = self.templates["before_after"]
        elif ("timeline" in t) or ("로드맵" in t) or ("roadmap" in t):
            chosen = self.templates["timeline"]
        elif "kpi" in t:
            chosen = self.templates["kpi_cards"]
        elif ("section" in t) or ("구분" in t):
            chosen = self.templates["section_divider"]
        elif ("swot" in t):
            chosen = self.templates.get("swot", self.templates["comparison_table"])
        elif ("bcg" in t):
            chosen = self.templates.get("bcg", self.templates["comparison_table"])
        elif ("dashboard" in t):
            chosen = self.templates.get("dashboard", self.templates["two_column"])
        elif ("heatmap" in t):
            chosen = self.templates.get("heatmap", self.templates["matrix"])
        elif ("text chart bottom" in t) or ("chart bottom" in t):
            chosen = self.templates.get("text_chart_bottom", self.templates["text_chart_right"])
        elif ("dual chart" in t) or ("dual_chart" in t):
            chosen = self.templates.get("dual_chart", self.templates["two_column"])
        elif ("waterfall" in t):
            chosen = self.templates.get("waterfall", self.templates["text_chart_right"])
        else:
            items = [str(c)] if c else []
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.5), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        if not items:
            tf.text = spec.get("headline", "") or spec.get("title", "")
            return
        tf.text = items[0]
        for it in items[1:]:
            p = tf.add_paragraph()
            p.text = it
            p.level = 0
            p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(83, 86, 90)


class TwoColumnTemplate(_TemplateBase):
    name = "two_column"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        left = spec.get("left") or spec.get("content_left") or []
        right = spec.get("right") or spec.get("content_right") or []
        # Left
        lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.0), Inches(5.0))
        ltf = lb.text_frame
        ltf.word_wrap = True
        if left:
            ltf.text = str(left[0])
            for it in left[1:7]:
                p = ltf.add_paragraph()
                p.text = str(it)
        # Right
        rb = slide.shapes.add_textbox(Inches(6.0), Inches(1.2), Inches(5.0), Inches(5.0))
        rtf = rb.text_frame
        rtf.word_wrap = True
        if right:
            rtf.text = str(right[0])
            for it in right[1:7]:
                p = rtf.add_paragraph()
                p.text = str(it)


class ThreeColumnTemplate(_TemplateBase):
    name = "three_column"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        cols = spec.get("columns") or spec.get("three_column") or []
        if not cols or not isinstance(cols, list):
            # Fallback from generic content list into 3 chunks
            base = spec.get("content") or []
            if isinstance(base, list) and base:
                n = max(1, len(base) // 3)
                cols = [base[:n], base[n:2*n], base[2*n:]]
            else:
                cols = [[], [], []]
        # positions
        zones = [(Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.0)),
                 (Inches(4.5), Inches(1.2), Inches(3.8), Inches(5.0)),
                 (Inches(8.5), Inches(1.2), Inches(3.8), Inches(5.0))]
        for i, z in enumerate(zones):
            items = cols[i] if i < len(cols) else []
            tb = slide.shapes.add_textbox(*z)
            tf = tb.text_frame
            tf.word_wrap = True
            if items:
                tf.text = str(items[0])
                for it in items[1:7]:
                    p = tf.add_paragraph()
                    p.text = str(it)


class TextChartRightTemplate(_TemplateBase):
    name = "text_chart_right"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        # left text
        items = spec.get("bullets") or spec.get("content") or []
        lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(5.0))
        ltf = lb.text_frame
        ltf.word_wrap = True
        if isinstance(items, list) and items:
            ltf.text = str(items[0])
            for it in items[1:6]:
                p = ltf.add_paragraph()
                p.text = str(it)
        # right chart (placeholder text if chart generator not used here)
        rb = slide.shapes.add_textbox(Inches(6.2), Inches(1.2), Inches(4.5), Inches(5.0))
        rtf = rb.text_frame
        rtf.text = spec.get("chart_title", "[Chart]")


class ComparisonTableTemplate(_TemplateBase):
    name = "comparison_table"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        # simple 2-column table from lists
        left = spec.get("left") or []
        right = spec.get("right") or []
        rows = max(len(left), len(right), 1)
        cols = 2
        table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.2), Inches(10.5), Inches(5.0)).table
        table.cell(0,0).text = spec.get("left_header", "Option A")
        table.cell(0,1).text = spec.get("right_header", "Option B")
        for r in range(rows):
            table.cell(r+1,0).text = str(left[r]) if r < len(left) else ""
            table.cell(r+1,1).text = str(right[r]) if r < len(right) else ""


class ActionListTemplate(_TemplateBase):
    name = "action_list"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        title = spec.get("title") or spec.get("headline") or "Action Items"
        items = spec.get("bullets") or spec.get("content") or []
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(10.5), Inches(1.0))
        tf = tb.text_frame
        tf.text = str(title)
        for p in tf.paragraphs:
            p.font.name = 'Arial'
            p.font.size = Pt(22)
            p.font.bold = True
        lb = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(10.0), Inches(5.0))
        ltf = lb.text_frame
        ltf.word_wrap = True
        if isinstance(items, list) and items:
            for idx, it in enumerate(items[:8]):
                p = ltf.paragraphs[0] if idx == 0 else ltf.add_paragraph()
                p.text = f"✓ {it}"
                p.level = 0
                p.font.name = 'Arial'
                p.font.size = Pt(14)
        else:
            ltf.text = "✓ 액션 항목을 정의하세요"

class Matrix2x2Template(_TemplateBase):
    name = "matrix"

    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        table = slide.shapes.add_table(2, 2, Inches(2.0), Inches(1.5), Inches(6.0), Inches(4.5)).table
        data = spec.get("matrix") or [["", ""], ["", ""]]
        for r in range(2):
            for c in range(2):
                try:
                    table.cell(r, c).text = str(data[r][c])
                except Exception:
                    table.cell(r, c).text = ""


class TemplateOrchestrator:
    def __init__(self) -> None:
        self.templates = self._load_templates()
        # register extended analysis templates
        try:
            from app.services.templates_ext import (SWOTTemplate, BCGMatrixTemplate, DashboardTemplate, HeatmapTemplate, TextChartBottomTemplate, DualChartTemplate, WaterfallTemplate)
            self.templates.update({
                "swot": SWOTTemplate(),
                "bcg": BCGMatrixTemplate(),
                "dashboard": DashboardTemplate(),
                "heatmap": HeatmapTemplate(),
                "text_chart_bottom": TextChartBottomTemplate(),
                "dual_chart": DualChartTemplate(),
                "waterfall": WaterfallTemplate(),
            })
        except Exception:
            pass

    def _load_templates(self) -> Dict[str, _TemplateBase]:
        return {
            "title_slide": TitleSlideTemplate(),
            "full_text": FullTextTemplate(),
            "two_column": TwoColumnTemplate(),
            "three_column": ThreeColumnTemplate(),
            "text_chart_right": TextChartRightTemplate(),
            "comparison_table": ComparisonTableTemplate(),
            "action_list": ActionListTemplate(),
            "matrix": Matrix2x2Template(),
            # extended (lightweight)
            "title_body": TitleBodyTemplate(),
            "section_divider": SectionDividerTemplate(),
            "pros_cons": ProsConsTemplate(),
            "before_after": BeforeAfterTemplate(),
            "timeline": TimelineTemplate(),
            "kpi_cards": KPICardsTemplate(),
            "title_body": TitleBodyTemplate(),
            "section_divider": SectionDividerTemplate(),
            "pros_cons": ProsConsTemplate(),
            "before_after": BeforeAfterTemplate(),
            "timeline": TimelineTemplate(),
            "kpi_cards": KPICardsTemplate(),
        }

    def select_and_apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        """Select best template and apply to the given slide.
        Heuristics only; AI customization not enabled by default here.
        """
        t = (spec.get("title") or spec.get("headline") or "").lower()
        stype = (spec.get("type") or "").lower()
        # quick numeric check
        has_numbers = any(ch.isdigit() for ch in (spec.get("headline") or ""))
        # categories
        num_categories = len(spec.get("columns") or [])
        # comparison
        is_comparison = bool(spec.get("left") or spec.get("right")) or ("comparison" in t or "비교" in t)

        # rules
        if spec.get("slide_number", 1) == 1 or (stype == "title"):
            chosen = self.templates["title_slide"]
        elif stype == "executive summary":
            chosen = self.templates.get("title_body", self.templates["full_text"])
        elif has_numbers and (spec.get("chart_data") or "chart" in t or "차트" in t or "data" in t or "분석" in t):
            chosen = self.templates["text_chart_right"]
        elif is_comparison:
            chosen = self.templates["comparison_table"]
        elif (stype == "market analysis") or (num_categories == 3) or (spec.get("columns") is not None):
            chosen = self.templates["three_column"]
        elif (stype in ["strategic options", "impact analysis"]) or ("matrix" in t or "매트릭스" in t) or spec.get("matrix"):
            chosen = self.templates["matrix"]
        elif stype == "recommendations":
            chosen = self.templates.get("action_list", self.templates.get("kpi_cards", self.templates["full_text"]))
        else:
            chosen = self.templates["full_text"]

        chosen.apply(slide, spec)




