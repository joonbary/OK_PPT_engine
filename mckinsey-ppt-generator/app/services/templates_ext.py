"""
Extended slide templates (UTF-8)
Lightweight implementations for additional layouts.
"""
from typing import Dict, Any
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class TitleBodyTemplate:
    name = "title_body"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        title = spec.get("title") or spec.get("headline") or ""
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.5), Inches(0.7))
        tf = tbox.text_frame; tf.text = title; tf.word_wrap = True
        b = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9.5), Inches(5.6))
        bf = b.text_frame; bf.word_wrap = True
        for i, it in enumerate(spec.get("content") or []):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph(); p.text = str(it)

class SectionDividerTemplate:
    name = "section_divider"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        title = spec.get("title") or "Section"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9.5), Inches(1.2))
        tf = box.text_frame; tf.text = title
        for p in tf.paragraphs:
            p.font.name = 'Arial'; p.font.size = Pt(32); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

class ProsConsTemplate:
    name = "pros_cons"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        pros = spec.get('pros') or []
        cons = spec.get('cons') or []
        hb = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(5.0), Inches(0.6))
        hb.text_frame.text = 'Pros'
        hb2 = slide.shapes.add_textbox(Inches(5.5), Inches(0.8), Inches(5.0), Inches(0.6))
        hb2.text_frame.text = 'Cons'
        lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.0), Inches(5.0))
        ltf = lb.text_frame; ltf.word_wrap = True
        if pros:
            ltf.text = str(pros[0])
            for it in pros[1:7]:
                p = ltf.add_paragraph(); p.text = str(it)
        rb = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(5.0), Inches(5.0))
        rtf = rb.text_frame; rtf.word_wrap = True
        if cons:
            rtf.text = str(cons[0])
            for it in cons[1:7]:
                p = rtf.add_paragraph(); p.text = str(it)

class BeforeAfterTemplate:
    name = "before_after"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        before = spec.get('before') or []
        after = spec.get('after') or []
        hb = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(5.0), Inches(0.6))
        hb.text_frame.text = 'Before'
        hb2 = slide.shapes.add_textbox(Inches(5.5), Inches(0.8), Inches(5.0), Inches(0.6))
        hb2.text_frame.text = 'After'
        lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.0), Inches(5.0))
        ltf = lb.text_frame; ltf.word_wrap = True
        if before:
            ltf.text = str(before[0])
            for it in before[1:7]:
                p = ltf.add_paragraph(); p.text = str(it)
        rb = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(5.0), Inches(5.0))
        rtf = rb.text_frame; rtf.word_wrap = True
        if after:
            rtf.text = str(after[0])
            for it in after[1:7]:
                p = rtf.add_paragraph(); p.text = str(it)

class TimelineTemplate:
    name = "timeline"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        milestones = spec.get('milestones') or []
        b = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10.5), Inches(5.4))
        tf = b.text_frame; tf.word_wrap = True
        tf.text = ' → '.join([str(m) for m in milestones[:8]])

class KPICardsTemplate:
    name = "kpi_cards"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        kpis = spec.get('kpis') or []
        x = 0.5
        for i, k in enumerate(kpis[:4]):
            box = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(2.3), Inches(2.0))
            box.text_frame.text = f"{k.get('name','KPI')}\n{k.get('value','-')}"
            x += 2.6

# Analysis templates (placeholders)
class SWOTTemplate:
    name = "swot"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        table = slide.shapes.add_table(2, 2, Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.0)).table
        labels = [["Strengths","Weaknesses"],["Opportunities","Threats"]]
        for r in range(2):
            for c in range(2):
                table.cell(r,c).text = labels[r][c]

class BCGMatrixTemplate:
    name = "bcg"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(0.6)).text_frame.text = "BCG Matrix (placeholder)"
        slide.shapes.add_table(2, 2, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))

class DashboardTemplate:
    name = "dashboard"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(0.6)).text_frame.text = "Dashboard (2x2)"
        xs = [0.5, 5.5]; ys = [1.5, 4.0]
        for y in ys:
            for x in xs:
                slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(2.0)).text_frame.text = "[Chart]"

class HeatmapTemplate:
    name = "heatmap"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(0.6)).text_frame.text = "Heatmap (placeholder)"

class TextChartBottomTemplate:
    name = "text_chart_bottom"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10.5), Inches(3.0))
        tb.text_frame.text = (spec.get('headline') or '')
        slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(10.5), Inches(2.0)).text_frame.text = "[Chart]"

class DualChartTemplate:
    name = "dual_chart"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(0.6)).text_frame.text = "Dual Charts"
        slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5), Inches(4.0)).text_frame.text = "[Chart A]"
        slide.shapes.add_textbox(Inches(5.6), Inches(1.6), Inches(5), Inches(4.0)).text_frame.text = "[Chart B]"

class WaterfallTemplate:
    name = "waterfall"
    def apply(self, slide: Slide, spec: Dict[str, Any]) -> None:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(0.6)).text_frame.text = "Waterfall (placeholder)"
