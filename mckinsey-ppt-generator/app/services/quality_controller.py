# -*- coding: utf-8 -*-
"""
?덉쭏 蹂댁쬆 而⑦듃濡ㅻ윭 (媛뺥솕 踰꾩쟾)
- 5媛 吏  ?덉쭏 湲곗? ?뺣? ?됯?
- HeadlineGenerator + InsightLadder ?듯빀
- 媛 以??됯퇏 ?먯닔 怨꾩궛
"""

from typing import Dict, List, Optional
from dataclasses import dataclass
from pptx import Presentation
import re
import logging

from app.services.slide_validator import SlideValidator, ValidationResult
from app.services.headline_generator import SoWhatTester
from app.services.insight_ladder import InsightLevel, InsightQualityEvaluator
from app.models.workflow_models import QualityScore as WorkflowQualityScore # Added

logger = logging.getLogger(__name__)


@dataclass
class QualityScore:
    """?덉쭏 ?먯닔"""
    clarity: float       # 紐낇솗??(20%)
    insight: float       # ?몄궗?댄듃 (25%)
    structure: float     # 援ъ“ (20%)
    visual: float        # ?쒓컖 (15%)
    actionability: float # ?ㅽ뻾媛 ?μ꽦 (20%)
    total: float         # 媛 以??됯퇏
    passed: bool         # 紐⑺몴 ?ъ꽦 ?щ?
    details: Dict        # ?몃? ?먯닔


class QualityController:
    """
    PPT ?덉쭏 ?됯? 諛?蹂댁쬆 (媛뺥솕 踰꾩쟾)
    
    Criteria:
    1. Clarity (紐낇솗?? 20%
       - So What ?뚯뒪??(40%)
       - ?ㅻ뱶?쇱씤 ?덉쭏 (30%)
       - 硫붿떆吏  ?쇨???(20%)
       - ?⑹뼱 ?쇨???(10%)
    
    2. Insight (?몄궗?댄듃) 25%
       - 4?④퀎 ?섎뜑 ?꾨떖 (40%)
       - ?곗씠??湲곕컲 (30%)
       - 鍮꾧탳 遺꾩꽍 (20%)
       - ?꾨왂???⑥쓽 (10%)
    
    3. Structure (援ъ“) 20%
       - MECE ?먯튃 (40%)
       - ?쇰━???먮쫫 (30%)
       - ?쇰씪誘몃뱶 援ъ“ (30%)
    
    4. Visual (?쒓컖) 15%
       - ?붿옄???쇨???(40%)
       - 媛 ?낆꽦 (30%)
       - McKinsey ?쒖? (30%)
    
    5. Actionability (?ㅽ뻾媛 ?μ꽦) 20%
       - 援ъ껜??沅뚭퀬 (50%)
       - ?뺣웾??(30%)
       - ?곗꽑?쒖쐞 (20%)
    """
    
    WEIGHTS = {
        "clarity": 0.20,
        "insight": 0.25,
        "structure": 0.20,
        "visual": 0.15,
        "actionability": 0.20
    }
    
    def __init__(self, target_score: float = 0.85):
        """
        珥덇린??        
        Args:
            target_score: 紐⑺몴 ?덉쭏 ?먯닔 (湲곕낯: 0.85)
        """
        self.target_score = target_score
        self.validator = SlideValidator()
        self.so_what_tester = SoWhatTester()
        self.insight_evaluator = InsightQualityEvaluator()
        self.logger = logging.getLogger(self.__class__.__name__)
    
    def evaluate(self, prs: Presentation) -> QualityScore:
        """
        ?꾨젅?좏뀒?댁뀡 ?덉쭏 ?됯?
        
        Args:
            prs: PowerPoint ?꾨젅?좏뀒?댁뀡
        
        Returns:
            QualityScore: 媛?湲곗?蹂??먯닔 諛?珥앹젏
        """
        try:
            scores = {}
            details = {}
            
            # 1. Clarity ?됯? (紐낇솗??
            scores["clarity"], details["clarity"] = self._evaluate_clarity(prs)
            
            # 2. Insight ?됯? (?몄궗?댄듃)
            scores["insight"], details["insight"] = self._evaluate_insight(prs)
            
            # 3. Structure ?됯? (援ъ“)
            scores["structure"], details["structure"] = self._evaluate_structure(prs)
            
            # 4. Visual ?됯? (?쒓컖)
            scores["visual"], details["visual"] = self._evaluate_visual(prs)
            
            # 5. Actionability ?됯? (?ㅽ뻾媛 ?μ꽦)
            scores["actionability"], details["actionability"] = self._evaluate_actionability(prs)
            
            # 媛 以??됯퇏 怨꾩궛
            total = sum(
                scores[criterion] * weight 
                for criterion, weight in self.WEIGHTS.items()
            )
            
            result = QualityScore(
                clarity=scores["clarity"],
                insight=scores["insight"],
                structure=scores["structure"],
                visual=scores["visual"],
                actionability=scores["actionability"],
                total=total,
                passed=total >= self.target_score,
                details=details
            )
            
            self.logger.info(f"Quality evaluation complete: {total:.3f}")
            return result
            
        except Exception as e:
            self.logger.error(f"Quality evaluation failed: {e}")
            # ?대갚: ??? ?먯닔 諛섑솚
            return QualityScore(
                clarity=0.5,
                insight=0.5,
                structure=0.5,
                visual=0.5,
                actionability=0.5,
                total=0.5,
                passed=False,
                details={"error": str(e)}
            )
    
    def _evaluate_clarity(self, prs: Presentation) -> tuple[float, Dict]:
        """
        紐낇솗???됯? (媛뺥솕)
        
        ?됯? ??ぉ:
        1. So What ?뚯뒪??(40%)
        2. ?ㅻ뱶?쇱씤 ?덉쭏 (30%)
        3. 硫붿떆吏  ?쇨???(20%)
        4. ?⑹뼱 ?쇨???(10%)
        """
        slide_scores = []
        details = {
            "so_what_passed": 0,
            "headline_quality": 0,
            "message_consistency": 0,
            "terminology_consistency": 0,
            "total_slides": len(prs.slides)
        }
        
        for slide in prs.slides:
            slide_score = 0.0
            
            # 1. So What ?뚯뒪??(40%)
            if self._has_title(slide):
                title = self._get_title_text(slide)
                so_what_result = self.so_what_tester.test(title)
                slide_score += so_what_result["score"] * 0.4
                
                if so_what_result["passed"]:
                    details["so_what_passed"] += 1
            
            # 2. ?ㅻ뱶?쇱씤 ?덉쭏 (30%)
            headline_score = self._evaluate_headline_quality(slide)
            slide_score += headline_score * 0.3
            details["headline_quality"] += headline_score
            
            # 3. 硫붿떆吏  ?쇨???(20%)
            consistency_score = self._evaluate_message_consistency(slide)
            slide_score += consistency_score * 0.2
            details["message_consistency"] += consistency_score
            
            # 4. ?⑹뼱 ?쇨???(10%)
            terminology_score = self._evaluate_terminology(slide)
            slide_score += terminology_score * 0.1
            details["terminology_consistency"] += terminology_score
            
            slide_scores.append(slide_score)
        
        # ?됯퇏 怨꾩궛
        avg_score = sum(slide_scores) / len(slide_scores) if slide_scores else 0.0
        
        # ?곸꽭 ?뺣낫 ?됯퇏??        if details["total_slides"] > 0:
            details["so_what_pass_rate"] = details["so_what_passed"] / details["total_slides"]
            details["avg_headline_quality"] = details["headline_quality"] / details["total_slides"]
            details["avg_message_consistency"] = details["message_consistency"] / details["total_slides"]
            details["avg_terminology"] = details["terminology_consistency"] / details["total_slides"]
        
        self.logger.info(f"Clarity score: {avg_score:.3f}")
        return avg_score, details
    
    def _evaluate_headline_quality(self, slide: 'Slide') -> float:
        """
        ?ㅻ뱶?쇱씤 ?덉쭏 ?됯?
        
        湲곗?:
        - ?숈궗 ?ы븿: +0.3
        - ?レ옄 ?ы븿: +0.3
        - 20???댁긽: +0.2
        - ?⑥쓽 ?ㅼ썙?? +0.2
        """
        if not self._has_title(slide):
            return 0.0
        
        title = self._get_title_text(slide)
        score = 0.0
        
        # ?숈궗 寃 ??        action_verbs = ["?쒓났", "?뺣낫", "?ъ꽦", "?ㅽ쁽", "媛 ??, "?꾩슂", "媛쒖꽑", "利앷?", "媛먯냼"]
        if any(verb in title for verb in action_verbs):
            score += 0.3
        
        # ?レ옄 寃 ??        if re.search(r'\d+', title):
            score += 0.3
        
        # 湲몄씠 寃 ??        if len(title) >= 20:
            score += 0.2
        
        # ?⑥쓽 ?ㅼ썙??寃 ??        implication_keywords = ["媛 ??, "?꾩슂", "?ㅽ쁽", "?뺣낫", "湲고쉶", "?꾪삊", "以묒슂", "?듭떖"]
        if any(keyword in title for keyword in implication_keywords):
            score += 0.2
        
        return min(1.0, score)
    
    def _evaluate_message_consistency(self, slide: 'Slide') -> float:
        """
        硫붿떆吏  ?쇨????됯?
        
        ?쒕ぉ怨?蹂몃Ц???쇨???寃 ??        """
        if not self._has_title(slide):
            return 0.5
        
        title = self._get_title_text(slide)
        content = self._extract_slide_content(slide)
        
        # ?ㅼ썙???쇱튂??寃 ??        title_keywords = set(re.findall(r'\w+', title.lower()))
        content_keywords = set(re.findall(r'\w+', content.lower()))
        
        if not title_keywords or not content_keywords:
            return 0.5
        
        # 援먯쭛??鍮꾩쑉
        overlap = len(title_keywords & content_keywords)
        union = len(title_keywords | content_keywords)
        
        consistency = overlap / union if union > 0 else 0.0
        
        # 0.3 ~ 1.0 踰붿쐞濡??뺢퇋??(?덈Т ??? ?먯닔 諛⑹?)
        return max(0.3, min(1.0, consistency * 2))
    
    def _evaluate_terminology(self, slide: 'Slide') -> float:
        """
        ?⑹뼱 ?쇨????됯?
        
        鍮꾩쫰?덉뒪/?꾨Ц ?⑹뼱 ?ъ슜 ?щ?
        """
        content = self._extract_slide_content(slide)
        
        # McKinsey 鍮꾩쫰?덉뒪 ?⑹뼱
        business_terms = [
            "?꾨왂", "?깆옣", "?쒖옣", "寃쎌웳", "媛 移?, "?⑥쑉", "理쒖쟻??,
            "?곸떊", "李⑤퀎??, "?ъ??붾떇", "?ㅽ뻾", "ROI", "KPI"
        ]
        
        term_count = sum(1 for term in business_terms if term in content)
        
        # 2媛??댁긽?대㈃ 留뚯젏
        if term_count >= 2:
            return 1.0
        elif term_count == 1:
            return 0.7
        else:
            return 0.5
    
    def _evaluate_insight(self, prs: Presentation) -> tuple[float, Dict]:
        """
        ?몄궗?댄듃 源딆씠 ?됯? (媛뺥솕)
        
        ?됯? 湲곗?:
        1. 4?④퀎 ?섎뜑 ?꾨떖 (40%)
        2. ?곗씠??湲곕컲 (30%)
        3. 鍮꾧탳 遺꾩꽍 ?ы븿 (20%)
        4. ?꾨왂???⑥쓽 (10%)
        """
        slide_scores = []
        details = {
            "avg_insight_level": 0,
            "data_based_slides": 0,
            "comparison_slides": 0,
            "strategic_slides": 0,
            "total_slides": len(prs.slides)
        }
        
        total_insight_level = 0
        
        for slide in prs.slides:
            slide_score = 0.0
            content = self._extract_slide_content(slide)
            
            # 1. ?몄궗?댄듃 ?덈꺼 寃 ??(40%)
            insight_level = self._detect_insight_level(content)
            slide_score += (insight_level / 4.0) * 0.4
            total_insight_level += insight_level
            
            # 2. ?곗씠??湲곕컲 寃 ??(30%)
            if self._has_quantification(content):
                slide_score += 0.3
                details["data_based_slides"] += 1
            
            # 3. 鍮꾧탳 遺꾩꽍 寃 ??(20%)
            comparison_keywords = ["? 鍮?, "鍮꾧탳", "諛?, "李⑥씠", "利앷?", "媛먯냼", "?믪?", "???"]
            if any(keyword in content for keyword in comparison_keywords):
                slide_score += 0.2
                details["comparison_slides"] += 1
            
            # 4. ?꾨왂???⑥쓽 寃 ??(10%)
            strategy_keywords = ["?꾨왂", "?꾩슂", "媛 ??, "沅뚭퀬", "?쒖븞", "?ㅽ뻾", "?ъ옄", "?뺣?"]
            if any(keyword in content for keyword in strategy_keywords):
                slide_score += 0.1
                details["strategic_slides"] += 1
            
            slide_scores.append(slide_score)
        
        # ?됯퇏 怨꾩궛
        avg_score = sum(slide_scores) / len(slide_scores) if slide_scores else 0.0
        
        # ?곸꽭 ?뺣낫
        if details["total_slides"] > 0:
            details["avg_insight_level"] = total_insight_level / details["total_slides"]
            details["data_based_rate"] = details["data_based_slides"] / details["total_slides"]
            details["comparison_rate"] = details["comparison_slides"] / details["total_slides"]
            details["strategic_rate"] = details["strategic_slides"] / details["total_slides"]
        
        self.logger.info(f"Insight score: {avg_score:.3f}, Avg level: {details.get('avg_insight_level', 0):.1f}")
        return avg_score, details
    
    def _detect_insight_level(self, content: str) -> int:
        """
        ?몄궗?댄듃 ?덈꺼 媛먯? (1-4)
        
        Level 1: ?⑥닚 ?쒖닠 ("留ㅼ텧??1000??)
        Level 2: 鍮꾧탳 ("?꾨뀈 ? 鍮?10% 利앷?")
        Level 3: ?먯씤 ("?좎젣?덉씠 70% 湲곗뿬")
        Level 4: ?꾨왂 ("?쇱씤 ?뺣? ?꾩슂")
        """
        level = 1
        
        # Level 2: 鍮꾧탳 ?ㅼ썙??        comparison_keywords = ["? 鍮?, "鍮꾧탳", "諛?, "利앷?", "媛먯냼", "?믪?", "???"]
        if any(word in content for word in comparison_keywords):
            level = 2
        
        # Level 3: ?먯씤 ?ㅼ썙??        implication_keywords = ["?먯씤", "湲곗뿬", "?곹뼢", "寃곌낵", "?④낵", "?뚮Ц", "?뺣텇", "?붿씤"]
        if any(word in content for word in implication_keywords):
            level = 3
        
        # Level 4: ?꾨왂 ?ㅼ썙??        action_keywords = ["?꾨왂", "?꾩슂", "沅뚭퀬", "?쒖븞", "?ㅽ뻾", "媛 ??, "?ъ옄", "?뺣?", "媛쒖꽑"]
        if any(word in content for word in action_keywords):
            level = 4
        
        return level
    
    def _evaluate_structure(self, prs: Presentation) -> tuple[float, Dict]:
        """
        援ъ“ ?쇰━???됯?
        
        湲곗?:
        1. MECE ?먯튃 (40%)
        2. ?쇰━???먮쫫 (30%)
        3. ?쇰씪誘몃뱶 援ъ“ (30%)
        """
        details = {
            "mece_score": 0,
    def _evaluate_structure(self, prs: Presentation) -> tuple[float, Dict]:
        """구조 평가 (MECE 40%, Flow 35%, Pyramid 25%)"""
        try:
            slide_specs: List[Dict] = []
            for s in prs.slides:
                title = ''
                try:
                    if s.shapes.title:
                        title = s.shapes.title.text or ''
                except Exception:
                    pass
                body_texts = []
                for sh in s.shapes:
                    if hasattr(sh, 'text_frame') and sh.text_frame and sh.text_frame.text:
                        body_texts.append(sh.text_frame.text)
                slide_specs.append({'title': title, 'content': body_texts, 'headline': body_texts[0] if body_texts else title})

            from app.services.mece_validator import MECEValidator
            from app.services.logic_flow_analyzer import LogicFlowAnalyzer
            mece = MECEValidator().validate_mece(slide_specs)
            flow = LogicFlowAnalyzer().analyze_flow(slide_specs)
            pyramid = self._evaluate_pyramid(slide_specs)

            structure_score = (mece.score * 0.40 + flow.transition_quality * 0.35 + pyramid * 0.25)
            details = {
                'mece_score': mece.score,
                'flow_transition_quality': flow.transition_quality,
                'pyramid_score': pyramid,
                'mece_overlaps': [(o.slide1, o.slide2, o.similarity) for o in mece.overlaps],
                'mece_gaps': [g.area for g in mece.gaps],
                'flow_gaps': flow.logic_gaps,
                'suggestions': mece.suggestions + flow.improvement_suggestions,
            }
            return max(0.0, min(1.0, structure_score)), details
        except Exception as e:
            logger.error(f"Structure evaluation failed: {e}")
            return 0.6, {"error": str(e)}

    def _evaluate_pyramid(self, slides: List[Dict]) -> float:
        """Pyramid Principle: conclusion first, supporting layers present."""
        if not slides:
            return 0.0
        first = slides[0].get('headline', '') or slides[0].get('title', '')
        import re as _re
        has_conclusion = bool(_re.search(r"결론|권고|요약|summary|recommend", first.lower()))
        support = any(len((s.get('content') or [])) >= 2 for s in slides[1:])
        score = 0.0
        if has_conclusion:
            score += 0.6
        if support:
            score += 0.4
        return score
        # ?ㅼ썙??湲곕컲 ?④퀎 媛먯?
        flow_score = 0.0
        
        # 泥??щ씪?대뱶: ?꾩엯/媛쒖슂
        intro_keywords = ["媛쒖슂", "?뚭컻", "諛곌꼍", "紐⑹쟻", "?붿빟"]
        if any(keyword in titles[0] for keyword in intro_keywords):
            flow_score += 0.3
        
        # 以묎컙 ?щ씪?대뱶: 遺꾩꽍/?곸꽭
        middle = titles[1:-1]
        analysis_keywords = ["遺꾩꽍", "?꾪솴", "臾몄젣", "?댁뒋", "湲고쉶", "?꾪삊"]
        if any(any(keyword in title for keyword in analysis_keywords) for title in middle):
            flow_score += 0.4
        
        # 留덉?留??щ씪?대뱶: 寃곕줎/沅뚭퀬
        conclusion_keywords = ["寃곕줎", "沅뚭퀬", "?쒖븞", "?ㅽ뻾", "?ㅼ쓬?④퀎", "?붿빟"]
        if any(keyword in titles[-1] for keyword in conclusion_keywords):
            flow_score += 0.3
        
        return max(0.7, flow_score)  # 理쒖냼 0.7
    
    def _evaluate_visual(self, prs: Presentation) -> tuple[float, Dict]:
        """
        ?쒓컖???덉쭏 ?됯? (SlideValidator ?쒖슜)
        
        湲곗?:
        1. ?붿옄???쇨???(40%)
        2. 媛 ?낆꽦 (30%)
        3. McKinsey ?쒖? (30%)
        """
        visual_issues = []
        details = {
            "total_issues": 0,
            "critical_issues": 0,
            "warning_issues": 0,
            "total_slides": len(prs.slides)
        }
        
        for slide in prs.slides:
            result = self.validator.validate_slide(slide)
            visual_issues.extend(result.issues)
            
            # ?댁뒋 遺꾨쪟
            for issue in result.issues:
                if issue.severity == "critical":
                    details["critical_issues"] += 1
                elif issue.severity == "warning":
                    details["warning_issues"] += 1
        
        details["total_issues"] = len(visual_issues)
        
        # ?먯닔 怨꾩궛 (?댁뒋 ?섏뿉 諛섎퉬濡 )
        if details["total_slides"] > 0:
            # ?щ씪?대뱶???됯퇏 ?댁뒋 ??            avg_issues_per_slide = details["total_issues"] / details["total_slides"]
            
            # 0 ?댁뒋 = 1.0, 10 ?댁뒋 = 0.0
            visual_score = max(0.0, 1.0 - (avg_issues_per_slide / 10.0))
        else:
            visual_score = 0.5
        
        self.logger.info(f"Visual score: {visual_score:.3f}, Issues: {details['total_issues']}")
        return visual_score, details
    
    def _evaluate_actionability(self, prs: Presentation) -> tuple[float, Dict]:
        """
        ?ㅽ뻾 媛 ?μ꽦 ?됯? (媛쒖꽑 踰꾩쟾)
        """
        slide_scores = []
        details = {
            "actionable_slides": 0,
            "quantified_slides": 0,
            "prioritized_slides": 0,
            "total_slides": len(prs.slides)
        }
        
        for slide in prs.slides:
            slide_score = 0.0
            content = self._extract_slide_content(slide)
            
            # ?붾쾭源? 肄섑뀗痢??뺤씤
            self.logger.debug(f"Evaluating slide content: {content[:100]}...")
            
            # 1. 援ъ껜??沅뚭퀬 (50%)
            action_keywords = ["沅뚭퀬", "?쒖븞", "?ㅽ뻾", "異붿쭊", "?꾩슂", "?댁빞", "?쒗뻾", "?ъ옄", "?뺣?", "媛쒖꽑"]
            if any(keyword in content for keyword in action_keywords):
                slide_score += 0.5
                details["actionable_slides"] += 1
                self.logger.debug("??Actionable keywords found")
            
            # 2. ?뺣웾??(30%)
            if self._has_quantification(content):
                slide_score += 0.3
                details["quantified_slides"] += 1
                self.logger.debug("??Quantification found")
            
            # 3. ?곗꽑?쒖쐞 (20%) - 媛뺥솕??媛먯?
            priority_patterns = [
                r'\[理쒖슦??]', r'\[?듭떖\]', r'\[以묒슂\]',  # ? 愿꾪샇 ?⑦꽩
                r'?곗꽑?쒖쐞\s*[1-3]', r'[1-3]?쒖쐞',        # ?レ옄 ?⑦꽩
                '?곗꽑', '?듭떖', '以묒슂', '湲닿툒', '理쒖슦??    # ?ㅼ썙??            ]
            
            has_priority = False
            for pattern in priority_patterns:
                if isinstance(pattern, str):
                    if pattern in content:
                        has_priority = True
                        break
                else:
                    if re.search(pattern, content):
                        has_priority = True
                        break
            
            if has_priority:
                slide_score += 0.2
                details["prioritized_slides"] += 1
                self.logger.debug("??Priority found")
            
            slide_scores.append(slide_score)
        
        # ?됯퇏 怨꾩궛
        avg_score = sum(slide_scores) / len(slide_scores) if slide_scores else 0.0
        
        # ?곸꽭 ?뺣낫
        if details["total_slides"] > 0:
            details["actionable_rate"] = details["actionable_slides"] / details["total_slides"]
            details["quantified_rate"] = details["quantified_slides"] / details["total_slides"]
            details["prioritized_rate"] = details["prioritized_slides"] / details["total_slides"]
        
        self.logger.info(
            f"Actionability: {avg_score:.3f} "
            f"(actionable: {details['actionable_slides']}, "
            f"quantified: {details['quantified_slides']}, "
            f"prioritized: {details['prioritized_slides']})"
        )
        
        return avg_score, details
    
    # === ?ы띁 硫붿꽌??===
    
    def _has_title(self, slide: 'Slide') -> bool:
        """?щ씪?대뱶???쒕ぉ???덈뒗吏  ?뺤씤"""
        try:
            return slide.shapes.title is not None and slide.shapes.title.has_text_frame
        except:
            return False
    
    def _get_title_text(self, slide: 'Slide') -> str:
        """?щ씪?대뱶 ?쒕ぉ ?띿뒪??異붿텧"""
        try:
            if self._has_title(slide):
                return slide.shapes.title.text
        except:
            pass
        return ""
    
    def _extract_slide_content(self, slide: 'Slide') -> str:
        """
        ?щ씪?대뱶 ?꾩껜 ?띿뒪??異붿텧"""
        content = []
        
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            content.append(text)
        except:
            pass
        
        return " ".join(content)
    
    def _has_quantification(self, text: str) -> bool:
        """?뺣웾???ы븿 ?щ?"""
        return bool(re.search(r'\d+', text))

    # WorkflowQualityScore濡?蹂 ?섑븯??硫붿꽌??異붽?
    def evaluate_to_workflow_score(self, prs: Presentation) -> WorkflowQualityScore:
        """
        WorkflowQualityScore ?뺤떇?쇰줈 ?됯? 寃곌낵 諛섑솚
        
        Args:
            prs: PowerPoint ?꾨젅?좏뀒?댁뀡
        
        Returns:
            WorkflowQualityScore: workflow_models??QualityScore 媛앹껜
        """
        result = self.evaluate(prs)
        
        workflow_score = WorkflowQualityScore(
            clarity=result.clarity,
            insight=result.insight,
            structure=result.structure,
            visual=result.visual,
            actionability=result.actionability,
            total=result.total,
            passed=result.passed,
            target_score=self.target_score
        )
        
        # 媛 以??됯퇏 ?ш퀎??        workflow_score.calculate_total(self.WEIGHTS)
        
        return workflow_score
