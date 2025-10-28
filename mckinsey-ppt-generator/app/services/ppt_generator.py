"""
Main PowerPoint generation engine using python-pptx.

This module provides comprehensive PPT generation capabilities including
slide creation, content formatting, chart generation, and file management
with McKinsey-style professional formatting.
"""

import os
import io
import logging
from typing import Dict, List, Optional, Any, Tuple, Union
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from sqlalchemy.orm import Session
from PIL import Image

from app.core.config import settings
from app.models.presentation import Presentation as PresentationModel, Slide as SlideModel
from app.services.content_analyzer import ContentAnalyzer
from app.services.layout_library import LayoutLibrary
from app.services.layout_applier import LayoutApplier
from app.services.text_fitter import TextFitter
from app.services.mckinsey_styles import McKinseyStyles, SlideLayoutType
from app.services.template_manager import TemplateManager

# Configure logging
logger = logging.getLogger(__name__)


@dataclass
class SlideContent:
    """Data structure for slide content"""
    title: str = ""
    subtitle: str = ""
    content: List[str] = None
    layout_type: str = "content"
    images: List[str] = None
    charts: List[Dict[str, Any]] = None
    tables: List[Dict[str, Any]] = None
    speaker_notes: str = ""
    
    def __post_init__(self):
        if self.content is None:
            self.content = []
        if self.images is None:
            self.images = []
        if self.charts is None:
            self.charts = []
        if self.tables is None:
            self.tables = []


@dataclass
class PresentationMetadata:
    """Metadata for generated presentations"""
    title: str
    author: str = ""
    subject: str = ""
    category: str = ""
    creation_time: datetime = None
    slide_count: int = 0
    word_count: int = 0
    generation_time: float = 0.0
    ai_model: str = ""
    ai_tokens_used: int = 0
    
    def __post_init__(self):
        if self.creation_time is None:
            self.creation_time = datetime.now()


class PPTGenerator:
    """Main PowerPoint generation engine"""
    
    def __init__(self, db: Session = None, output_dir: str = None, template_manager: TemplateManager = None):
        """
        Initialize PPT generator with layout optimization components
        
        Args:
            db: Database session
            output_dir: Directory for saving generated presentations
            template_manager: Template manager instance
        """
        self.db = db
        self.mckinsey_styles = McKinseyStyles()
        self.content_analyzer = ContentAnalyzer()
        self.layout_library = LayoutLibrary()
        self.layout_applier = LayoutApplier()
        self.text_fitter = TextFitter()
        
        # Set up output directory
        if output_dir:
            self.output_dir = Path(output_dir)
        else:
            self.output_dir = Path("output/presentations")
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Template manager
        if template_manager:
            self.template_manager = template_manager
        elif db:
            self.template_manager = TemplateManager(db)
        else:
            self.template_manager = None
    
    async def generate(self, result: Dict) -> Presentation:
        """
        AI 워크플로우 결과를 실제 PPT로 변환
        
        Args:
            result: WorkflowOrchestrator의 결과
                {
                    "slides": [
                        {
                            "type": "title",
                            "headline": "AI 생성 헤드라인", 
                            "subtitle": "부제목",
                            "content": ["포인트1", "포인트2"],
                            "chart": {...}
                        },
                        ...
                    ],
                    "quality_score": 0.85
                }
        
        Returns:
            Presentation: 생성된 PPT 객체
        """
        logger.info("🎨 Generating PPT from AI result...")
        
        prs = Presentation()
        
        # AI 결과에서 슬라이드 데이터 추출
        slides_data = result.get("slides", [])
        
        if not slides_data:
            logger.warning("⚠️ No slides data in result, generating fallback")
            return self._generate_fallback_ppt()
        
        logger.info(f"📊 Creating {len(slides_data)} slides from AI content")
        
        # McKinsey 색상 정의
        mckinsey_blue = RGBColor(0, 118, 168)  # #0076A8
        text_gray = RGBColor(83, 86, 90)       # #53565A
        
        # 각 슬라이드를 AI 데이터로 생성
        for idx, slide_data in enumerate(slides_data, 1):
            try:
                slide_type = slide_data.get("type", "content")
                logger.info(f"  Creating slide {idx}/{len(slides_data)}: {slide_type}")
                
                # 슬라이드 타입에 따른 레이아웃 선택
                if slide_type == "title":
                    layout = prs.slide_layouts[0]  # Title Slide
                else:
                    layout = prs.slide_layouts[1]  # Title and Content
                
                slide = prs.slides.add_slide(layout)
                
                # AI 생성 헤드라인 적용 (핵심!)
                if slide.shapes.title:
                    headline = slide_data.get("headline", "")
                    slide.shapes.title.text = headline
                    logger.info(f"    ✨ AI Headline: {headline}")
                    
                    # McKinsey 스타일 적용
                    title_frame = slide.shapes.title.text_frame
                    for paragraph in title_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = mckinsey_blue
                
                # 부제목 (표지 슬라이드용)
                if slide_type == "title" and len(slide.shapes) > 1:
                    subtitle = slide_data.get("subtitle", "")
                    slide.shapes[1].text = subtitle
                
                # AI 생성 콘텐츠 적용
                content_list = slide_data.get("content", [])
                if content_list and len(slide.shapes) > 1:
                    # Content placeholder 찾기
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            text_frame = shape.text_frame
                            text_frame.clear()
                            
                            for item in content_list:
                                p = text_frame.add_paragraph()
                                p.text = item
                                p.level = 0
                                p.font.size = Pt(14)
                                p.font.color.rgb = text_gray
                                p.space_after = Pt(12)
                            break
                
            except Exception as e:
                logger.error(f"  ❌ Failed to create slide {idx}: {e}")
                # 실패해도 계속 진행
        
        logger.info(f"✅ PPT generation complete: {len(prs.slides)} slides")
        return prs
    
    def _generate_fallback_ppt(self) -> Presentation:
        """긴급 폴백 (AI 데이터 없을 때만)"""
        logger.warning("🚨 Generating emergency fallback PPT")
        prs = Presentation()
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Error: No AI Content"
        
        return prs
    
    def create_presentation(self, slides_content: List[SlideContent], 
                          metadata: PresentationMetadata,
                          template_name: str = "McKinsey Professional",
                          save_to_db: bool = True) -> Tuple[Optional[str], Optional[PresentationModel]]:
        """
        Create a complete PowerPoint presentation
        
        Args:
            slides_content: List of slide content data
            metadata: Presentation metadata
            template_name: Name of template to use
            save_to_db: Whether to save to database
            
        Returns:
            Tuple of (file_path, presentation_model)
        """
        start_time = datetime.now()
        
        try:
            # Get template
            template = self.template_manager.get_template_by_name(template_name)
            if not template:
                logger.warning(f"Template '{template_name}' not found, using default")
                template = self.template_manager.get_template_by_name("McKinsey Professional")
            
            # Create presentation
            prs = Presentation()
            
            # Apply template styling
            if template:
                self.template_manager.apply_template_to_presentation(prs, template)
                self.template_manager.increment_usage_count(template.id)
            
            # Generate slides
            for i, slide_content in enumerate(slides_content):
                try:
                    slide = self._create_slide(prs, slide_content, template, i + 1)
                    if not slide:
                        logger.warning(f"Failed to create slide {i + 1}")
                except Exception as e:
                    logger.error(f"Error creating slide {i + 1}: {e}")
                    continue
            
            # Set presentation properties
            self._set_presentation_properties(prs, metadata)
            
            # Generate filename and save
            filename = self._generate_filename(metadata.title)
            file_path = self.output_dir / filename
            
            prs.save(str(file_path))
            
            # Calculate metrics
            generation_time = (datetime.now() - start_time).total_seconds()
            metadata.generation_time = generation_time
            metadata.slide_count = len(prs.slides)
            metadata.word_count = self._count_words_in_presentation(prs)
            
            # Save to database if requested
            presentation_model = None
            if save_to_db:
                presentation_model = self._save_to_database(
                    slides_content, metadata, str(file_path), template
                )
            
            logger.info(f"Created presentation: {file_path}")
            logger.info(f"Generation time: {generation_time:.2f}s, Slides: {metadata.slide_count}")
            
            return str(file_path), presentation_model
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return None, None
    
    def _create_slide(self, presentation: Presentation, content: SlideContent, 
                     template, slide_number: int) -> Optional[Slide]:
        """Create a single slide with content"""
        try:
            # Determine layout based on content type
            layout_index = self._get_layout_index(content.layout_type)
            slide_layout = presentation.slide_layouts[layout_index]
            slide = presentation.slides.add_slide(slide_layout)
            
            # Get layout configuration from template
            if template:
                layout_config = self.template_manager.get_layout_config(
                    template, content.layout_type
                )
            else:
                layout_config = self.mckinsey_styles.get_layout_config(
                    SlideLayoutType(content.layout_type)
                )
            
            # Add content based on layout type
            if content.layout_type == "title":
                self._create_title_slide(slide, content, layout_config)
            elif content.layout_type == "two_column":
                self._create_two_column_slide(slide, content, layout_config)
            elif content.layout_type == "chart":
                self._create_chart_slide(slide, content, layout_config)
            elif content.layout_type == "image":
                self._create_image_slide(slide, content, layout_config)
            elif content.layout_type == "table":
                self._create_table_slide(slide, content, layout_config)
            else:
                self._create_content_slide(slide, content, layout_config)
            
            # Add speaker notes
            if content.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = content.speaker_notes
            
            # Add footer with slide number
            self.mckinsey_styles.create_header_footer(
                slide, 
                footer_text=f"Slide {slide_number}",
                page_number=True
            )
            
            return slide
            
        except Exception as e:
            logger.error(f"Error creating slide: {e}")
            return None
    
    def _create_title_slide(self, slide: Slide, content: SlideContent, 
                          layout_config: Dict[str, Any]) -> None:
        """Create title slide"""
        try:
            # Add title
            if content.title:
                title_pos = layout_config.get("title_position")
                if title_pos:
                    title_box = slide.shapes.add_textbox(*title_pos)
                    title_frame = title_box.text_frame
                    title_frame.text = content.title
                    self.mckinsey_styles.apply_text_formatting(title_frame, "title")
                else:
                    # Use default title placeholder
                    if slide.shapes.title:
                        slide.shapes.title.text = content.title
            
            # Add subtitle
            if content.subtitle:
                subtitle_pos = layout_config.get("subtitle_position")
                if subtitle_pos:
                    subtitle_box = slide.shapes.add_textbox(*subtitle_pos)
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_frame.text = content.subtitle
                    self.mckinsey_styles.apply_text_formatting(subtitle_frame, "subtitle")
                else:
                    # Use placeholders if available
                    placeholders = [shape for shape in slide.shapes if hasattr(shape, 'placeholder_format')]
                    if len(placeholders) > 1:
                        placeholders[1].text = content.subtitle
                        
        except Exception as e:
            logger.error(f"Error creating title slide: {e}")
    
    def _create_content_slide(self, slide: Slide, content: SlideContent, 
                            layout_config: Dict[str, Any]) -> None:
        """Create content slide with bullet points"""
        try:
            # Add title
            if content.title and slide.shapes.title:
                slide.shapes.title.text = content.title
                if hasattr(slide.shapes.title, 'text_frame'):
                    self.mckinsey_styles.apply_text_formatting(
                        slide.shapes.title.text_frame, "heading"
                    )
            
            # Add content
            if content.content:
                content_pos = layout_config.get("content_position", 
                    self.mckinsey_styles.layouts.get_content_position())
                
                content_box = slide.shapes.add_textbox(*content_pos)
                content_frame = content_box.text_frame
                
                # Create bullet points
                self.mckinsey_styles.create_bullet_points(content_frame, content.content)
                
        except Exception as e:
            logger.error(f"Error creating content slide: {e}")
    
    def _create_two_column_slide(self, slide: Slide, content: SlideContent, 
                               layout_config: Dict[str, Any]) -> None:
        """Create two-column layout slide with overflow prevention"""
        try:
            # Add title (limit length for Korean text)
            if content.title and slide.shapes.title:
                title_text = content.title[:50] if len(content.title) > 50 else content.title
                slide.shapes.title.text = title_text
                if hasattr(slide.shapes.title, 'text_frame'):
                    self.mckinsey_styles.apply_text_formatting(
                        slide.shapes.title.text_frame, "heading"
                    )
            
            # Split content between columns (limit to prevent overflow)
            max_items_per_column = 4
            if content.content:
                # Process content to ensure it fits
                processed_content = []
                for item in content.content[:max_items_per_column * 2]:
                    # Truncate long items
                    if len(item) > 60:
                        item = item[:57] + "..."
                    processed_content.append(item)
                
                mid_point = len(processed_content) // 2
                left_content = processed_content[:mid_point]
                right_content = processed_content[mid_point:]
            else:
                left_content = []
                right_content = []
            
            # Left column
            if left_content:
                left_pos = layout_config.get("left_column_position",
                    self.mckinsey_styles.layouts.get_left_column_position())
                
                left_box = slide.shapes.add_textbox(*left_pos)
                left_frame = left_box.text_frame
                self.mckinsey_styles.create_bullet_points(left_frame, left_content)
            
            # Right column
            if right_content:
                right_pos = layout_config.get("right_column_position",
                    self.mckinsey_styles.layouts.get_right_column_position())
                
                right_box = slide.shapes.add_textbox(*right_pos)
                right_frame = right_box.text_frame
                self.mckinsey_styles.create_bullet_points(right_frame, right_content)
                
        except Exception as e:
            logger.error(f"Error creating two-column slide: {e}")
    
    def _create_chart_slide(self, slide: Slide, content: SlideContent, 
                          layout_config: Dict[str, Any]) -> None:
        """Create slide with charts"""
        try:
            # Add title
            if content.title and slide.shapes.title:
                slide.shapes.title.text = content.title
                if hasattr(slide.shapes.title, 'text_frame'):
                    self.mckinsey_styles.apply_text_formatting(
                        slide.shapes.title.text_frame, "heading"
                    )
            
            # Add charts
            if content.charts:
                chart_pos = layout_config.get("chart_position")
                if not chart_pos:
                    chart_pos = (
                        self.mckinsey_styles.layouts.MARGIN_LEFT + Inches(1.0),
                        self.mckinsey_styles.layouts.CONTENT_TOP,
                        self.mckinsey_styles.layouts.CHART_WIDTH,
                        self.mckinsey_styles.layouts.CHART_HEIGHT
                    )
                
                # Create first chart (support for multiple charts can be added)
                chart_data = content.charts[0]
                self._add_chart_to_slide(slide, chart_data, chart_pos)
                
        except Exception as e:
            logger.error(f"Error creating chart slide: {e}")
    
    def _create_image_slide(self, slide: Slide, content: SlideContent, 
                          layout_config: Dict[str, Any]) -> None:
        """Create slide with images"""
        try:
            # Add title
            if content.title and slide.shapes.title:
                slide.shapes.title.text = content.title
                if hasattr(slide.shapes.title, 'text_frame'):
                    self.mckinsey_styles.apply_text_formatting(
                        slide.shapes.title.text_frame, "heading"
                    )
            
            # Add images
            if content.images:
                image_path = content.images[0]  # Use first image
                if os.path.exists(image_path):
                    # Calculate image position
                    img_left = self.mckinsey_styles.layouts.MARGIN_LEFT + Inches(1.0)
                    img_top = self.mckinsey_styles.layouts.CONTENT_TOP
                    img_width = Inches(6.0)
                    img_height = Inches(4.0)
                    
                    # Add image
                    slide.shapes.add_picture(
                        image_path, img_left, img_top, img_width, img_height
                    )
            
            # Add caption if content available
            if content.content:
                caption_top = self.mckinsey_styles.layouts.CONTENT_TOP + Inches(4.5)
                caption_box = slide.shapes.add_textbox(
                    self.mckinsey_styles.layouts.MARGIN_LEFT,
                    caption_top,
                    self.mckinsey_styles.layouts.CONTENT_WIDTH,
                    Inches(1.0)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = " ".join(content.content)
                self.mckinsey_styles.apply_text_formatting(caption_frame, "body")
                
        except Exception as e:
            logger.error(f"Error creating image slide: {e}")
    
    def _create_table_slide(self, slide: Slide, content: SlideContent, 
                          layout_config: Dict[str, Any]) -> None:
        """Create slide with tables"""
        try:
            # Add title
            if content.title and slide.shapes.title:
                slide.shapes.title.text = content.title
                if hasattr(slide.shapes.title, 'text_frame'):
                    self.mckinsey_styles.apply_text_formatting(
                        slide.shapes.title.text_frame, "heading"
                    )
            
            # Add table
            if content.tables:
                table_data = content.tables[0]  # Use first table
                self._add_table_to_slide(slide, table_data, layout_config)
                
        except Exception as e:
            logger.error(f"Error creating table slide: {e}")
    
    def _add_chart_to_slide(self, slide: Slide, chart_data: Dict[str, Any], 
                          position: Tuple[float, float, float, float]) -> None:
        """Add chart to slide"""
        try:
            chart_type = chart_data.get("type", "column")
            chart_title = chart_data.get("title", "")
            categories = chart_data.get("categories", [])
            series_data = chart_data.get("series", [])
            
            if not categories or not series_data:
                logger.warning("Chart data incomplete, skipping chart creation")
                return
            
            # Create chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = categories
            
            for series in series_data:
                series_name = series.get("name", "Series")
                values = series.get("values", [])
                chart_data_obj.add_series(series_name, values)
            
            # Determine chart type
            pptx_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            if chart_type.lower() == "line":
                pptx_chart_type = XL_CHART_TYPE.LINE
            elif chart_type.lower() == "pie":
                pptx_chart_type = XL_CHART_TYPE.PIE
            elif chart_type.lower() == "bar":
                pptx_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
            
            # Add chart to slide
            chart = slide.shapes.add_chart(
                pptx_chart_type, 
                position[0], position[1], position[2], position[3],
                chart_data_obj
            ).chart
            
            # Apply McKinsey styling
            if chart_title:
                chart.chart_title.text_frame.text = chart_title
            
            # Style chart (simplified)
            self._style_chart(chart)
            
        except Exception as e:
            logger.error(f"Error adding chart to slide: {e}")
    
    def _add_table_to_slide(self, slide: Slide, table_data: Dict[str, Any], 
                          layout_config: Dict[str, Any]) -> None:
        """Add table to slide"""
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            
            if not headers and not rows:
                logger.warning("Table data incomplete, skipping table creation")
                return
            
            # Calculate table dimensions
            table_rows = len(rows) + (1 if headers else 0)
            table_cols = max(len(headers), max(len(row) for row in rows) if rows else 0)
            
            if table_rows == 0 or table_cols == 0:
                return
            
            # Position table
            table_left = self.mckinsey_styles.layouts.MARGIN_LEFT
            table_top = self.mckinsey_styles.layouts.CONTENT_TOP
            table_width = self.mckinsey_styles.layouts.CONTENT_WIDTH
            table_height = Inches(3.0)
            
            # Add table
            table = slide.shapes.add_table(
                table_rows, table_cols,
                table_left, table_top, table_width, table_height
            ).table
            
            # Add headers
            if headers:
                for col_idx, header in enumerate(headers):
                    if col_idx < table_cols:
                        cell = table.cell(0, col_idx)
                        cell.text = str(header)
                        # Apply header styling
                        self._style_table_cell(cell, is_header=True)
            
            # Add data rows
            for row_idx, row_data in enumerate(rows):
                table_row_idx = row_idx + (1 if headers else 0)
                if table_row_idx < table_rows:
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < table_cols:
                            cell = table.cell(table_row_idx, col_idx)
                            cell.text = str(cell_data)
                            # Apply data cell styling
                            self._style_table_cell(cell, is_header=False)
                            
        except Exception as e:
            logger.error(f"Error adding table to slide: {e}")
    
    def _style_chart(self, chart) -> None:
        """Apply McKinsey styling to chart"""
        try:
            chart_style = self.mckinsey_styles.get_chart_style_config()
            
            # Apply colors to series
            for i, series in enumerate(chart.series):
                color = self.mckinsey_styles.colors.get_chart_color(i)
                # Apply color (simplified implementation)
                
        except Exception as e:
            logger.error(f"Error styling chart: {e}")
    
    def _style_table_cell(self, cell, is_header: bool = False) -> None:
        """Apply McKinsey styling to table cell"""
        try:
            table_style = self.mckinsey_styles.get_table_style_config()
            
            # Set background color
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = table_style["header_background"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = table_style["row_background_1"]
            
            # Style text
            if cell.text_frame and cell.text_frame.paragraphs:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        font.name = table_style["font_name"]
                        font.size = table_style["font_size"]
                        
                        if is_header:
                            font.color.rgb = table_style["header_text_color"]
                            font.bold = True
                        else:
                            font.color.rgb = table_style["text_color"]
                            
        except Exception as e:
            logger.error(f"Error styling table cell: {e}")
    
    def _get_layout_index(self, layout_type: str) -> int:
        """Get slide layout index based on layout type"""
        layout_mapping = {
            "title": 0,
            "content": 1,
            "two_column": 3,
            "blank": 6,
            "chart": 1,
            "image": 1,
            "table": 1
        }
        return layout_mapping.get(layout_type, 1)
    
    def _set_presentation_properties(self, presentation: Presentation, 
                                   metadata: PresentationMetadata) -> None:
        """Set presentation properties and metadata"""
        try:
            core_props = presentation.core_properties
            core_props.title = metadata.title
            core_props.author = metadata.author
            core_props.subject = metadata.subject
            core_props.category = metadata.category
            core_props.created = metadata.creation_time
            
        except Exception as e:
            logger.error(f"Error setting presentation properties: {e}")
    
    def _count_words_in_presentation(self, presentation: Presentation) -> int:
        """Count total words in presentation"""
        try:
            word_count = 0
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        words = shape.text.split()
                        word_count += len(words)
                        
                # Count speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame.text:
                    words = slide.notes_slide.notes_text_frame.text.split()
                    word_count += len(words)
            
            return word_count
            
        except Exception as e:
            logger.error(f"Error counting words: {e}")
            return 0
    
    def _generate_filename(self, title: str) -> str:
        """Generate filename for presentation"""
        try:
            # Clean title for filename
            clean_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
            clean_title = clean_title.replace(' ', '_')
            
            # Add timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            return f"{clean_title}_{timestamp}.pptx"
            
        except Exception as e:
            logger.error(f"Error generating filename: {e}")
            return f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    
    def _save_to_database(self, slides_content: List[SlideContent], 
                         metadata: PresentationMetadata,
                         file_path: str, template) -> Optional[PresentationModel]:
        """Save presentation to database"""
        try:
            # Create presentation record
            presentation = PresentationModel(
                title=metadata.title,
                description=metadata.subject,
                user_email=metadata.author,
                presentation_metadata={
                    "category": metadata.category,
                    "creation_time": metadata.creation_time.isoformat(),
                    "template_used": template.name if template else "None"
                },
                settings={
                    "template_id": template.id if template else None,
                    "ai_model": metadata.ai_model,
                    "ai_tokens_used": metadata.ai_tokens_used
                },
                outline=[{
                    "slide_number": i + 1,
                    "title": slide.title,
                    "layout_type": slide.layout_type
                } for i, slide in enumerate(slides_content)],
                content=[{
                    "slide_number": i + 1,
                    "title": slide.title,
                    "subtitle": slide.subtitle,
                    "content": slide.content,
                    "layout_type": slide.layout_type,
                    "speaker_notes": slide.speaker_notes
                } for i, slide in enumerate(slides_content)],
                file_path=file_path,
                slide_count=metadata.slide_count,
                word_count=metadata.word_count,
                ai_model=metadata.ai_model,
                ai_tokens_used=metadata.ai_tokens_used,
                generation_time=int(metadata.generation_time)
            )
            
            # Save to database
            self.db.add(presentation)
            self.db.commit()
            self.db.refresh(presentation)
            
            # Create slide records
            for i, slide_content in enumerate(slides_content):
                slide_record = SlideModel(
                    presentation_id=presentation.id,
                    slide_number=i + 1,
                    title=slide_content.title,
                    subtitle=slide_content.subtitle,
                    content={
                        "text_content": slide_content.content,
                        "charts": slide_content.charts,
                        "tables": slide_content.tables
                    },
                    layout_type=slide_content.layout_type,
                    images=slide_content.images,
                    charts=slide_content.charts,
                    speaker_notes=slide_content.speaker_notes
                )
                
                self.db.add(slide_record)
            
            self.db.commit()
            
            logger.info(f"Saved presentation to database: {presentation.id}")
            return presentation
            
        except Exception as e:
            logger.error(f"Error saving to database: {e}")
            self.db.rollback()
            return None
    
    def load_presentation_from_file(self, file_path: str) -> Optional[Presentation]:
        """Load existing presentation from file"""
        try:
            if not os.path.exists(file_path):
                logger.error(f"Presentation file not found: {file_path}")
                return None
            
            presentation = Presentation(file_path)
            logger.info(f"Loaded presentation from: {file_path}")
            return presentation
            
        except Exception as e:
            logger.error(f"Error loading presentation: {e}")
            return None
    
    def update_slide_content(self, presentation: Presentation, 
                           slide_index: int, new_content: SlideContent) -> bool:
        """Update content of a specific slide"""
        try:
            if slide_index >= len(presentation.slides):
                logger.error(f"Slide index {slide_index} out of range")
                return False
            
            slide = presentation.slides[slide_index]
            
            # Update title
            if new_content.title and slide.shapes.title:
                slide.shapes.title.text = new_content.title
            
            # Update content (simplified implementation)
            # This would need more sophisticated logic to handle different content types
            
            logger.info(f"Updated slide {slide_index + 1}")
            return True
            
        except Exception as e:
            logger.error(f"Error updating slide content: {e}")
            return False
    
    def export_to_pdf(self, presentation_path: str, pdf_path: str = None) -> Optional[str]:
        """Export presentation to PDF (requires additional dependencies)"""
        try:
            # Note: PDF export requires additional tools like LibreOffice or comtypes
            # This is a placeholder implementation
            logger.warning("PDF export not implemented yet")
            return None
            
        except Exception as e:
            logger.error(f"Error exporting to PDF: {e}")
            return None
    
    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about a presentation file"""
        try:
            if not os.path.exists(file_path):
                return {}
            
            presentation = Presentation(file_path)
            
            info = {
                "file_path": file_path,
                "file_size": os.path.getsize(file_path),
                "slide_count": len(presentation.slides),
                "created": datetime.fromtimestamp(os.path.getctime(file_path)),
                "modified": datetime.fromtimestamp(os.path.getmtime(file_path))
            }
            
            # Get core properties if available
            if hasattr(presentation, 'core_properties'):
                core_props = presentation.core_properties
                info.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "category": core_props.category or ""
                })
            
            return info
            
        except Exception as e:
            logger.error(f"Error getting presentation info: {e}")
            return {}


def create_sample_presentation() -> List[SlideContent]:
    """Create sample presentation content for testing"""
    slides = [
        SlideContent(
            title="Strategic Initiative Overview",
            subtitle="McKinsey & Company Professional Template",
            layout_type="title"
        ),
        SlideContent(
            title="Executive Summary",
            content=[
                "Key findings from market analysis",
                "Strategic recommendations for growth",
                "Implementation roadmap and timeline",
                "Expected outcomes and ROI projections"
            ],
            layout_type="content",
            speaker_notes="Present key findings with confidence and data backing"
        ),
        SlideContent(
            title="Market Analysis Results",
            content=[
                "Market size: $2.5B globally",
                "Growth rate: 15% annually",
                "Competitive landscape analysis",
                "Customer segmentation insights"
            ],
            layout_type="two_column",
            charts=[{
                "type": "column",
                "title": "Market Growth Projection",
                "categories": ["2023", "2024", "2025", "2026"],
                "series": [{
                    "name": "Market Size (B$)",
                    "values": [2.1, 2.4, 2.8, 3.2]
                }]
            }]
        )
    ]
    
    return slides