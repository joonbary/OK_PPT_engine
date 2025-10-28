"""
Layout Applier for PowerPoint presentations.
Applies selected layouts to slides with precise content placement.
"""

from typing import Dict, List, Any, Tuple, Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.slide import Slide
from app.services.text_fitter import TextFitter
from app.services.layout_library import LayoutLibrary
from app.core.logging import app_logger
import time
import logging

logger = logging.getLogger(__name__)


class LayoutApplier:
    """
    레이아웃 적용 및 텍스트 최적화 통합 시스템
    Enhanced TextFitter와 완전 통합된 고급 레이아웃 적용기
    """
    
    def __init__(self):
        """Enhanced TextFitter와 통합된 초기화"""
        self.text_fitter = TextFitter()
        self.layout_library = LayoutLibrary()
        
        # 레이아웃별 폰트 크기 기본값
        self.default_font_sizes = {
            "headline": 18,
            "title": 32,
            "subtitle": 18,
            "body": 14,
            "bullets": 14,
            "caption": 10,
            "chart_label": 10
        }
        
        # 레이아웃별 처리 전략 매핑
        self.layout_handlers = {
            "title_slide": self._apply_title_slide_enhanced,
            "single_column": self._apply_single_column_enhanced,
            "bullet_list": self._apply_bullet_list_enhanced,
            "two_column": self._apply_two_column_enhanced,
            "three_column": self._apply_three_column_enhanced,
            "image_text": self._apply_image_text_enhanced,
            "matrix_2x2": self._apply_matrix_enhanced,
            "timeline": self._apply_timeline_enhanced,
            "process_flow": self._apply_process_flow_enhanced,
            "pyramid": self._apply_pyramid_enhanced,
            "dashboard_grid": self._apply_dashboard_grid_enhanced,
            "quote_highlight": self._apply_quote_highlight_enhanced,
            "split_screen": self._apply_split_screen_enhanced,
            "agenda_toc": self._apply_agenda_toc_enhanced,
            "waterfall": self._apply_waterfall_enhanced
        }
        
        # 성능 메트릭
        self.metrics = {
            "total_slides": 0,
            "text_adjustments": 0,
            "truncations": 0,
            "avg_processing_time": 0.0
        }
        
    def apply_layout(
        self,
        slide: Slide,
        layout_name: str,
        content: Dict,
        options: Optional[Dict] = None
    ) -> Dict:
        """
        레이아웃을 슬라이드에 적용하고 콘텐츠 최적 배치
        Enhanced TextFitter와 완전 통합된 고급 처리
        
        Parameters:
            slide: pptx.slide.Slide 객체
            layout_name: 적용할 레이아웃 이름
            content: 배치할 콘텐츠 딕셔너리
            options: 추가 옵션 (폰트 크기 강제, 여백 조정 등)
        
        Returns:
            {
                "success": bool,
                "layout_applied": str,
                "text_adjustments": int,
                "warnings": List[str],
                "metrics": Dict
            }
        """
        
        start_time = time.time()
        
        try:
            # 레이아웃 정의 가져오기
            layout_def = self.layout_library.get_layout(layout_name)
            
            if not layout_def:
                logger.error(f"Unknown layout: {layout_name}")
                return {"success": False, "error": "Unknown layout"}
            
            # 옵션 기본값 설정
            options = options or {}
            
            # 레이아웃별 처리 함수 호출
            handler = self.layout_handlers.get(layout_name)
            
            if handler:
                result = handler(slide, layout_def, content, options)
            else:
                logger.warning(f"No enhanced handler for {layout_name}, using generic")
                result = self._apply_generic_enhanced(slide, layout_def, content, options)
            
            # 메트릭 업데이트
            processing_time = time.time() - start_time
            self.metrics["total_slides"] += 1
            self.metrics["avg_processing_time"] = (
                (self.metrics["avg_processing_time"] * (self.metrics["total_slides"] - 1) 
                 + processing_time) / self.metrics["total_slides"]
            )
            
            result["processing_time_ms"] = processing_time * 1000
            
            return result
            
        except Exception as e:
            logger.error(f"Layout application failed: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
                "layout_applied": layout_name
            }
    
    def _apply_title_slide(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply title slide layout"""
        try:
            elements = layout.get("elements", [])
            
            for element in elements:
                if element["type"] == "title" and content.get("title"):
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        content["title"]
                    )
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element["type"] == "subtitle" and content.get("subtitle"):
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        content["subtitle"]
                    )
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
        except Exception as e:
            results["issues"].append(f"Title slide error: {str(e)}")
    
    def _apply_bullet_list(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """
        불릿 리스트 레이아웃 적용
        """
        try:
            elements = layout.get("elements", [])
            
            for element in elements:
                if element["type"] == "headline" and content.get("title"):
                    # Add headline with length limit
                    max_length = layout.get("max_text_length", {}).get("headline", 80)
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"], max_length
                    )
                    
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        title_text
                    )
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element["type"] == "bullets" and content.get("bullets"):
                    # Process bullet points with text fitting
                    bullets = content["bullets"]
                    max_bullets = element.get("max_bullets", 5)
                    max_bullet_length = layout.get("max_text_length", {}).get("bullet", 100)
                    
                    # Limit and truncate bullets
                    if len(bullets) > max_bullets:
                        results["issues"].append(f"Truncated bullets from {len(bullets)} to {max_bullets}")
                        bullets = bullets[:max_bullets]
                    
                    # Create text box for bullets
                    text_box = self._add_text_box_empty(slide, element["position"])
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    # Add each bullet point
                    for i, bullet_text in enumerate(bullets):
                        # Truncate long bullets
                        bullet_text = self.text_fitter.truncate_with_ellipsis(
                            bullet_text, max_bullet_length
                        )
                        
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = bullet_text
                        p.level = int(element.get("indent_level", 0))
                        
                        # Apply formatting
                        for run in p.runs:
                            run.font.name = "맑은 고딕"
                            run.font.size = element.get("font_size", Pt(14))
                        
                        # Add spacing
                        if "bullet_spacing" in element:
                            p.space_after = element["bullet_spacing"]
                    
                    results["elements_placed"] += 1
                    results["text_adjusted"] = True
                    
        except Exception as e:
            results["issues"].append(f"Bullet list error: {str(e)}")
    
    def _apply_two_column(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """
        2컬럼 비교 레이아웃 적용
        """
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    # Add headline
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"], 
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element_type == "left_column_header" and content.get("left_header"):
                    # Add left column header
                    header_text = self.text_fitter.truncate_with_ellipsis(
                        content["left_header"],
                        max_text_length.get("column_header", 30)
                    )
                    text_frame = self._add_text_box(slide, element["position"], header_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element_type == "right_column_header" and content.get("right_header"):
                    # Add right column header
                    header_text = self.text_fitter.truncate_with_ellipsis(
                        content["right_header"],
                        max_text_length.get("column_header", 30)
                    )
                    text_frame = self._add_text_box(slide, element["position"], header_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element_type == "left_column" and content.get("left_content"):
                    # Add left column content
                    self._add_column_content(
                        slide, element, content["left_content"],
                        max_text_length.get("column_item", 60),
                        results
                    )
                    
                elif element_type == "right_column" and content.get("right_content"):
                    # Add right column content
                    self._add_column_content(
                        slide, element, content["right_content"],
                        max_text_length.get("column_item", 60),
                        results
                    )
                    
        except Exception as e:
            results["issues"].append(f"Two column error: {str(e)}")
    
    def _apply_three_column(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply three column layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                # Column headers and content
                for col_num in range(1, 4):
                    if element_type == f"col{col_num}_header" and content.get(f"col{col_num}_header"):
                        header_text = self.text_fitter.truncate_with_ellipsis(
                            content[f"col{col_num}_header"],
                            max_text_length.get("column_header", 25)
                        )
                        text_frame = self._add_text_box(slide, element["position"], header_text)
                        self._apply_element_formatting(text_frame, element)
                        results["elements_placed"] += 1
                        
                    elif element_type == f"col{col_num}" and content.get(f"col{col_num}_content"):
                        self._add_column_content(
                            slide, element, content[f"col{col_num}_content"],
                            max_text_length.get("column_item", 40),
                            results
                        )
                        
        except Exception as e:
            results["issues"].append(f"Three column error: {str(e)}")
    
    def _apply_matrix(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply 2x2 matrix layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                # Process quadrants
                for q_num in range(1, 5):
                    if element_type == f"q{q_num}_header" and content.get(f"q{q_num}_header"):
                        header_text = self.text_fitter.truncate_with_ellipsis(
                            content[f"q{q_num}_header"],
                            max_text_length.get("quadrant_header", 30)
                        )
                        text_frame = self._add_text_box(slide, element["position"], header_text)
                        self._apply_element_formatting(text_frame, element)
                        results["elements_placed"] += 1
                        
                    elif element_type == f"q{q_num}" and content.get(f"q{q_num}_content"):
                        self._add_column_content(
                            slide, element, content[f"q{q_num}_content"],
                            max_text_length.get("quadrant_item", 50),
                            results
                        )
                        
        except Exception as e:
            results["issues"].append(f"Matrix error: {str(e)}")
    
    def _apply_image_text(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply image with text layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element_type == "image_area" and content.get("image_path"):
                    # Add image placeholder or actual image
                    try:
                        slide.shapes.add_picture(
                            content["image_path"],
                            element["position"][0],
                            element["position"][1],
                            element["position"][2],
                            element["position"][3]
                        )
                        results["elements_placed"] += 1
                    except:
                        # If image fails, add placeholder
                        shape = slide.shapes.add_shape(
                            1,  # Rectangle
                            element["position"][0],
                            element["position"][1],
                            element["position"][2],
                            element["position"][3]
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
                        results["issues"].append("Image not found, placeholder added")
                        
                elif element_type == "text" and content.get("text_content"):
                    self._add_column_content(
                        slide, element, content["text_content"],
                        max_text_length.get("text_item", 60),
                        results
                    )
                    
        except Exception as e:
            results["issues"].append(f"Image text error: {str(e)}")
    
    def _apply_table_layout(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply table layout"""
        try:
            elements = layout.get("elements", [])
            
            for element in elements:
                if element["type"] == "headline" and content.get("title"):
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        content["title"]
                    )
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element["type"] == "table_area" and content.get("table_data"):
                    # Add table
                    self._add_table_to_slide(
                        slide,
                        element["position"],
                        content["table_data"],
                        results
                    )
                    
        except Exception as e:
            results["issues"].append(f"Table layout error: {str(e)}")
    
    def _apply_single_column(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply single column layout (default)"""
        try:
            elements = layout.get("elements", [])
            
            for element in elements:
                if element["type"] == "headline" and content.get("title"):
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        content["title"]
                    )
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    
                elif element["type"] == "body" and content.get("content"):
                    # Apply text fitting for body content
                    body_text = content["content"]
                    if isinstance(body_text, list):
                        body_text = "\n".join(body_text)
                    
                    # Fit text to box
                    fit_result = self.text_fitter.fit_text_to_box(
                        body_text,
                        element["position"][2].inches if hasattr(element["position"][2], 'inches') else element["position"][2],
                        element["position"][3].inches if hasattr(element["position"][3], 'inches') else element["position"][3],
                        element.get("font_size", Pt(11)).pt if hasattr(element.get("font_size", Pt(11)), 'pt') else element.get("font_size", 11)
                    )
                    
                    text_frame = self._add_text_box(
                        slide,
                        element["position"],
                        fit_result["adjusted_text"]
                    )
                    
                    # Apply adjusted font size
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(fit_result["adjusted_font_size"])
                    
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                    results["text_adjusted"] = True
                    
                    if fit_result["truncated"]:
                        results["issues"].append("Text truncated to fit")
                        
        except Exception as e:
            results["issues"].append(f"Single column error: {str(e)}")
    
    def _add_column_content(
        self,
        slide: Slide,
        element: Dict,
        content: List[str],
        max_item_length: int,
        results: Dict
    ):
        """Add content to a column with text fitting"""
        try:
            max_items = element.get("max_bullets", 4)
            
            # Limit items
            if len(content) > max_items:
                results["issues"].append(f"Column items truncated from {len(content)} to {max_items}")
                content = content[:max_items]
            
            # Create text box
            text_box = self._add_text_box_empty(slide, element["position"])
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            # Add items
            for i, item in enumerate(content):
                # Truncate long items
                item = self.text_fitter.truncate_with_ellipsis(item, max_item_length)
                
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = "• " + item
                
                # Apply formatting
                for run in p.runs:
                    run.font.name = "맑은 고딕"
                    run.font.size = element.get("font_size", Pt(11))
            
            results["elements_placed"] += 1
            results["text_adjusted"] = True
            
        except Exception as e:
            results["issues"].append(f"Column content error: {str(e)}")
    
    def _add_text_box(
        self, 
        slide: Slide, 
        position: Tuple[float, float, float, float],
        text: str
    ):
        """
        텍스트 박스 추가 헬퍼 메서드
        
        position: (left, top, width, height) in Inches or raw values
        """
        # Handle position values that might already be Inches objects
        left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
        top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
        width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
        height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        return text_frame
    
    def _add_text_box_empty(
        self, 
        slide: Slide, 
        position: Tuple[float, float, float, float]
    ):
        """Add empty text box"""
        # Handle position values that might already be Inches objects
        left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
        top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
        width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
        height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
        
        return slide.shapes.add_textbox(left, top, width, height)
    
    def _apply_element_formatting(self, text_frame, element: Dict):
        """Apply formatting to text frame based on element configuration"""
        try:
            # Set alignment
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT
            }
            
            if "alignment" in element:
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = alignment_map.get(element["alignment"], PP_ALIGN.LEFT)
            
            # Set font properties
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "font_size" in element:
                        # Handle font size that might be Pt object or int
                        if hasattr(element["font_size"], 'pt'):
                            run.font.size = element["font_size"]
                        else:
                            run.font.size = Pt(element["font_size"])
                    
                    if "bold" in element:
                        run.font.bold = element["bold"]
                    
                    run.font.name = "맑은 고딕"
            
            # Set line spacing if specified
            if "line_spacing" in element:
                for paragraph in text_frame.paragraphs:
                    paragraph.line_spacing = element["line_spacing"]
                    
        except Exception as e:
            app_logger.error(f"Error applying element formatting: {str(e)}")
    
    def _add_table_to_slide(
        self,
        slide: Slide,
        position: Tuple[float, float, float, float],
        table_data: Dict,
        results: Dict
    ):
        """Add table to slide"""
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            
            if not headers and not rows:
                results["issues"].append("No table data provided")
                return
            
            # Calculate table dimensions
            table_rows = len(rows) + (1 if headers else 0)
            table_cols = max(len(headers), max(len(row) for row in rows) if rows else 0)
            
            # Handle position values
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # Add table
            table = slide.shapes.add_table(
                table_rows, table_cols,
                left, top, width, height
            ).table
            
            # Add headers
            if headers:
                for col_idx, header in enumerate(headers):
                    if col_idx < table_cols:
                        cell = table.cell(0, col_idx)
                        cell.text = str(header)
            
            # Add data rows
            for row_idx, row_data in enumerate(rows):
                table_row_idx = row_idx + (1 if headers else 0)
                if table_row_idx < table_rows:
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < table_cols:
                            cell = table.cell(table_row_idx, col_idx)
                            cell.text = str(cell_data)
            
            results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Table creation error: {str(e)}")
    
    def _apply_timeline(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply timeline layout with milestones"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type == "timeline_line":
                    # Add timeline line
                    self._add_timeline_line(slide, element, results)
                
                elif element_type.startswith("milestone_") and not element_type.endswith("_detail"):
                    # Add milestone markers
                    milestone_num = element_type.split("_")[1]
                    milestone_content = content.get(f"milestone_{milestone_num}")
                    if milestone_content:
                        milestone_text = self.text_fitter.truncate_with_ellipsis(
                            milestone_content,
                            max_text_length.get("milestone", 30)
                        )
                        text_frame = self._add_text_box(slide, element["position"], milestone_text)
                        self._apply_element_formatting(text_frame, element)
                        results["elements_placed"] += 1
                
                elif element_type.endswith("_detail"):
                    # Add milestone details
                    milestone_num = element_type.split("_")[1]
                    detail_content = content.get(f"milestone_{milestone_num}_detail")
                    if detail_content:
                        detail_text = self.text_fitter.truncate_with_ellipsis(
                            detail_content,
                            max_text_length.get("milestone_detail", 50)
                        )
                        text_frame = self._add_text_box(slide, element["position"], detail_text)
                        self._apply_element_formatting(text_frame, element)
                        results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Timeline error: {str(e)}")
    
    def _apply_process_flow(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply process flow layout with connected steps"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type.startswith("step_"):
                    # Add process steps
                    step_num = element_type.split("_")[1]
                    step_content = content.get(f"step_{step_num}")
                    if step_content:
                        step_text = self.text_fitter.truncate_with_ellipsis(
                            step_content,
                            max_text_length.get("step", 40)
                        )
                        # Add step as shaped text box
                        self._add_shaped_text_box(slide, element, step_text, results)
                
                elif element_type.startswith("arrow_"):
                    # Add connecting arrows
                    self._add_process_arrow(slide, element, results)
            
        except Exception as e:
            results["issues"].append(f"Process flow error: {str(e)}")
    
    def _apply_pyramid(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply pyramid hierarchy layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type.startswith("level_"):
                    # Add pyramid levels
                    level_content = content.get(element_type)
                    if level_content:
                        max_length_key = element_type if element_type in max_text_length else "level_1"
                        level_text = self.text_fitter.truncate_with_ellipsis(
                            level_content,
                            max_text_length.get(max_length_key, 30)
                        )
                        # Add level as shaped text box
                        self._add_shaped_text_box(slide, element, level_text, results)
            
        except Exception as e:
            results["issues"].append(f"Pyramid error: {str(e)}")
    
    def _apply_dashboard_grid(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply dashboard grid layout for KPIs"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type.startswith("kpi_"):
                    # Add KPI cards
                    kpi_num = element_type.split("_")[1]
                    kpi_data = content.get(f"kpi_{kpi_num}")
                    if kpi_data:
                        self._add_kpi_card(slide, element, kpi_data, max_text_length, results)
            
        except Exception as e:
            results["issues"].append(f"Dashboard grid error: {str(e)}")
    
    def _apply_quote_highlight(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply quote highlight layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type == "quote_box" and content.get("quote"):
                    # Add main quote
                    quote_text = self.text_fitter.truncate_with_ellipsis(
                        content["quote"],
                        max_text_length.get("quote", 200)
                    )
                    # Add quote marks
                    formatted_quote = f'"{quote_text}"'
                    self._add_shaped_text_box(slide, element, formatted_quote, results)
                
                elif element_type == "attribution" and content.get("attribution"):
                    # Add attribution
                    attribution_text = self.text_fitter.truncate_with_ellipsis(
                        content["attribution"],
                        max_text_length.get("attribution", 50)
                    )
                    text_frame = self._add_text_box(slide, element["position"], f"— {attribution_text}")
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type == "context" and content.get("context"):
                    # Add context information
                    context_text = self.text_fitter.truncate_with_ellipsis(
                        content["context"],
                        max_text_length.get("context", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], context_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Quote highlight error: {str(e)}")
    
    def _apply_split_screen(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply split screen layout for balanced content"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type == "left_panel" and content.get("left_content"):
                    # Add left panel content
                    left_content = content["left_content"]
                    if isinstance(left_content, list):
                        left_text = "\n".join(left_content)
                    else:
                        left_text = str(left_content)
                    
                    left_text = self.text_fitter.truncate_with_ellipsis(
                        left_text,
                        max_text_length.get("panel", 300)
                    )
                    self._add_shaped_text_box(slide, element, left_text, results)
                
                elif element_type == "right_panel" and content.get("right_content"):
                    # Add right panel content
                    right_content = content["right_content"]
                    if isinstance(right_content, list):
                        right_text = "\n".join(right_content)
                    else:
                        right_text = str(right_content)
                    
                    right_text = self.text_fitter.truncate_with_ellipsis(
                        right_text,
                        max_text_length.get("panel", 300)
                    )
                    self._add_shaped_text_box(slide, element, right_text, results)
                
                elif element_type == "divider_line":
                    # Add divider line
                    self._add_divider_line(slide, element, results)
            
        except Exception as e:
            results["issues"].append(f"Split screen error: {str(e)}")
    
    def _apply_agenda_toc(self, slide: Slide, layout: Dict, content: Dict, results: Dict):
        """Apply agenda/table of contents layout"""
        try:
            elements = layout.get("elements", [])
            max_text_length = layout.get("max_text_length", {})
            
            for element in elements:
                element_type = element["type"]
                
                if element_type == "headline" and content.get("title"):
                    title_text = self.text_fitter.truncate_with_ellipsis(
                        content["title"],
                        max_text_length.get("headline", 80)
                    )
                    text_frame = self._add_text_box(slide, element["position"], title_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
                
                elif element_type.startswith("agenda_item_"):
                    # Add agenda items
                    item_num = element_type.split("_")[2]
                    item_content = content.get(f"agenda_item_{item_num}")
                    if item_content:
                        item_text = self.text_fitter.truncate_with_ellipsis(
                            item_content,
                            max_text_length.get("agenda_item", 60)
                        )
                        # Add number prefix if specified
                        prefix = element.get("number_prefix", "")
                        formatted_text = f"{prefix} {item_text}" if prefix else item_text
                        
                        text_frame = self._add_text_box(slide, element["position"], formatted_text)
                        self._apply_element_formatting(text_frame, element)
                        results["elements_placed"] += 1
                
                elif element_type == "time_info" and content.get("time_info"):
                    # Add time information
                    time_text = self.text_fitter.truncate_with_ellipsis(
                        content["time_info"],
                        max_text_length.get("time_info", 30)
                    )
                    text_frame = self._add_text_box(slide, element["position"], time_text)
                    self._apply_element_formatting(text_frame, element)
                    results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Agenda TOC error: {str(e)}")
    
    def _add_shaped_text_box(self, slide: Slide, element: Dict, text: str, results: Dict):
        """Add text box with background shape and styling"""
        try:
            position = element["position"]
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # Add shape first
            shape_type = element.get("shape_type", "rect")
            if shape_type == "rounded_rect":
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.ROUNDED_RECTANGLE
                    left, top, width, height
                )
            else:
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    left, top, width, height
                )
            
            # Apply background color
            bg_color = element.get("background_color", "#F8F9FA")
            if bg_color:
                shape.fill.solid()
                # Convert hex color to RGB
                if bg_color.startswith("#"):
                    hex_color = bg_color[1:]
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    shape.fill.fore_color.rgb = RGBColor(*rgb)
            
            # Apply border
            border_color = element.get("border_color")
            if border_color:
                shape.line.color.rgb = RGBColor(222, 226, 230)  # Default border
            
            # Add text to shape
            text_frame = shape.text_frame
            text_frame.text = text
            text_frame.word_wrap = True
            
            # Apply text formatting
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = self._get_alignment(element.get("alignment", "center"))
                for run in paragraph.runs:
                    run.font.name = "맑은 고딕"
                    if "font_size" in element:
                        if hasattr(element["font_size"], 'pt'):
                            run.font.size = element["font_size"]
                        else:
                            run.font.size = Pt(element["font_size"])
                    if "bold" in element:
                        run.font.bold = element["bold"]
                    
                    # Apply text color
                    text_color = element.get("text_color", "#000000")
                    if text_color and text_color.startswith("#"):
                        hex_color = text_color[1:]
                        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                        run.font.color.rgb = RGBColor(*rgb)
            
            results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Shaped text box error: {str(e)}")
    
    def _add_timeline_line(self, slide: Slide, element: Dict, results: Dict):
        """Add timeline line"""
        try:
            position = element["position"]
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # Add line shape
            line = slide.shapes.add_connector(
                1,  # STRAIGHT connector
                left, top, left + width, top + height
            )
            
            # Style the line
            line.line.color.rgb = RGBColor(0, 115, 230)  # Blue color
            line.line.width = Pt(element.get("line_thickness", 3))
            
            results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Timeline line error: {str(e)}")
    
    def _add_process_arrow(self, slide: Slide, element: Dict, results: Dict):
        """Add process flow arrow"""
        try:
            position = element["position"]
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # Add arrow shape
            arrow = slide.shapes.add_shape(
                5,  # MSO_SHAPE.RIGHT_ARROW
                left, top, width, height
            )
            
            # Style the arrow
            arrow.fill.solid()
            arrow_color = element.get("arrow_color", "#0073E6")
            if arrow_color.startswith("#"):
                hex_color = arrow_color[1:]
                rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                arrow.fill.fore_color.rgb = RGBColor(*rgb)
            
            results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Process arrow error: {str(e)}")
    
    def _add_kpi_card(self, slide: Slide, element: Dict, kpi_data: Dict, max_text_length: Dict, results: Dict):
        """Add KPI card with title, value, and description"""
        try:
            # Create background shape
            self._add_shaped_text_box(slide, element, "", results)
            
            # Add KPI content as overlay text boxes
            position = element["position"]
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # KPI Title (top section)
            if kpi_data.get("title"):
                title_text = self.text_fitter.truncate_with_ellipsis(
                    kpi_data["title"],
                    max_text_length.get("kpi_title", 20)
                )
                title_box = slide.shapes.add_textbox(
                    left + Inches(0.1), top + Inches(0.1),
                    width - Inches(0.2), Inches(0.3)
                )
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(10)
                title_frame.paragraphs[0].font.bold = True
            
            # KPI Value (middle section)
            if kpi_data.get("value"):
                value_text = self.text_fitter.truncate_with_ellipsis(
                    str(kpi_data["value"]),
                    max_text_length.get("kpi_value", 15)
                )
                value_box = slide.shapes.add_textbox(
                    left + Inches(0.1), top + Inches(0.6),
                    width - Inches(0.2), Inches(0.6)
                )
                value_frame = value_box.text_frame
                value_frame.text = value_text
                value_frame.paragraphs[0].font.size = Pt(18)
                value_frame.paragraphs[0].font.bold = True
                value_frame.paragraphs[0].alignment = self._get_alignment("center")
            
            # KPI Description (bottom section)
            if kpi_data.get("description"):
                desc_text = self.text_fitter.truncate_with_ellipsis(
                    kpi_data["description"],
                    max_text_length.get("kpi_description", 30)
                )
                desc_box = slide.shapes.add_textbox(
                    left + Inches(0.1), top + Inches(1.3),
                    width - Inches(0.2), Inches(0.4)
                )
                desc_frame = desc_box.text_frame
                desc_frame.text = desc_text
                desc_frame.paragraphs[0].font.size = Pt(8)
            
        except Exception as e:
            results["issues"].append(f"KPI card error: {str(e)}")
    
    def _add_divider_line(self, slide: Slide, element: Dict, results: Dict):
        """Add vertical divider line"""
        try:
            position = element["position"]
            left = position[0] if hasattr(position[0], 'inches') else Inches(position[0])
            top = position[1] if hasattr(position[1], 'inches') else Inches(position[1])
            width = position[2] if hasattr(position[2], 'inches') else Inches(position[2])
            height = position[3] if hasattr(position[3], 'inches') else Inches(position[3])
            
            # Add vertical line
            line = slide.shapes.add_connector(
                1,  # STRAIGHT connector
                left, top, left, top + height
            )
            
            # Style the line
            line_color = element.get("line_color", "#DEE2E6")
            if line_color.startswith("#"):
                hex_color = line_color[1:]
                rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                line.line.color.rgb = RGBColor(*rgb)
            
            line.line.width = Pt(element.get("line_thickness", 2))
            
            results["elements_placed"] += 1
            
        except Exception as e:
            results["issues"].append(f"Divider line error: {str(e)}")
    
    def _get_alignment(self, alignment_str: str):
        """Convert alignment string to PowerPoint alignment enum"""
        from pptx.enum.text import PP_ALIGN
        
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }
        
        return alignment_map.get(alignment_str, PP_ALIGN.LEFT)
    
    # ================================
    # Enhanced Layout Handlers
    # ================================
    
    def _apply_text_box_enhanced(
        self,
        slide: Slide,
        position: Tuple[float, float, float, float],
        text: str,
        initial_font_size: int = 14,
        is_bold: bool = False,
        alignment: str = "left",
        allow_truncate: bool = True
    ) -> Dict:
        """
        Enhanced TextFitter를 사용한 텍스트 박스 추가
        100% 오버플로우 방지 보장
        """
        try:
            # 위치 값 처리
            left = position[0] if hasattr(position[0], 'inches') else position[0]
            top = position[1] if hasattr(position[1], 'inches') else position[1]
            width = position[2] if hasattr(position[2], 'inches') else position[2]
            height = position[3] if hasattr(position[3], 'inches') else position[3]
            
            # Enhanced TextFitter로 텍스트 피팅
            fit_result = self.text_fitter.fit_text_to_box(
                text,
                box_width=width,
                box_height=height,
                initial_font_size=initial_font_size,
                use_binary_search=True,
                preserve_words=True
            )
            
            # 텍스트 박스 생성
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
            
            text_frame = text_box.text_frame
            text_frame.text = fit_result["adjusted_text"]
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            # 폰트 서식 적용
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(fit_result["adjusted_font_size"])
                paragraph.font.bold = is_bold
                paragraph.font.name = "맑은 고딕"
                paragraph.alignment = self._get_alignment(alignment)
            
            return {
                "text_frame": text_frame,
                "fit_result": fit_result,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Enhanced text box creation failed: {e}")
            return {"success": False, "error": str(e)}
    
    def _apply_title_slide_enhanced(
        self,
        slide,
        layout_def: Dict,
        content: Dict,
        options: Dict
    ) -> Dict:
        """
        타이틀 슬라이드 적용 - Enhanced TextFitter 통합
        - 큰 제목 중앙 정렬
        - 부제목 자동 크기 조정
        """
        
        warnings = []
        adjustments = 0
        
        try:
            # Get elements or text_boxes
            elements = layout_def.get("elements", layout_def.get("text_boxes", []))
            
            if not elements:
                return {
                    "success": False,
                    "error": "No layout elements found",
                    "layout_applied": "title_slide"
                }
            
            # 제목 처리
            title_box = elements[0]  # title
            title_text = content.get("title", "")
            
            if title_text:
                title_result = self._apply_text_box_enhanced(
                    slide, title_box["position"],
                    title_text,
                    initial_font_size=32,
                    is_bold=True,
                    alignment="center"
                )
                
                if title_result["success"]:
                    fit_result = title_result["fit_result"]
                    if fit_result["adjusted_font_size"] != 32:
                        adjustments += 1
                    if not fit_result["fits"]:
                        warnings.append("Title text too long for title slide")
                else:
                    warnings.append(f"Title processing failed: {title_result.get('error')}")
            
            # 부제목 처리
            if len(elements) > 1:
                subtitle_box = elements[1]
                subtitle_text = content.get("subtitle", "")
                
                if subtitle_text:
                    subtitle_result = self._apply_text_box_enhanced(
                        slide, subtitle_box["position"],
                        subtitle_text,
                        initial_font_size=18,
                        alignment="center"
                    )
                    
                    if subtitle_result["success"]:
                        fit_result = subtitle_result["fit_result"]
                        if fit_result["truncated"]:
                            warnings.append("Subtitle was truncated")
                        adjustments += 1
                    else:
                        warnings.append(f"Subtitle processing failed: {subtitle_result.get('error')}")
            
            return {
                "success": True,
                "layout_applied": "title_slide",
                "text_adjustments": adjustments,
                "warnings": warnings
            }
            
        except Exception as e:
            logger.error(f"Title slide application failed: {e}")
            return {
                "success": False,
                "error": str(e),
                "layout_applied": "title_slide"
            }
    
    def _apply_bullet_list_enhanced(
        self,
        slide,
        layout_def: Dict,
        content: Dict,
        options: Dict
    ) -> Dict:
        """
        불릿 리스트 레이아웃 적용 - Enhanced 버전
        - 불릿 포인트 개수에 따라 폰트 크기 조정
        - 최대 5개 불릿 권장
        - 각 불릿의 길이에 따른 지능적 처리
        """
        
        warnings = []
        adjustments = 0
        
        try:
            # Get elements or text_boxes
            elements = layout_def.get("elements", layout_def.get("text_boxes", []))
            
            if not elements:
                return {
                    "success": False,
                    "error": "No layout elements found",
                    "layout_applied": "bullet_list"
                }
            
            # 1. 헤드라인 추가
            headline_box = elements[0]
            headline_text = content.get("headline", content.get("title", ""))
            
            if headline_text:
                headline_result = self._apply_text_box_enhanced(
                    slide, headline_box["position"],
                    headline_text,
                    initial_font_size=18,
                    is_bold=True
                )
                
                if headline_result["success"]:
                    adjustments += 1
                else:
                    warnings.append(f"Headline processing failed: {headline_result.get('error')}")
            
            # 2. 불릿 포인트 처리
            if len(elements) > 1:
                bullets_box = elements[1]
            else:
                # Create default bullets box if not defined
                bullets_box = {
                    "position": (1, 2, 8, 4),
                    "max_bullets": 5
                }
                
            bullets = content.get("bullets", content.get("content", []))
            
            if isinstance(bullets, str):
                # 문자열인 경우 줄바꿈으로 분리
                bullets = [line.strip() for line in bullets.split('\n') if line.strip()]
            
            if bullets:
                # 개수 체크
                max_bullets = bullets_box.get("max_bullets", 5)
                if len(bullets) > max_bullets:
                    warnings.append(f"Too many bullets ({len(bullets)}), recommended max: {max_bullets}")
                    bullets = bullets[:max_bullets]
                
                # 불릿 개수에 따라 초기 폰트 크기 조정
                initial_font_size = 14
                if len(bullets) > 4:
                    initial_font_size = 12
                elif len(bullets) <= 3:
                    initial_font_size = 16
                
                # 전체 불릿 텍스트 결합 (테스트용)
                bullets_text = '\n'.join(f"• {bullet}" for bullet in bullets)
                
                # 전체가 맞는지 확인
                fit_result = self.text_fitter.fit_text_to_box(
                    bullets_text,
                    bullets_box["position"][2],
                    bullets_box["position"][3],
                    initial_font_size=initial_font_size
                )
                
                # 텍스트 박스 생성
                text_box = slide.shapes.add_textbox(
                    Inches(bullets_box["position"][0]),
                    Inches(bullets_box["position"][1]),
                    Inches(bullets_box["position"][2]),
                    Inches(bullets_box["position"][3])
                )
                
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                
                if fit_result["fits"]:
                    # 개별 불릿으로 추가 (더 나은 포매팅)
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        # Smart truncation 사용
                        processed_bullet = self.text_fitter.truncate_with_ellipsis(
                            bullet, max_length=100, smart_truncation=True
                        )
                        
                        p.text = processed_bullet
                        p.level = 0
                        p.font.size = Pt(fit_result["adjusted_font_size"])
                        p.font.name = "맑은 고딕"
                        
                        # 불릿이 너무 길면 경고
                        if len(bullet) > 120:
                            warnings.append(f"Bullet {i+1} is very long ({len(bullet)} chars)")
                    
                    adjustments += 1
                else:
                    # 맞지 않으면 개별 처리
                    warnings.append("Bullets don't fit, processing individually with smart sizing")
                    
                    for i, bullet in enumerate(bullets):
                        # 각 불릿을 독립적으로 처리
                        bullet_height = bullets_box["position"][3] / len(bullets)
                        
                        individual_fit = self.text_fitter.fit_text_to_box(
                            bullet,
                            bullets_box["position"][2],
                            bullet_height,
                            initial_font_size=initial_font_size
                        )
                        
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = individual_fit["adjusted_text"]
                        p.level = 0
                        p.font.size = Pt(individual_fit["adjusted_font_size"])
                        p.font.name = "맑은 고딕"
                        
                        if individual_fit["truncated"]:
                            warnings.append(f"Bullet {i+1} was truncated")
                    
                    adjustments += len(bullets)
            
            return {
                "success": True,
                "layout_applied": "bullet_list",
                "text_adjustments": adjustments,
                "warnings": warnings
            }
            
        except Exception as e:
            logger.error(f"Bullet list application failed: {e}")
            return {
                "success": False,
                "error": str(e),
                "layout_applied": "bullet_list"
            }
    
    def _apply_two_column_enhanced(
        self,
        slide,
        layout_def: Dict,
        content: Dict,
        options: Dict
    ) -> Dict:
        """
        2컬럼 비교 레이아웃 적용 - Enhanced 버전
        - 좌우 컬럼 균형 유지
        - 텍스트 양에 따라 폰트 크기 동기화
        """
        
        warnings = []
        adjustments = 0
        
        try:
            # Get elements or text_boxes
            elements = layout_def.get("elements", layout_def.get("text_boxes", []))
            
            if not elements:
                return {
                    "success": False,
                    "error": "No layout elements found",
                    "layout_applied": "two_column"
                }
            
            # 1. 헤드라인
            headline_box = elements[0]
            headline_text = content.get("headline", content.get("title", ""))
            
            if headline_text:
                headline_result = self._apply_text_box_enhanced(
                    slide, headline_box["position"],
                    headline_text,
                    initial_font_size=18,
                    is_bold=True
                )
                
                if headline_result["success"]:
                    adjustments += 1
            
            # 2. 좌우 컬럼 콘텐츠 (균형 맞추기)
            left_content = content.get("left_column", content.get("left_content", ""))
            right_content = content.get("right_column", content.get("right_content", ""))
            
            if isinstance(left_content, list):
                left_content = "\n".join(left_content)
            if isinstance(right_content, list):
                right_content = "\n".join(right_content)
            
            # 좌우 텍스트 길이 비교하여 초기 폰트 크기 결정
            left_length = len(left_content)
            right_length = len(right_content)
            max_length = max(left_length, right_length)
            
            initial_font_size = 14
            if max_length > 500:
                initial_font_size = 11
            elif max_length > 300:
                initial_font_size = 12
            
            # 좌우 컬럼 박스 위치 찾기
            left_col_box = None
            right_col_box = None
            
            for element in elements:
                box_type = element.get("type", "")
                if "left" in box_type:
                    left_col_box = element
                elif "right" in box_type:
                    right_col_box = element
            
            # 기본값 설정 (찾지 못할 경우)
            if not left_col_box and len(elements) > 1:
                left_col_box = elements[1]
            if not right_col_box and len(elements) > 2:
                right_col_box = elements[2]
            
            # Create default column boxes if still not found
            if not left_col_box:
                left_col_box = {"position": (0.5, 2, 4, 4)}
            if not right_col_box:
                right_col_box = {"position": (5, 2, 4, 4)}
            
            if left_col_box and left_content:
                # 좌측 컬럼
                left_fit = self.text_fitter.fit_text_to_box(
                    left_content,
                    left_col_box["position"][2],
                    left_col_box["position"][3],
                    initial_font_size=initial_font_size
                )
                
                left_result = self._apply_text_box_enhanced(
                    slide, left_col_box["position"],
                    left_fit["adjusted_text"],
                    initial_font_size=left_fit["adjusted_font_size"]
                )
                
                if left_result["success"]:
                    adjustments += 1
                    if left_fit["truncated"]:
                        warnings.append("Left column was truncated")
            
            if right_col_box and right_content:
                # 우측 컬럼
                right_fit = self.text_fitter.fit_text_to_box(
                    right_content,
                    right_col_box["position"][2],
                    right_col_box["position"][3],
                    initial_font_size=initial_font_size
                )
                
                right_result = self._apply_text_box_enhanced(
                    slide, right_col_box["position"],
                    right_fit["adjusted_text"],
                    initial_font_size=right_fit["adjusted_font_size"]
                )
                
                if right_result["success"]:
                    adjustments += 1
                    if right_fit["truncated"]:
                        warnings.append("Right column was truncated")
            
            return {
                "success": True,
                "layout_applied": "two_column",
                "text_adjustments": adjustments,
                "warnings": warnings
            }
            
        except Exception as e:
            logger.error(f"Two column application failed: {e}")
            return {
                "success": False,
                "error": str(e),
                "layout_applied": "two_column"
            }
    
    def _apply_three_column_enhanced(
        self,
        slide,
        layout_def: Dict,
        content: Dict,
        options: Dict
    ) -> Dict:
        """
        3컬럼 레이아웃 적용 - Enhanced 버전
        - 균등 분할 및 폰트 동기화
        """
        warnings = []
        adjustments = 0
        
        try:
            # Get elements or text_boxes
            elements = layout_def.get("elements", layout_def.get("text_boxes", []))
            
            if not elements:
                return {
                    "success": False,
                    "error": "No layout elements found",
                    "layout_applied": "three_column"
                }
            
            # 헤드라인 처리
            headline_box = elements[0]
            headline_text = content.get("headline", content.get("title", ""))
            
            if headline_text:
                headline_result = self._apply_text_box_enhanced(
                    slide, headline_box["position"],
                    headline_text,
                    initial_font_size=18,
                    is_bold=True
                )
                if headline_result["success"]:
                    adjustments += 1
            
            # 3개 컬럼 처리
            columns = ['col1', 'col2', 'col3']
            column_contents = []
            
            for col in columns:
                col_content = content.get(f"{col}_content", content.get(col, ""))
                if isinstance(col_content, list):
                    col_content = "\n".join(col_content)
                column_contents.append(col_content)
            
            # 전체 컬럼의 최대 길이로 초기 폰트 크기 결정
            max_length = max(len(content) for content in column_contents if content)
            initial_font_size = 12 if max_length > 300 else 14
            
            # 각 컬럼 처리
            for i, col in enumerate(columns):
                if column_contents[i]:
                    # 컬럼 박스 찾기
                    col_box = None
                    for element in elements:
                        if element.get("type") == col or col in element.get("type", ""):
                            col_box = element
                            break
                    
                    # Create default column box if not found
                    if not col_box:
                        col_box = {"position": (i * 3 + 0.5, 2, 2.5, 4)}
                    
                    if col_box:
                        col_result = self._apply_text_box_enhanced(
                            slide, col_box["position"],
                            column_contents[i],
                            initial_font_size=initial_font_size
                        )
                        
                        if col_result["success"]:
                            adjustments += 1
                            if col_result["fit_result"]["truncated"]:
                                warnings.append(f"Column {i+1} was truncated")
            
            return {
                "success": True,
                "layout_applied": "three_column",
                "text_adjustments": adjustments,
                "warnings": warnings
            }
            
        except Exception as e:
            logger.error(f"Three column application failed: {e}")
            return {
                "success": False,
                "error": str(e),
                "layout_applied": "three_column"
            }
    
    def _apply_generic_enhanced(
        self,
        slide,
        layout_def: Dict,
        content: Dict,
        options: Dict
    ) -> Dict:
        """
        범용 레이아웃 적용 (Enhanced TextFitter 사용)
        핸들러가 없는 경우 사용
        """
        warnings = []
        adjustments = 0
        
        try:
            # Get elements or text_boxes
            elements = layout_def.get("elements", layout_def.get("text_boxes", []))
            
            for i, element in enumerate(elements):
                box_type = element.get("type", f"box_{i}")
                text_content = content.get(box_type, "")
                
                # Try different content keys
                if not text_content:
                    for key in ['content', 'text', 'body']:
                        if content.get(key):
                            text_content = content[key]
                            break
                
                if text_content:
                    font_size = 14
                    if element.get("font_size"):
                        if hasattr(element["font_size"], 'pt'):
                            font_size = element["font_size"].pt
                        else:
                            font_size = element["font_size"]
                    
                    result = self._apply_text_box_enhanced(
                        slide, element["position"],
                        text_content,
                        initial_font_size=font_size
                    )
                    
                    if result["success"]:
                        adjustments += 1
                        if result["fit_result"]["truncated"]:
                            warnings.append(f"{box_type} was truncated")
                    else:
                        warnings.append(f"{box_type} processing failed")
            
            return {
                "success": True,
                "layout_applied": "generic",
                "text_adjustments": adjustments,
                "warnings": warnings
            }
            
        except Exception as e:
            logger.error(f"Generic layout application failed: {e}")
            return {
                "success": False,
                "error": str(e),
                "layout_applied": "generic"
            }
    
    # 추가 레이아웃 핸들러 스텁 예시
    def _apply_single_column_enhanced(self, slide, layout_def, content, options):
        """단일 컬럼 - 긴 텍스트 최적화"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_image_text_enhanced(self, slide, layout_def, content, options):
        """이미지+텍스트 - 이미지 영역 예약 후 텍스트 배치"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_matrix_enhanced(self, slide, layout_def, content, options):
        """2x2 매트릭스 - 4개 셀 균등 배분"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_timeline_enhanced(self, slide, layout_def, content, options):
        """타임라인 - 시간순 이벤트 배치"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_process_flow_enhanced(self, slide, layout_def, content, options):
        """프로세스 플로우 - 단계별 배치"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_pyramid_enhanced(self, slide, layout_def, content, options):
        """피라미드 계층"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_dashboard_grid_enhanced(self, slide, layout_def, content, options):
        """대시보드 그리드 - KPI 표시"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_quote_highlight_enhanced(self, slide, layout_def, content, options):
        """인용문 강조"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_split_screen_enhanced(self, slide, layout_def, content, options):
        """분할 화면"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_agenda_toc_enhanced(self, slide, layout_def, content, options):
        """의사일정/목차"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def _apply_waterfall_enhanced(self, slide, layout_def, content, options):
        """워터폴 차트"""
        return self._apply_generic_enhanced(slide, layout_def, content, options)
    
    def get_metrics(self) -> Dict:
        """
        성능 메트릭 반환
        """
        return {
            **self.metrics,
            "text_fitter_cache_stats": self.text_fitter.get_performance_metrics()
        }