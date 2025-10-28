"""
SlideFixer Fallback (UTF-8)

Minimal auto-fix implementation used when the full SlideFixer is unavailable.
Performs safe adjustments: margins, simple text truncation, spacing tweaks.
"""

from typing import Optional, Dict, Any
from pptx.slide import Slide
from pptx.util import Inches, Pt


class SlideFixer:
    def __init__(self) -> None:
        # Minimal init
        self.safe_margin = Inches(0.3)

    def fix_slide(
        self,
        slide: Slide,
        validation_result: Optional[Any] = None,
        aggressive_mode: bool = False,
        slide_number: Optional[int] = None,
    ) -> Dict[str, Any]:
        fixes_applied: list[str] = []
        fixes_failed: list[str] = []

        # Interpret validation_result into simple issue categories
        issues = []
        try:
            if validation_result is None:
                issues = []
            elif hasattr(validation_result, "issues"):
                # ValidationResult-like
                issues = [getattr(i, "category", None) for i in getattr(validation_result, "issues", [])]
            elif isinstance(validation_result, dict):
                issues = validation_result.get("issues", []) or []
        except Exception:
            issues = []

        # Always run safe margin/text passes first
        try:
            self._fix_margins(slide)
            fixes_applied.append("margins_adjusted")
        except Exception as e:
            fixes_failed.append(f"margins_adjust_failed: {e}")

        try:
            self._reduce_text_density(slide)
            fixes_applied.append("text_density_reduced")
        except Exception as e:
            fixes_failed.append(f"text_density_failed: {e}")

        # Targeted fixes by category (best-effort)
        try:
            self._fix_text_overflow(slide)
            fixes_applied.append("text_overflow_pass")
        except Exception as e:
            fixes_failed.append(f"text_overflow_failed: {e}")

        try:
            self._fix_out_of_bounds(slide)
            fixes_applied.append("out_of_bounds_pass")
        except Exception as e:
            fixes_failed.append(f"out_of_bounds_failed: {e}")

        try:
            self._fix_shape_overlap(slide)
            fixes_applied.append("shape_overlap_pass")
        except Exception as e:
            fixes_failed.append(f"shape_overlap_failed: {e}")

        try:
            self._fix_font_consistency(slide)
            fixes_applied.append("font_consistency_pass")
        except Exception as e:
            fixes_failed.append(f"font_consistency_failed: {e}")

        try:
            self._apply_mckinsey_style(slide)
            fixes_applied.append("mckinsey_style_pass")
        except Exception as e:
            fixes_failed.append(f"mckinsey_style_failed: {e}")

        return {
            "slide_number": slide_number or 0,
            "fixes_applied": fixes_applied,
            "fixes_failed": fixes_failed,
        }

    def _fix_margins(self, slide: Slide) -> None:
        margin = int(self.safe_margin)
        # assume 16:9 default; use slide width/height from shapes' parent
        sw = getattr(slide.part.slide_layout.part.presentation.slide_width, "__int__", lambda: Inches(13.33))()
        sh = getattr(slide.part.slide_layout.part.presentation.slide_height, "__int__", lambda: Inches(7.5))()
        try:
            sw = int(slide.part.slide_layout.part.presentation.slide_width)
            sh = int(slide.part.slide_layout.part.presentation.slide_height)
        except Exception:
            pass
        for shape in slide.shapes:
            try:
                if shape.left < margin:
                    shape.left = margin
                if shape.top < margin:
                    shape.top = margin
                if shape.left + shape.width > sw - margin:
                    new_w = sw - margin - shape.left
                    if new_w > 0:
                        shape.width = int(new_w)
                    else:
                        shape.left = int(sw - margin - shape.width)
                if shape.top + shape.height > sh - margin:
                    new_h = sh - margin - shape.top
                    if new_h > 0:
                        shape.height = int(new_h)
                    else:
                        shape.top = int(sh - margin - shape.height)
            except Exception:
                continue

    def _reduce_text_density(self, slide: Slide) -> None:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if not tf or not tf.text:
                continue
            # truncate overly long text
            if len(tf.text) > 500:
                txt = tf.text[:500]
                tf.clear()
                tf.text = txt + "…"
            # limit paragraphs/bullets
            count = 0
            for p in list(tf.paragraphs):
                if p.text.strip():
                    count += 1
                if count > 7:
                    p.text = ""
            # spacing
            for p in tf.paragraphs:
                try:
                    p.space_after = Pt(8)
                    p.line_spacing = 1.2
                except Exception:
                    pass

    def _fix_text_overflow(self, slide: Slide) -> None:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if not tf or not tf.text:
                continue
            # Reduce font size if too long
            if len(tf.text) > 300:
                for p in tf.paragraphs:
                    try:
                        # best-effort smaller font
                        for r in p.runs:
                            if hasattr(r.font, "size") and r.font.size and r.font.size.pt > 12:
                                r.font.size = Pt(12)
                    except Exception:
                        pass

    def _fix_out_of_bounds(self, slide: Slide) -> None:
        # ensure shapes are within slide area with safe margin
        margin = int(self.safe_margin)
        try:
            sw = int(slide.part.slide_layout.part.presentation.slide_width)
            sh = int(slide.part.slide_layout.part.presentation.slide_height)
        except Exception:
            sw = int(Inches(13.33))
            sh = int(Inches(7.5))
        for shape in slide.shapes:
            try:
                if shape.left < margin:
                    shape.left = margin
                if shape.top < margin:
                    shape.top = margin
                if shape.left + shape.width > sw - margin:
                    shape.left = int(sw - margin - shape.width)
                if shape.top + shape.height > sh - margin:
                    shape.top = int(sh - margin - shape.height)
            except Exception:
                continue

    def _fix_shape_overlap(self, slide: Slide) -> None:
        # naive pass: nudge overlapping shapes downwards
        shapes = [s for s in slide.shapes]
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                s1, s2 = shapes[i], shapes[j]
                try:
                    if (s1.left < s2.left + s2.width and s1.left + s1.width > s2.left and
                        s1.top < s2.top + s2.height and s1.top + s1.height > s2.top):
                        # overlap detected → nudge s2
                        s2.top = int(s1.top + s1.height + int(self.safe_margin))
                except Exception:
                    continue

    def _fix_font_consistency(self, slide: Slide) -> None:
        # Set Arial 12pt for body paragraphs
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                try:
                    for r in p.runs:
                        if hasattr(r.font, "name") and (not r.font.name or r.font.name.strip() == ""):
                            r.font.name = "Arial"
                        if hasattr(r.font, "size") and (not r.font.size or (hasattr(r.font.size, 'pt') and r.font.size.pt > 14)):
                            r.font.size = Pt(12)
                except Exception:
                    continue

    def _apply_mckinsey_style(self, slide: Slide) -> None:
        # Title-like paragraphs bold/blue, body gray
        from pptx.dml.color import RGBColor
        BLUE = RGBColor(0, 118, 168)
        GRAY = RGBColor(83, 86, 90)
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for idx, p in enumerate(tf.paragraphs):
                try:
                    for r in p.runs:
                        r.font.name = r.font.name or "Arial"
                        if idx == 0:
                            r.font.bold = True
                            r.font.color.rgb = BLUE
                        else:
                            r.font.color.rgb = GRAY
                except Exception:
                    continue
