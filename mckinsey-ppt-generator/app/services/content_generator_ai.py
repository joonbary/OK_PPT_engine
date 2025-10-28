"""
ContentGeneratorAI (UTF-8, minimal dependencies)

Generates slides from a document using an AI JSON instruction pattern, with
fallback content when AI is unavailable. Designed to avoid heavy optional
imports so it works in constrained environments.
"""

import asyncio
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from app.services.simple_chart_generator import SimpleChartGenerator


class ContentGeneratorAI:
    def __init__(self, language: str = "ko") -> None:
        self.language = language
        try:
            import os
            if os.getenv("OPENAI_API_KEY"):
                from openai import AsyncOpenAI  # type: ignore
                self.llm = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            else:
                self.llm = None
        except Exception:
            self.llm = None
        try:
            from app.utils.text_formatter import EumsumStyleConverter
            self._formatter = EumsumStyleConverter()
        except Exception:
            self._formatter = None

    async def generate_from_document_with_ai(self, document: str, num_slides: int = 6, target_audience: str = "executive") -> Presentation:
        # Optional agent-based enrichment
        sections = []
        data_insights_seed = []
        strategies_seed = []
        try:
            from app.agents.strategist_agent import StrategistAgent
            from app.agents.data_analyst_agent import DataAnalystAgent
            strat = StrategistAgent()
            analyst = DataAnalystAgent()
            sres = await strat.process(input_data={'document': document, 'num_slides': num_slides}, context={'language': self.language})
            ares = await analyst.process(input_data={'document': document, 'outline': sres.get('outline', []), 'pyramid': sres.get('pyramid', {})}, context={'language': self.language})
            sections = [item.get('title','Section') for item in (sres.get('outline') or []) if isinstance(item, dict) and item.get('title')]
            data_insights_seed = ares.get('insights') or []
            strategies_seed = ares.get('recommendations') or []
        except Exception:
            pass
        # Prefer MECE-based slide outline when available
        try:
            slide_outline = sres.get('slide_outline') if isinstance(sres, dict) else None
        except Exception:
            slide_outline = None
        if slide_outline:
            slides: List[Dict] = []
            # Localize title for display if needed
            def _localize(title: str) -> str:
                if (self.language or 'ko').lower().startswith('ko'):
                    m = {
                        'Title': '표지',
                        'Executive Summary': '핵심 요약',
                        'Market Analysis': '시장 분석',
                        'Strategic Options': '전략 옵션',
                        'Recommendations': '권고안',
                        'Internal Analysis': '내부 분석',
                        'Impact Analysis': '영향 분석',
                        'Challenges': '과제',
                    }
                    return m.get(title, title)
                return title

            slides.append({
                'slide_number': 0,
                'title': _localize(slide_outline[0].get('type', 'Title')),
                'type': 'title',
                'enriched_content': {'source_content': document[:600], 'key_messages': []},
                'generation_instruction': self._make_instruction(
                    title=slide_outline[0].get('content_focus', 'Presentation'),
                    source=document,
                    data_insights=[],
                    strategies=[],
                    industry_context='',
                )
            })
            for item in slide_outline[1:]:
                # type은 템플릿 선택에 사용(영문 표준 유지), title은 표시용(언어 현지화)
                raw_type = (item.get('type') or item.get('content_focus') or 'content')
                title = _localize(raw_type)
                focus = item.get('mece_segment')
                extra = f"다음 콘텐츠에만 집중:\n{focus}\n" if (focus and focus != 'ALL') else ''
                slides.append({
                    'slide_number': item.get('slide_number', len(slides)),
                    'title': title,
                    'type': raw_type,  # 템플릿 오케스트레이터가 활용
                    'enriched_content': {'source_content': (focus or document)[:600], 'key_messages': []},
                    'generation_instruction': self._make_instruction(
                        title=title,
                        source=(focus or document),
                        data_insights=[],
                        strategies=[],
                        industry_context=extra,
                    )
                })
            return await self.generate_slides_with_ai(slides)
        if not sections:
            sections = [
                "Executive Summary",
                "Market Analysis",
                "Strategic Options",
                "Recommendations",
            ]
        # 요청 슬라이드 수를 존중: 제목(1) + 본문(num_slides-1)
        desired = max(1, num_slides - 1)
        # 섹션명을 순환하여 원하는 개수 채우기
        seq: List[str] = []
        while len(seq) < desired:
            for s in sections:
                if len(seq) < desired:
                    seq.append(s)
                else:
                    break

        slides: List[Dict] = []
        slides.append({
            'slide_number': 0,
            'title': "Strategic Analysis & Recommendations",
            'type': 'title',
            'enriched_content': {
                'data_insights': [],
                'visualizations': [],
                'strategic_recommendations': ["Define priorities", "Phase execution"],
                'industry_context': '',
                'source_content': document[:600],
                'key_messages': []
            },
            'generation_instruction': self._make_instruction(
                title="Strategic Analysis & Recommendations",
                source=document,
                data_insights=[],
                strategies=["Define priorities", "Phase execution"],
                industry_context='',
            )
        })

        for i, title in enumerate(seq):
            title_keys = set([w.lower() for w in title.split() if w])
            # normalize insights
            norm_ins = []
            for it in (data_insights_seed if isinstance(data_insights_seed, list) else []):
                try:
                    s = it.get('insight') if isinstance(it, dict) else str(it)
                except Exception:
                    s = str(it)
                norm_ins.append(s)
            picked_ins = [s for s in norm_ins if any(k in s.lower() for k in title_keys)] or norm_ins[:2]
            # strategies
            strat_list = strategies_seed if isinstance(strategies_seed, list) else [str(strategies_seed)]
            picked_strat = [s for s in strat_list if any(k in str(s).lower() for k in title_keys)] or strat_list[:2]

            slides.append({
                'slide_number': i + 1,
                'title': title,
                'type': 'content',
                'enriched_content': {
                    'data_insights': picked_ins,
                    'visualizations': [],
                    'strategic_recommendations': picked_strat,
                    'industry_context': '',
                    'source_content': document[:600],
                    'key_messages': []
                },
                'generation_instruction': self._make_instruction(
                    title=title,
                    source=document,
                    data_insights=picked_ins,
                    strategies=picked_strat,
                    industry_context='',
                )
            })

        return await self.generate_slides_with_ai(slides)

    async def generate_slides_with_ai(self, integrated_insights: List[Dict]) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for idx, slide_data in enumerate(integrated_insights):
            content = await self._generate_content(slide_data)
            if idx == 0:
                self._add_title_slide(prs, content)
            else:
                self._add_content_slide(prs, slide_data, content)
        return prs

    async def _generate_content(self, slide_data: Dict) -> Dict:
        instruction = slide_data.get('generation_instruction', '')
        if not self.llm:
            return self._post_process(self._fallback(slide_data))
        try:
            resp = await asyncio.wait_for(
                self.llm.chat.completions.create(
                    model="gpt-4-turbo-preview",
                    messages=[
                        {"role": "system", "content": self._system_instruction_json_only()},
                        {"role": "user", "content": instruction},
                    ],
                    temperature=0.7,
                    max_tokens=900,
                    response_format={"type": "json_object"},
                ),
                timeout=60,
            )
            import json
            data = json.loads((resp.choices[0].message.content or '').strip())
            return self._post_process(data)
        except Exception:
            return self._post_process(self._fallback(slide_data))

    def _fallback(self, slide_data: Dict) -> Dict:
        enriched = slide_data.get('enriched_content', {}) or {}
        title = slide_data.get('title', 'Slide')
        def first_sentence(s: str) -> str:
            s = (s or '').strip()
            return (s.split('.')[0] + '.') if s else ''
        points: List[str] = []
        for it in (enriched.get('data_insights') or [])[:2]:
            t = it.get('insight') if isinstance(it, dict) else str(it)
            fs = first_sentence(str(t))
            if fs:
                points.append(fs[:150])
        for it in (enriched.get('strategic_recommendations') or [])[:2]:
            fs = first_sentence(str(it))
            if fs:
                points.append(fs[:150])
        # 지역화: 한국어일 때 기본 문구도 한국어로
        is_ko = (getattr(self, 'language', 'ko') or 'ko').lower().startswith('ko')
        headline_suffix = '핵심 요약' if is_ko else 'Key Findings'
        pending_msg = '분석 결과 준비 중' if is_ko else 'Analysis results pending'
        return {
            'headline': f"{title}: {headline_suffix}",
            'key_points': points or [pending_msg],
            'supporting_detail': (enriched.get('source_content') or '')[:300],
        }

    def _post_process(self, data: Dict) -> Dict:
        if (getattr(self, 'language', 'ko') or "").lower().startswith("ko") and getattr(self, '_formatter', None):
            try:
                if 'headline' in data:
                    data['headline'] = self._formatter.convert_headline(str(data.get('headline') or ''))
                if 'key_points' in data and isinstance(data.get('key_points'), list):
                    data['key_points'] = self._formatter.convert_bullet_points([str(x) for x in data['key_points']])
            except Exception:
                pass
        return data

    def _add_title_slide(self, prs: Presentation, content: Dict) -> None:
        blank = prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else 0]
        slide = prs.slides.add_slide(blank)
        tbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = tbox.text_frame
        head = content.get('headline', 'Presentation Title')
        try:
            head = str(head)
        except Exception:
            head = 'Presentation Title'
        tf.text = head
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Arial'
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 118, 168)
        sub = content.get('supporting_detail')
        if sub is not None:
            sbox = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.8))
            sf = sbox.text_frame
            try:
                sub_text = sub if isinstance(sub, str) else str(sub)
            except Exception:
                sub_text = ''
            sf.text = sub_text[:100]
            sf.word_wrap = True
            for p in sf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(83, 86, 90)

    def _add_content_slide(self, prs: Presentation, slide_data: Dict, content: Dict) -> None:
        blank = prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else 0]
        slide = prs.slides.add_slide(blank)
        # Build spec and normalize
        title_lower = (slide_data.get('title', '') or '').lower()
        spec = {
            'slide_number': slide_data.get('slide_number', 1),
            'title': slide_data.get('title', ''),
            'type': slide_data.get('type', ''),
            'headline': content.get('headline', ''),
            'content': (slide_data.get('enriched_content') or {}).get('key_messages') or content.get('key_points') or [],
            'bullets': content.get('key_points') or [],
        }
        # Layout intent based on slide type/title
        bullets = list(spec.get('bullets') or [])
        def split_three_columns(items):
            n = len(items)
            if n <= 0:
                return []
            k = max(1, (n + 2) // 3)
            return [items[:k], items[k:2*k], items[2*k:]]
        if any(key in title_lower for key in ['market analysis','market','고객','시장']) and bullets:
            spec['columns'] = split_three_columns(bullets)
        if any(key in title_lower for key in ['strategic options','matrix','impact','매트릭스','전략']) and bullets:
            # Create a simple 2x2 matrix from bullets
            m = [["", ""], ["", ""]]
            for i, it in enumerate(bullets[:4]):
                m[i//2][i%2] = str(it)
            spec['matrix'] = m
        ec = slide_data.get('enriched_content') or {}
        for k in ('columns','matrix','left','right','pros','cons','before','after','milestones','kpis','chart_data','chart_title'):
            if k in ec:
                spec[k] = ec.get(k)
        if slide_data.get('layout_type'):
            spec['layout_type'] = slide_data.get('layout_type')
        if slide_data.get('content_type'):
            spec['content_type'] = slide_data.get('content_type')
        if isinstance(ec.get('columns'), list):
            spec['columns'] = ec.get('columns')
        if isinstance(ec.get('matrix'), list):
            spec['matrix'] = ec.get('matrix')
        try:
            from app.services.layout_validator import LayoutValidator
            spec = LayoutValidator().normalize(spec)
        except Exception:
            pass
        try:
            from app.services.template_orchestrator import TemplateOrchestrator
            TemplateOrchestrator().select_and_apply(slide, spec)
            return
        except Exception:
            pass
        # Fallback simple rendering
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        tf = tb.text_frame; tf.text = slide_data.get('title', ''); tf.word_wrap = True
        hb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.6))
        hf = hb.text_frame; hf.text = content.get('headline', ''); hf.word_wrap = True
        pb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4))
        ptf = pb.text_frame; ptf.word_wrap = True
        for i, point in enumerate(content.get('key_points') or []):
            p = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph(); p.text = f"• {point}"

    def _make_instruction(self, title: str, source: str, data_insights: List, strategies: List, industry_context: str) -> str:
        lang = (getattr(self, 'language', 'ko') or 'ko').lower()
        lang_msg = {
            'ko': '모든 출력을 한국어로 작성하세요. 제목과 불릿포인트, 본문 모두 한국어로 작성.',
            'en': 'Write all output in English. Title, bullets, and body in English.',
            'ja': 'すべて日本語で出力してください。タイトルや箇条書き、本文は日本語で。',
        }.get(lang, 'Write all output in the specified language.')
        fi = "\n".join([f"- {str(x)[:200]}" for x in data_insights]) or "(none)"
        fs = "\n".join([f"- {str(x)[:200]}" for x in strategies]) or "(none)"
        return (
            f"Slide: {title}\n\n"
            f"{lang_msg}\n"
            "Create McKinsey-style slide content. Respond with JSON only.\n"
            "Fields: headline (1-2 sentences), key_points (3-5), supporting_detail.\n\n"
            f"Source:\n{source[:800]}\n\n"
            f"Data insights:\n{fi}\n\n"
            f"Strategy recommendations:\n{fs}\n\n"
            f"Industry context:\n{(industry_context or '')[:400]}\n"
        )

    def _system_instruction_json_only(self) -> str:
        lang = (getattr(self, 'language', 'ko') or 'ko').lower()
        if lang.startswith('ko'):
            return 'JSON만 반환하세요. 모든 텍스트는 한국어로 작성.'
        if lang.startswith('ja'):
            return 'JSONのみ返してください。すべて日本語で記述。'
        return 'Respond with pure JSON only. Use the requested language.'



