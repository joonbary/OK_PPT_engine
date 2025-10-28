"""
McKinsey 표준 레이아웃 시스템
Task 4 - McKinsey 표준 레이아웃 구현

5가지 표준 레이아웃:
1. Title Slide - 제목 슬라이드
2. Executive Summary - 요약 슬라이드  
3. Content with Chart - 차트 포함 콘텐츠
4. Bullet Points - 불릿 포인트
5. Comparison - 비교 분석
"""

import logging
from typing import Dict, List, Optional, Any, Tuple
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class LayoutManager:
    """
    McKinsey 표준 레이아웃 관리자
    
    모든 슬라이드를 5가지 표준 레이아웃 중 하나로 구성
    일관성 있는 위치, 크기, 정렬 보장
    """
    
    def __init__(self):
        """LayoutManager 초기화"""
        self.logger = logging.getLogger(__name__)
        
        # 슬라이드 크기 (16:9)
        self.SLIDE_WIDTH = Inches(13.33)
        self.SLIDE_HEIGHT = Inches(7.5)
        
        # McKinsey 표준 여백
        self.MARGIN_TOP = Inches(1.0)      # 상단 여백
        self.MARGIN_BOTTOM = Inches(0.5)   # 하단 여백
        self.MARGIN_LEFT = Inches(0.75)    # 좌측 여백
        self.MARGIN_RIGHT = Inches(0.75)   # 우측 여백
        
        # 콘텐츠 영역
        self.CONTENT_WIDTH = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
        self.CONTENT_HEIGHT = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM
        
        # 표준 위치
        self.TITLE_TOP = Inches(0.5)       # 제목 상단 위치
        self.TITLE_HEIGHT = Inches(0.75)   # 제목 높이
        self.CONTENT_TOP = Inches(1.5)     # 콘텐츠 시작 위치
        
        # McKinsey 색상
        self.MCKINSEY_BLUE = RGBColor(0, 118, 168)
        self.MCKINSEY_GRAY = RGBColor(83, 86, 90)
        
        self.logger.info("LayoutManager initialized with McKinsey standards")
    
    def apply_layout(self, slide: Slide, layout_type: str, content: Dict[str, Any]) -> None:
        """
        슬라이드에 표준 레이아웃 적용
        
        Args:
            slide: 레이아웃을 적용할 슬라이드
            layout_type: 레이아웃 타입 (title, summary, chart, bullets, comparison)
            content: 슬라이드 콘텐츠
        """
        self.logger.info(f"Applying {layout_type} layout to slide")
        
        try:
            # 기존 콘텐츠 정리
            self._clear_slide(slide)
            
            # 레이아웃별 적용
            if layout_type == "title":
                self._apply_title_layout(slide, content)
            elif layout_type == "summary":
                self._apply_summary_layout(slide, content)
            elif layout_type == "chart":
                self._apply_chart_layout(slide, content)
            elif layout_type == "bullets":
                self._apply_bullets_layout(slide, content)
            elif layout_type == "comparison":
                self._apply_comparison_layout(slide, content)
            else:
                # 기본 레이아웃 (bullets)
                self._apply_bullets_layout(slide, content)
            
            self.logger.info(f"Successfully applied {layout_type} layout")
            
        except Exception as e:
            self.logger.error(f"Failed to apply layout: {e}")
            raise
    
    def _clear_slide(self, slide: Slide) -> None:
        """슬라이드 초기화 (빈 텍스트 박스 제거)"""
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                pass
    
    def _apply_title_layout(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        레이아웃 1: Title Slide
        - 중앙 정렬된 제목과 부제목
        - McKinsey Blue 강조
        """
        # 제목
        title_left = self.MARGIN_LEFT
        title_top = Inches(2.5)  # 중앙 상단
        title_width = self.CONTENT_WIDTH
        title_height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "Title")
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 제목 스타일
        for run in title_frame.paragraphs[0].runs:
            run.font.name = "Arial"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.MCKINSEY_BLUE
        
        # 부제목
        if content.get("subtitle"):
            subtitle_top = title_top + title_height + Inches(0.3)
            subtitle_box = slide.shapes.add_textbox(
                title_left, subtitle_top, title_width, Inches(1.0)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content["subtitle"]
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            for run in subtitle_frame.paragraphs[0].runs:
                run.font.name = "Arial"
                run.font.size = Pt(24)
                run.font.color.rgb = self.MCKINSEY_GRAY
        
        # 날짜/저자 (하단)
        if content.get("date") or content.get("author"):
            footer_text = []
            if content.get("date"):
                footer_text.append(content["date"])
            if content.get("author"):
                footer_text.append(content["author"])
            
            footer_top = self.SLIDE_HEIGHT - Inches(1.0)
            footer_box = slide.shapes.add_textbox(
                title_left, footer_top, title_width, Inches(0.5)
            )
            footer_frame = footer_box.text_frame
            footer_frame.text = " | ".join(footer_text)
            footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            for run in footer_frame.paragraphs[0].runs:
                run.font.name = "Arial"
                run.font.size = Pt(12)
                run.font.color.rgb = self.MCKINSEY_GRAY
    
    def _apply_summary_layout(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        레이아웃 2: Executive Summary
        - 상단 제목
        - 3-4개 핵심 포인트 박스
        """
        # 제목
        self._add_slide_title(slide, content.get("title", "Executive Summary"))
        
        # 핵심 포인트 박스들
        points = content.get("points", [])[:4]  # 최대 4개
        
        if len(points) <= 2:
            # 2개 이하: 가로 배치
            box_width = (self.CONTENT_WIDTH - Inches(0.5)) / len(points)
            box_height = Inches(3.5)
            box_top = self.CONTENT_TOP + Inches(0.5)
            
            for i, point in enumerate(points):
                box_left = self.MARGIN_LEFT + (box_width + Inches(0.5)) * i
                self._add_summary_box(slide, box_left, box_top, box_width, box_height, point)
        else:
            # 3-4개: 2x2 그리드
            box_width = (self.CONTENT_WIDTH - Inches(0.5)) / 2
            box_height = (self.CONTENT_HEIGHT - Inches(1.0)) / 2
            
            for i, point in enumerate(points):
                row = i // 2
                col = i % 2
                box_left = self.MARGIN_LEFT + (box_width + Inches(0.5)) * col
                box_top = self.CONTENT_TOP + (box_height + Inches(0.5)) * row
                self._add_summary_box(slide, box_left, box_top, box_width, box_height, point)
    
    def _add_summary_box(self, slide: Slide, left: float, top: float, 
                        width: float, height: float, content: Dict) -> None:
        """요약 박스 추가"""
        # 박스 배경
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
        box.line.color.rgb = self.MCKINSEY_BLUE
        box.line.width = Pt(1.5)
        
        # 텍스트 추가
        text_frame = box.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        
        # 제목
        p = text_frame.add_paragraph()
        p.text = content.get("title", "Key Point")
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.MCKINSEY_BLUE
        
        # 내용
        p = text_frame.add_paragraph()
        p.text = content.get("text", "")
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = self.MCKINSEY_GRAY
        p.space_after = Pt(6)
    
    def _apply_chart_layout(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        레이아웃 3: Content with Chart
        - 상단 제목
        - 좌측 텍스트 (40%)
        - 우측 차트 (60%)
        """
        # 제목
        self._add_slide_title(slide, content.get("title", "Analysis"))
        
        # 좌측 텍스트 영역 (40%)
        text_width = self.CONTENT_WIDTH * 0.4 - Inches(0.25)
        text_left = self.MARGIN_LEFT
        text_top = self.CONTENT_TOP
        text_height = self.CONTENT_HEIGHT - Inches(0.5)
        
        text_box = slide.shapes.add_textbox(
            text_left, text_top, text_width, text_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # 텍스트 포인트 추가
        points = content.get("points", [])
        for i, point in enumerate(points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {point}"
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.color.rgb = self.MCKINSEY_GRAY
            p.space_after = Pt(12)
            p.level = 0
        
        # 우측 차트 영역 (60%)
        chart_width = self.CONTENT_WIDTH * 0.6 - Inches(0.25)
        chart_left = text_left + text_width + Inches(0.5)
        chart_top = self.CONTENT_TOP
        chart_height = self.CONTENT_HEIGHT - Inches(0.5)
        
        # 차트 플레이스홀더 (실제 차트는 chart_generator가 처리)
        if content.get("chart_path"):
            try:
                slide.shapes.add_picture(
                    content["chart_path"],
                    chart_left, chart_top,
                    width=chart_width, height=chart_height
                )
            except Exception as e:
                self.logger.warning(f"Could not add chart image: {e}")
                # 차트 대체 박스
                chart_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, chart_left, chart_top, chart_width, chart_height
                )
                chart_box.fill.solid()
                chart_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
                chart_box.text_frame.text = "[Chart]"
                chart_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _apply_bullets_layout(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        레이아웃 4: Bullet Points
        - 상단 제목
        - 구조화된 불릿 포인트
        """
        # 제목
        self._add_slide_title(slide, content.get("title", "Key Points"))
        
        # 불릿 포인트 영역
        bullets_left = self.MARGIN_LEFT
        bullets_top = self.CONTENT_TOP
        bullets_width = self.CONTENT_WIDTH
        bullets_height = self.CONTENT_HEIGHT - Inches(0.5)
        
        bullets_box = slide.shapes.add_textbox(
            bullets_left, bullets_top, bullets_width, bullets_height
        )
        text_frame = bullets_box.text_frame
        text_frame.word_wrap = True
        
        # 불릿 포인트 추가
        points = content.get("points", [])
        for i, point in enumerate(points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            
            if isinstance(point, dict):
                # 메인 포인트
                p.text = f"• {point.get('main', '')}"
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self.MCKINSEY_BLUE
                p.space_after = Pt(6)
                p.level = 0
                
                # 서브 포인트
                for sub in point.get("sub", []):
                    sub_p = text_frame.add_paragraph()
                    sub_p.text = f"– {sub}"
                    sub_p.font.name = "Arial"
                    sub_p.font.size = Pt(14)
                    sub_p.font.color.rgb = self.MCKINSEY_GRAY
                    sub_p.space_after = Pt(4)
                    sub_p.level = 1
            else:
                # 단순 텍스트
                p.text = f"• {point}"
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = self.MCKINSEY_GRAY
                p.space_after = Pt(12)
                p.level = 0
    
    def _apply_comparison_layout(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        레이아웃 5: Comparison
        - 상단 제목
        - 2-3개 컬럼 비교
        """
        # 제목
        self._add_slide_title(slide, content.get("title", "Comparison"))
        
        # 비교 항목들
        items = content.get("items", [])[:3]  # 최대 3개
        num_items = len(items)
        
        if num_items == 0:
            return
        
        # 컬럼 계산
        col_width = (self.CONTENT_WIDTH - Inches(0.3) * (num_items - 1)) / num_items
        col_height = self.CONTENT_HEIGHT - Inches(0.5)
        col_top = self.CONTENT_TOP
        
        for i, item in enumerate(items):
            col_left = self.MARGIN_LEFT + (col_width + Inches(0.3)) * i
            
            # 컬럼 박스
            col_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, col_left, col_top, col_width, col_height
            )
            col_box.fill.solid()
            col_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
            col_box.line.color.rgb = self.MCKINSEY_BLUE
            col_box.line.width = Pt(1)
            
            # 텍스트
            text_frame = col_box.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 헤더
            p = text_frame.add_paragraph()
            p.text = item.get("header", f"Option {i+1}")
            p.font.name = "Arial"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.MCKINSEY_BLUE
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(12)
            
            # 포인트들
            for point in item.get("points", []):
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.name = "Arial"
                p.font.size = Pt(12)
                p.font.color.rgb = self.MCKINSEY_GRAY
                p.space_after = Pt(6)
    
    def _add_slide_title(self, slide: Slide, title: str) -> None:
        """슬라이드 제목 추가 (표준 위치)"""
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, self.TITLE_TOP, 
            self.CONTENT_WIDTH, self.TITLE_HEIGHT
        )
        text_frame = title_box.text_frame
        text_frame.text = title
        
        p = text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.MCKINSEY_BLUE
        p.alignment = PP_ALIGN.LEFT
    
    def detect_layout_type(self, content: Dict[str, Any]) -> str:
        """
        콘텐츠 기반 레이아웃 타입 자동 감지
        
        Args:
            content: 슬라이드 콘텐츠
            
        Returns:
            레이아웃 타입 문자열
        """
        # 제목 슬라이드 감지
        if content.get("is_title") or content.get("slide_number") == 1:
            return "title"
        
        # 요약 슬라이드 감지
        if "executive" in content.get("title", "").lower() or \
           "summary" in content.get("title", "").lower():
            return "summary"
        
        # 차트 포함 감지
        if content.get("chart") or content.get("chart_path"):
            return "chart"
        
        # 비교 분석 감지
        if "comparison" in content.get("title", "").lower() or \
           "vs" in content.get("title", "").lower() or \
           content.get("comparison_items"):
            return "comparison"
        
        # 기본: 불릿 포인트
        return "bullets"
    
    def validate_layout(self, slide: Slide) -> Dict[str, Any]:
        """
        레이아웃 검증
        
        Args:
            slide: 검증할 슬라이드
            
        Returns:
            검증 결과 딕셔너리
        """
        issues = []
        warnings = []
        
        for shape in slide.shapes:
            # 경계 검사
            if shape.left < 0 or shape.top < 0:
                issues.append(f"Shape outside slide boundaries")
            
            if shape.left + shape.width > self.SLIDE_WIDTH:
                warnings.append(f"Shape exceeds right boundary")
            
            if shape.top + shape.height > self.SLIDE_HEIGHT:
                warnings.append(f"Shape exceeds bottom boundary")
            
            # 텍스트 오버플로우 검사
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text and len(text_frame.text) > 500:
                    warnings.append(f"Text might be too long")
        
        return {
            "valid": len(issues) == 0,
            "issues": issues,
            "warnings": warnings,
            "shape_count": len(slide.shapes)
        }
    
    def get_layout_stats(self) -> Dict[str, Any]:
        """레이아웃 통계 반환"""
        return {
            "layouts_available": ["title", "summary", "chart", "bullets", "comparison"],
            "slide_dimensions": {
                "width": float(self.SLIDE_WIDTH / Inches(1)),
                "height": float(self.SLIDE_HEIGHT / Inches(1))
            },
            "margins": {
                "top": float(self.MARGIN_TOP / Inches(1)),
                "bottom": float(self.MARGIN_BOTTOM / Inches(1)),
                "left": float(self.MARGIN_LEFT / Inches(1)),
                "right": float(self.MARGIN_RIGHT / Inches(1))
            },
            "content_area": {
                "width": float(self.CONTENT_WIDTH / Inches(1)),
                "height": float(self.CONTENT_HEIGHT / Inches(1))
            }
        }