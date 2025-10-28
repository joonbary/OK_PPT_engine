"""
Enhanced Text Fitter for PowerPoint presentations.
Provides pixel-perfect text measurement and intelligent overflow prevention.

Features:
- PIL integration for accurate text measurement
- Multi-language support (Korean, English, Chinese)  
- Smart line breaking and truncation
- Performance optimization with caching
- Advanced text density analysis
- Multi-column support
"""

from pptx.util import Inches, Pt
from typing import Dict, Tuple, Optional, List, Union, Any
import re
import os
import platform
import time
from functools import lru_cache
from collections import defaultdict
import threading
from pathlib import Path

# PIL imports for accurate text measurement
from PIL import Image, ImageDraw, ImageFont

# Optional fonttools import for advanced font detection
try:
    import fonttools.ttLib as ttLib
    FONTTOOLS_AVAILABLE = True
except ImportError:
    ttLib = None
    FONTTOOLS_AVAILABLE = False

from app.core.logging import app_logger


class TextFitter:
    """
    Enhanced Text Fitter with pixel-perfect measurement and intelligent overflow prevention.
    
    Features:
    - Accurate PIL-based text measurement
    - Multi-language font support
    - Smart line breaking and truncation
    - Performance caching
    - Advanced text analysis
    """
    
    def __init__(self):
        """Initialize enhanced text fitter with advanced features"""
        # Basic settings
        self.min_font_size = Pt(10)
        self.max_font_size = Pt(18)
        self.default_font_size = Pt(14)
        self.line_spacing = 1.2
        
        # Font management
        self._font_cache = {}
        self._measurement_cache = {}
        self._cache_lock = threading.Lock()
        self._cache_hits = 0
        self._cache_misses = 0
        
        # Font configuration
        self.font_families = {
            'korean': ['맑은 고딕', 'Malgun Gothic', 'NanumGothic', 'Arial Unicode MS'],
            'english': ['Calibri', 'Arial', 'Helvetica', 'Segoe UI'],
            'chinese': ['SimSun', 'Microsoft YaHei', 'PingFang SC', 'Arial Unicode MS'],
            'fallback': ['Arial', 'DejaVu Sans', 'Liberation Sans']
        }
        
        # Performance settings
        self.enable_caching = True
        self.cache_max_size = 1000
        self.binary_search_precision = 0.5  # Font size precision for binary search
        
        # Text analysis settings
        self.korean_width_multiplier = 1.4  # Korean chars are wider
        self.chinese_width_multiplier = 1.3  # Chinese chars width
        self.punctuation_width_multiplier = 0.6  # Punctuation is narrower
        
        # Initialize font system
        self._initialize_fonts()
        
        app_logger.info("Enhanced TextFitter initialized with PIL integration")
        
    def _initialize_fonts(self):
        """Initialize font system and detect available fonts"""
        try:
            self.system_fonts = self._detect_system_fonts()
            self.available_fonts = self._validate_fonts()
            app_logger.info(f"Font system initialized with {len(self.available_fonts)} fonts")
        except Exception as e:
            app_logger.error(f"Font initialization failed: {e}")
            self.available_fonts = {}
            
    def _detect_system_fonts(self) -> Dict[str, List[str]]:
        """Detect available system fonts by platform"""
        system = platform.system()
        font_paths = []
        
        if system == "Windows":
            font_paths = [
                "C:/Windows/Fonts",
                "C:/Windows/System32/Fonts",
                os.path.expanduser("~/AppData/Local/Microsoft/Windows/Fonts")
            ]
        elif system == "Darwin":  # macOS
            font_paths = [
                "/System/Library/Fonts",
                "/Library/Fonts",
                os.path.expanduser("~/Library/Fonts")
            ]
        else:  # Linux
            font_paths = [
                "/usr/share/fonts",
                "/usr/local/share/fonts",
                os.path.expanduser("~/.fonts"),
                os.path.expanduser("~/.local/share/fonts")
            ]
        
        fonts = defaultdict(list)
        
        if FONTTOOLS_AVAILABLE and ttLib:
            # Use fonttools for advanced font detection
            for path in font_paths:
                if os.path.exists(path):
                    for file in Path(path).rglob("*.ttf"):
                        try:
                            font = ttLib.TTFont(str(file))
                            name = font['name'].getDebugName(1)  # Font family name
                            if name:
                                fonts[name.lower()].append(str(file))
                        except:
                            continue
        else:
            # Fallback to simple file-based detection
            app_logger.warning("fonttools not available, using simplified font detection")
            for path in font_paths:
                if os.path.exists(path):
                    for file in Path(path).rglob("*.ttf"):
                        # Extract font name from filename
                        font_name = file.stem.lower()
                        fonts[font_name].append(str(file))
                    
                    # Add common font aliases for Windows
                    if system == "Windows":
                        common_fonts = {
                            'calibri': ['calibri.ttf', 'calibrib.ttf', 'calibrii.ttf', 'calibriz.ttf'],
                            'arial': ['arial.ttf', 'arialbd.ttf', 'ariali.ttf', 'arialbi.ttf'],
                            'malgun gothic': ['malgun.ttf', 'malgunbd.ttf'],
                            '맑은 고딕': ['malgun.ttf', 'malgunbd.ttf']
                        }
                        
                        for font_name, files in common_fonts.items():
                            for font_file in files:
                                font_path = os.path.join(path, font_file)
                                if os.path.exists(font_path):
                                    fonts[font_name].append(font_path)
        
        return dict(fonts)
        
    def _validate_fonts(self) -> Dict[str, str]:
        """Validate and map preferred fonts to actual font files"""
        validated = {}
        
        for category, font_names in self.font_families.items():
            for font_name in font_names:
                font_key = font_name.lower()
                if font_key in self.system_fonts:
                    validated[category] = self.system_fonts[font_key][0]
                    break
            
            # Fallback to default if no preferred font found
            if category not in validated and 'arial' in self.system_fonts:
                validated[category] = self.system_fonts['arial'][0]
        
        return validated
        
    @lru_cache(maxsize=100)
    def _get_font_object(self, font_name: str, font_size: int) -> Optional[ImageFont.FreeTypeFont]:
        """Get cached PIL font object"""
        cache_key = f"{font_name}_{font_size}"
        
        with self._cache_lock:
            if cache_key in self._font_cache:
                self._cache_hits += 1
                return self._font_cache[cache_key]
            
            self._cache_misses += 1
            
        try:
            # Determine font category
            font_category = self._detect_text_language("")
            font_path = self.available_fonts.get(font_category) or self.available_fonts.get('fallback')
            
            if font_path:
                font = ImageFont.truetype(font_path, font_size)
            else:
                # Fallback to default font
                font = ImageFont.load_default()
                
            with self._cache_lock:
                if len(self._font_cache) >= self.cache_max_size:
                    # Remove oldest entries
                    oldest_keys = list(self._font_cache.keys())[:10]
                    for key in oldest_keys:
                        del self._font_cache[key]
                        
                self._font_cache[cache_key] = font
                
            return font
            
        except Exception as e:
            app_logger.error(f"Failed to load font {font_name}: {e}")
            return ImageFont.load_default()
            
    def _detect_text_language(self, text: str) -> str:
        """Detect primary language of text for font selection"""
        if not text:
            return 'english'
            
        korean_count = len(re.findall(r'[\u3131-\ucb4c]', text))
        chinese_count = len(re.findall(r'[\u4e00-\u9fff]', text))
        total_chars = len(text)
        
        if korean_count / total_chars > 0.3:
            return 'korean'
        elif chinese_count / total_chars > 0.3:
            return 'chinese'
        else:
            return 'english'
            
    def measure_text_precise(self, text: str, font_size: int, font_name: str = 'Calibri') -> Tuple[float, float]:
        """Measure text dimensions with pixel-perfect accuracy using PIL"""
        cache_key = f"{text}_{font_size}_{font_name}"
        
        if self.enable_caching:
            with self._cache_lock:
                if cache_key in self._measurement_cache:
                    self._cache_hits += 1
                    return self._measurement_cache[cache_key]
                    
                self._cache_misses += 1
        
        try:
            # Get font object
            font = self._get_font_object(font_name, font_size)
            if not font:
                return self._fallback_measurement(text, font_size)
            
            # Create temporary image for measurement
            img = Image.new('RGB', (1, 1), 'white')
            draw = ImageDraw.Draw(img)
            
            # Get text bounding box
            bbox = draw.textbbox((0, 0), text, font=font)
            width = bbox[2] - bbox[0]
            height = bbox[3] - bbox[1]
            
            # Convert pixels to inches (assuming 96 DPI)
            width_inches = width / 96.0
            height_inches = height / 96.0
            
            result = (width_inches, height_inches)
            
            # Cache result
            if self.enable_caching:
                with self._cache_lock:
                    if len(self._measurement_cache) >= self.cache_max_size:
                        # Remove oldest entries
                        oldest_keys = list(self._measurement_cache.keys())[:10]
                        for key in oldest_keys:
                            del self._measurement_cache[key]
                            
                    self._measurement_cache[cache_key] = result
            
            return result
            
        except Exception as e:
            app_logger.error(f"PIL text measurement failed: {e}")
            return self._fallback_measurement(text, font_size)
            
    def _fallback_measurement(self, text: str, font_size: int) -> Tuple[float, float]:
        """Fallback measurement using character-based calculation"""
        # Enhanced fallback with better language support
        korean_chars = len(re.findall(r'[\u3131-\ucb4c]', text))
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        punctuation_chars = len(re.findall(r'[.,;:!?\-()\[\]{}]', text))
        latin_chars = len(text) - korean_chars - chinese_chars - punctuation_chars
        
        # Calculate width with language-specific multipliers
        base_char_width = font_size * 0.5 / 72  # Convert to inches
        total_width = (
            korean_chars * base_char_width * self.korean_width_multiplier +
            chinese_chars * base_char_width * self.chinese_width_multiplier +
            punctuation_chars * base_char_width * self.punctuation_width_multiplier +
            latin_chars * base_char_width
        )
        
        # Calculate height
        line_count = text.count('\n') + 1
        line_height = font_size * self.line_spacing / 72  # Convert to inches
        total_height = line_count * line_height
        
        return (total_width, total_height)
        
    def fit_text_to_box(
        self, 
        text: str, 
        box_width: float, 
        box_height: float,
        initial_font_size: int = 14,
        font_name: str = 'Calibri',
        use_binary_search: bool = True,
        preserve_words: bool = True,
        max_lines: Optional[int] = None
    ) -> Dict:
        """
        Enhanced text fitting with PIL integration and intelligent algorithms.
        
        Args:
            text: Text to fit
            box_width: Width in inches
            box_height: Height in inches
            initial_font_size: Starting font size in points
            font_name: Font family name
            use_binary_search: Use binary search for optimal font size
            preserve_words: Avoid breaking words when wrapping
            max_lines: Maximum number of lines allowed
            
        Returns:
        {
            "adjusted_font_size": int,
            "adjusted_text": str,
            "fits": bool,
            "overflow_lines": int,
            "truncated": bool,
            "actual_width": float,
            "actual_height": float,
            "measurement_method": str,
            "cache_hit": bool,
            "processing_time": float
        }
        """
        start_time = time.time()
        
        try:
            # Prepare text for processing
            processed_text = self._prepare_text(text)
            
            # Determine optimal font size
            if use_binary_search:
                optimal_size = self._binary_search_font_size(
                    processed_text, box_width, box_height, font_name, initial_font_size
                )
            else:
                optimal_size = self._linear_search_font_size(
                    processed_text, box_width, box_height, font_name, initial_font_size
                )
            
            # Apply intelligent line breaking
            wrapped_text = self._intelligent_wrap(
                processed_text, box_width, optimal_size, font_name, preserve_words
            )
            
            # Check if wrapped text fits
            actual_width, actual_height = self.measure_text_precise(
                wrapped_text, optimal_size, font_name
            )
            
            fits = actual_height <= box_height
            truncated = False
            overflow_lines = 0
            
            # Handle overflow with smart truncation
            if not fits:
                if max_lines:
                    # Respect max lines constraint
                    wrapped_text = self._truncate_to_lines(wrapped_text, max_lines)
                else:
                    # Use intelligent truncation
                    wrapped_text = self._smart_truncate(
                        wrapped_text, box_width, box_height, optimal_size, font_name
                    )
                
                truncated = True
                
                # Recalculate dimensions
                actual_width, actual_height = self.measure_text_precise(
                    wrapped_text, optimal_size, font_name
                )
                
                # Calculate overflow
                if actual_height > box_height:
                    line_height = optimal_size * self.line_spacing / 72
                    overflow_lines = int((actual_height - box_height) / line_height)
                    fits = False
                else:
                    fits = True
                    overflow_lines = 0
            
            processing_time = time.time() - start_time
            
            return {
                "adjusted_font_size": optimal_size,
                "adjusted_text": wrapped_text,
                "fits": fits,
                "overflow_lines": overflow_lines,
                "truncated": truncated,
                "actual_width": actual_width,
                "actual_height": actual_height,
                "measurement_method": "PIL" if self.available_fonts else "fallback",
                "cache_hit": self._cache_hits > self._cache_misses,
                "processing_time": processing_time
            }
            
        except Exception as e:
            app_logger.error(f"Error in enhanced text fitting: {str(e)}")
            processing_time = time.time() - start_time
            
            # Fallback to basic fitting
            return {
                "adjusted_font_size": initial_font_size,
                "adjusted_text": text,
                "fits": False,
                "overflow_lines": 0,
                "truncated": False,
                "actual_width": 0,
                "actual_height": 0,
                "measurement_method": "error_fallback",
                "cache_hit": False,
                "processing_time": processing_time
            }
    
    def _prepare_text(self, text: str) -> str:
        """Prepare text for processing (normalize whitespace, etc.)"""
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text.strip())
        # Preserve intentional line breaks
        text = text.replace('\n', '\n')
        return text
        
    def _binary_search_font_size(
        self, text: str, box_width: float, box_height: float, 
        font_name: str, initial_size: int
    ) -> int:
        """Use binary search to find optimal font size"""
        min_size = self.min_font_size.pt
        max_size = min(initial_size, self.max_font_size.pt)
        
        best_size = min_size
        
        while max_size - min_size > self.binary_search_precision:
            mid_size = (min_size + max_size) / 2
            
            # Test if text fits at this size
            test_text = self._intelligent_wrap(text, box_width, mid_size, font_name, True)
            _, height = self.measure_text_precise(test_text, int(mid_size), font_name)
            
            if height <= box_height:
                best_size = int(mid_size)
                min_size = mid_size
            else:
                max_size = mid_size
                
        return best_size
        
    def _linear_search_font_size(
        self, text: str, box_width: float, box_height: float,
        font_name: str, initial_size: int
    ) -> int:
        """Linear search for font size (fallback method)"""
        current_size = min(initial_size, self.max_font_size.pt)
        
        while current_size >= self.min_font_size.pt:
            test_text = self._intelligent_wrap(text, box_width, current_size, font_name, True)
            _, height = self.measure_text_precise(test_text, current_size, font_name)
            
            if height <= box_height:
                return current_size
                
            current_size -= 1
            
        return self.min_font_size.pt
        
    def _intelligent_wrap(
        self, text: str, box_width: float, font_size: int, 
        font_name: str, preserve_words: bool = True
    ) -> str:
        """Intelligent line breaking with language awareness"""
        if not text:
            return text
            
        # Calculate characters per line using precise measurement
        sample_char = 'A'  # Use average character for width estimation
        char_width, _ = self.measure_text_precise(sample_char, font_size, font_name)
        max_chars_per_line = int(box_width / char_width) if char_width > 0 else 50
        
        # Handle existing line breaks
        paragraphs = text.split('\n')
        wrapped_paragraphs = []
        
        for paragraph in paragraphs:
            if len(paragraph) <= max_chars_per_line:
                wrapped_paragraphs.append(paragraph)
                continue
                
            if preserve_words:
                wrapped = self._wrap_preserving_words(paragraph, max_chars_per_line)
            else:
                wrapped = self._wrap_hard_break(paragraph, max_chars_per_line)
                
            wrapped_paragraphs.append(wrapped)
            
        return '\n'.join(wrapped_paragraphs)
        
    def _wrap_preserving_words(self, text: str, max_chars: int) -> str:
        """Wrap text preserving word boundaries with Korean particle awareness"""
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            word_length = self._calculate_display_length(word)
            space_length = 1 if current_line else 0
            
            # Check if adding this word would exceed the limit
            if current_length + word_length + space_length > max_chars:
                if current_line:
                    lines.append(' '.join(current_line))
                    current_line = [word]
                    current_length = word_length
                else:
                    # Word is too long, must break it intelligently
                    lines.append(self._break_long_word(word, max_chars))
                    current_line = []
                    current_length = 0
            else:
                current_line.append(word)
                current_length += word_length + space_length
        
        # Add remaining words
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)
        
    def _wrap_hard_break(self, text: str, max_chars: int) -> str:
        """Hard wrap at character limit"""
        lines = []
        current_pos = 0
        
        while current_pos < len(text):
            end_pos = min(current_pos + max_chars, len(text))
            lines.append(text[current_pos:end_pos])
            current_pos = end_pos
            
        return '\n'.join(lines)
        
    def _calculate_display_length(self, text: str) -> int:
        """Calculate display length considering multi-byte characters"""
        korean_chars = len(re.findall(r'[\u3131-\ucb4c]', text))
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        other_chars = len(text) - korean_chars - chinese_chars
        
        # Weight different character types
        return int(
            korean_chars * self.korean_width_multiplier +
            chinese_chars * self.chinese_width_multiplier +
            other_chars
        )
        
    def _break_long_word(self, word: str, max_chars: int) -> str:
        """Intelligently break long words, considering language patterns"""
        if len(word) <= max_chars:
            return word
            
        # For Korean text, try to break at syllable boundaries
        if re.search(r'[\u3131-\ucb4c]', word):
            return self._break_korean_word(word, max_chars)
        else:
            # For Latin text, break at character limit
            return word[:max_chars]
            
    def _break_korean_word(self, word: str, max_chars: int) -> str:
        """Break Korean words at syllable boundaries"""
        # Find good break points (after complete syllables)
        break_point = max_chars
        for i in range(max_chars-1, 0, -1):
            if not re.match(r'[\u3131-\u314F]', word[i]):  # Not a jamo
                break_point = i
                break
                
        return word[:break_point]
        
    def _smart_truncate(
        self, text: str, box_width: float, box_height: float,
        font_size: int, font_name: str
    ) -> str:
        """Smart truncation with sentence and phrase awareness"""
        lines = text.split('\n')
        line_height = font_size * self.line_spacing / 72
        max_lines = int(box_height / line_height)
        
        if len(lines) <= max_lines:
            return text
            
        # Truncate to fit available lines
        truncated_lines = lines[:max_lines-1]  # Leave room for ellipsis
        
        if truncated_lines:
            # Add smart ellipsis to last line
            last_line = truncated_lines[-1]
            if len(last_line) > 3:
                # Try to end at sentence boundary
                sentence_end = max(
                    last_line.rfind('.'),
                    last_line.rfind('!'),
                    last_line.rfind('?'),
                    last_line.rfind('\u3002')  # Korean period
                )
                
                if sentence_end > len(last_line) * 0.7:
                    truncated_lines[-1] = last_line[:sentence_end+1]
                else:
                    # End at word boundary
                    space_pos = last_line.rfind(' ')
                    if space_pos > len(last_line) * 0.7:
                        truncated_lines[-1] = last_line[:space_pos] + "..."
                    else:
                        truncated_lines[-1] = last_line[:-3] + "..."
        
        return '\n'.join(truncated_lines)
        
    def _truncate_to_lines(self, text: str, max_lines: int) -> str:
        """Truncate text to specified number of lines"""
        lines = text.split('\n')
        if len(lines) <= max_lines:
            return text
            
        truncated = lines[:max_lines-1]
        if len(lines) > max_lines:
            last_line = lines[max_lines-1]
            if len(last_line) > 3:
                truncated.append(last_line[:-3] + "...")
            else:
                truncated.append("...")
                
        return '\n'.join(truncated)
    
    def calculate_text_dimensions(
        self, 
        text: str, 
        font_size: int, 
        box_width: float,
        font_name: str = 'Calibri'
    ) -> Tuple[float, float]:
        """
        Calculate text dimensions with enhanced accuracy.
        
        Args:
            text: Text to measure
            font_size: Font size in points
            box_width: Box width in inches
            font_name: Font family name
            
        Returns: (required_width, required_height) in inches
        """
        try:
            # Use enhanced measurement if available
            if self.available_fonts:
                # Apply intelligent wrapping first
                wrapped_text = self._intelligent_wrap(text, box_width, font_size, font_name)
                return self.measure_text_precise(wrapped_text, font_size, font_name)
            else:
                # Fallback to improved estimation
                return self._enhanced_fallback_calculation(text, font_size, box_width)
            
        except Exception as e:
            app_logger.error(f"Error in enhanced text dimensions calculation: {str(e)}")
            return self._enhanced_fallback_calculation(text, font_size, box_width)
            
    def _enhanced_fallback_calculation(self, text: str, font_size: int, box_width: float) -> Tuple[float, float]:
        """Enhanced fallback calculation with better language support"""
        # Apply intelligent wrapping using character-based estimation
        avg_char_width = font_size * 0.5 / 72  # Convert to inches
        max_chars_per_line = int(box_width / avg_char_width) if avg_char_width > 0 else 50
        
        # Wrap text
        wrapped_text = self._intelligent_wrap(text, box_width, font_size, 'Calibri', True)
        lines = wrapped_text.split('\n')
        
        # Calculate dimensions with language-specific character widths
        max_width = 0
        for line in lines:
            line_width = self._calculate_line_width_fallback(line, font_size)
            max_width = max(max_width, line_width)
        
        # Calculate height
        line_height = font_size * self.line_spacing / 72
        total_height = len(lines) * line_height
        
        return (max_width, total_height)
        
    def _calculate_line_width_fallback(self, line: str, font_size: int) -> float:
        """Calculate line width using character-based estimation"""
        korean_chars = len(re.findall(r'[\u3131-\ucb4c]', line))
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', line))
        punctuation_chars = len(re.findall(r'[.,;:!?\-()\[\]{}]', line))
        latin_chars = len(line) - korean_chars - chinese_chars - punctuation_chars
        
        # Calculate width with language-specific multipliers
        base_char_width = font_size * 0.5 / 72  # Convert to inches
        total_width = (
            korean_chars * base_char_width * self.korean_width_multiplier +
            chinese_chars * base_char_width * self.chinese_width_multiplier +
            punctuation_chars * base_char_width * self.punctuation_width_multiplier +
            latin_chars * base_char_width
        )
        
        return total_width
    
    def auto_adjust_font_size(
        self, 
        text: str, 
        box_width: float, 
        box_height: float,
        start_size: int = 14,
        font_name: str = "Arial",
        use_binary_search: bool = False
    ) -> int:
        """
        텍스트가 박스에 맞는 최대 폰트 크기 찾기
        
        알고리즘:
        1. 시작 크기에서 텍스트 측정
        2. 넘치면 1pt씩 줄임
        3. min_font_size까지 반복
        4. 여전히 넘치면 텍스트 잘라내기 또는 경고
        """
        if use_binary_search:
            return self._binary_search_font_size(text, box_width, box_height, font_name, start_size)
        else:
            # Linear search fallback
            current_size = min(start_size, self.max_font_size.pt)
            
            while current_size >= self.min_font_size.pt:
                width, height = self.calculate_text_dimensions(text, current_size, box_width, font_name)
                
                if height <= box_height:
                    return current_size
                
                current_size -= 1
            
            return self.min_font_size.pt
    
    def truncate_with_ellipsis(
        self, 
        text: str, 
        max_length: int,
        smart_truncation: bool = True,
        preserve_sentences: bool = True
    ) -> str:
        """
        Enhanced text truncation with intelligent sentence and word boundary detection.
        
        Args:
            text: Text to truncate
            max_length: Maximum character length
            smart_truncation: Use intelligent truncation algorithms
            preserve_sentences: Try to preserve complete sentences
            
        Returns:
            Truncated text with appropriate ellipsis
        """
        if len(text) <= max_length:
            return text
            
        if not smart_truncation:
            # Simple truncation
            return text[:max_length-3] + "..." if max_length > 3 else text[:max_length]
            
        # Smart truncation with multiple strategies
        return self._apply_smart_truncation(text, max_length, preserve_sentences)
        
    def _apply_smart_truncation(self, text: str, max_length: int, preserve_sentences: bool) -> str:
        """Apply intelligent truncation strategies"""
        if max_length <= 3:
            return text[:max_length]
            
        truncated = text[:max_length-3]
        
        if preserve_sentences:
            # Try to end at sentence boundary
            sentence_boundaries = [
                truncated.rfind('.'),
                truncated.rfind('!'),
                truncated.rfind('?'),
                truncated.rfind('\u3002'),  # Korean period
                truncated.rfind('\uff01'),  # Korean exclamation
                truncated.rfind('\uff1f'),  # Korean question mark
            ]
            
            best_boundary = max(sentence_boundaries)
            if best_boundary > max_length * 0.6:  # Keep if we don't lose too much
                return text[:best_boundary+1]
        
        # Try to break at word boundaries
        word_boundaries = [
            truncated.rfind(' '),
            truncated.rfind('\t'),
            # Korean particles
            truncated.rfind('을'), truncated.rfind('를'),
            truncated.rfind('이'), truncated.rfind('가'),
            truncated.rfind('의'), truncated.rfind('에'),
            truncated.rfind('에서'), truncated.rfind('부터'),
        ]
        
        best_boundary = max(word_boundaries)
        if best_boundary > max_length * 0.7:
            return text[:best_boundary] + "..."
            
        # Fallback to character boundary
        return truncated + "..."
    
    def smart_line_break(
        self, 
        text: str, 
        max_chars_per_line: int,
        preserve_words: bool = True,
        language_aware: bool = True
    ) -> str:
        """
        Enhanced smart line breaking with language awareness and word preservation.
        
        Args:
            text: Text to break into lines
            max_chars_per_line: Maximum characters per line
            preserve_words: Avoid breaking words
            language_aware: Use language-specific breaking rules
            
        Returns:
            Text with intelligent line breaks
        """
        if len(text) <= max_chars_per_line:
            return text
            
        if language_aware:
            return self._language_aware_line_break(text, max_chars_per_line, preserve_words)
        else:
            return self._simple_line_break(text, max_chars_per_line, preserve_words)
            
    def _language_aware_line_break(self, text: str, max_chars: int, preserve_words: bool) -> str:
        """Language-aware line breaking"""
        # Detect primary language
        language = self._detect_text_language(text)
        
        if language == 'korean':
            return self._korean_line_break(text, max_chars, preserve_words)
        elif language == 'chinese':
            return self._chinese_line_break(text, max_chars)
        else:
            return self._latin_line_break(text, max_chars, preserve_words)
            
    def _korean_line_break(self, text: str, max_chars: int, preserve_words: bool) -> str:
        """Korean-specific line breaking with particle awareness"""
        lines = []
        current_line = ""
        
        i = 0
        while i < len(text):
            char = text[i]
            
            # Check if adding this character would exceed limit
            char_width = self.korean_width_multiplier if re.match(r'[\u3131-\ucb4c]', char) else 1
            current_width = self._calculate_display_length(current_line)
            
            if current_width + char_width > max_chars and current_line:
                # Try to find a good break point
                break_point = self._find_korean_break_point(current_line)
                if break_point > 0:
                    lines.append(current_line[:break_point])
                    current_line = current_line[break_point:] + char
                else:
                    lines.append(current_line)
                    current_line = char
            else:
                current_line += char
                
            i += 1
            
        if current_line:
            lines.append(current_line)
            
        return '\n'.join(lines)
        
    def _find_korean_break_point(self, line: str) -> int:
        """Find optimal break point in Korean text"""
        # Korean particles that are good break points
        particles = ['을', '를', '이', '가', '의', '에', '에서', '부터', '까지']
        
        best_break = 0
        for particle in particles:
            pos = line.rfind(particle)
            if pos > best_break:
                best_break = pos + len(particle)
                
        # If no particles found, break at last space
        if best_break == 0:
            space_pos = line.rfind(' ')
            if space_pos > 0:
                best_break = space_pos
                
        return best_break
        
    def _chinese_line_break(self, text: str, max_chars: int) -> str:
        """Chinese-specific line breaking"""
        lines = []
        current_line = ""
        
        for char in text:
            char_width = self.chinese_width_multiplier if re.match(r'[\u4e00-\u9fff]', char) else 1
            current_width = self._calculate_display_length(current_line)
            
            if current_width + char_width > max_chars and current_line:
                lines.append(current_line)
                current_line = char
            else:
                current_line += char
                
        if current_line:
            lines.append(current_line)
            
        return '\n'.join(lines)
        
    def _latin_line_break(self, text: str, max_chars: int, preserve_words: bool) -> str:
        """Latin text line breaking"""
        if preserve_words:
            return self._wrap_preserving_words(text, max_chars)
        else:
            return self._wrap_hard_break(text, max_chars)
            
    def _simple_line_break(self, text: str, max_chars: int, preserve_words: bool) -> str:
        """Simple line breaking without language awareness"""
        if preserve_words:
            words = text.split()
            lines = []
            current_line = []
            current_length = 0
            
            for word in words:
                word_length = len(word)
                space_length = 1 if current_line else 0
                
                if current_length + word_length + space_length > max_chars:
                    if current_line:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                        current_length = word_length
                    else:
                        lines.append(word[:max_chars])
                        current_line = []
                        current_length = 0
                else:
                    current_line.append(word)
                    current_length += word_length + space_length
                    
            if current_line:
                lines.append(' '.join(current_line))
                
            return '\n'.join(lines)
        else:
            return self._wrap_hard_break(text, max_chars)
    
    def _wrap_text(self, text: str, box_width_pts: float, char_width: float) -> list:
        """
        Legacy wrap text method (maintained for backward compatibility).
        
        Args:
            text: Text to wrap
            box_width_pts: Box width in points
            char_width: Average character width in points
            
        Returns:
            List of wrapped lines
        """
        max_chars_per_line = int(box_width_pts / char_width) if char_width > 0 else 50
        
        # Use enhanced line breaking
        wrapped_text = self.smart_line_break(text, max_chars_per_line, True, True)
        return wrapped_text.split('\n')
    
    def _truncate_to_fit(
        self, 
        text: str, 
        box_width: float, 
        box_height: float, 
        font_size: int,
        font_name: str = 'Calibri'
    ) -> str:
        """
        Enhanced truncate text to fit within box dimensions.
        
        Args:
            text: Text to truncate
            box_width: Box width in inches
            box_height: Box height in inches
            font_size: Font size in points
            font_name: Font family name
            
        Returns:
            Intelligently truncated text with ellipsis
        """
        # Calculate maximum lines that can fit
        line_height = font_size * self.line_spacing / 72  # Convert to inches
        max_lines = int(box_height / line_height)
        
        # Apply intelligent wrapping first
        wrapped_text = self._intelligent_wrap(text, box_width, font_size, font_name, True)
        
        # Use smart truncation
        return self._smart_truncate(wrapped_text, box_width, box_height, font_size, font_name)
    
    # Enhanced methods need to be added after the main enhanced functionality is complete
    
    def calculate_bullet_spacing(
        self,
        bullet_count: int,
        box_height: float,
        font_size: int
    ) -> Dict[str, float]:
        """
        Calculate optimal spacing for bullet points with enhanced algorithms.
        
        Args:
            bullet_count: Number of bullet points
            box_height: Available height in inches
            font_size: Font size in points
            
        Returns:
            Dictionary with enhanced spacing recommendations
        """
        # Calculate total height needed with line spacing
        line_height = font_size * self.line_spacing / 72  # Convert to inches
        total_text_height = bullet_count * line_height
        
        # Calculate optimal spacing distribution
        if total_text_height >= box_height:
            # Too tight, use minimum spacing
            return {
                "line_spacing": 1.0,
                "paragraph_spacing": 0,
                "bullet_indent": 0.2,
                "fits": False,
                "recommended_font_size": max(8, int(font_size * 0.8))
            }
        else:
            # Calculate optimal spacing for visual appeal
            available_space = box_height - total_text_height
            
            # Distribute extra space between bullets
            if bullet_count > 1:
                paragraph_spacing = min(12, available_space * 72 / (bullet_count - 1))
            else:
                paragraph_spacing = 6
                
            # Ensure reasonable spacing bounds
            paragraph_spacing = max(3, min(paragraph_spacing, 18))
            
            return {
                "line_spacing": 1.2,
                "paragraph_spacing": paragraph_spacing,
                "bullet_indent": 0.25,
                "fits": True,
                "recommended_font_size": font_size,
                "vertical_centering": available_space / 2 if available_space > 0.5 else 0
            }
            
    def get_performance_metrics(self) -> Dict[str, Any]:
        """Get performance metrics for the text fitter"""
        total_requests = self._cache_hits + self._cache_misses
        cache_hit_rate = (self._cache_hits / total_requests * 100) if total_requests > 0 else 0
        
        return {
            "cache_hit_rate": cache_hit_rate,
            "total_cache_hits": self._cache_hits,
            "total_cache_misses": self._cache_misses,
            "font_cache_size": len(self._font_cache),
            "measurement_cache_size": len(self._measurement_cache),
            "available_fonts": len(self.available_fonts),
            "font_families_supported": list(self.font_families.keys())
        }
        
    def clear_cache(self):
        """Clear all caches"""
        with self._cache_lock:
            self._font_cache.clear()
            self._measurement_cache.clear()
            self._cache_hits = 0
            self._cache_misses = 0
            
        app_logger.info("TextFitter caches cleared")
        
    def benchmark_performance(self, test_texts: List[str], iterations: int = 100) -> Dict[str, float]:
        """Benchmark text fitter performance"""
        import time
        
        times = {
            "fit_text_total": 0,
            "measure_text_total": 0,
            "wrap_text_total": 0,
            "truncate_text_total": 0
        }
        
        for _ in range(iterations):
            for text in test_texts:
                # Benchmark fit_text_to_box
                start = time.time()
                self.fit_text_to_box(text, 4.0, 2.0)
                times["fit_text_total"] += time.time() - start
                
                # Benchmark measure_text_precise
                start = time.time()
                self.measure_text_precise(text, 14)
                times["measure_text_total"] += time.time() - start
                
                # Benchmark smart_line_break
                start = time.time()
                self.smart_line_break(text, 50)
                times["wrap_text_total"] += time.time() - start
                
                # Benchmark truncate_with_ellipsis
                start = time.time()
                self.truncate_with_ellipsis(text, 100)
                times["truncate_text_total"] += time.time() - start
                
        # Calculate averages
        total_operations = iterations * len(test_texts)
        return {
            "avg_fit_text_ms": (times["fit_text_total"] / total_operations) * 1000,
            "avg_measure_text_ms": (times["measure_text_total"] / total_operations) * 1000,
            "avg_wrap_text_ms": (times["wrap_text_total"] / total_operations) * 1000,
            "avg_truncate_text_ms": (times["truncate_text_total"] / total_operations) * 1000,
            "total_operations": total_operations,
            "cache_hit_rate": self.get_performance_metrics()["cache_hit_rate"]
        }
    
    def validate_text_boundaries(
        self,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 14,
        slide_width: float = 10.0,
        slide_height: float = 7.5
    ) -> Dict[str, Any]:
        """
        슬라이드 경계 내 텍스트 위치 검증
        
        Args:
            text: 검증할 텍스트
            left: 왼쪽 위치 (인치)
            top: 상단 위치 (인치)
            width: 텍스트 박스 너비 (인치)
            height: 텍스트 박스 높이 (인치)
            font_size: 폰트 크기 (pt)
            slide_width: 슬라이드 너비 (인치, 기본 10")
            slide_height: 슬라이드 높이 (인치, 기본 7.5")
            
        Returns:
            {
                'is_valid': bool,
                'issues': List[str],
                'adjusted_position': Dict,
                'adjusted_font_size': int
            }
        """
        issues = []
        adjusted = {}
        
        # 최소 여백 정의
        MIN_MARGIN = 0.5  # 0.5인치
        
        # 1. 경계 검증
        if left < MIN_MARGIN:
            issues.append(f"Left margin too small: {left:.2f}")
            adjusted['left'] = MIN_MARGIN
        else:
            adjusted['left'] = left
            
        if top < MIN_MARGIN:
            issues.append(f"Top margin too small: {top:.2f}")
            adjusted['top'] = MIN_MARGIN
        else:
            adjusted['top'] = top
            
        if left + width > slide_width - MIN_MARGIN:
            issues.append(f"Right boundary exceeded: {left + width:.2f}")
            adjusted['width'] = slide_width - adjusted['left'] - MIN_MARGIN
        else:
            adjusted['width'] = width
            
        if top + height > slide_height - MIN_MARGIN:
            issues.append(f"Bottom boundary exceeded: {top + height:.2f}")
            adjusted['height'] = slide_height - adjusted['top'] - MIN_MARGIN
        else:
            adjusted['height'] = height
            
        # 2. 텍스트 크기 조정
        text_width, text_height = self.calculate_text_dimensions(
            text, font_size, adjusted['width']
        )
        
        adjusted_font_size = font_size
        
        # 텍스트가 박스보다 크면 폰트 크기 자동 조정
        if text_height > adjusted['height']:
            # 최적 폰트 크기 찾기
            adjusted_font_size = self.auto_adjust_font_size(
                text, 
                adjusted['width'], 
                adjusted['height'],
                font_size
            )
            issues.append(f"Font size adjusted from {font_size}pt to {adjusted_font_size}pt")
            
        # 3. 결과 반환
        is_valid = len(issues) == 0
        
        return {
            'is_valid': is_valid,
            'issues': issues,
            'adjusted_position': adjusted,
            'adjusted_font_size': adjusted_font_size,
            'recommendations': self._get_boundary_recommendations(issues)
        }
    
    def _get_boundary_recommendations(self, issues: List[str]) -> List[str]:
        """
        경계 문제에 대한 권장사항 제공
        """
        recommendations = []
        
        if any('margin too small' in issue for issue in issues):
            recommendations.append("최소 여백 0.5인치 확보 필요")
            
        if any('boundary exceeded' in issue for issue in issues):
            recommendations.append("텍스트 박스 크기 조정 필요")
            
        if any('Font size adjusted' in issue for issue in issues):
            recommendations.append("콘텐츠 양 축소 또는 슬라이드 분할 검토")
            
        return recommendations
    
    def fit_text_to_shape(
        self,
        shape,
        text: str,
        min_font_size: int = 8,
        max_font_size: int = 18
    ) -> bool:
        """
        python-pptx shape에 텍스트 맞추기
        
        Args:
            shape: pptx shape 객체
            text: 삽입할 텍스트
            min_font_size: 최소 폰트 크기
            max_font_size: 최대 폰트 크기
            
        Returns:
            bool: 성공 여부
        """
        try:
            from pptx.util import Inches, Pt
            
            # Shape 크기 가져오기
            width_inches = shape.width / 914400  # EMU to inches
            height_inches = shape.height / 914400
            
            # 최적 폰트 크기 찾기
            optimal_font_size = self.auto_adjust_font_size(
                text, width_inches, height_inches, max_font_size
            )
            
            # 폰트 크기가 너무 작으면 텍스트 자르기
            if optimal_font_size < min_font_size:
                # 텍스트 자르기
                max_chars = int(width_inches * height_inches * 10)  # 대략적인 계산
                text = self.truncate_with_ellipsis(text, max_chars)
                optimal_font_size = min_font_size
            
            # Shape에 텍스트 적용
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                p = text_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(optimal_font_size)
                p.font.name = "Arial"
                
                app_logger.info(f"Text fitted to shape with font size {optimal_font_size}pt")
                return True
            
            return False
            
        except Exception as e:
            app_logger.error(f"Failed to fit text to shape: {e}")
            return False