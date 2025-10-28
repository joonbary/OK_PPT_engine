from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from datetime import datetime
import os
from loguru import logger
from app.templates.mckinsey_layouts import McKinseyLayoutManager

class PptxGeneratorService:
    """
    styled_slides 데이터를 기반으로 실제 PPTX 파일을 생성하는 서비스
    """
    def __init__(self):
        self.output_dir = "output/generated_presentations"
        os.makedirs(self.output_dir, exist_ok=True)
        # Fix #3: McKinsey Layout Manager 통합
        self.layout_manager = McKinseyLayoutManager()
        logger.info(f"PptxGeneratorService initialized. Output directory: {self.output_dir}")

    def generate_pptx(self, styled_slides: List[Dict], output_filename: str, chart_images: List[Dict] = None) -> str:
        """
        스타일링된 슬라이드 데이터를 기반으로 PPTX 파일을 생성하고 저장합니다.

        Args:
            styled_slides (List[Dict]): DesignAgent에서 생성된 스타일링된 슬라이드 데이터 목록
            output_filename (str): 저장할 PPTX 파일의 이름 (예: "my_presentation.pptx")
            chart_images (List[Dict], optional): 차트 이미지 정보 목록. Defaults to None.

        Returns:
            str: 생성된 PPTX 파일의 절대 경로
        """
        prs = Presentation()
        self.apply_global_styles(prs)

        for i, slide_data in enumerate(styled_slides):
            slide = self._add_slide(prs, slide_data)
            
            if chart_images:
                for chart in chart_images:
                    if chart.get('slide_index') == i:
                        self.add_chart_to_slide(
                            slide,
                            chart['path'],
                            position=chart.get('position', 'center')
                        )

        output_path = os.path.join(self.output_dir, output_filename)
        prs.save(output_path)
        # 절대 경로 반환
        absolute_path = os.path.abspath(output_path)
        logger.info(f"PPTX file generated and saved to: {absolute_path}")
        return absolute_path

    def apply_global_styles(self, prs: Presentation):
        """
        프레젠테이션 전체에 적용될 기본 스타일을 설정합니다.
        """
        # 기본 폰트 설정 (Arial)
        for slide_layout in prs.slide_layouts:
            for shape in slide_layout.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font = run.font
                            font.name = 'Arial'
                            font.size = Pt(18) # 기본 크기
                            font.color.theme_color = MSO_THEME_COLOR.TEXT_1

    def _add_slide(self, prs: Presentation, slide_data: Dict):
        """
        단일 슬라이드를 프레젠테이션에 추가하고 내용을 채웁니다.
        Fix #3: McKinsey Layout Manager 적용
        """
        try:
            # 빈 슬라이드 생성
            blank_slide_layout = prs.slide_layouts[6] # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 슬라이드 인덱스 계산
            slide_idx = len(prs.slides) - 1
            
            # 슬라이드 데이터를 McKinsey Layout 형식으로 변환
            layout_content = self._convert_to_layout_format(slide_data)
            
            # 레이아웃 타입 결정
            layout_type = self.layout_manager.determine_layout_type(layout_content, slide_idx)
            logger.info(f"🎨 Applying McKinsey layout: {layout_type} for slide {slide_idx}")
            
            # McKinsey 레이아웃 적용
            result = self.layout_manager.apply_layout(slide, layout_type, layout_content)
            
            # 차트 슬라이드의 경우 추가 처리
            if layout_type == 'chart_slide' and isinstance(result, tuple):
                slide, chart_area = result
                # 차트 데이터가 있으면 처리
                if slide_data.get('chart_data'):
                    self._add_chart_to_layout_area(slide, slide_data, chart_area)
            
            # speaker_notes 추가
            if 'speaker_notes' in slide_data and slide_data['speaker_notes']:
                slide.notes_slide.notes_text_frame.text = slide_data['speaker_notes']
            
            # AI 헤드라인 로그 (기존 형식 유지)
            if 'headline' in slide_data and slide_data['headline']:
                logger.info(f"    ✨ AI Headline: {slide_data['headline']}")
            elif 'title' in slide_data and slide_data['title']:
                logger.info(f"    📝 Using title: {slide_data['title']}")
            
            return slide
            
        except Exception as e:
            logger.error(f"❌ McKinsey layout application failed: {e}")
            # Fallback to original method
            return self._add_slide_fallback(prs, slide_data)
    
    def _convert_to_layout_format(self, slide_data: Dict) -> Dict:
        """슬라이드 데이터를 McKinsey Layout 형식으로 변환"""
        layout_content = {}
        
        # 제목 처리
        if 'headline' in slide_data and slide_data['headline']:
            layout_content['title'] = slide_data['headline']
        elif 'title' in slide_data and slide_data['title']:
            layout_content['title'] = slide_data['title']
        
        # 콘텐츠 처리
        if 'key_points' in slide_data and slide_data['key_points']:
            layout_content['bullets'] = slide_data['key_points']
            layout_content['body'] = slide_data['key_points']
        elif 'content' in slide_data and slide_data['content']:
            layout_content['content'] = slide_data['content']
            layout_content['body'] = slide_data['content']
        
        # 차트 데이터 처리
        if 'chart_data' in slide_data:
            layout_content['chart_data'] = slide_data['chart_data']
            layout_content['has_chart'] = True
        
        # 2단 레이아웃 처리
        if 'left_content' in slide_data and 'right_content' in slide_data:
            layout_content['left_content'] = slide_data['left_content']
            layout_content['right_content'] = slide_data['right_content']
            layout_content['has_columns'] = True
        
        return layout_content
    
    def _add_chart_to_layout_area(self, slide, slide_data: Dict, chart_area: Dict):
        """차트 영역에 차트 추가"""
        try:
            from app.services.simple_chart_generator import SimpleChartGenerator
            chart_generator = SimpleChartGenerator()
            
            chart_data = slide_data.get('chart_data')
            if not chart_data:
                return
            
            chart_type = chart_data.get('type', 'bar')
            
            if chart_type == 'bar':
                success = chart_generator.create_bar_chart(chart_data, slide)
            elif chart_type == 'line':
                success = chart_generator.create_line_chart(chart_data, slide)
            elif chart_type == 'pie':
                success = chart_generator.create_pie_chart(chart_data, slide)
            else:
                success = False
            
            if success:
                logger.info(f"✅ Chart added to McKinsey layout: {chart_type}")
            else:
                logger.error(f"❌ Failed to add chart to McKinsey layout")
                
        except Exception as e:
            logger.error(f"Error adding chart to layout area: {e}")
    
    def _add_slide_fallback(self, prs: Presentation, slide_data: Dict):
        """Fallback 슬라이드 생성 (기존 방식)"""
        blank_slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # 제목 추가 - AI 생성 헤드라인 우선 사용
        title_text = None
        if 'headline' in slide_data and slide_data['headline']:
            title_text = slide_data['headline']
            logger.info(f"    ✨ AI Headline: {title_text}")
        elif 'title' in slide_data and slide_data['title']:
            title_text = slide_data['title']
            logger.info(f"    📝 Using title: {title_text}")
        
        if title_text:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = title_text
            p.font.name = slide_data.get('fonts', {}).get('title', {}).get('family', 'Arial')
            p.font.size = Pt(slide_data.get('fonts', {}).get('title', {}).get('size', 24))
            p.font.bold = slide_data.get('fonts', {}).get('title', {}).get('weight', 'bold') == 'bold'
            p.font.color.rgb = self._hex_to_rgb(slide_data.get('colors', {}).get('text', '#000000'))

        # 주요 내용 추가 (key_points 또는 content)
        content_text = ""
        if 'key_points' in slide_data and slide_data['key_points']:
            content_text = "\n".join([f"- {point}" for point in slide_data['key_points']])
        elif 'content' in slide_data and slide_data['content']:
            content_text = slide_data['content']

        if content_text:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            body_shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = body_shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p = text_frame.paragraphs[0]
            p.text = content_text
            p.font.name = slide_data.get('fonts', {}).get('body', {}).get('family', 'Arial')
            p.font.size = Pt(slide_data.get('fonts', {}).get('body', {}).get('size', 14))
            p.font.color.rgb = self._hex_to_rgb(slide_data.get('colors', {}).get('text', '#000000'))

        # speaker_notes 추가
        if 'speaker_notes' in slide_data and slide_data['speaker_notes']:
            slide.notes_slide.notes_text_frame.text = slide_data['speaker_notes']
        
        return slide

    def add_chart_to_slide(self, slide, image_path: str, position: str = "center"):
        """차트 이미지를 슬라이드에 추가"""
        from pptx.util import Inches
        
        if not os.path.exists(image_path):
            logger.error(f"차트 이미지 없음: {image_path}")
            return False
        
        # 위치 계산
        if position == "center":
            left = Inches(1.5)
            top = Inches(2.5)
            width = Inches(7.0)
            height = Inches(4.0)
        elif position == "right":
            left = Inches(5.5)
            top = Inches(2.0)
            width = Inches(4.0)
            height = Inches(4.5)
        else:  # left
            left = Inches(0.5)
            top = Inches(2.0)
            width = Inches(4.0)
            height = Inches(4.5)
        
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            logger.info(f"차트 추가 성공: {image_path}")
            return True
        except Exception as e:
            logger.error(f"차트 추가 실패: {e}")
            return False

    def _hex_to_rgb(self, hex_color: str):
        """16진수 색상 코드를 RGBColor 객체로 변환"""
        from pptx.dml.color import RGBColor
        hex_color = hex_color.lstrip('#')
        return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))


if __name__ == '__main__':
    # 테스트용 styled_slides 데이터
    test_slides = [
        {
            'title': '테스트 슬라이드 1: 시장 분석',
            'key_points': ['시장 규모 1000억원 달성', '모바일 쇼핑 비중 70%'],
            'colors': {'primary': '#0076A8', 'text': '#000000'},
            'fonts': {'title': {'family': 'Arial', 'size': 28, 'weight': 'bold'}, 'body': {'family': 'Arial', 'size': 18}},
            'speaker_notes': '이 슬라이드는 2024년 전자상거래 시장의 주요 현황을 보여줍니다.'
        },
        {
            'title': '테스트 슬라이드 2: 전략적 권고사항',
            'content': '모바일 앱 UI/UX 전면 개편 및 AI 기반 추천 알고리즘 도입이 필요합니다.',
            'chart_specs': [{'type': 'bar', 'data': {'labels': ['A', 'B'], 'values': [10, 20]}}],
            'colors': {'primary': '#0076A8', 'text': '#000000'},
            'fonts': {'title': {'family': 'Arial', 'size': 28, 'weight': 'bold'}, 'body': {'family': 'Arial', 'size': 18}},
            'speaker_notes': '핵심 과제를 해결하기 위한 구체적인 전략을 제시합니다.'
        }
    ]

    generator = PptxGeneratorService()
    output_file = f"test_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    generated_path = generator.generate_pptx(test_slides, output_file)
    print(f"Generated test PPTX at: {generated_path}")