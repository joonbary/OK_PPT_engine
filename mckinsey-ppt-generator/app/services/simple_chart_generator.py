"""
Simple Chart Generator for Direct PPT Charts
차트를 PPT에 직접 생성하는 모듈
"""

import logging
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class SimpleChartGenerator:
    """PPT 네이티브 차트 생성기"""
    
    # McKinsey 색상
    MCKINSEY_COLORS = [
        '0076A8',  # McKinsey Blue
        'F47621',  # Orange  
        '6BA644',  # Green
        'E31B23',  # Red
        '53565A',  # Gray
        '808285',  # Light Gray
    ]
    
    def create_bar_chart(self, chart_data: dict, slide):
        """막대 차트 생성"""
        try:
            # 차트 데이터 준비
            chart_data_obj = ChartData()
            
            # 카테고리 (레이블)
            categories = chart_data.get('labels', ['A', 'B', 'C', 'D'])
            chart_data_obj.categories = categories
            
            # 데이터 시리즈
            values = chart_data.get('values', [30, 40, 25, 50])
            title = chart_data.get('title', 'Data Analysis')
            
            # 값이 리스트가 아니면 리스트로 변환
            if not isinstance(values, list):
                values = [values]
            
            # 카테고리와 값 개수 맞추기
            if len(values) < len(categories):
                values.extend([0] * (len(categories) - len(values)))
            elif len(values) > len(categories):
                values = values[:len(categories)]
            
            chart_data_obj.add_series('Values', values)
            
            # 차트 추가 (위치와 크기)
            x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(4.5)
            
            # 차트 생성
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                x, y, cx, cy,
                chart_data_obj
            ).chart
            
            # 차트 제목
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            
            # 색상 설정 (McKinsey Blue)
            if chart.plots and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'series') and len(plot.series) > 0:
                    series = plot.series[0]
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0, 118, 168)  # McKinsey Blue
            
            logger.info(f"✅ Bar chart created: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Bar chart creation failed: {e}")
            return False
    
    def create_line_chart(self, chart_data: dict, slide):
        """선 차트 생성"""
        try:
            # 차트 데이터 준비
            chart_data_obj = ChartData()
            
            # 카테고리 (X축)
            categories = chart_data.get('labels', ['Q1', 'Q2', 'Q3', 'Q4'])
            chart_data_obj.categories = categories
            
            # 데이터 시리즈 (Y축)
            values = chart_data.get('values', [20, 35, 30, 42])
            title = chart_data.get('title', 'Trend Analysis')
            
            # 값이 리스트가 아니면 리스트로 변환
            if not isinstance(values, list):
                values = [values]
                
            # 카테고리와 값 개수 맞추기
            if len(values) < len(categories):
                values.extend([values[-1] if values else 0] * (len(categories) - len(values)))
            elif len(values) > len(categories):
                values = values[:len(categories)]
            
            chart_data_obj.add_series('Trend', values)
            
            # 차트 추가
            x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(4.5)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS,
                x, y, cx, cy,
                chart_data_obj
            ).chart
            
            # 차트 제목
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            
            # 선 색상 설정
            if chart.plots and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'series') and len(plot.series) > 0:
                    series = plot.series[0]
                    line = series.format.line
                    line.color.rgb = RGBColor(0, 118, 168)  # McKinsey Blue
                    line.width = Pt(2.5)
            
            logger.info(f"✅ Line chart created: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Line chart creation failed: {e}")
            return False
    
    def create_pie_chart(self, chart_data: dict, slide):
        """파이 차트 생성"""
        try:
            # 차트 데이터 준비
            chart_data_obj = ChartData()
            
            # 카테고리
            categories = chart_data.get('labels', ['A', 'B', 'C', 'D'])
            values = chart_data.get('values', [30, 25, 20, 25])
            title = chart_data.get('title', 'Distribution')
            
            # 값이 리스트가 아니면 리스트로 변환
            if not isinstance(values, list):
                values = [values]
            
            # 카테고리와 값 개수 맞추기
            if len(values) < len(categories):
                values.extend([10] * (len(categories) - len(values)))
            elif len(values) > len(categories):
                values = values[:len(categories)]
            
            chart_data_obj.categories = categories
            chart_data_obj.add_series('Distribution', values)
            
            # 차트 추가
            x, y, cx, cy = Inches(2.5), Inches(2), Inches(8), Inches(4.5)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                x, y, cx, cy,
                chart_data_obj
            ).chart
            
            # 차트 제목
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            
            # 데이터 레이블 표시
            if chart.plots and len(chart.plots) > 0:
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.show_percentage = True
                data_labels.show_category_name = True
                data_labels.position = -30  # 내부 표시
            
            logger.info(f"✅ Pie chart created: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Pie chart creation failed: {e}")
            return False
    
    def create_table(self, table_data: dict, slide):
        """테이블 생성"""
        try:
            rows = table_data.get('rows', 3)
            cols = table_data.get('cols', 3)
            data = table_data.get('data', [])
            title = table_data.get('title', '')
            
            # 테이블 추가
            x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(4)
            
            table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
            
            # 헤더 스타일
            for col_idx in range(cols):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 118, 168)  # McKinsey Blue
                
                # 헤더 텍스트
                if col_idx < len(data[0]) if data else False:
                    cell.text = str(data[0][col_idx])
                    # 헤더 텍스트 색상
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
                            run.font.size = Pt(11)
            
            # 데이터 행
            for row_idx in range(1, rows):
                for col_idx in range(cols):
                    if row_idx < len(data) and col_idx < len(data[row_idx]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(data[row_idx][col_idx])
                        # 데이터 텍스트 스타일
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(10)
            
            logger.info(f"✅ Table created: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Table creation failed: {e}")
            return False