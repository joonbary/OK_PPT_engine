"""
Slide Validator for PowerPoint presentations.
Validates slide quality and detects issues.
"""

from typing import Dict, List, Optional, Any, Tuple, Union
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import time
from dataclasses import dataclass, field
from enum import Enum
from app.core.logging import app_logger


class IssueSeverity(Enum):
    """검증 이슈의 심각도 레벨"""
    CRITICAL = "critical"  # 반드시 수정해야 하는 오류
    WARNING = "warning"   # 권장 사항
    SUGGESTION = "suggestion"  # 개선 제안
    INFO = "info"  # 정보 제공


class IssueCategory(Enum):
    """검증 이슈의 카테고리"""
    TEXT_OVERFLOW = "text_overflow"
    SHAPE_OVERLAP = "shape_overlap"
    OUT_OF_BOUNDS = "out_of_bounds"
    FONT_CONSISTENCY = "font_consistency"
    MARGIN_ISSUES = "margin_issues"
    READABILITY = "readability"
    CONTENT_DENSITY = "content_density"
    EMPTY_CONTENT = "empty_content"
    MCKINSEY_STYLE = "mckinsey_style"
    PERFORMANCE = "performance"


@dataclass
class ValidationIssue:
    """개별 검증 이슈를 나타내는 데이터 클래스"""
    severity: IssueSeverity
    category: IssueCategory
    message: str
    shape_id: Optional[str] = None
    position: Optional[Tuple[float, float]] = None
    suggested_fix: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def __str__(self) -> str:
        return f"[{self.severity.value.upper()}] {self.category.value}: {self.message}"


@dataclass
class ValidationResult:
    """전체 슬라이드 검증 결과를 나타내는 데이터 클래스"""
    is_valid: bool
    slide_number: Optional[int] = None
    issues: List[ValidationIssue] = field(default_factory=list)
    metrics: Dict[str, Any] = field(default_factory=dict)
    processing_time_ms: float = 0.0
    
    @property
    def critical_issues(self) -> List[ValidationIssue]:
        return [issue for issue in self.issues if issue.severity == IssueSeverity.CRITICAL]
    
    @property
    def warnings(self) -> List[ValidationIssue]:
        return [issue for issue in self.issues if issue.severity == IssueSeverity.WARNING]
    
    @property
    def suggestions(self) -> List[ValidationIssue]:
        return [issue for issue in self.issues if issue.severity == IssueSeverity.SUGGESTION]
    
    @property
    def info_items(self) -> List[ValidationIssue]:
        return [issue for issue in self.issues if issue.severity == IssueSeverity.INFO]
    
    def get_issues_by_category(self, category: IssueCategory) -> List[ValidationIssue]:
        return [issue for issue in self.issues if issue.category == category]
    
    def to_dict(self) -> Dict[str, Any]:
        """레거시 호환성을 위한 딕셔너리 변환"""
        return {
            "is_valid": self.is_valid,
            "issues": [issue.message for issue in self.critical_issues],
            "warnings": [issue.message for issue in self.warnings],
            "suggestions": [issue.message for issue in self.suggestions],
            "metrics": self.metrics,
            "processing_time_ms": self.processing_time_ms
        }


class SlideValidator:
    """
    생성된 슬라이드의 품질을 검증하고 문제점 감지
    """
    
    def __init__(self):
        """Initialize slide validator with validation rules"""
        # Import TextFitter for enhanced validation
        try:
            from .text_fitter import TextFitter
            self.text_fitter = TextFitter()
            self.text_fitter_available = True
            app_logger.info("Enhanced TextFitter integration enabled for validation")
        except ImportError as e:
            self.text_fitter = None
            self.text_fitter_available = False
            app_logger.warning(f"TextFitter not available for validation: {e}")
        
        # Slide dimensions (16:9 format)
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        
        # Margin limits (McKinsey standards) - 강화된 기준
        self.min_margin = Inches(0.5)  # PDF 변환 안전 마진
        self.max_margin = Inches(1.0)
        self.recommended_margin = Inches(0.75)
        self.pdf_safe_margin = Inches(0.6)  # PDF 변환시 최소 안전 마진
        
        # Font size limits (McKinsey standards) - PDF 안전 기준
        self.min_font_size = Pt(11)  # PDF 최소 가독 크기
        self.max_font_size = Pt(32)  # 적절한 최대 크기
        self.pdf_min_font_size = Pt(10)  # PDF 불가능 한계
        self.recommended_min_font = Pt(12)
        self.recommended_title_font = Pt(24)
        
        # Text length limits (McKinsey best practices)
        self.max_title_length = 60  # 더 간결한 제목
        self.max_bullet_length = 100  # 더 간결한 불릿
        self.max_bullets_per_slide = 5  # 더 적은 불릿
        self.max_text_lines_per_box = 8
        
        # Color contrast ratio for readability
        self.min_contrast_ratio = 4.5  # WCAG AA standard
        
        # McKinsey style compliance rules
        self.mckinsey_style_rules = {
            "preferred_fonts": ["Arial", "Calibri", "Helvetica"],
            "max_font_variants": 2,
            "title_font_min": Pt(20),
            "body_font_range": (Pt(11), Pt(16)),
            "bullet_indent_standard": Inches(0.25)
        }
        
        # Performance thresholds
        self.performance_thresholds = {
            "max_validation_time_ms": 200,  # PDF 검증을 위해 시간 여유
            "max_shapes_per_slide": 12,  # PDF 성능을 위해 제한
            "max_text_boxes_per_slide": 6  # PDF 변환 안정성
        }
        
        # PDF 변환 경고 임계값
        self.pdf_warning_thresholds = {
            "min_text_box_width": Inches(1.0),  # 최소 텍스트 박스 너비
            "min_text_box_height": Inches(0.5),  # 최소 텍스트 박스 높이
            "max_text_density": 0.8,  # 최대 텍스트 밀도
            "min_shape_spacing": Inches(0.1)  # 최소 도형 간격
        }
        
    def validate_slide(self, slide: Slide, slide_number: Optional[int] = None) -> Union[ValidationResult, Dict[str, Any]]:
        """
        Enhanced 슬라이드 검증 (새로운 ValidationResult 구조 사용)
        
        Args:
            slide: 검증할 슬라이드
            slide_number: 슬라이드 번호 (옵션)
            
        Returns:
            ValidationResult 객체 또는 레거시 호환 딕셔너리
        """
        start_time = time.perf_counter()
        
        try:
            result = ValidationResult(
                is_valid=True,
                slide_number=slide_number,
                issues=[],
                metrics={},
                processing_time_ms=0.0
            )
            
            # 1. Enhanced Text overflow check with TextFitter integration
            self._check_text_overflow_enhanced(slide, result)
            
            # 2. Precise shape overlap detection
            self._check_overlapping_shapes_enhanced(slide, result)
            
            # 3. Boundary violation check
            self._check_out_of_bounds_enhanced(slide, result)
            
            # 4. Advanced font consistency validation
            self._check_font_consistency_enhanced(slide, result)
            
            # 5. McKinsey style compliance check
            self._check_mckinsey_style_compliance(slide, result)
            
            # 6. Margin and spacing validation
            self._check_margins_enhanced(slide, result)
            
            # 7. Readability and content quality
            self._check_readability_enhanced(slide, result)
            
            # 8. Content density and structure
            self._check_content_density_enhanced(slide, result)
            
            # 9. Performance metrics
            self._check_performance_metrics(slide, result)
            
            # 10. Empty content validation
            self._check_empty_content(slide, result)
            
            # 11. PDF 변환 호환성 검증 (NEW)
            self._check_pdf_compatibility(slide, result)
            
            # 12. 레이아웃 무결성 검증 (NEW)
            self._check_layout_integrity(slide, result)
            
            # Determine overall validity
            result.is_valid = len(result.critical_issues) == 0
            
            # Calculate processing time
            end_time = time.perf_counter()
            result.processing_time_ms = (end_time - start_time) * 1000
            
            # Add performance info
            if result.processing_time_ms > self.performance_thresholds["max_validation_time_ms"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.PERFORMANCE,
                    message=f"검증 시간 초과: {result.processing_time_ms:.1f}ms (목표: {self.performance_thresholds['max_validation_time_ms']}ms)",
                    suggested_fix="슬라이드 복잡도를 줄이거나 요소 수를 감소시키세요"
                ))
            
            app_logger.info(f"Enhanced slide validation: valid={result.is_valid}, "
                          f"critical={len(result.critical_issues)}, "
                          f"warnings={len(result.warnings)}, "
                          f"time={result.processing_time_ms:.1f}ms")
            
            return result
            
        except Exception as e:
            end_time = time.perf_counter()
            processing_time = (end_time - start_time) * 1000
            
            app_logger.error(f"Error in enhanced slide validation: {str(e)}")
            
            error_result = ValidationResult(
                is_valid=False,
                slide_number=slide_number,
                issues=[ValidationIssue(
                    severity=IssueSeverity.CRITICAL,
                    category=IssueCategory.PERFORMANCE,
                    message=f"검증 시스템 오류: {str(e)}",
                    suggested_fix="슬라이드 구조를 확인하고 다시 시도하세요"
                )],
                metrics={"validation_error": True},
                processing_time_ms=processing_time
            )
            
            return error_result
            
    def validate_slide_legacy(self, slide: Slide) -> Dict[str, Any]:
        """
        레거시 호환성을 위한 기존 검증 메서드
        기존 코드와의 호환성을 유지하기 위해 유지
        """
        result = self.validate_slide(slide)
        if isinstance(result, ValidationResult):
            return result.to_dict()
        return result
    
    def _check_text_overflow_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """
        Enhanced TextFitter와 통합된 정밀한 텍스트 오버플로우 검증
        """
        try:
            overflow_count = 0
            shapes_analyzed = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    continue
                    
                if not shape.text_frame.text.strip():
                    continue
                    
                shapes_analyzed += 1
                text_frame = shape.text_frame
                text = text_frame.text
                
                # Shape 정보 추출
                shape_width = shape.width / 914400  # EMU to inches
                shape_height = shape.height / 914400
                shape_id = f"shape_{shape_idx}"
                
                # Enhanced TextFitter를 사용한 정밀한 검증
                if self.text_fitter_available:
                    try:
                        # 현재 텍스트가 박스에 맞는지 정밀 검증
                        fit_result = self.text_fitter.fit_text_to_box(
                            text, 
                            shape_width, 
                            shape_height, 
                            initial_font_size=12,
                            use_binary_search=True
                        )
                        
                        # 실제 폰트 크기 확인
                        current_font_size = 12  # 기본값
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    current_font_size = run.font.size.pt
                                    break
                            if current_font_size != 12:
                                break
                        
                        # 오버플로우 검증
                        adjusted_font_size = fit_result.get("adjusted_font_size", current_font_size)
                        
                        if adjusted_font_size < current_font_size * 0.9:  # 10% 여유
                            overflow_count += 1
                            
                            severity = IssueSeverity.CRITICAL if adjusted_font_size < current_font_size * 0.7 else IssueSeverity.WARNING
                            
                            result.issues.append(ValidationIssue(
                                severity=severity,
                                category=IssueCategory.TEXT_OVERFLOW,
                                message=f"텍스트 오버플로우: 현재 {current_font_size}pt → 권장 {adjusted_font_size:.1f}pt",
                                shape_id=shape_id,
                                position=(shape.left / 914400, shape.top / 914400),
                                suggested_fix=f"폰트 크기를 {adjusted_font_size:.1f}pt로 조정하거나 텍스트를 줄이세요",
                                metadata={
                                    "current_font_size": current_font_size,
                                    "recommended_font_size": adjusted_font_size,
                                    "text_length": len(text),
                                    "shape_size": (shape_width, shape_height),
                                    "overflow_ratio": current_font_size / adjusted_font_size if adjusted_font_size > 0 else 1
                                }
                            ))
                            
                        # 텍스트가 완전히 축약된 경우
                        if fit_result.get("truncated", False):
                            adjusted_text = fit_result.get("adjusted_text", text)
                            result.issues.append(ValidationIssue(
                                severity=IssueSeverity.CRITICAL,
                                category=IssueCategory.TEXT_OVERFLOW,
                                message=f"텍스트 강제 축약: {len(text)}자 → {len(adjusted_text)}자",
                                shape_id=shape_id,
                                suggested_fix="텍스트 박스 크기를 늘리거나 텍스트를 줄이세요",
                                metadata={
                                    "original_length": len(text),
                                    "truncated_length": len(adjusted_text),
                                    "truncation_ratio": len(adjusted_text) / len(text) if len(text) > 0 else 1
                                }
                            ))
                        
                    except Exception as e:
                        app_logger.warning(f"TextFitter validation failed for shape {shape_idx}: {e}")
                        # Fallback to basic validation
                        self._basic_overflow_check(shape, shape_idx, result)
                else:
                    # Fallback to basic validation when TextFitter is not available
                    self._basic_overflow_check(shape, shape_idx, result)
            
            # 메트릭 업데이트
            result.metrics.update({
                "text_overflow_count": overflow_count,
                "shapes_analyzed": shapes_analyzed,
                "text_fitter_enabled": self.text_fitter_available
            })
            
        except Exception as e:
            app_logger.error(f"Error in enhanced text overflow check: {e}")
            result.issues.append(ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.TEXT_OVERFLOW,
                message=f"텍스트 오버플로우 검증 오류: {str(e)}",
                suggested_fix="슬라이드를 수동으로 확인하세요"
            ))
    
    def _basic_overflow_check(self, shape, shape_idx: int, result: ValidationResult) -> None:
        """
        TextFitter가 사용 불가능할 때 사용하는 기본 오버플로우 검증
        """
        try:
            text_frame = shape.text_frame
            text = text_frame.text
            shape_area = (shape.width / 914400) * (shape.height / 914400)
            
            # 기본적인 추정 (20자/인치² at 12pt)
            estimated_capacity = shape_area * 20
            
            if len(text) > estimated_capacity * 1.5:  # 50% 초과시 경고
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.TEXT_OVERFLOW,
                    message=f"텍스트 과밀 가능성: {len(text)}자 (박스 크기: {shape_area:.2f}인치²)",
                    shape_id=f"shape_{shape_idx}",
                    suggested_fix="텍스트를 줄이거나 박스 크기를 늘리세요",
                    metadata={
                        "text_length": len(text),
                        "estimated_capacity": estimated_capacity,
                        "overflow_ratio": len(text) / estimated_capacity
                    }
                ))
                
        except Exception as e:
            app_logger.warning(f"Basic overflow check failed: {e}")
    
    def _check_pdf_compatibility(self, slide: Slide, result: ValidationResult) -> None:
        """
        PDF 변환 호환성 사전 검증
        PDF 변환 시 발생할 수 있는 문제를 미리 감지
        """
        try:
            pdf_issues = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_id = f"shape_{shape_idx}"
                
                # 1. 텍스트 박스 크기 검증
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    width = shape.width / 914400  # EMU to inches
                    height = shape.height / 914400
                    
                    # 너무 작은 텍스트 박스
                    if width < self.pdf_warning_thresholds["min_text_box_width"] or \
                       height < self.pdf_warning_thresholds["min_text_box_height"]:
                        pdf_issues += 1
                        result.issues.append(ValidationIssue(
                            severity=IssueSeverity.WARNING,
                            category=IssueCategory.OUT_OF_BOUNDS,
                            message=f"PDF 변환 위험: 텍스트 박스가 너무 작음 ({width:.1f}x{height:.1f}인치)",
                            shape_id=shape_id,
                            suggested_fix=f"텍스트 박스 크기를 {self.pdf_warning_thresholds['min_text_box_width']:.1f}x{self.pdf_warning_thresholds['min_text_box_height']:.1f}인치 이상으로 조정",
                            metadata={
                                "current_size": (width, height),
                                "min_safe_size": (
                                    self.pdf_warning_thresholds["min_text_box_width"] / Inches(1),
                                    self.pdf_warning_thresholds["min_text_box_height"] / Inches(1)
                                )
                            }
                        ))
                    
                    # 2. 폰트 크기 PDF 호환성
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size < self.pdf_min_font_size:
                                pdf_issues += 1
                                result.issues.append(ValidationIssue(
                                    severity=IssueSeverity.CRITICAL,
                                    category=IssueCategory.READABILITY,
                                    message=f"PDF 변환 실패 위험: 폰트 크기 {run.font.size.pt}pt는 PDF에서 읽기 어려움",
                                    shape_id=shape_id,
                                    suggested_fix=f"폰트 크기를 최소 {self.pdf_min_font_size.pt}pt 이상으로 조정",
                                    metadata={
                                        "current_font_size": run.font.size.pt,
                                        "min_pdf_font_size": self.pdf_min_font_size.pt
                                    }
                                ))
                                break
                        if pdf_issues > 0:
                            break
                
                # 3. 도형 경계 검증
                left = shape.left / 914400
                top = shape.top / 914400
                right = left + (shape.width / 914400)
                bottom = top + (shape.height / 914400)
                
                # PDF 안전 마진 검사
                if left < self.pdf_safe_margin / Inches(1) or \
                   top < self.pdf_safe_margin / Inches(1) or \
                   right > (self.slide_width - self.pdf_safe_margin) / Inches(1) or \
                   bottom > (self.slide_height - self.pdf_safe_margin) / Inches(1):
                    pdf_issues += 1
                    result.issues.append(ValidationIssue(
                        severity=IssueSeverity.WARNING,
                        category=IssueCategory.MARGIN_ISSUES,
                        message=f"PDF 마진 경고: 요소가 안전 마진({self.pdf_safe_margin / Inches(1):.1f}인치)에 너무 가깝게 위치",
                        shape_id=shape_id,
                        position=(left, top),
                        suggested_fix="PDF 변환 시 잘릴 수 있으니 여백을 늘리세요",
                        metadata={
                            "shape_bounds": (left, top, right, bottom),
                            "safe_margin": self.pdf_safe_margin / Inches(1)
                        }
                    ))
            
            # 메트릭 업데이트
            result.metrics.update({
                "pdf_compatibility_issues": pdf_issues,
                "pdf_ready": pdf_issues == 0
            })
            
        except Exception as e:
            app_logger.error(f"Error in PDF compatibility check: {e}")
            result.issues.append(ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.PERFORMANCE,
                message="PDF 호환성 검증 실패",
                suggested_fix="PDF 변환 전 수동 검토 필요"
            ))
    
    def _check_layout_integrity(self, slide: Slide, result: ValidationResult) -> None:
        """
        레이아웃 무결성 검증
        슬라이드 레이아웃이 McKinsey 표준을 따르는지 확인
        """
        try:
            # 레이아웃 표준 검사
            text_boxes = 0
            images = 0
            charts = 0
            tables = 0
            total_area_used = 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_boxes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    charts += 1
                elif hasattr(shape, 'table'):
                    tables += 1
                
                # 사용 면적 계산
                shape_area = (shape.width / 914400) * (shape.height / 914400)
                total_area_used += shape_area
            
            # 슬라이드 면적
            slide_area = (self.slide_width / Inches(1)) * (self.slide_height / Inches(1))
            usage_ratio = total_area_used / slide_area if slide_area > 0 else 0
            
            # 1. 과밀도 검사
            if usage_ratio > self.pdf_warning_thresholds["max_text_density"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.CONTENT_DENSITY,
                    message=f"레이아웃 과밀: 슬라이드 면적의 {usage_ratio*100:.1f}% 사용 중",
                    suggested_fix="콘텐츠를 여러 슬라이드로 분할하거나 여백을 늘리세요",
                    metadata={
                        "usage_ratio": usage_ratio,
                        "element_counts": {
                            "text_boxes": text_boxes,
                            "images": images,
                            "charts": charts,
                            "tables": tables
                        }
                    }
                ))
            
            # 2. 레이아웃 균형 검사
            if text_boxes > self.performance_thresholds["max_text_boxes_per_slide"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.CONTENT_DENSITY,
                    message=f"텍스트 박스 과다: {text_boxes}개 (최대 권장: {self.performance_thresholds['max_text_boxes_per_slide']}개)",
                    suggested_fix="텍스트를 통합하거나 슬라이드를 분할하세요",
                    metadata={"text_box_count": text_boxes}
                ))
            
            # 3. McKinsey 표준 레이아웃 패턴 검사
            if text_boxes == 0 and charts == 0 and tables == 0:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.EMPTY_CONTENT,
                    message="비표준 레이아웃: 의미 있는 콘텐츠 없음",
                    suggested_fix="McKinsey 표준 레이아웃(title, summary, chart, bullets, comparison) 중 하나를 선택하세요"
                ))
            
            # 메트릭 업데이트
            result.metrics.update({
                "layout_integrity": {
                    "text_boxes": text_boxes,
                    "images": images,
                    "charts": charts,
                    "tables": tables,
                    "area_usage_ratio": usage_ratio,
                    "balanced": text_boxes <= self.performance_thresholds["max_text_boxes_per_slide"] and usage_ratio <= 0.8
                }
            })
            
        except Exception as e:
            app_logger.error(f"Error in layout integrity check: {e}")
            result.issues.append(ValidationIssue(
                severity=IssueSeverity.WARNING,
                category=IssueCategory.PERFORMANCE,
                message="레이아웃 무결성 검증 실패",
                suggested_fix="레이아웃을 수동으로 확인하세요"
            ))
    
    def _check_text_overflow(self, slide: Slide) -> List:
        """
        텍스트가 박스를 넘어가는지 체크
        """
        overflowing = []
        
        try:
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                
                if not shape.text_frame or not shape.text_frame.text:
                    continue
                    
                text_frame = shape.text_frame
                
                # Check if text frame has auto-size set to overflow
                if hasattr(text_frame, 'auto_size'):
                    if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
                        # Text might overflow if it's too long
                        text_length = len(text_frame.text)
                        
                        # Estimate if text fits based on shape size and text length
                        shape_area = (shape.width / 914400) * (shape.height / 914400)  # Convert EMU to inches²
                        
                        # Rough estimate: 20 chars per square inch at 12pt font
                        estimated_capacity = shape_area * 20
                        
                        if text_length > estimated_capacity * 1.5:  # 50% overflow threshold
                            overflowing.append({
                                "shape": shape,
                                "text_length": text_length,
                                "estimated_overflow": text_length - estimated_capacity
                            })
                
                # Check for very long text in small shapes
                if shape.width < Inches(2) and len(text_frame.text) > 100:
                    overflowing.append({
                        "shape": shape,
                        "reason": "Long text in small shape"
                    })
                    
        except Exception as e:
            app_logger.error(f"Error checking text overflow: {str(e)}")
        
        return overflowing
    
    def _check_overlapping_shapes(self, slide: Slide) -> List:
        """
        요소들이 겹치는지 체크
        """
        shapes = []
        overlapping = []
        
        try:
            # Collect all visible shapes
            for shape in slide.shapes:
                # Skip if shape is likely a background or decoration
                if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                    if shape.fill.type == 6:  # Background fill
                        continue
                
                shapes.append(shape)
            
            # Check for overlaps
            for i, shape1 in enumerate(shapes):
                for shape2 in shapes[i+1:]:
                    if self._shapes_overlap(shape1, shape2):
                        overlapping.append((shape1, shape2))
                        
        except Exception as e:
            app_logger.error(f"Error checking overlapping shapes: {str(e)}")
        
        return overlapping
    
    def _shapes_overlap(self, shape1, shape2) -> bool:
        """
        두 도형이 겹치는지 확인
        """
        try:
            # Get bounding boxes
            box1 = (
                shape1.left,
                shape1.top,
                shape1.left + shape1.width,
                shape1.top + shape1.height
            )
            box2 = (
                shape2.left,
                shape2.top,
                shape2.left + shape2.width,
                shape2.top + shape2.height
            )
            
            # Check if boxes overlap
            overlap = not (
                box1[2] < box2[0] or  # shape1 is left of shape2
                box1[0] > box2[2] or  # shape1 is right of shape2
                box1[3] < box2[1] or  # shape1 is above shape2
                box1[1] > box2[3]     # shape1 is below shape2
            )
            
            # Allow small overlaps (less than 0.1 inch)
            if overlap:
                overlap_width = min(box1[2], box2[2]) - max(box1[0], box2[0])
                overlap_height = min(box1[3], box2[3]) - max(box1[1], box2[1])
                
                # Convert to inches (EMU to inches: divide by 914400)
                overlap_area = (overlap_width / 914400) * (overlap_height / 914400)
                
                # Ignore very small overlaps
                if overlap_area < 0.01:  # Less than 0.01 square inches
                    return False
            
            return overlap
            
        except Exception as e:
            app_logger.error(f"Error checking shape overlap: {str(e)}")
            return False
    
    def _check_out_of_bounds(self, slide: Slide) -> List:
        """
        슬라이드 경계를 넘어간 요소 체크
        """
        out_of_bounds = []
        
        try:
            # Get slide dimensions in EMU
            slide_width = self.slide_width.inches * 914400
            slide_height = self.slide_height.inches * 914400
            
            # Allow small margin for edge cases
            margin = Inches(0.1).inches * 914400
            
            for shape in slide.shapes:
                violations = []
                
                # Check right boundary
                if shape.left + shape.width > slide_width + margin:
                    violations.append("right")
                
                # Check bottom boundary
                if shape.top + shape.height > slide_height + margin:
                    violations.append("bottom")
                
                # Check left boundary
                if shape.left < -margin:
                    violations.append("left")
                
                # Check top boundary
                if shape.top < -margin:
                    violations.append("top")
                
                if violations:
                    out_of_bounds.append({
                        "shape": shape,
                        "violations": violations
                    })
                    
        except Exception as e:
            app_logger.error(f"Error checking out of bounds: {str(e)}")
        
        return out_of_bounds
    
    def _check_font_consistency(self, slide: Slide) -> str:
        """
        폰트 일관성 체크
        """
        font_sizes = set()
        font_names = set()
        
        try:
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                
                if not shape.text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.add(run.font.size.pt)
                        if run.font.name:
                            font_names.add(run.font.name)
            
            # Check for too many different font sizes
            if len(font_sizes) > 4:
                return f"{len(font_sizes)}개의 다른 폰트 크기 사용"
            
            # Check for multiple font families
            if len(font_names) > 2:
                return f"{len(font_names)}개의 다른 폰트 사용"
            
            # Check for very small fonts
            if font_sizes and min(font_sizes) < self.recommended_min_font.pt:
                return f"너무 작은 폰트 크기 ({min(font_sizes)}pt)"
            
            return ""
            
        except Exception as e:
            app_logger.error(f"Error checking font consistency: {str(e)}")
            return ""
    
    def _check_margins(self, slide: Slide) -> str:
        """
        여백 체크
        """
        try:
            min_left = None
            min_top = None
            max_right = None
            max_bottom = None
            
            for shape in slide.shapes:
                if min_left is None or shape.left < min_left:
                    min_left = shape.left
                if min_top is None or shape.top < min_top:
                    min_top = shape.top
                    
                right = shape.left + shape.width
                bottom = shape.top + shape.height
                
                if max_right is None or right > max_right:
                    max_right = right
                if max_bottom is None or bottom > max_bottom:
                    max_bottom = bottom
            
            if min_left is not None:
                # Convert EMU to inches
                left_margin = min_left / 914400
                top_margin = min_top / 914400
                right_margin = (self.slide_width.inches * 914400 - max_right) / 914400
                bottom_margin = (self.slide_height.inches * 914400 - max_bottom) / 914400
                
                issues = []
                if left_margin < self.min_margin.inches:
                    issues.append(f"왼쪽 여백 부족 ({left_margin:.2f}\")")
                if top_margin < self.min_margin.inches:
                    issues.append(f"상단 여백 부족 ({top_margin:.2f}\")")
                if right_margin < self.min_margin.inches:
                    issues.append(f"오른쪽 여백 부족 ({right_margin:.2f}\")")
                if bottom_margin < self.min_margin.inches:
                    issues.append(f"하단 여백 부족 ({bottom_margin:.2f}\")")
                
                return ", ".join(issues) if issues else ""
            
            return ""
            
        except Exception as e:
            app_logger.error(f"Error checking margins: {str(e)}")
            return ""
    
    def _check_text_readability(self, slide: Slide) -> List[str]:
        """
        텍스트 가독성 체크
        """
        issues = []
        
        try:
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                
                if not shape.text_frame or not shape.text_frame.text:
                    continue
                
                text_frame = shape.text_frame
                
                # Check for all caps text (harder to read)
                if text_frame.text.isupper() and len(text_frame.text) > 20:
                    issues.append("긴 대문자 텍스트 발견 (가독성 저하)")
                
                # Check for very long paragraphs
                for paragraph in text_frame.paragraphs:
                    if len(paragraph.text) > 200:
                        issues.append("너무 긴 단락 발견 (200자 초과)")
                        break
                
                # Check line count in text frame
                line_count = len(text_frame.text.split('\n'))
                if line_count > 10:
                    issues.append(f"텍스트 박스에 너무 많은 줄 ({line_count}줄)")
                    
        except Exception as e:
            app_logger.error(f"Error checking text readability: {str(e)}")
        
        return issues
    
    def _check_content_density(self, slide: Slide) -> str:
        """
        콘텐츠 밀도 체크
        """
        try:
            total_text_length = 0
            text_box_count = 0
            bullet_count = 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if shape.text_frame.text:
                        total_text_length += len(shape.text_frame.text)
                        text_box_count += 1
                        
                        # Count bullets
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.level >= 0 and paragraph.text.strip():
                                bullet_count += 1
            
            # Check for too dense content
            if total_text_length > 500:
                return f"콘텐츠 과밀 (총 {total_text_length}자)"
            
            if bullet_count > self.max_bullets_per_slide:
                return f"불릿 포인트 과다 ({bullet_count}개)"
            
            if text_box_count > 6:
                return f"텍스트 박스 과다 ({text_box_count}개)"
            
            return ""
            
        except Exception as e:
            app_logger.error(f"Error checking content density: {str(e)}")
            return ""
    
    def _is_empty_slide(self, slide: Slide) -> bool:
        """
        빈 슬라이드인지 확인
        """
        try:
            has_content = False
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if shape.text_frame.text.strip():
                        has_content = True
                        break
                
                # Check for images, charts, tables
                if hasattr(shape, 'image'):
                    has_content = True
                    break
                    
                if hasattr(shape, 'chart'):
                    has_content = True
                    break
                    
                if hasattr(shape, 'table'):
                    has_content = True
                    break
            
            return not has_content
            
        except Exception as e:
            app_logger.error(f"Error checking empty slide: {str(e)}")
            return False
    
    def _generate_suggestions(self, slide: Slide, metrics: Dict) -> List[str]:
        """
        개선 제안 생성
        """
        suggestions = []
        
        try:
            # Based on metrics, provide suggestions
            if metrics.get("text_overflow_count", 0) > 0:
                suggestions.append("텍스트 양을 줄이거나 폰트 크기를 조정하세요")
            
            if metrics.get("overlapping_shapes", 0) > 0:
                suggestions.append("겹치는 요소들을 재배치하세요")
            
            if metrics.get("content_density") == "high":
                suggestions.append("콘텐츠를 여러 슬라이드로 나누는 것을 고려하세요")
            
            if not metrics.get("font_consistency", True):
                suggestions.append("폰트 스타일을 통일하여 일관성을 높이세요")
            
            if metrics.get("readability_issues", 0) > 0:
                suggestions.append("텍스트를 더 짧고 명확하게 작성하세요")
            
            if metrics.get("is_empty", False):
                suggestions.append("슬라이드에 콘텐츠를 추가하세요")
            
            return suggestions
            
        except Exception as e:
            app_logger.error(f"Error generating suggestions: {str(e)}")
            return []
    
    def validate_presentation(self, presentation) -> Dict[str, Any]:
        """
        전체 프레젠테이션 검증
        """
        try:
            total_issues = 0
            total_warnings = 0
            slide_validations = []
            
            for i, slide in enumerate(presentation.slides):
                validation = self.validate_slide(slide)
                validation["slide_number"] = i + 1
                slide_validations.append(validation)
                
                total_issues += len(validation["issues"])
                total_warnings += len(validation["warnings"])
            
            # Overall presentation metrics
            valid_slides = sum(1 for v in slide_validations if v["is_valid"])
            
            return {
                "is_valid": total_issues == 0,
                "total_slides": len(presentation.slides),
                "valid_slides": valid_slides,
                "total_issues": total_issues,
                "total_warnings": total_warnings,
                "slide_validations": slide_validations,
                "validation_rate": valid_slides / len(presentation.slides) if presentation.slides else 0
            }
            
        except Exception as e:
            app_logger.error(f"Error validating presentation: {str(e)}")
            return {
                "is_valid": False,
                "error": str(e)
            }
    
    # ========== Enhanced Validation Methods ==========
    
    def _check_overlapping_shapes_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """정밀한 요소 겹침 검증"""
        try:
            shapes = [shape for shape in slide.shapes if hasattr(shape, 'left')]
            overlap_count = 0
            
            for i, shape1 in enumerate(shapes):
                for j, shape2 in enumerate(shapes[i+1:], i+1):
                    if self._shapes_overlap_precise(shape1, shape2):
                        overlap_count += 1
                        overlap_area = self._calculate_overlap_area(shape1, shape2)
                        
                        # 겹침 심각도 판정
                        shape1_area = (shape1.width * shape1.height) / (914400 * 914400)
                        overlap_percentage = overlap_area / shape1_area * 100
                        
                        severity = IssueSeverity.CRITICAL if overlap_percentage > 20 else IssueSeverity.WARNING
                        
                        result.issues.append(ValidationIssue(
                            severity=severity,
                            category=IssueCategory.SHAPE_OVERLAP,
                            message=f"요소 겹침: {overlap_percentage:.1f}% 겹침 (shape_{i} ↔ shape_{j})",
                            shape_id=f"shape_{i}_shape_{j}",
                            suggested_fix="겹치는 요소들을 재배치하거나 크기를 조정하세요",
                            metadata={
                                "overlap_area_sq_inches": overlap_area,
                                "overlap_percentage": overlap_percentage,
                                "shape1_index": i,
                                "shape2_index": j
                            }
                        ))
            
            result.metrics["overlapping_shapes"] = overlap_count
            
        except Exception as e:
            app_logger.error(f"Error in enhanced shape overlap check: {e}")
    
    def _shapes_overlap_precise(self, shape1, shape2) -> bool:
        """정밀한 겹침 검증"""
        try:
            box1 = (shape1.left, shape1.top, shape1.left + shape1.width, shape1.top + shape1.height)
            box2 = (shape2.left, shape2.top, shape2.left + shape2.width, shape2.top + shape2.height)
            
            # 겹치는지 확인
            if (box1[2] <= box2[0] or box1[0] >= box2[2] or 
                box1[3] <= box2[1] or box1[1] >= box2[3]):
                return False
            
            # 겹침 면적 계산
            overlap_area = self._calculate_overlap_area(shape1, shape2)
            return overlap_area > 0.001  # 0.001 평방인치 이상
            
        except Exception:
            return False
    
    def _calculate_overlap_area(self, shape1, shape2) -> float:
        """겹침 면적 계산 (평방인치)"""
        try:
            left = max(shape1.left, shape2.left)
            top = max(shape1.top, shape2.top)
            right = min(shape1.left + shape1.width, shape2.left + shape2.width)
            bottom = min(shape1.top + shape1.height, shape2.top + shape2.height)
            
            if right > left and bottom > top:
                overlap_width = (right - left) / 914400
                overlap_height = (bottom - top) / 914400
                return overlap_width * overlap_height
            return 0.0
            
        except Exception:
            return 0.0
    
    def _check_out_of_bounds_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """정밀한 경계 초과 검증"""
        try:
            slide_width_emu = self.slide_width.inches * 914400
            slide_height_emu = self.slide_height.inches * 914400
            margin_emu = Inches(0.05).inches * 914400  # 5% 여유
            
            out_of_bounds_count = 0
            
            for idx, shape in enumerate(slide.shapes):
                violations = []
                shape_id = f"shape_{idx}"
                
                # 경계 검사
                if shape.left < -margin_emu:
                    violations.append("left")
                if shape.top < -margin_emu:
                    violations.append("top")
                if shape.left + shape.width > slide_width_emu + margin_emu:
                    violations.append("right")
                if shape.top + shape.height > slide_height_emu + margin_emu:
                    violations.append("bottom")
                
                if violations:
                    out_of_bounds_count += 1
                    
                    # 얼마나 벗어났는지 계산
                    excess_amounts = []
                    if "left" in violations:
                        excess_amounts.append(f"왼쪽으로 {abs(shape.left) / 914400:.2f}인치")
                    if "right" in violations:
                        excess = (shape.left + shape.width - slide_width_emu) / 914400
                        excess_amounts.append(f"오른쪽으로 {excess:.2f}인치")
                    if "top" in violations:
                        excess_amounts.append(f"위로 {abs(shape.top) / 914400:.2f}인치")
                    if "bottom" in violations:
                        excess = (shape.top + shape.height - slide_height_emu) / 914400
                        excess_amounts.append(f"아래로 {excess:.2f}인치")
                    
                    result.issues.append(ValidationIssue(
                        severity=IssueSeverity.CRITICAL,
                        category=IssueCategory.OUT_OF_BOUNDS,
                        message=f"슬라이드 경계 초과: {', '.join(excess_amounts)}",
                        shape_id=shape_id,
                        position=(shape.left / 914400, shape.top / 914400),
                        suggested_fix="요소를 슬라이드 경계 내로 이동시키세요",
                        metadata={
                            "violations": violations,
                            "excess_amounts": excess_amounts
                        }
                    ))
            
            result.metrics["out_of_bounds_count"] = out_of_bounds_count
            
        except Exception as e:
            app_logger.error(f"Error in enhanced out of bounds check: {e}")
    
    def _check_font_consistency_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """고급 폰트 일관성 검증"""
        try:
            font_sizes = set()
            font_names = set()
            font_details = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    continue
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.font.size:
                            font_sizes.add(run.font.size.pt)
                        if run.font.name:
                            font_names.add(run.font.name)
                            
                        font_details.append({
                            "shape": shape_idx,
                            "paragraph": para_idx,
                            "run": run_idx,
                            "font_name": run.font.name,
                            "font_size": run.font.size.pt if run.font.size else None,
                            "text": run.text[:20] + "..." if len(run.text) > 20 else run.text
                        })
            
            # 폰트 다양성 검사
            if len(font_names) > self.mckinsey_style_rules["max_font_variants"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.FONT_CONSISTENCY,
                    message=f"폰트 종류 과다: {len(font_names)}개 (권장: {self.mckinsey_style_rules['max_font_variants']}개 이하)",
                    suggested_fix=f"폰트를 {self.mckinsey_style_rules['preferred_fonts']} 중에서 선택하세요",
                    metadata={"font_names": list(font_names)}
                ))
            
            # 폰트 크기 다양성 검사
            if len(font_sizes) > 4:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.FONT_CONSISTENCY,
                    message=f"폰트 크기 과다: {len(font_sizes)}개 (권장: 4개 이하)",
                    suggested_fix="폰트 크기를 통일하여 일관성을 높이세요",
                    metadata={"font_sizes": sorted(list(font_sizes))}
                ))
            
            # McKinsey 추천 폰트 검사
            non_preferred_fonts = font_names - set(self.mckinsey_style_rules["preferred_fonts"])
            if non_preferred_fonts:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.SUGGESTION,
                    category=IssueCategory.MCKINSEY_STYLE,
                    message=f"비추천 폰트 사용: {list(non_preferred_fonts)}",
                    suggested_fix=f"McKinsey 추천 폰트 사용: {self.mckinsey_style_rules['preferred_fonts']}",
                    metadata={"non_preferred_fonts": list(non_preferred_fonts)}
                ))
            
            # 너무 작은 폰트 검사
            small_fonts = [size for size in font_sizes if size < self.recommended_min_font.pt]
            if small_fonts:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.READABILITY,
                    message=f"너무 작은 폰트: {min(small_fonts)}pt (권장: {self.recommended_min_font.pt}pt 이상)",
                    suggested_fix=f"폰트 크기를 {self.recommended_min_font.pt}pt 이상으로 조정하세요",
                    metadata={"small_font_sizes": small_fonts}
                ))
            
            result.metrics.update({
                "font_consistency": len(font_names) <= self.mckinsey_style_rules["max_font_variants"],
                "font_variety_count": len(font_names),
                "font_size_variety_count": len(font_sizes),
                "font_details": font_details[:10]  # 처음 10개만 저장
            })
            
        except Exception as e:
            app_logger.error(f"Error in enhanced font consistency check: {e}")
    
    def _check_mckinsey_style_compliance(self, slide: Slide, result: ValidationResult) -> None:
        """McKinsey 스타일 가이드 준수 검증"""
        try:
            compliance_score = 100
            violations = []
            
            # 제목 폰트 크기 검사
            title_shapes = self._identify_title_shapes(slide)
            for shape_idx, shape in title_shapes:
                avg_font_size = self._get_average_font_size(shape)
                if avg_font_size and avg_font_size < self.mckinsey_style_rules["title_font_min"].pt:
                    violations.append(f"제목 폰트 크기 부족: {avg_font_size}pt")
                    compliance_score -= 15
            
            # 본문 폰트 크기 검사
            body_shapes = self._identify_body_shapes(slide)
            for shape_idx, shape in body_shapes:
                avg_font_size = self._get_average_font_size(shape)
                min_body, max_body = self.mckinsey_style_rules["body_font_range"]
                if avg_font_size and not (min_body.pt <= avg_font_size <= max_body.pt):
                    violations.append(f"본문 폰트 크기 부적절: {avg_font_size}pt")
                    compliance_score -= 10
            
            # 불릿 포인트 개수 검사
            bullet_count = self._count_bullets(slide)
            if bullet_count > self.max_bullets_per_slide:
                violations.append(f"불릿 포인트 과다: {bullet_count}개")
                compliance_score -= 20
            
            # 슬라이드 복잡도 검사
            shape_count = len([s for s in slide.shapes if hasattr(s, 'text_frame')])
            if shape_count > self.performance_thresholds["max_text_boxes_per_slide"]:
                violations.append(f"텍스트 박스 과다: {shape_count}개")
                compliance_score -= 15
            
            # 결과 기록
            if violations:
                severity = IssueSeverity.WARNING if compliance_score > 70 else IssueSeverity.CRITICAL
                result.issues.append(ValidationIssue(
                    severity=severity,
                    category=IssueCategory.MCKINSEY_STYLE,
                    message=f"McKinsey 스타일 준수율: {compliance_score}%",
                    suggested_fix="McKinsey 스타일 가이드를 참조하여 수정하세요",
                    metadata={
                        "compliance_score": compliance_score,
                        "violations": violations,
                        "bullet_count": bullet_count,
                        "text_box_count": shape_count
                    }
                ))
            
            result.metrics["mckinsey_compliance_score"] = compliance_score
            
        except Exception as e:
            app_logger.error(f"Error in McKinsey style compliance check: {e}")
    
    def _identify_title_shapes(self, slide: Slide) -> List[Tuple[int, Any]]:
        """제목으로 추정되는 요소 식별"""
        title_shapes = []
        try:
            for idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                    # 위치와 크기로 제목 판단
                    top_position = shape.top / 914400
                    font_size = self._get_average_font_size(shape)
                    
                    if top_position < 2.0 and font_size and font_size >= 18:
                        title_shapes.append((idx, shape))
        except Exception:
            pass
        return title_shapes
    
    def _identify_body_shapes(self, slide: Slide) -> List[Tuple[int, Any]]:
        """본문으로 추정되는 요소 식별"""
        body_shapes = []
        try:
            for idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                    top_position = shape.top / 914400
                    if top_position >= 2.0:  # 제목 아래 영역
                        body_shapes.append((idx, shape))
        except Exception:
            pass
        return body_shapes
    
    def _get_average_font_size(self, shape) -> Optional[float]:
        """텍스트 박스의 평균 폰트 크기 계산"""
        try:
            font_sizes = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
            return sum(font_sizes) / len(font_sizes) if font_sizes else None
        except Exception:
            return None
    
    def _count_bullets(self, slide: Slide) -> int:
        """불릿 포인트 개수 계산"""
        try:
            bullet_count = 0
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level >= 0 and paragraph.text.strip():
                            bullet_count += 1
            return bullet_count
        except Exception:
            return 0
    
    def _check_margins_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """고급 여백 검증"""
        try:
            if not slide.shapes:
                return
            
            # 실제 사용 영역 계산
            min_left = min(shape.left for shape in slide.shapes)
            min_top = min(shape.top for shape in slide.shapes)
            max_right = max(shape.left + shape.width for shape in slide.shapes)
            max_bottom = max(shape.top + shape.height for shape in slide.shapes)
            
            # 여백 계산 (인치)
            margins = {
                "left": min_left / 914400,
                "top": min_top / 914400,
                "right": (self.slide_width.inches * 914400 - max_right) / 914400,
                "bottom": (self.slide_height.inches * 914400 - max_bottom) / 914400
            }
            
            # 여백 문제 검사
            margin_issues = []
            for side, margin in margins.items():
                if margin < self.min_margin.inches:
                    severity = IssueSeverity.CRITICAL if margin < self.min_margin.inches * 0.5 else IssueSeverity.WARNING
                    margin_issues.append(f"{side}: {margin:.2f}인치")
                    
                    result.issues.append(ValidationIssue(
                        severity=severity,
                        category=IssueCategory.MARGIN_ISSUES,
                        message=f"{side} 여백 부족: {margin:.2f}인치 (최소: {self.min_margin.inches}인치)",
                        suggested_fix=f"{side} 여백을 {self.recommended_margin.inches}인치 이상으로 조정하세요",
                        metadata={"current_margin": margin, "required_margin": self.min_margin.inches}
                    ))
            
            result.metrics.update({
                "margins": margins,
                "margin_issues_count": len(margin_issues)
            })
            
        except Exception as e:
            app_logger.error(f"Error in enhanced margin check: {e}")
    
    def _check_readability_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """고급 가독성 검증"""
        try:
            readability_issues = 0
            total_text_length = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    continue
                
                text_frame = shape.text_frame
                text = text_frame.text
                total_text_length += len(text)
                
                # 대문자 텍스트 검사
                if text.isupper() and len(text) > 20:
                    readability_issues += 1
                    result.issues.append(ValidationIssue(
                        severity=IssueSeverity.WARNING,
                        category=IssueCategory.READABILITY,
                        message=f"긴 대문자 텍스트: {len(text)}자",
                        shape_id=f"shape_{shape_idx}",
                        suggested_fix="대소문자를 적절히 혼용하세요",
                        metadata={"text_length": len(text)}
                    ))
                
                # 한 줄당 글자 수 검사
                lines = text.split('\n')
                for line_idx, line in enumerate(lines):
                    if len(line) > 60:  # 한 줄 60자 초과
                        readability_issues += 1
                        result.issues.append(ValidationIssue(
                            severity=IssueSeverity.SUGGESTION,
                            category=IssueCategory.READABILITY,
                            message=f"긴 텍스트 줄: {len(line)}자 (권장: 60자 이하)",
                            shape_id=f"shape_{shape_idx}_line_{line_idx}",
                            suggested_fix="텍스트를 여러 줄로 나누세요",
                            metadata={"line_length": len(line), "line_text": line[:30] + "..."}
                        ))
                
                # 줄 수 검사
                if len(lines) > self.max_text_lines_per_box:
                    readability_issues += 1
                    result.issues.append(ValidationIssue(
                        severity=IssueSeverity.WARNING,
                        category=IssueCategory.READABILITY,
                        message=f"텍스트 박스에 줄 수 과다: {len(lines)}줄 (권장: {self.max_text_lines_per_box}줄 이하)",
                        shape_id=f"shape_{shape_idx}",
                        suggested_fix="텍스트를 여러 박스로 나누거나 줄이세요",
                        metadata={"line_count": len(lines)}
                    ))
            
            result.metrics.update({
                "readability_issues_count": readability_issues,
                "total_text_length": total_text_length
            })
            
        except Exception as e:
            app_logger.error(f"Error in enhanced readability check: {e}")
    
    def _check_content_density_enhanced(self, slide: Slide, result: ValidationResult) -> None:
        """고급 콘텐츠 밀도 검증"""
        try:
            metrics = {
                "total_text_length": 0,
                "text_box_count": 0,
                "bullet_count": 0,
                "shape_count": len(slide.shapes),
                "content_density_score": 0
            }
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                    metrics["text_box_count"] += 1
                    metrics["total_text_length"] += len(shape.text_frame.text)
                    
                    # 불릿 포인트 계산
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level >= 0 and paragraph.text.strip():
                            metrics["bullet_count"] += 1
            
            # 밀도 점수 계산 (0-100)
            density_factors = [
                min(metrics["total_text_length"] / 500 * 30, 30),  # 텍스트 길이 (최대 30점)
                min(metrics["bullet_count"] / 5 * 25, 25),  # 불릿 수 (최대 25점)
                min(metrics["text_box_count"] / 8 * 25, 25),  # 텍스트박스 수 (최대 25점)
                min(metrics["shape_count"] / 15 * 20, 20)  # 전체 요소 수 (최대 20점)
            ]
            metrics["content_density_score"] = sum(density_factors)
            
            # 밀도 검사
            if metrics["content_density_score"] > 80:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.CONTENT_DENSITY,
                    message=f"콘텐츠 과밀: 밀도 점수 {metrics['content_density_score']:.0f}/100",
                    suggested_fix="콘텐츠를 여러 슬라이드로 나누거나 텍스트를 줄이세요",
                    metadata=metrics
                ))
            elif metrics["content_density_score"] > 60:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.SUGGESTION,
                    category=IssueCategory.CONTENT_DENSITY,
                    message=f"콘텐츠 밀도 높음: {metrics['content_density_score']:.0f}/100",
                    suggested_fix="가독성 향상을 위해 콘텐츠를 정리하는 것을 고려하세요",
                    metadata=metrics
                ))
            
            result.metrics.update(metrics)
            
        except Exception as e:
            app_logger.error(f"Error in enhanced content density check: {e}")
    
    def _check_performance_metrics(self, slide: Slide, result: ValidationResult) -> None:
        """성능 메트릭 검증"""
        try:
            shape_count = len(slide.shapes)
            text_boxes = len([s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame])
            
            # 성능 임계치 검사
            if shape_count > self.performance_thresholds["max_shapes_per_slide"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.WARNING,
                    category=IssueCategory.PERFORMANCE,
                    message=f"요소 과다: {shape_count}개 (권장: {self.performance_thresholds['max_shapes_per_slide']}개 이하)",
                    suggested_fix="불필요한 요소를 제거하여 성능을 개선하세요",
                    metadata={"shape_count": shape_count}
                ))
            
            if text_boxes > self.performance_thresholds["max_text_boxes_per_slide"]:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.SUGGESTION,
                    category=IssueCategory.PERFORMANCE,
                    message=f"텍스트 박스 과다: {text_boxes}개 (권장: {self.performance_thresholds['max_text_boxes_per_slide']}개 이하)",
                    suggested_fix="텍스트 박스를 통합하거나 줄이세요",
                    metadata={"text_box_count": text_boxes}
                ))
            
            result.metrics.update({
                "performance_shape_count": shape_count,
                "performance_text_boxes": text_boxes
            })
            
        except Exception as e:
            app_logger.error(f"Error in performance metrics check: {e}")
    
    def _check_empty_content(self, slide: Slide, result: ValidationResult) -> None:
        """빈 콘텐츠 검증"""
        try:
            has_meaningful_content = False
            empty_shapes = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if shape.text_frame.text.strip():
                        has_meaningful_content = True
                    else:
                        empty_shapes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.CHART or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_meaningful_content = True
            
            if not has_meaningful_content:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.CRITICAL,
                    category=IssueCategory.EMPTY_CONTENT,
                    message="빈 슬라이드: 의미있는 콘텐츠가 없습니다",
                    suggested_fix="텍스트, 이미지, 차트 등의 콘텐츠를 추가하세요",
                    metadata={"is_empty": True}
                ))
            elif empty_shapes > 0:
                result.issues.append(ValidationIssue(
                    severity=IssueSeverity.SUGGESTION,
                    category=IssueCategory.EMPTY_CONTENT,
                    message=f"빈 텍스트 박스: {empty_shapes}개",
                    suggested_fix="빈 텍스트 박스를 제거하거나 내용을 추가하세요",
                    metadata={"empty_shapes": empty_shapes}
                ))
            
            result.metrics.update({
                "has_meaningful_content": has_meaningful_content,
                "empty_shapes_count": empty_shapes
            })
            
        except Exception as e:
            app_logger.error(f"Error in empty content check: {e}")