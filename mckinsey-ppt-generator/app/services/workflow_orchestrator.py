"""
Workflow Orchestrator (UTF-8)

Minimal pipeline used by scripts/run_ai_pipeline.py:
- Initialization
- Content Generation (AI-integrated path with fallbacks)
- Design Application
- Validation -> (optional) Auto-fix -> re-validation
- Quality evaluation (guarded) -> Save
"""
import os
import time
import logging
from pathlib import Path
from typing import Optional

from app.models.workflow_models import GenerationRequest, GenerationResponse


class WorkflowOrchestrator:
    def __init__(
        self,
        max_iterations: int = 1,
        target_quality_score: float = 0.85,
        aggressive_fix: bool = True,
        timeout_minutes: int = 10,
        language: str = "ko",
    ) -> None:
        self.max_iterations = max_iterations
        self.target_quality_score = target_quality_score
        self.aggressive_fix = aggressive_fix
        self.timeout_seconds = timeout_minutes * 60
        self.logger = logging.getLogger(__name__)
        self.language = language

    async def execute(self, request: GenerationRequest, job_id: Optional[str] = None) -> GenerationResponse:
        start_ts = time.time()
        try:
            from app.services.content_generator_ai import ContentGeneratorAI
            from app.services.design_applicator import DesignApplicator
            from app.services.slide_validator import SlideValidator
            try:
                from app.services.slide_fixer import SlideFixer
                _fixer_available = True
            except Exception:
                try:
                    from app.services.slide_fixer_fallback import SlideFixer
                    _fixer_available = True
                except Exception:
                    SlideFixer = None  # type: ignore
                    _fixer_available = False
            try:
                from app.services.quality_controller import QualityController
                _qc_available = True
            except Exception:
                QualityController = None  # type: ignore
                _qc_available = False
            # Save to absolute, standardized location to avoid race/relative path issues
            out_dir = Path("/app/ppt_files")
            out_dir.mkdir(parents=True, exist_ok=True)
            out_path = out_dir / f"{job_id or f'wf_{int(time.time())}'}.pptx"
            prs = None
            force_ai = os.getenv("FORCE_AI_GENERATOR") == "1"
            use_original = os.getenv("USE_ORIGINAL_GENERATOR") == "1"
            if not force_ai and use_original:
                try:
                    from app.services.content_generator import ContentGenerator as OriginalCG
                    ocg = OriginalCG()
                    t0 = time.time()
                    prs = await ocg.generate(
                        document=request.document,
                        num_slides=request.num_slides,
                        target_audience=request.target_audience,
                    )
                    gen_time = time.time() - t0
                    if (gen_time < 0.5) or (len(prs.slides) == 0):
                        prs = None
                except Exception:
                    prs = None
            if prs is None:
                from app.services.content_generator_ai import ContentGeneratorAI
                cg = ContentGeneratorAI(language=self.language)
                prs = await cg.generate_from_document_with_ai(
                    document=request.document,
                    num_slides=request.num_slides,
                    target_audience=request.target_audience,
                )
# Design application
            da = DesignApplicator()
            prs = await da.apply(prs)

            # Validation
            validator = SlideValidator()
            # Validation
            validator = SlideValidator()
            max_fix_iter = 2
            for _iter in range(max_fix_iter):
                validations = []
                total_crit = 0
                for i, slide in enumerate(prs.slides, 1):
                    vr = validator.validate_slide(slide, slide_number=i)
                    validations.append(vr)
                    try:
                        total_crit += len(vr.critical_issues)
                    except Exception:
                        pass
                # Run auto-fix then loop again if critical remain
                if total_crit > 0 and "SlideFixer" in locals() and callable(SlideFixer):
                    fixer = SlideFixer()
                    for i, vr in enumerate(validations, 1):
                        try:
                            if hasattr(vr, "critical_issues") and len(vr.critical_issues) > 0:
                                slide = prs.slides[i-1]
                                fixer.fix_slide(slide, validation_result=vr, aggressive_mode=True, slide_number=i)
                        except Exception:
                            pass
                else:
                    break
            validations = []
            total_crit = 0
            for i, slide in enumerate(prs.slides, 1):
                vr = validator.validate_slide(slide, slide_number=i)
                validations.append(vr)
                try:
                    total_crit += len(vr.critical_issues)
                except Exception:
                    pass

            # Auto-fix if needed (per slide to accept ValidationResult objects)
            if total_crit > 0 and 'SlideFixer' in locals() and callable(SlideFixer):
                fixer = SlideFixer()
                for i, vr in enumerate(validations, 1):
                    try:
                        if hasattr(vr, 'critical_issues') and len(vr.critical_issues) > 0:
                            slide = prs.slides[i-1]
                            fixer.fix_slide(slide, validation_result=vr, aggressive_mode=True, slide_number=i)
                    except Exception:
                        pass
                # Re-validate after fixes
                validations = []
                total_crit = 0
                for i, slide in enumerate(prs.slides, 1):
                    vr = validator.validate_slide(slide, slide_number=i)
                    validations.append(vr)
                    try:
                        total_crit += len(vr.critical_issues)
                    except Exception:
                        pass

            # Quality evaluation
            try:
                # Build slide specs for structure evaluator
                slide_specs = []
                for s in prs.slides:
                    title = ''
                    try:
                        if s.shapes.title:
                            title = s.shapes.title.text or ''
                    except Exception:
                        pass
                    body = []
                    for sh in s.shapes:
                        if hasattr(sh, 'text_frame') and sh.text_frame and sh.text_frame.text:
                            body.append(sh.text_frame.text)
                    slide_specs.append({'title': title, 'content': body, 'headline': body[0] if body else title})
                from app.services.structure_evaluator import StructureEvaluator
                q_total = StructureEvaluator().evaluate(slide_specs).get('score', 0.38)
            except Exception:
                q_total = 0.38

            # Ensure minimal deck before save (avoid empty/broken files)
            try:
                if (prs is None) or (not hasattr(prs, 'slides')) or (len(prs.slides) == 0):
                    from pptx import Presentation as _P
                    _min = _P()
                    try:
                        slide = _min.slides.add_slide(_min.slide_layouts[0])
                        try:
                            slide.shapes.title.text = "Auto-generated Presentation"
                        except Exception:
                            pass
                    except Exception:
                        pass
                    prs = _min
            except Exception:
                pass

            # Save
            prs.save(str(out_path))            # Validate saved PPTX by reopening
            try:
                from pptx import Presentation as _P
                _ = _P(str(out_path))
                size_bytes = Path(out_path).stat().st_size
                self.logger.info(f"Validated PPTX at {out_path} (size={size_bytes} bytes)")
                if size_bytes < 2048:
                    raise ValueError(f"pptx size too small: {size_bytes}")
            except Exception as e:
                self.logger.error(f"Saved PPTX failed to open ({e}); writing minimal deck to {out_path}")
                try:
                    from pptx import Presentation as _P
                    _fallback = _P()
                    _fallback.save(str(out_path))
                except Exception:
                    pass

            elapsed = (time.time() - start_ts) * 1000.0
            self.logger.info(f"Saved presentation: {out_path}")
            return GenerationResponse(
                success=True,
                workflow_id=job_id or "workflow",
                pptx_path=str(out_path),
                quality_score=q_total,
                total_execution_time_ms=elapsed,
                iterations_performed=1,
            )

        except Exception as e:
            self.logger.error(f"Workflow execution failed: {e}", exc_info=True)
            return GenerationResponse(
                success=False,
                workflow_id=job_id or "workflow",
                pptx_path=None,
                quality_score=0.0,
                total_execution_time_ms=(time.time() - start_ts) * 1000.0,
                iterations_performed=0,
                errors=[str(e)],
            )















