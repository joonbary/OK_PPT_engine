"""
McKinsey Style Validator
스타일 규정 준수 검증 모듈
"""

import logging
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class StyleValidator:
    """McKinsey 스타일 검증기"""
    
    # McKinsey 허용 색상
    ALLOWED_COLORS = [
        RGBColor(0, 118, 168),       # McKinsey Blue
        RGBColor(244, 118, 33),      # Orange
        RGBColor(83, 86, 90),        # Gray
        RGBColor(217, 217, 217),     # Light Gray
        RGBColor(255, 255, 255),     # White
        RGBColor(0, 0, 0),           # Black
        RGBColor(107, 166, 68),      # Green (accent)
        RGBColor(227, 27, 35),       # Red (warning)
        RGBColor(128, 130, 133),     # Light Gray variant
    ]
    
    # 허용 폰트
    ALLOWED_FONTS = ['Arial', 'Calibri']  # Calibri는 fallback
    
    async def is_mckinsey_compliant(self, slide) -> bool:
        """
        슬라이드가 McKinsey 스타일 규정을 준수하는지 검증
        
        Returns:
            bool: 규정 준수 여부
        """
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            # 폰트 검증
                            if paragraph.font.name not in self.ALLOWED_FONTS:
                                logger.warning(f"Non-compliant font: {paragraph.font.name}")
                                return False
                            
                            # 색상 검증
                            if paragraph.font.color and paragraph.font.color.rgb:
                                color_rgb = paragraph.font.color.rgb
                                if not self._is_allowed_color(color_rgb):
                                    logger.warning(f"Non-compliant color detected")
                                    return False
            
            return True
            
        except Exception as e:
            logger.debug(f"Validation error: {e}")
            return False
    
    def _is_allowed_color(self, color_rgb) -> bool:
        """색상이 허용된 McKinsey 색상인지 확인"""
        for allowed_color in self.ALLOWED_COLORS:
            if self._colors_match(color_rgb, allowed_color):
                return True
        return False
    
    def _colors_match(self, color1, color2) -> bool:
        """두 색상이 일치하는지 확인"""
        try:
            # RGBColor 객체 비교
            if hasattr(color1, 'rgb') and hasattr(color2, 'rgb'):
                return color1.rgb == color2.rgb
            elif hasattr(color1, '__iter__') and hasattr(color2, '__iter__'):
                return tuple(color1) == tuple(color2)
            else:
                return color1 == color2
        except:
            return False