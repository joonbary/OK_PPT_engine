"""
McKinsey-specific styles and formatting for PowerPoint presentations.

This module defines the visual identity and styling guidelines used by McKinsey & Company
for professional presentations, including colors, fonts, layouts, and design patterns.
"""

from typing import Dict, List, Tuple, Optional, Any
from enum import Enum
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE


class McKinseyColors:
    """McKinsey brand colors and professional color palette"""
    
    # Primary McKinsey Colors
    PRIMARY_BLUE = RGBColor(0, 40, 86)      # #002856 - Deep McKinsey Blue
    SECONDARY_BLUE = RGBColor(0, 119, 200)  # #0077C8 - Bright Blue
    DARK_BLUE = RGBColor(0, 32, 91)         # Dark Blue
    LIGHT_BLUE = RGBColor(0, 162, 225)      # #00A2E1 - Light Blue
    
    # Supporting Colors
    LIGHT_GRAY = RGBColor(242, 242, 242)
    MEDIUM_GRAY = RGBColor(128, 128, 128)
    DARK_GRAY = RGBColor(64, 64, 64)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    
    # Accent Colors for Charts and Highlights
    GREEN = RGBColor(0, 164, 153)     # #00A499 - Teal Green
    RED = RGBColor(220, 38, 38)       # #DC2626 - Emphasis Red
    ORANGE = RGBColor(255, 147, 0)    # #FF9300 - McKinsey Orange
    PURPLE = RGBColor(102, 45, 145)   # Purple
    TEAL = RGBColor(0, 164, 153)      # Teal
    
    # Chart Color Palette
    CHART_COLORS = [
        PRIMARY_BLUE,
        GREEN,
        ORANGE,
        RED,
        PURPLE,
        TEAL,
        SECONDARY_BLUE,
        MEDIUM_GRAY
    ]
    
    @classmethod
    def get_chart_color(cls, index: int) -> RGBColor:
        """Get chart color by index with cycling"""
        return cls.CHART_COLORS[index % len(cls.CHART_COLORS)]


class McKinseyFonts:
    """McKinsey font specifications and text formatting"""
    
    # Primary Fonts
    PRIMARY_FONT = "Calibri"
    SECONDARY_FONT = "Arial"
    MONOSPACE_FONT = "Consolas"
    
    # Font Sizes - Professional McKinsey Standards
    TITLE_SIZE = Pt(36)
    SUBTITLE_SIZE = Pt(24)
    HEADING_SIZE = Pt(28)
    SUBHEADING_SIZE = Pt(20)
    BODY_SIZE = Pt(14)
    SMALL_SIZE = Pt(12)
    FOOTNOTE_SIZE = Pt(9)
    
    # Font Weights and Styles
    BOLD = True
    ITALIC = True
    
    @classmethod
    def get_title_format(cls) -> Dict[str, Any]:
        """Get title text formatting"""
        return {
            'font_name': cls.PRIMARY_FONT,
            'font_size': cls.TITLE_SIZE,
            'bold': cls.BOLD,
            'color': McKinseyColors.PRIMARY_BLUE
        }
    
    @classmethod
    def get_subtitle_format(cls) -> Dict[str, Any]:
        """Get subtitle text formatting"""
        return {
            'font_name': cls.PRIMARY_FONT,
            'font_size': cls.SUBTITLE_SIZE,
            'bold': False,
            'color': McKinseyColors.DARK_GRAY
        }
    
    @classmethod
    def get_body_format(cls) -> Dict[str, Any]:
        """Get body text formatting"""
        return {
            'font_name': cls.PRIMARY_FONT,
            'font_size': cls.BODY_SIZE,
            'bold': False,
            'color': McKinseyColors.BLACK
        }
    
    @classmethod
    def get_heading_format(cls) -> Dict[str, Any]:
        """Get heading text formatting"""
        return {
            'font_name': cls.PRIMARY_FONT,
            'font_size': cls.HEADING_SIZE,
            'bold': cls.BOLD,
            'color': McKinseyColors.PRIMARY_BLUE
        }


class McKinseyLayouts:
    """McKinsey slide layout specifications and dimensions"""
    
    # Standard Slide Dimensions (16:9)
    SLIDE_WIDTH = Inches(13.33)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Margins and Spacing
    MARGIN_LEFT = Inches(0.75)
    MARGIN_RIGHT = Inches(0.75)
    MARGIN_TOP = Inches(1.0)
    MARGIN_BOTTOM = Inches(0.5)
    
    # Content Areas
    CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM
    
    # Title Area
    TITLE_HEIGHT = Inches(1.2)
    TITLE_TOP = MARGIN_TOP
    
    # Main Content Area
    CONTENT_TOP = TITLE_TOP + TITLE_HEIGHT + Inches(0.3)
    CONTENT_MAIN_HEIGHT = CONTENT_HEIGHT - TITLE_HEIGHT - Inches(0.3)
    
    # Two-Column Layout
    COLUMN_GAP = Inches(0.5)
    COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2
    
    # Chart and Image Dimensions
    CHART_HEIGHT = Inches(4.5)
    CHART_WIDTH = Inches(6.0)
    
    @classmethod
    def get_title_position(cls) -> Tuple[float, float, float, float]:
        """Get title text box position (left, top, width, height)"""
        return (
            cls.MARGIN_LEFT,
            cls.TITLE_TOP,
            cls.CONTENT_WIDTH,
            cls.TITLE_HEIGHT
        )
    
    @classmethod
    def get_content_position(cls) -> Tuple[float, float, float, float]:
        """Get main content area position"""
        return (
            cls.MARGIN_LEFT,
            cls.CONTENT_TOP,
            cls.CONTENT_WIDTH,
            cls.CONTENT_MAIN_HEIGHT
        )
    
    @classmethod
    def get_left_column_position(cls) -> Tuple[float, float, float, float]:
        """Get left column position for two-column layout"""
        return (
            cls.MARGIN_LEFT,
            cls.CONTENT_TOP,
            cls.COLUMN_WIDTH,
            cls.CONTENT_MAIN_HEIGHT
        )
    
    @classmethod
    def get_right_column_position(cls) -> Tuple[float, float, float, float]:
        """Get right column position for two-column layout"""
        return (
            cls.MARGIN_LEFT + cls.COLUMN_WIDTH + cls.COLUMN_GAP,
            cls.CONTENT_TOP,
            cls.COLUMN_WIDTH,
            cls.CONTENT_MAIN_HEIGHT
        )


class SlideLayoutType(Enum):
    """Enumeration of available slide layout types"""
    TITLE = "title"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    IMAGE = "image"
    CHART = "chart"
    TABLE = "table"
    BLANK = "blank"
    SECTION_HEADER = "section_header"
    COMPARISON = "comparison"
    TIMELINE = "timeline"


class McKinseyStyles:
    """Main class for applying McKinsey styles to PowerPoint presentations"""
    
    def __init__(self):
        self.colors = McKinseyColors()
        self.fonts = McKinseyFonts()
        self.layouts = McKinseyLayouts()
    
    def apply_text_formatting(self, text_frame, format_type: str = "body") -> None:
        """Apply McKinsey text formatting to a text frame"""
        try:
            if not text_frame or not text_frame.text:
                return
            
            # Get format specifications
            if format_type == "title":
                format_spec = self.fonts.get_title_format()
            elif format_type == "subtitle":
                format_spec = self.fonts.get_subtitle_format()
            elif format_type == "heading":
                format_spec = self.fonts.get_heading_format()
            else:
                format_spec = self.fonts.get_body_format()
            
            # Apply formatting to all paragraphs
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                
                for run in paragraph.runs:
                    font = run.font
                    font.name = format_spec['font_name']
                    font.size = format_spec['font_size']
                    font.bold = format_spec.get('bold', False)
                    font.italic = format_spec.get('italic', False)
                    font.color.rgb = format_spec['color']
                    
        except Exception as e:
            print(f"Error applying text formatting: {e}")
    
    def create_bullet_points(self, text_frame, bullet_points: List[str], 
                           level: int = 0) -> None:
        """Create formatted bullet points with overflow prevention"""
        try:
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Limit number of bullet points to prevent overflow
            max_points = 7
            processed_points = []
            
            for point in bullet_points[:max_points]:
                # Truncate long text for Korean/English mixed content
                if len(point) > 80:
                    point = point[:77] + "..."
                processed_points.append(point)
            
            for i, point in enumerate(processed_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = point
                p.level = level
                p.space_after = Pt(6)  # Add spacing between points
                p.line_spacing = 1.2   # Better line spacing for readability
                
                # Format the paragraph
                for run in p.runs:
                    font = run.font
                    font.name = "맑은 고딕"  # Better Korean font support
                    font.size = Pt(11)  # Slightly smaller for better fit
                    font.color.rgb = self.fonts.get_body_format()['color']
                    
        except Exception as e:
            print(f"Error creating bullet points: {e}")
    
    def apply_shape_style(self, shape, style_type: str = "default") -> None:
        """Apply McKinsey styling to shapes"""
        try:
            if style_type == "primary":
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.colors.PRIMARY_BLUE
                shape.line.color.rgb = self.colors.PRIMARY_BLUE
            elif style_type == "secondary":
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.colors.LIGHT_GRAY
                shape.line.color.rgb = self.colors.MEDIUM_GRAY
            elif style_type == "accent":
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.colors.GREEN
                shape.line.color.rgb = self.colors.GREEN
            else:
                # Default style
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.colors.WHITE
                shape.line.color.rgb = self.colors.DARK_GRAY
                
        except Exception as e:
            print(f"Error applying shape style: {e}")
    
    def get_layout_config(self, layout_type: SlideLayoutType) -> Dict[str, Any]:
        """Get layout configuration for a specific slide type"""
        configs = {
            SlideLayoutType.TITLE: {
                "title_position": self.layouts.get_title_position(),
                "subtitle_position": (
                    self.layouts.MARGIN_LEFT,
                    self.layouts.TITLE_TOP + self.layouts.TITLE_HEIGHT + Inches(0.5),
                    self.layouts.CONTENT_WIDTH,
                    Inches(1.0)
                ),
                "background_color": self.colors.WHITE
            },
            SlideLayoutType.CONTENT: {
                "title_position": self.layouts.get_title_position(),
                "content_position": self.layouts.get_content_position(),
                "background_color": self.colors.WHITE
            },
            SlideLayoutType.TWO_COLUMN: {
                "title_position": self.layouts.get_title_position(),
                "left_column_position": self.layouts.get_left_column_position(),
                "right_column_position": self.layouts.get_right_column_position(),
                "background_color": self.colors.WHITE
            },
            SlideLayoutType.CHART: {
                "title_position": self.layouts.get_title_position(),
                "chart_position": (
                    self.layouts.MARGIN_LEFT + Inches(1.0),
                    self.layouts.CONTENT_TOP,
                    self.layouts.CHART_WIDTH,
                    self.layouts.CHART_HEIGHT
                ),
                "background_color": self.colors.WHITE
            },
            SlideLayoutType.SECTION_HEADER: {
                "title_position": (
                    self.layouts.MARGIN_LEFT,
                    Inches(2.5),
                    self.layouts.CONTENT_WIDTH,
                    Inches(2.0)
                ),
                "background_color": self.colors.PRIMARY_BLUE
            }
        }
        
        return configs.get(layout_type, configs[SlideLayoutType.CONTENT])
    
    def create_header_footer(self, slide, header_text: str = "", 
                           footer_text: str = "", page_number: bool = True) -> None:
        """Add header and footer to slide"""
        try:
            # Add footer
            if footer_text or page_number:
                footer_height = Inches(0.3)
                footer_top = self.layouts.SLIDE_HEIGHT - footer_height - Inches(0.1)
                
                footer_box = slide.shapes.add_textbox(
                    self.layouts.MARGIN_LEFT,
                    footer_top,
                    self.layouts.CONTENT_WIDTH,
                    footer_height
                )
                
                footer_frame = footer_box.text_frame
                footer_frame.margin_left = 0
                footer_frame.margin_right = 0
                footer_frame.margin_top = 0
                footer_frame.margin_bottom = 0
                
                p = footer_frame.paragraphs[0]
                p.text = footer_text
                p.alignment = PP_ALIGN.LEFT
                
                # Format footer text
                for run in p.runs:
                    font = run.font
                    font.name = self.fonts.PRIMARY_FONT
                    font.size = self.fonts.FOOTNOTE_SIZE
                    font.color.rgb = self.colors.MEDIUM_GRAY
                    
        except Exception as e:
            print(f"Error creating header/footer: {e}")
    
    def get_chart_style_config(self) -> Dict[str, Any]:
        """Get chart styling configuration"""
        return {
            "colors": [color for color in self.colors.CHART_COLORS],
            "font_name": self.fonts.PRIMARY_FONT,
            "title_size": self.fonts.HEADING_SIZE,
            "label_size": self.fonts.SMALL_SIZE,
            "background_color": self.colors.WHITE,
            "grid_color": self.colors.LIGHT_GRAY,
            "border_color": self.colors.MEDIUM_GRAY
        }
    
    def get_table_style_config(self) -> Dict[str, Any]:
        """Get table styling configuration"""
        return {
            "header_background": self.colors.PRIMARY_BLUE,
            "header_text_color": self.colors.WHITE,
            "row_background_1": self.colors.WHITE,
            "row_background_2": self.colors.LIGHT_GRAY,
            "border_color": self.colors.MEDIUM_GRAY,
            "text_color": self.colors.BLACK,
            "font_name": self.fonts.PRIMARY_FONT,
            "font_size": self.fonts.SMALL_SIZE
        }


def create_mckinsey_template() -> Dict[str, Any]:
    """Create a complete McKinsey template configuration"""
    styles = McKinseyStyles()
    
    return {
        "name": "McKinsey Professional",
        "description": "Official McKinsey & Company presentation template",
        "colors": {
            "primary": styles.colors.PRIMARY_BLUE,
            "secondary": styles.colors.SECONDARY_BLUE,
            "accent": styles.colors.GREEN,
            "background": styles.colors.WHITE,
            "text": styles.colors.BLACK
        },
        "fonts": {
            "primary": styles.fonts.PRIMARY_FONT,
            "secondary": styles.fonts.SECONDARY_FONT,
            "title_size": styles.fonts.TITLE_SIZE,
            "body_size": styles.fonts.BODY_SIZE
        },
        "layouts": {
            layout_type.value: styles.get_layout_config(layout_type)
            for layout_type in SlideLayoutType
        },
        "chart_style": styles.get_chart_style_config(),
        "table_style": styles.get_table_style_config()
    }