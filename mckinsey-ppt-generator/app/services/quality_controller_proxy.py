"""
QualityController Proxy (UTF-8)

Provides a QualityController-compatible interface that computes the
structure-related quality score using StructureEvaluator (MECE/Flow/Pyramid).
This helps integrate with code paths expecting a QualityController.
"""

from typing import Optional
from dataclasses import dataclass
from pptx import Presentation

from app.services.structure_evaluator import StructureEvaluator


@dataclass
class QualityScore:
    clarity: float = 0.0
    insight: float = 0.0
    structure: float = 0.0
    visual: float = 0.0
    actionability: float = 0.0
    total: float = 0.0
    passed: bool = False
    details: dict = None


class QualityController:
    def __init__(self, target_score: float = 0.85) -> None:
        self.target = target_score

    def evaluate(self, prs: Presentation) -> QualityScore:
        # Build slide specs
        slide_specs = []
        for s in prs.slides:
            title = ''
            try:
                if s.shapes.title:
                    title = s.shapes.title.text or ''
            except Exception:
                pass
            body = []
            for sh in s.shapes:
                if hasattr(sh, 'text_frame') and sh.text_frame and sh.text_frame.text:
                    body.append(sh.text_frame.text)
            slide_specs.append({'title': title, 'content': body, 'headline': body[0] if body else title})

        result = StructureEvaluator().evaluate(slide_specs)
        total = result.get('score', 0.0)
        qs = QualityScore(
            clarity=0.0,
            insight=0.0,
            structure=total,
            visual=0.0,
            actionability=0.0,
            total=total,
            passed=total >= self.target,
            details=result,
        )
        return qs
