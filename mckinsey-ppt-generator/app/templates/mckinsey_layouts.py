"""
McKinsey Layout Template System
Fix #3: 표준화된 레이아웃 템플릿
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, Any, List, Optional, Tuple
import logging

logger = logging.getLogger(__name__)


class McKinseyLayoutManager:
    """McKinsey 표준 레이아웃 관리자"""
    
    # McKinsey 색상 팔레트
    COLORS = {
        'primary': RGBColor(0, 118, 168),      # McKinsey Blue
        'secondary': RGBColor(244, 118, 33),   # Orange
        'text': RGBColor(83, 86, 90),          # Dark Gray
        'light_gray': RGBColor(217, 217, 217), # Light Gray
        'white': RGBColor(255, 255, 255),      # White
    }
    
    # 표준 여백
    MARGINS = {
        'top': Inches(0.5),
        'bottom': Inches(0.5),
        'left': Inches(0.5),
        'right': Inches(0.5),
        'spacing': Inches(0.3)
    }
    
    @staticmethod
    def apply_layout(slide, layout_type: str, content: Dict[str, Any]) -> Any:
        """
        레이아웃 타입에 따라 콘텐츠 배치
        
        Layout Types:
        - title_slide: 제목 슬라이드
        - content_slide: 일반 콘텐츠
        - two_column: 2단 레이아웃
        - three_column: 3단 레이아웃
        - chart_slide: 차트 + 인사이트
        - bullet_slide: 불릿 포인트
        - matrix_slide: 2x2 매트릭스
        """
        
        logger.info(f"Applying layout: {layout_type}")
        
        try:
            # Normalize aliases to supported internal ids
            aliases = {
                'title_and_content': 'content_slide',
                'content': 'content_slide',
                'bullet_points': 'bullet_slide',
                'matrix': 'matrix_slide',
                'comparison': 'three_column',
                'split_text_chart': 'chart_slide',
                'dual_header': 'content_slide',
                'conclusion_slide': 'content_slide',
                'image_slide': 'image',
            }
            lt = aliases.get((layout_type or '').strip(), (layout_type or '').strip())
            if lt == 'title_slide':
                return McKinseyLayoutManager._layout_title_slide(slide, content)
            elif lt == 'chart_slide':
                return McKinseyLayoutManager._layout_chart_slide(slide, content)
            elif lt == 'two_column':
                return McKinseyLayoutManager._layout_two_column(slide, content)
            elif lt == 'three_column':
                return McKinseyLayoutManager._layout_three_column(slide, content)
            elif lt == 'matrix_slide':
                return McKinseyLayoutManager._layout_matrix_slide(slide, content)
            elif lt == 'bullet_slide':
                return McKinseyLayoutManager._layout_bullet_slide(slide, content)
            elif lt == 'image':
                return McKinseyLayoutManager._layout_content_slide(slide, content)
            else:
                return McKinseyLayoutManager._layout_content_slide(slide, content)
        except Exception as e:
            logger.error(f"Layout application failed: {e}")
            # Fallback to content slide
            return McKinseyLayoutManager._layout_content_slide(slide, content)
    
    @staticmethod
    def _layout_title_slide(slide, content: Dict[str, Any]):
        """제목 슬라이드 레이아웃"""
        # 제목 (중앙)
        title_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2.8),
            Inches(11.33),  # 16:9 기준
            Inches(1.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = content.get('title', 'Presentation Title')
        
        # 제목 스타일
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['primary']
        
        # 부제 (중앙 하단)
        if content.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(4.5),
                Inches(11.33),
                Inches(1)
            )
            sub_frame = subtitle_box.text_frame
            sub_frame.text = content['subtitle']
            
            for paragraph in sub_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 날짜/작성자 (하단)
        if content.get('date') or content.get('author'):
            footer_text = []
            if content.get('date'):
                footer_text.append(content['date'])
            if content.get('author'):
                footer_text.append(content['author'])
            
            footer_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(6.5),
                Inches(11.33),
                Inches(0.5)
            )
            footer_frame = footer_box.text_frame
            footer_frame.text = ' | '.join(footer_text)
            
            for paragraph in footer_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['light_gray']
        
        return slide
    
    @staticmethod
    def _layout_chart_slide(slide, content: Dict[str, Any]) -> Tuple[Any, Dict]:
        """차트 슬라이드 레이아웃 (차트 70% + 인사이트 30%)"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),  # 전체 너비에서 좌우 여백 제외
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Analysis')
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 차트 영역 (좌측 70%)
        chart_area = {
            'left': Inches(0.5),
            'top': Inches(1.5),
            'width': Inches(8.5),
            'height': Inches(5.0)
        }
        
        # 인사이트 영역 (우측 30%)
        insight_box = slide.shapes.add_textbox(
            Inches(9.2),
            Inches(1.5),
            Inches(3.6),
            Inches(5.0)
        )
        insight_frame = insight_box.text_frame
        insight_frame.word_wrap = True
        
        # 인사이트 제목
        p = insight_frame.paragraphs[0]
        p.text = "Key Insights"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = McKinseyLayoutManager.COLORS['primary']
        
        # 인사이트 내용
        insights = content.get('insights', [])
        if isinstance(insights, str):
            insights = [insights]
        
        for insight in insights[:3]:  # 최대 3개 인사이트
            p = insight_frame.add_paragraph()
            p.text = f"• {insight}"
            p.font.name = 'Arial'
            p.font.size = Pt(11)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
            p.space_after = Pt(6)
        
        return slide, chart_area
    
    @staticmethod
    def _layout_two_column(slide, content: Dict[str, Any]):
        """2단 레이아웃"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '')
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 좌측 컬럼
        left_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(5.9),
            Inches(5.0)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        # 좌측 제목 (optional)
        if content.get('left_title'):
            p = left_frame.paragraphs[0]
            p.text = content['left_title']
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['primary']
            left_frame.add_paragraph()
        
        # 좌측 내용
        left_content = content.get('left_content', '')
        if isinstance(left_content, list):
            for item in left_content:
                p = left_frame.add_paragraph()
                p.text = f"• {item}"
                try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
                p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        else:
            p = left_frame.paragraphs[-1] if left_frame.paragraphs else left_frame.add_paragraph()
            p.text = left_content
            try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 우측 컬럼
        right_box = slide.shapes.add_textbox(
            Inches(6.9),
            Inches(1.5),
            Inches(5.9),
            Inches(5.0)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        # 우측 제목 (optional)
        if content.get('right_title'):
            p = right_frame.paragraphs[0]
            p.text = content['right_title']
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['primary']
            right_frame.add_paragraph()
        
        # 우측 내용
        right_content = content.get('right_content', '')
        if isinstance(right_content, list):
            for item in right_content:
                p = right_frame.add_paragraph()
                p.text = f"• {item}"
                try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
                p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        else:
            p = right_frame.paragraphs[-1] if right_frame.paragraphs else right_frame.add_paragraph()
            p.text = right_content
            try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        return slide

    @staticmethod
    def _layout_three_column(slide, content: Dict[str, Any]):
        """3단 레이아웃"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '')
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']

        # 세 컬럼 영역
        lefts = [0.5, 4.6, 8.7]
        width = 3.6
        col_keys = content.get('columns') or []
        bullets = content.get('bullets') or content.get('body') or []
        # 분배
        if not col_keys and isinstance(bullets, list):
            # 균등 분배
            n = len(bullets)
            chunk = max(1, (n + 2) // 3)
            col_keys = [bullets[0:chunk], bullets[chunk:2*chunk], bullets[2*chunk:]]

        for idx in range(3):
            box = slide.shapes.add_textbox(
                Inches(lefts[idx]),
                Inches(1.5),
                Inches(width),
                Inches(5.0)
            )
            tf = box.text_frame
            tf.word_wrap = True
            items = col_keys[idx] if idx < len(col_keys) else []
            if isinstance(items, str):
                p = tf.paragraphs[0]
                p.text = items
                p.font.name = 'Arial'
                try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
                p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
            elif isinstance(items, list):
                for j, it in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"• {it}"
                    p.level = 0
                    p.font.name = 'Arial'
                    try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
                    p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        return slide
    
    @staticmethod
    def _layout_matrix_slide(slide, content: Dict[str, Any]):
        """2x2 매트릭스 레이아웃"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '2x2 Matrix')
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 매트릭스 4개 영역
        quadrants = [
            {'left': 0.5, 'top': 1.5, 'title': 'High/High'},  # 좌상
            {'left': 6.9, 'top': 1.5, 'title': 'Low/High'},   # 우상
            {'left': 0.5, 'top': 4.0, 'title': 'High/Low'},   # 좌하
            {'left': 6.9, 'top': 4.0, 'title': 'Low/Low'}     # 우하
        ]
        
        matrix_data = content.get('matrix', {})
        
        for i, quad in enumerate(quadrants):
            box = slide.shapes.add_textbox(
                Inches(quad['left']),
                Inches(quad['top']),
                Inches(5.4),
                Inches(2.0)
            )
            text_frame = box.text_frame
            text_frame.word_wrap = True
            
            # Quadrant 제목
            p = text_frame.paragraphs[0]
            quad_key = f"q{i+1}"
            quad_data = matrix_data.get(quad_key, {})
            p.text = quad_data.get('title', quad['title'])
            p.font.bold = True
            try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['primary']
            
            # Quadrant 내용
            quad_content = quad_data.get('content', [])
            if isinstance(quad_content, str):
                quad_content = [quad_content]
            
            for item in quad_content[:2]:  # 최대 2개 항목
                p = text_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        return slide
    
    @staticmethod
    def _layout_bullet_slide(slide, content: Dict[str, Any]):
        """불릿 포인트 슬라이드"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '')
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 불릿 포인트 영역
        body_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(12.33),
            Inches(5.0)
        )
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        
        # 불릿 포인트 추가
        bullets = content.get('bullets', content.get('body', []))
        if isinstance(bullets, str):
            bullets = [bullets]
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
            p.level = 0  # 불릿 레벨
            p.space_after = Pt(12)
        
        return slide
    
    @staticmethod
    def _layout_content_slide(slide, content: Dict[str, Any]):
        """일반 콘텐츠 슬라이드 (기본)"""
        # 제목
        title_box = slide.shapes.add_textbox(
            McKinseyLayoutManager.MARGINS['left'],
            McKinseyLayoutManager.MARGINS['top'],
            Inches(12.33),
            Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '')
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        # 본문
        body_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(12.33),
            Inches(5.0)
        )
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        
        body_content = content.get('body', content.get('content', ''))
        
        # 리스트인 경우 불릿 포인트로
        if isinstance(body_content, list):
            for i, item in enumerate(body_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {item}"
                try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
                p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
                p.space_after = Pt(6)
        else:
            p = text_frame.paragraphs[0]
            p.text = str(body_content)
            try:
                p.font.name = '맑은 고딕'
            except Exception:
                p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = McKinseyLayoutManager.COLORS['text']
        
        return slide
    
    @staticmethod
    def determine_layout_type(slide_content: Dict[str, Any], slide_idx: int) -> str:
        """
        슬라이드 내용에 따라 적절한 레이아웃 타입 결정
        """
        # 첫 번째 슬라이드는 제목 슬라이드
        if slide_idx == 0:
            return 'title_slide'
        
        # 차트가 있으면 차트 슬라이드
        if slide_content.get('has_chart') or slide_content.get('chart_data'):
            return 'chart_slide'
        
        # 2단 구성이면 two_column
        if slide_content.get('has_columns') or (
            slide_content.get('left_content') and slide_content.get('right_content')
        ):
            return 'two_column'
        
        # 매트릭스 구성
        if slide_content.get('matrix') or '2x2' in slide_content.get('title', '').lower():
            return 'matrix_slide'
        
        # 불릿 포인트가 많으면 bullet_slide
        body = slide_content.get('body', slide_content.get('content', ''))
        if isinstance(body, list) and len(body) > 2:
            return 'bullet_slide'
        
        # 기본은 content_slide
        return 'content_slide'
