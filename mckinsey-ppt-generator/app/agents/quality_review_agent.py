"""
품질 검토 에이전트
So What 테스트와 McKinsey 품질 기준으로 검증합니다.
"""

from typing import Dict, List
from pptx import Presentation # Added
from app.agents.base_agent_v2 import BaseAgentV2
from app.db.models import AgentType
from app.services.quality_controller import QualityController, QualityScore # Added QualityScore
from app.services.headline_generator import HeadlineGenerator, SoWhatTester
from loguru import logger


class QualityReviewAgent(BaseAgentV2):
    """
    McKinsey 품질 검토 파트너
    
    역할:
    - So What 테스트 실행
    - MECE 원칙 검증
    - 정량화 체크
    - 개선 제안 생성
    """
    
    def __init__(self):
        super().__init__(
            name="Quality Review Agent",
            agent_type=AgentType.REVIEWER,
            model="gpt-4-turbo"
        )
        self.quality_controller = QualityController()
        self.headline_generator = HeadlineGenerator()
        self.so_what_tester = SoWhatTester()
        self.metrics = {}

    async def process(self, input_data: Dict, context: Dict) -> Dict:
        """
        품질 검토 파이프라인
        
        Args:
            input_data: {
                'pptx_file_path': str, # Path to the generated PPTX file
                'insights': List[Dict],
                'pyramid': Dict
            }
            context: {
                'job_id': str
            }
        
        Returns:
            {
                'quality_scores': Dict, # Will be converted from QualityScore object
                'so_what_results': List[Dict],
                'improvements': List[Dict],
                'passed': bool
            }
        """
        logger.info(f"Starting quality review for job {context.get('job_id')}")
        
        try:
            pptx_file_path = input_data.get('pptx_file_path')
            insights = input_data.get('insights', [])

            if not pptx_file_path:
                raise ValueError("pptx_file_path is required for QualityReviewAgent")

            # Load the presentation
            prs = Presentation(pptx_file_path)
            
            # Step 1: So What 테스트
            logger.info("Step 1: Running So What tests")
            so_what_results = await self._run_so_what_tests(prs) # Pass prs
            
            # Step 2: 품질 점수 계산
            logger.info("Step 2: Calculating quality scores")
            # _calculate_quality_scores will now return a QualityScore object
            quality_score_obj = await self._calculate_quality_scores(prs, insights) # Pass prs
            
            # Convert QualityScore object to dict for return consistency
            quality_scores_dict = {
                'clarity': quality_score_obj.clarity,
                'insight': quality_score_obj.insight,
                'structure': quality_score_obj.structure,
                'visual': quality_score_obj.visual,
                'actionability': quality_score_obj.actionability,
                'total_score': quality_score_obj.total,
                'details': quality_score_obj.details # Include details
            }
            
            # Step 3: 개선 제안
            logger.info("Step 3: Generating improvements")
            improvements = await self._generate_improvements(quality_score_obj, so_what_results) # Pass QualityScore object
            
            passed = quality_score_obj.passed # Use passed attribute from QualityScore object
            
            self.metrics['quality_score'] = quality_score_obj.total
            self.metrics['improvements_suggested'] = len(improvements)
            
            return {
                'quality_scores': quality_scores_dict, # Return dict version
                'so_what_results': so_what_results,
                'improvements': improvements,
                'passed': passed
            }
            
        except Exception as e:
            logger.error(f"Quality review failed: {e}")
            raise
    
    async def _run_so_what_tests(self, prs: Presentation) -> List[Dict]: # Changed signature
        """
        각 슬라이드 헤드라인에 So What 테스트 실행
        
        Returns:
            [{
                'slide_number': int,
                'headline': str,
                'score': float,
                'passed': bool,
                'issues': List[str]
            }]
        """
        results = []
        slide_number = 0 # Track slide number

        for slide in prs.slides: # Iterate through actual pptx slides
            slide_number += 1
            headline = self._get_title_text(slide) # Helper to extract title from pptx slide

            if headline: # Only test if a headline exists
                # SoWhatTester 사용
                test_result = self.so_what_tester.test(headline)

                result = {
                    'slide_number': slide_number,
                    'headline': headline,
                    'score': test_result.get('score', 0),
                    'passed': test_result.get('score', 0) >= 0.7,
                    'issues': test_result.get('issues', [])
                }
                results.append(result)

        return results
    
    async def _calculate_quality_scores(
        self,
        prs: Presentation, # Changed signature
        insights: List[Dict]
    ) -> QualityScore: # Changed return type
        """
        5가지 품질 기준 평가
        
        기준:
        1. Clarity (명확성): 20%
        2. Insight (인사이트): 25%
        3. Structure (구조): 20%
        4. Visual (시각): 15%
        5. Actionability (실행성): 20%
        
        Returns:
            {
                'clarity': float,
                'insight': float,
                'structure': float,
                'visual': float,
                'actionability': float,
                'total_score': float
            }
        """
        # QualityController 활용
        # Pass the Presentation object directly
        evaluation_result = self.quality_controller.evaluate(prs) # Call evaluate directly

        return evaluation_result # Return the QualityScore object directly
    
    async def _generate_improvements(
        self,
        quality_score_obj: QualityScore, # Changed type hint and variable name
        so_what_results: List[Dict]
    ) -> List[Dict]:
        """
        개선 제안 생성
        
        Returns:
            [{
                'category': str,  # clarity/insight/structure/visual/actionability
                'priority': str,  # high/medium/low
                'issue': str,
                'suggestion': str
            }]
        """
        improvements = []
        
        # 점수가 낮은 영역 식별
        threshold = 0.8
        
        # Access attributes directly from quality_score_obj
        # Iterate over the defined criteria
        criteria = {
            'clarity': quality_score_obj.clarity,
            'insight': quality_score_obj.insight,
            'structure': quality_score_obj.structure,
            'visual': quality_score_obj.visual,
            'actionability': quality_score_obj.actionability
        }

        for category, score in criteria.items(): # Iterate over criteria dict
            if score < threshold:
                improvement = {
                    'category': category,
                    'priority': 'high' if score < 0.7 else 'medium',
                    'issue': f'{category} 점수가 {score:.2f}로 낮음',
                    'suggestion': self._get_improvement_suggestion(category)
                }
                improvements.append(improvement)
        
        # So What 테스트 실패한 슬라이드
        failed_so_what = [r for r in so_what_results if not r['passed']]
        if failed_so_what:
            improvements.append({
                'category': 'clarity',
                'priority': 'high',
                'issue': f'{len(failed_so_what)}개 슬라이드가 So What 테스트 실패',
                'suggestion': 'HeadlineGenerator로 헤드라인 재생성 필요'
            })
        
        return improvements
    
    def _get_improvement_suggestion(self, category: str) -> str:
        """카테고리별 개선 제안"""
        suggestions = {
            'clarity': 'HeadlineGenerator를 사용하여 So What이 명확한 헤드라인으로 재작성',
            'insight': 'InsightLadder로 4단계 인사이트를 더 깊이 있게 도출',
            'structure': 'MECE 원칙을 재검토하고 논리적 흐름 개선',
            'visual': 'McKinsey 스타일 가이드를 엄격히 준수하도록 디자인 재적용',
            'actionability': '구체적 수치와 우선순위를 포함한 실행 권고 추가'
        }
        
        return suggestions.get(category, '전문가 검토 필요')

    def _has_title(self, slide) -> bool:
        """슬라이드에 제목이 있는지 확인"""
        try:
            return slide.shapes.title is not None and slide.shapes.title.has_text_frame
        except:
            return False

    def _get_title_text(self, slide) -> str:
        """
        슬라이드 제목 텍스트 추출
        """
        try:
            if self._has_title(slide):
                return slide.shapes.title.text
        except:
            pass
        return ""
